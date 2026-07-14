from __future__ import annotations

import io
from pathlib import Path

import pytest
from pptx import Presentation

from backend.analytics_module.pptx_builder.chart_resolver import PPTXChartResolver
from backend.analytics_module.pptx_builder.engine import PPTXEngine
from backend.analytics_module.pptx_builder.layout import REFERENCE_HEIGHT_EMU, REFERENCE_WIDTH_EMU
from backend.analytics_module.pptx_builder.presentation_planner import PresentationPlanner
from backend.analytics_module.pptx_builder.pptx_export_audit import audit_pptx_bytes
from backend.analytics_module.pptx_builder.pptx_geometry_audit import audit_presentation_geometry
from backend.analytics_module.pptx_builder.validation_gating import PPTXValidationMode
from backend.analytics_module.pptx_builder.validator import PPTXIntegrityValidator
from backend.tests.analytics.pptx_acceptance_checklist import (
    build_acceptance_checklist,
    build_golden_smoke_checklist,
)
from backend.tests.analytics.pptx_acceptance_contract import (
    KNOWN_BAD_DECK_FILENAME,
    KNOWN_REPORT_ID,
    PROTEIN_BAR_SCREEN_CHARTS,
    REPRESENTATIVE_SCREEN_CHARTS,
    build_acceptance_generation_report,
    build_known_report_fixture,
    build_protein_bar_screen_report,
    build_representative_screen_report,
    expected_backend_key,
)
from backend.analytics_module.pptx_builder.chart_contracts import TASTE_TEST_CONTRACTS

REPO_ROOT = Path(__file__).resolve().parents[3]


@pytest.fixture(scope="session")
def acceptance_generation_report() -> dict:
    return build_acceptance_generation_report()


@pytest.fixture(scope="session")
def representative_screen_report() -> dict:
    return build_representative_screen_report()


@pytest.fixture(scope="session")
def known_report_fixture() -> dict:
    return build_known_report_fixture()


@pytest.fixture(scope="session")
def protein_bar_screen_report() -> dict:
    return build_protein_bar_screen_report()


def _generate_acceptance_deck(report_doc: dict, marketeers_template_path: Path):
    engine = PPTXEngine(template_path=str(marketeers_template_path))
    intents = PresentationPlanner.define_slide_intents(report_doc)
    pptx_stream, slide_count = engine.generate_presentation(intents)
    presentation = Presentation(pptx_stream)
    return engine, intents, presentation, pptx_stream, slide_count


def test_marketeers_template_dimensions_match_reference_canvas(marketeers_presentation):
    assert marketeers_presentation.slide_width == REFERENCE_WIDTH_EMU
    assert marketeers_presentation.slide_height == REFERENCE_HEIGHT_EMU


def test_representative_chart_ids_match_frontend_dispatch_contract():
    resolver = PPTXChartResolver()

    for chart in REPRESENTATIVE_SCREEN_CHARTS:
        resolution = resolver.resolve(chart)
        assert resolution.registry_key == chart["expected_backend_key"]
        assert resolution.uses_fallback_table is False
        assert resolution.registry_key == expected_backend_key(chart)


@pytest.mark.asyncio
async def test_generated_acceptance_deck_records_bounds_and_duplicate_titles(
    marketeers_template_path,
    acceptance_generation_report,
):
    _engine, intents, presentation, pptx_stream, _slide_count = _generate_acceptance_deck(
        acceptance_generation_report,
        marketeers_template_path,
    )
    geometry = audit_presentation_geometry(presentation)
    assert geometry["layout_warning_count"] == 0
    assert geometry["duplicate_title_count"] == 0

    validator = PPTXIntegrityValidator(io.BytesIO(pptx_stream.getvalue()))
    validation = await validator.validate(
        acceptance_generation_report,
        intents,
        mode=PPTXValidationMode.PRODUCTION,
    )
    assert validation["layout_warning_count"] == geometry["layout_warning_count"]
    assert validation["duplicate_title_count"] == geometry["duplicate_title_count"]
    assert validation["unsupported_count"] == 0
    assert validation["error_placeholder_count"] == 0


@pytest.mark.asyncio
async def test_known_report_golden_smoke_checklist_matches_exported_pptx(
    marketeers_template_path,
    known_report_fixture,
):
    engine, intents, presentation, pptx_stream, _slide_count = _generate_acceptance_deck(
        known_report_fixture,
        marketeers_template_path,
    )
    validator = PPTXIntegrityValidator(io.BytesIO(pptx_stream.getvalue()))
    validation = await validator.validate(
        known_report_fixture,
        intents,
        mode=PPTXValidationMode.PRODUCTION,
        render_journal=engine.render_journal,
        narrative_journal=engine.narrative_render_journal,
    )
    smoke = build_golden_smoke_checklist(
        known_report_fixture,
        presentation,
        intents,
        validation,
        render_journal=engine.render_journal,
        narrative_journal=engine.narrative_render_journal,
    )

    assert smoke["passed"] is True, smoke["failed_checks"]
    assert validation["passes_gate"] is True


def test_representative_screen_report_fixture_matches_screen_chart_contract(
    representative_screen_report,
):
    fixture_ids = {chart["chart_id"] for chart in representative_screen_report["charts"]}
    contract_ids = {chart["chart_id"] for chart in REPRESENTATIVE_SCREEN_CHARTS}
    assert fixture_ids == contract_ids


@pytest.mark.asyncio
async def test_representative_screen_report_golden_export_passes_acceptance(
    marketeers_template_path,
    representative_screen_report,
):
    engine, intents, presentation, pptx_stream, _slide_count = _generate_acceptance_deck(
        representative_screen_report,
        marketeers_template_path,
    )
    validator = PPTXIntegrityValidator(io.BytesIO(pptx_stream.getvalue()))
    validation = await validator.validate(
        representative_screen_report,
        intents,
        mode=PPTXValidationMode.PRODUCTION,
        render_journal=engine.render_journal,
        narrative_journal=engine.narrative_render_journal,
    )
    checklist = build_acceptance_checklist(
        representative_screen_report,
        presentation,
        intents,
        validation,
        render_journal=engine.render_journal,
        narrative_journal=engine.narrative_render_journal,
    )

    assert checklist["passed"] is True, checklist["failed_checks"]
    assert validation["passes_gate"] is True


@pytest.mark.asyncio
async def test_representative_screen_report_golden_snapshot_fingerprint(
    marketeers_template_path,
    representative_screen_report,
):
    engine, intents, _presentation, pptx_stream, _slide_count = _generate_acceptance_deck(
        representative_screen_report,
        marketeers_template_path,
    )
    validator = PPTXIntegrityValidator(io.BytesIO(pptx_stream.getvalue()))
    validation = await validator.validate(
        representative_screen_report,
        intents,
        mode=PPTXValidationMode.PRODUCTION,
        render_journal=engine.render_journal,
        narrative_journal=engine.narrative_render_journal,
    )
    fingerprint = {
        "passes_gate": validation["passes_gate"],
        "unsupported_count": validation["unsupported_count"],
        "error_placeholder_count": validation["error_placeholder_count"],
        "layout_warning_count": validation["layout_warning_count"],
        "duplicate_title_count": validation["duplicate_title_count"],
        "markers": {
            "executive_summary": validation["export_audit"].get("text_markers", {}).get("executive_summary", 0),
            "strategic_positioning": validation["export_audit"].get("text_markers", {}).get("strategic_positioning", 0),
            "competitive_swot": validation["export_audit"].get("text_markers", {}).get("competitive_swot", 0),
            "recommendations_4p": validation["export_audit"].get("text_markers", {}).get("recommendations_4p", 0),
        },
    }

    assert fingerprint["passes_gate"] is True
    assert fingerprint["unsupported_count"] == 0
    assert fingerprint["error_placeholder_count"] == 0
    assert fingerprint["layout_warning_count"] == 0
    assert fingerprint["duplicate_title_count"] == 0
    assert fingerprint["markers"]["executive_summary"] >= 1
    assert fingerprint["markers"]["strategic_positioning"] >= 1
    assert fingerprint["markers"]["competitive_swot"] >= 1
    assert fingerprint["markers"]["recommendations_4p"] >= 1


def test_protein_bar_screen_report_fixture_matches_screen_chart_contract(
    protein_bar_screen_report,
):
    fixture_ids = {chart["chart_id"] for chart in protein_bar_screen_report["charts"]}
    contract_ids = {chart["chart_id"] for chart in PROTEIN_BAR_SCREEN_CHARTS}
    assert fixture_ids == contract_ids


@pytest.mark.asyncio
async def test_protein_bar_screen_report_golden_export_passes_acceptance(
    marketeers_template_path,
    protein_bar_screen_report,
):
    engine, intents, presentation, pptx_stream, _slide_count = _generate_acceptance_deck(
        protein_bar_screen_report,
        marketeers_template_path,
    )
    validator = PPTXIntegrityValidator(io.BytesIO(pptx_stream.getvalue()))
    validation = await validator.validate(
        protein_bar_screen_report,
        intents,
        mode=PPTXValidationMode.PRODUCTION,
        render_journal=engine.render_journal,
        narrative_journal=engine.narrative_render_journal,
    )
    checklist = build_acceptance_checklist(
        protein_bar_screen_report,
        presentation,
        intents,
        validation,
        render_journal=engine.render_journal,
        narrative_journal=engine.narrative_render_journal,
    )

    assert checklist["passed"] is True, checklist["failed_checks"]
    assert validation["passes_gate"] is True


def test_known_bad_deck_regression_fixture_fails_golden_expectations():
    bad_deck_path = REPO_ROOT / KNOWN_BAD_DECK_FILENAME
    if not bad_deck_path.is_file():
        bad_deck_path = REPO_ROOT / "report_69ce229eeed39ea9d5282afa (15).pptx"
    if not bad_deck_path.is_file():
        pytest.skip(f"Known bad deck missing: {bad_deck_path}")

    audit = audit_pptx_bytes(bad_deck_path.read_bytes())
    assert audit["text_markers"]["analysis_interrupted"] > 0
    assert audit["text_markers"]["executive_summary"] == 0
    assert audit["error_placeholder_count"] > 0
    assert audit["unsupported_placeholder_count"] >= 0
    assert KNOWN_REPORT_ID in KNOWN_BAD_DECK_FILENAME


@pytest.mark.asyncio
@pytest.mark.parametrize(
    "contract_pattern",
    [contract.chart_id_pattern for contract in TASTE_TEST_CONTRACTS],
)
async def test_each_contract_builds_one_chart_deck_with_production_validation(
    marketeers_template_path,
    contract_pattern,
):
    contract = next(c for c in TASTE_TEST_CONTRACTS if c.chart_id_pattern == contract_pattern)
    chart_id = contract_pattern.replace("*", "sample")
    chart = {
        "chart_id": chart_id,
        "chart_type": contract.chart_type,
        "title": f"Contract {chart_id}",
        "data": contract.golden_sample_factory(),
    }
    report_doc = {
        "project_name": "Contract Validation Deck",
        "brand": "Hero Brand",
        "metadata": {
            "title": "Contract Validation Deck",
            "brand": "Hero Brand",
            "brands": ["Hero Brand", "Competitor A"],
            "company_name": "Marketeers",
            "date": "July 2026",
            "research_type": "Taste Test",
        },
        "charts": [chart],
        "insights": {},
    }

    engine, intents, presentation, pptx_stream, _slide_count = _generate_acceptance_deck(
        report_doc,
        marketeers_template_path,
    )
    validation = await PPTXIntegrityValidator(io.BytesIO(pptx_stream.getvalue())).validate(
        report_doc,
        intents,
        mode=PPTXValidationMode.PRODUCTION,
        render_journal=engine.render_journal,
        narrative_journal=engine.narrative_render_journal,
    )
    assert validation["is_corrupt"] is False
    assert validation["unsupported_count"] == 0
    assert validation["error_placeholder_count"] == 0
    entry = next((item for item in engine.render_journal if item.get("chart_id") == chart_id), None)
    if contract.chart_id_pattern.startswith("brand_card_"):
        # Brand cards are intentionally narrative-driven, not content-slide render_journal entries.
        assert any(
            str(item.get("section_id", "")).startswith("brand_profile::")
            for item in engine.narrative_render_journal
        )
    else:
        assert entry is not None, f"Missing render journal entry for {contract_pattern}"
        assert entry.get("registry_key") == contract.builder_registry_key
        assert entry.get("resolution_source") != "fallback_table"
        assert entry.get("render_status") == "rendered"

        chart_shapes = []
        table_shapes = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_chart", False):
                    chart_shapes.append(shape)
                if getattr(shape, "has_table", False):
                    table_shapes.append(shape)
        if contract.chart_type in {"wordcloud", "verbatim_analysis", "scorecard", "funnel_ratio_cards"}:
            assert len(presentation.slides) > 0
        else:
            assert chart_shapes or table_shapes
