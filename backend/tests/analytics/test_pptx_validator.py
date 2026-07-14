import io

import pytest
from pptx import Presentation
from pptx.util import Inches

from backend.analytics_module.pptx_builder.presentation_planner import SlideIntent, SlideType
from backend.analytics_module.pptx_builder.validation_gating import PPTXValidationMode
from backend.analytics_module.pptx_builder.validator import PPTXIntegrityValidator


def _build_minimal_pptx_bytes(*, title: str = "TITLE", body: str = "") -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    if body:
        textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
        textbox.text_frame.text = body
    stream = io.BytesIO()
    presentation.save(stream)
    return stream.getvalue()


@pytest.mark.asyncio
async def test_validator_reports_unsupported_count_from_payload():
    pptx_bytes = _build_minimal_pptx_bytes(title="Chart", body="Chart")
    validator = PPTXIntegrityValidator(io.BytesIO(pptx_bytes))
    report_doc = {
        "charts": [
            {"chart_type": "scatter_plot", "title": "Supported"},
            {"chart_type": "mystery_visual", "title": "Unsupported"},
        ]
    }

    result = await validator.validate(report_doc, intents=[], mode=PPTXValidationMode.QA)
    assert result["unsupported_count"] == 1


@pytest.mark.asyncio
async def test_validator_detects_unsupported_and_error_placeholders():
    pptx_bytes = _build_minimal_pptx_bytes(
        title="Broken Chart",
        body="DEVELOPMENT IN PROGRESS: NATIVE 'RADAR' RENDERER",
    )
    validator = PPTXIntegrityValidator(io.BytesIO(pptx_bytes))
    result = await validator.validate({"charts": []}, intents=[], mode=PPTXValidationMode.QA)

    assert result["unsupported_count"] >= 1
    assert result["validation_errors"]
    assert any("Unsupported renderer placeholder" in message for message in result["validation_errors"])


@pytest.mark.asyncio
async def test_qa_mode_allows_export_with_warnings(monkeypatch):
    monkeypatch.delenv("PPTX_VALIDATION_MODE", raising=False)
    pptx_bytes = _build_minimal_pptx_bytes(
        title="Broken Chart",
        body="DEVELOPMENT IN PROGRESS: NATIVE 'RADAR' RENDERER",
    )
    validator = PPTXIntegrityValidator(io.BytesIO(pptx_bytes))
    result = await validator.validate({"charts": []}, intents=[], mode=PPTXValidationMode.QA)

    assert result["passes_gate"] is True
    assert result["validation_errors"]


@pytest.mark.asyncio
async def test_production_mode_fails_gate_on_critical_issues():
    pptx_bytes = _build_minimal_pptx_bytes(
        title="Broken Chart",
        body="ANALYSIS INTERRUPTED",
    )
    validator = PPTXIntegrityValidator(io.BytesIO(pptx_bytes))
    result = await validator.validate({"charts": []}, intents=[], mode=PPTXValidationMode.PRODUCTION)

    assert result["passes_gate"] is False
    assert result["error_placeholder_count"] >= 1


@pytest.mark.asyncio
async def test_validator_flags_missing_required_executive_summary():
    pptx_bytes = _build_minimal_pptx_bytes(title="Charts Only", body="No narrative here.")
    validator = PPTXIntegrityValidator(io.BytesIO(pptx_bytes))
    intents = [SlideIntent(SlideType.EXECUTIVE_SUMMARY, data={"executive_summary": "Summary"})]
    result = await validator.validate(
        {"insights": {"executive_summary": "Summary"}},
        intents,
        mode=PPTXValidationMode.PRODUCTION,
    )

    assert result["passes_gate"] is False
    assert any("executive summary" in message.lower() for message in result["validation_errors"])


@pytest.mark.asyncio
async def test_validator_detects_out_of_bounds_shapes_in_qa_mode():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Overflow"
    slide.shapes.add_textbox(
        presentation.slide_width + Inches(1),
        Inches(1),
        Inches(2),
        Inches(1),
    )
    stream = io.BytesIO()
    presentation.save(stream)

    validator = PPTXIntegrityValidator(io.BytesIO(stream.getvalue()))
    result = await validator.validate({"charts": []}, intents=[], mode=PPTXValidationMode.QA)

    assert result["layout_warning_count"] >= 1
    assert result["passes_gate"] is True
    assert result["validation_warnings"]


@pytest.mark.asyncio
async def test_production_mode_fails_on_out_of_bounds_shapes():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Overflow"
    slide.shapes.add_textbox(
        presentation.slide_width + Inches(1),
        Inches(1),
        Inches(2),
        Inches(1),
    )
    stream = io.BytesIO()
    presentation.save(stream)

    validator = PPTXIntegrityValidator(io.BytesIO(stream.getvalue()))
    result = await validator.validate({"charts": []}, intents=[], mode=PPTXValidationMode.PRODUCTION)

    assert result["passes_gate"] is False
    assert any("extends outside the slide canvas" in message for message in result["validation_errors"])


@pytest.mark.asyncio
async def test_duplicate_titles_allow_section_divider_repetitions():
    presentation = Presentation()
    for _ in range(2):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Brand Awareness"
    stream = io.BytesIO()
    presentation.save(stream)

    validator = PPTXIntegrityValidator(io.BytesIO(stream.getvalue()))
    intents = [
        SlideIntent(SlideType.SECTION_DIVIDER, title="Brand Awareness"),
        SlideIntent(SlideType.SECTION_DIVIDER, title="Brand Awareness"),
    ]

    result = await validator.validate({"charts": []}, intents=intents, mode=PPTXValidationMode.PRODUCTION)

    assert result["passes_gate"] is True
    assert not any("Duplicate slide title" in message for message in result["validation_errors"])


@pytest.mark.asyncio
async def test_duplicate_titles_fail_without_section_divider_allowance():
    presentation = Presentation()
    for _ in range(2):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Brand Awareness"
    stream = io.BytesIO()
    presentation.save(stream)

    validator = PPTXIntegrityValidator(io.BytesIO(stream.getvalue()))
    result = await validator.validate({"charts": []}, intents=[], mode=PPTXValidationMode.PRODUCTION)

    assert result["passes_gate"] is False
    assert any("Duplicate slide title" in message for message in result["validation_errors"])


@pytest.mark.asyncio
async def test_production_mode_fails_when_rendered_chart_has_no_native_shape():
    pptx_bytes = _build_minimal_pptx_bytes(title="Empty Chart", body="No native chart here.")
    validator = PPTXIntegrityValidator(io.BytesIO(pptx_bytes))
    result = await validator.validate(
        {"charts": [{"chart_id": "overall_scatter", "chart_type": "scatter_plot"}]},
        intents=[SlideIntent(SlideType.CONTENT_SLIDE, data={"chart_type": "scatter_plot"})],
        mode=PPTXValidationMode.PRODUCTION,
        render_journal=[
            {
                "chart_id": "overall_scatter",
                "chart_type": "scatter_plot",
                "slide_index": 1,
                "render_status": "rendered",
            }
        ],
    )

    assert result["passes_gate"] is False
    assert any("rendered without a native chart shape" in message for message in result["validation_errors"])


@pytest.mark.asyncio
async def test_validator_exposes_render_tally_and_missing_narrative_sections():
    pptx_bytes = _build_minimal_pptx_bytes(title="Charts Only", body="No narrative here.")
    validator = PPTXIntegrityValidator(io.BytesIO(pptx_bytes))
    intents = [SlideIntent(SlideType.EXECUTIVE_SUMMARY, data={"executive_summary": "Summary"})]
    result = await validator.validate(
        {"insights": {"executive_summary": "Summary"}},
        intents,
        mode=PPTXValidationMode.PRODUCTION,
        render_journal=[{"render_status": "failed", "chart_id": "x", "slide_index": 1}],
    )

    assert result["render_tally"]["failed_chart_count"] == 1
    assert result["missing_narrative_sections"]
    assert any(section["section_id"] == "executive_summary" for section in result["missing_narrative_sections"])


@pytest.mark.asyncio
async def test_validator_blocks_unwhitelisted_fallback_table(monkeypatch):
    monkeypatch.setenv("PPTX_FALLBACK_TABLE_WHITELIST", "")
    pptx_bytes = _build_minimal_pptx_bytes(title="Chart", body="Chart")
    validator = PPTXIntegrityValidator(io.BytesIO(pptx_bytes))
    report_doc = {
        "charts": [
            {
                "chart_id": "mystery_chart",
                "chart_type": "unknown_type",
                "_resolution": {"uses_fallback_table": True},
            }
        ],
        "metadata": {"title": "Chart", "company_name": "Chart"},
    }
    result = await validator.validate(report_doc, intents=[], mode=PPTXValidationMode.PRODUCTION)
    assert result["passes_gate"] is False
    assert any("fallback table" in message.lower() for message in result["validation_errors"])


@pytest.mark.asyncio
async def test_validator_allows_whitelisted_fallback_table(monkeypatch):
    monkeypatch.setenv("PPTX_FALLBACK_TABLE_WHITELIST", "mystery_chart")
    pptx_bytes = _build_minimal_pptx_bytes(title="Chart", body="MYSTERY_CHART")
    validator = PPTXIntegrityValidator(io.BytesIO(pptx_bytes))
    report_doc = {
        "charts": [
            {
                "chart_id": "mystery_chart",
                "chart_type": "unknown_type",
                "_resolution": {"uses_fallback_table": True},
            }
        ],
        "metadata": {"title": "Chart", "company_name": "Chart"},
    }
    result = await validator.validate(report_doc, intents=[], mode=PPTXValidationMode.QA)
    assert not any("fallback table" in message.lower() for message in result["validation_errors"])


@pytest.mark.asyncio
async def test_validator_flags_missing_planned_chart_mapping():
    pptx_bytes = _build_minimal_pptx_bytes(title="Deck", body="DECK")
    validator = PPTXIntegrityValidator(io.BytesIO(pptx_bytes))
    report_doc = {
        "charts": [{"chart_id": "chart_present", "chart_type": "horizontal_bar"}],
        "metadata": {"title": "Deck", "company_name": "Deck"},
    }
    intents = [
        SlideIntent(SlideType.CONTENT_SLIDE, data={"chart_id": "chart_missing", "chart_type": "horizontal_bar"})
    ]
    result = await validator.validate(report_doc, intents=intents, mode=PPTXValidationMode.PRODUCTION)
    assert result["passes_gate"] is False
    assert any("missing from report payload" in message.lower() for message in result["validation_errors"])
