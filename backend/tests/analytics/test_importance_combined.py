import pytest
from backend.analytics_module.pptx_builder.theme import PPTXTheme
from backend.analytics_module.pptx_builder.pptx_importance_combined import PPTXImportanceCombined
from pptx.enum.chart import XL_CHART_TYPE

@pytest.fixture
def theme():
    return PPTXTheme()

def test_importance_combined_render(marketeers_presentation, marketeers_layout, theme):
    """
    Verification Phase 5:
    Ensures that the dual-panel scatter builder correctly renders two separate 
    XY scatter charts with the specified series data and panel titles.
    """
    builder = PPTXImportanceCombined(theme, marketeers_layout)
    
    # Mock unified payload structured as per Aggregator requirement
    sample_data = {
        "title": "Combined Importance Map",
        "subtitle": "Drill-down Analysis",
        "data": {
            "main_scatter": {
                "datasets": [
                    {
                        "brand": "Hero Brand",
                        "data": [{"x": 80, "y": 70, "attribute": "Quality"}]
                    },
                    {
                        "brand": "Competitor A",
                        "data": [{"x": 60, "y": 50, "attribute": "Price"}]
                    }
                ]
            },
            "sub_scatter": {
                "datasets": [
                    {
                        "brand": "Hero Brand",
                        "data": [
                            {"x": 90, "y": 85, "sub_attribute": "Durability"},
                            {"x": 40, "y": 30, "sub_attribute": "Reliability"}
                        ]
                    }
                ],
                "drill_attribute": "Quality"
            }
        }
    }
    
    # Prepare slide (using a standard layout from the session presentation)
    slide_layout = marketeers_presentation.slide_layouts[1] 
    slide = marketeers_presentation.slides.add_slide(slide_layout)
    
    # Execution: Build the slide
    builder.build(slide, sample_data)
    
    # Assertion 1: Slide contains exactly 2 native chart objects
    charts = [shape.chart for shape in slide.shapes if hasattr(shape, "chart")]
    assert len(charts) == 2, f"Expected 2 charts, found {len(charts)}"
    
    # Assertion 2: Both charts are XY_SCATTER type
    for i, chart in enumerate(charts):
        assert chart.chart_type == XL_CHART_TYPE.XY_SCATTER, f"Chart {i} is not XY_SCATTER"
        
    # Assertion 3: Correct number of series per chart (matching brand count)
    # Main chart has 2 brands
    assert len(charts[0].series) == 2, "Main scatter should have 2 brand series"
    # Sub chart has 1 brand
    assert len(charts[1].series) == 1, "Sub scatter should have 1 brand series"
    
    # Assertion 4: Verification of Panel Titles as text shapes
    title_texts = [
        shape.text_frame.text 
        for shape in slide.shapes 
        if hasattr(shape, "text_frame") and shape.text_frame
    ]
    assert "Overall Scatter" in title_texts, "Overall Scatter panel title missing"
    assert "Sub Scatter" in title_texts, "Sub Scatter panel title missing"

    # Assertion 5: Verify legend status (Shared simulation: only on left chart)
    assert charts[0].has_legend is True, "Left chart should have the shared legend"
    assert charts[1].has_legend is False, "Right chart should have legend suppressed for clarity"
