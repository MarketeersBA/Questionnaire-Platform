import pytest
from pptx.enum.chart import XL_DATA_LABEL_POSITION
from backend.analytics_module.pptx_builder.theme import PPTXTheme
from backend.analytics_module.pptx_builder.pptx_executive_summary import PPTXExecutiveSummary
from backend.analytics_module.pptx_builder.pptx_swot import PPTXSwot
from backend.analytics_module.pptx_builder.pptx_funnel_cards import PPTXFunnelCards
from backend.analytics_module.pptx_builder.pptx_radar import PPTXRadar
from backend.analytics_module.pptx_builder.pptx_grouped_bar import PPTXGroupedBar
from backend.analytics_module.pptx_builder.pptx_brand_profile_card import PPTXBrandProfileCard
from backend.analytics_module.pptx_builder.pptx_stacked_bar import PPTXStackedBar
from backend.analytics_module.pptx_builder.pptx_wordcloud import PPTXWordcloud
from backend.analytics_module.pptx_builder.pptx_verbatim import PPTXVerbatim
from backend.analytics_module.pptx_builder.pptx_brand_comparison import PPTXBrandComparison
from backend.analytics_module.pptx_builder.pptx_snake_line import PPTXSnakeLine
from backend.analytics_module.pptx_builder.pptx_reference_table import PPTXReferenceTable
from backend.analytics_module.pptx_builder.pptx_waterfall_bar import PPTXWaterfallBar
from backend.analytics_module.pptx_builder.pptx_generic_table import PPTXGenericTable


@pytest.fixture
def theme():
    return PPTXTheme()


def test_executive_summary_render(marketeers_presentation, marketeers_layout, theme):
    builder = PPTXExecutiveSummary(theme, marketeers_layout)
    context = {
        "executive_summary": "This is a summary",
        "key_findings": ["Finding 1", "Finding 2"],
        "opportunity_insights": [{"title": "O1", "insight": "D1"}],
    }
    slide = marketeers_presentation.slides.add_slide(marketeers_presentation.slide_layouts[1])
    builder.build(slide, context)
    assert len(slide.shapes) > 0


def test_swot_render(marketeers_presentation, marketeers_layout, theme):
    builder = PPTXSwot(theme, marketeers_layout)
    context = {
        "brand_name": "Test Brand",
        "swot": {
            "strengths": ["S1"],
            "weaknesses": ["W1"],
            "opportunities": ["O1"],
            "threats": ["T1"],
        },
    }
    slide = marketeers_presentation.slides.add_slide(marketeers_presentation.slide_layouts[1])
    builder.build(slide, context)
    assert len(slide.shapes) > 0


def test_funnel_cards_render(marketeers_presentation, marketeers_layout, theme):
    builder = PPTXFunnelCards(theme, marketeers_layout)
    context = {
        "data": {
            "brand_cards": [
                {
                    "brand": "Brand A",
                    "ratio_labels": ["A: 80%", "T: 50%"],
                }
            ]
        }
    }
    slide = marketeers_presentation.slides.add_slide(marketeers_presentation.slide_layouts[1])
    builder.build(slide, context)
    assert len(slide.shapes) > 0


def test_radar_render(marketeers_presentation, marketeers_layout, theme):
    builder = PPTXRadar(theme, marketeers_layout)
    context = {
        "data": {
            "labels": ["Attr 1", "Attr 2", "Attr 3"],
            "datasets": [
                {"label": "Brand A", "data": [4.5, 4.0, 4.2]},
                {"label": "Brand B", "data": [3.2, 3.5, 3.1]},
            ],
        },
    }
    slide = marketeers_presentation.slides.add_slide(marketeers_presentation.slide_layouts[1])
    builder.build(slide, context)
    charts = [shape for shape in slide.shapes if hasattr(shape, "chart")]
    assert len(charts) > 0


def test_grouped_bar_data_labels_enabled(marketeers_presentation, marketeers_layout, theme):
    builder = PPTXGroupedBar(theme, marketeers_layout)
    context = {
        "data": {
            "labels": ["Brand A", "Brand B"],
            "datasets": [{"label": "Overall", "data": [4.2, 3.8]}],
        }
    }
    slide = marketeers_presentation.slides.add_slide(marketeers_presentation.slide_layouts[1])
    builder.build(slide, context)
    chart = next(shape.chart for shape in slide.shapes if hasattr(shape, "chart"))
    labels = chart.series[0].data_labels
    assert labels.number_format == "0.0"
    assert labels.position == XL_DATA_LABEL_POSITION.OUTSIDE_END


def test_stacked_bar_data_labels_enabled(marketeers_presentation, marketeers_layout, theme):
    builder = PPTXStackedBar(theme, marketeers_layout)
    context = {
        "data": {
            "labels": ["Brand A", "Brand B"],
            "datasets": [
                {"label": "Aware", "data": [0.6, 0.5]},
                {"label": "Not Aware", "data": [0.4, 0.5]},
            ],
        }
    }
    slide = marketeers_presentation.slides.add_slide(marketeers_presentation.slide_layouts[1])
    builder.build(slide, context)
    chart = next(shape.chart for shape in slide.shapes if hasattr(shape, "chart"))
    labels = chart.series[0].data_labels
    assert labels.number_format == "0%"
    assert labels.position == XL_DATA_LABEL_POSITION.CENTER


def test_brand_profile_card_renders_with_layout_dimensions(marketeers_presentation, marketeers_layout, theme):
    builder = PPTXBrandProfileCard(theme, marketeers_layout)
    context = {
        "brand_data": {
            "data": {
                "profile": {
                    "Brand": "Hero Brand",
                    "Overall Score": 8.1,
                    "T2B %": 64,
                    "Evaluations": 412,
                },
                "strengths": [{"attribute": "Taste", "score": 8.5}],
            }
        },
        "ai_insight": "Hero brand leads on taste.",
        "brand_index": 1,
        "total_brands": 3,
    }
    slide = marketeers_presentation.slides.add_slide(marketeers_presentation.slide_layouts[1])
    builder.render(slide, context)
    assert len(slide.shapes) > 0


@pytest.mark.parametrize(
    "builder_cls, payload",
    [
        (PPTXWordcloud, {"data": {"words": []}}),
        (PPTXVerbatim, {"data": {"themes": [], "quotes": []}}),
        (PPTXBrandComparison, {"data": {"labels": [], "datasets": []}}),
        (PPTXSnakeLine, {"data": {"labels": [], "datasets": []}}),
        (PPTXReferenceTable, {"data": {"labels": [], "datasets": []}}),
        (PPTXWaterfallBar, {"data": {"labels": [], "datasets": []}}),
        (PPTXGenericTable, {"data": {"labels": ["A"], "datasets": []}}),
    ],
)
def test_soft_blank_builders_return_skipped_empty_data(
    builder_cls,
    payload,
    marketeers_presentation,
    marketeers_layout,
    theme,
):
    builder = builder_cls(theme, marketeers_layout)
    slide = marketeers_presentation.slides.add_slide(marketeers_presentation.slide_layouts[1])
    result = builder.build(slide, payload)
    assert result.status.value == "skipped_empty_data"


def test_snake_line_data_labels_enabled(marketeers_presentation, marketeers_layout, theme):
    builder = PPTXSnakeLine(theme, marketeers_layout)
    context = {
        "data": {
            "labels": ["Awareness", "Trial"],
            "datasets": [{"label": "Brand A", "data": [80, 55]}],
        }
    }
    slide = marketeers_presentation.slides.add_slide(marketeers_presentation.slide_layouts[1])
    builder.build(slide, context)
    chart = next(shape.chart for shape in slide.shapes if hasattr(shape, "chart"))
    labels = chart.series[0].data_labels
    assert labels.number_format == "0\"%\""
    assert labels.position == XL_DATA_LABEL_POSITION.ABOVE


def test_brand_comparison_data_labels_enabled(marketeers_presentation, marketeers_layout, theme):
    builder = PPTXBrandComparison(theme, marketeers_layout)
    context = {
        "data": {
            "labels": ["Brand A", "Brand B"],
            "datasets": [
                {"label": "Purchase Intent", "data": [72, 61]},
                {"label": "Overall Liking", "data": [81, 75]},
            ],
        }
    }
    slide = marketeers_presentation.slides.add_slide(marketeers_presentation.slide_layouts[1])
    builder.build(slide, context)
    chart = next(shape.chart for shape in slide.shapes if hasattr(shape, "chart"))
    labels = chart.series[0].data_labels
    assert labels.number_format == "0.0"
    assert labels.position == XL_DATA_LABEL_POSITION.OUTSIDE_END


def test_waterfall_data_labels_enabled(marketeers_presentation, marketeers_layout, theme):
    builder = PPTXWaterfallBar(theme, marketeers_layout)
    context = {
        "data": {
            "labels": ["Brand A", "Brand B"],
            "datasets": [
                {"label": "TOM", "data": [0.2, 0.15]},
                {"label": "Other_Unaided", "data": [0.35, 0.3]},
                {"label": "Aided", "data": [0.25, 0.4]},
            ],
        }
    }
    slide = marketeers_presentation.slides.add_slide(marketeers_presentation.slide_layouts[1])
    builder.build(slide, context)
    chart = next(shape.chart for shape in slide.shapes if hasattr(shape, "chart"))
    labels = chart.series[0].data_labels
    assert labels.number_format == "0%"
    assert labels.position == XL_DATA_LABEL_POSITION.OUTSIDE_END
