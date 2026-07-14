from __future__ import annotations

from pathlib import Path

import pytest
from PIL import Image

from backend.analytics_module.pptx_builder.theme import PPTXTheme

from backend.analytics_module.pptx_builder.chart_resolver import PPTXChartResolver
from backend.analytics_module.pptx_builder.engine import PPTXEngine
from backend.analytics_module.pptx_builder.hybrid_export.render_routing import (
  resolve_content_slide_builder,
  should_render_chart_as_image,
)
from backend.analytics_module.pptx_builder.pptx_image_chart import PPTXImageChart
from backend.analytics_module.pptx_builder.pptx_waterfall_bar import PPTXWaterfallBar


@pytest.fixture
def theme():
  return PPTXTheme()


def _write_png(path: Path) -> None:
  Image.new("RGB", (32, 32), color=(255, 0, 0)).save(path)


def test_should_render_chart_as_image_requires_existing_file(tmp_path: Path):
  image_path = tmp_path / "audience_affinity.png"
  chart = {"chart_id": "audience_affinity", "chart_type": "affinity_heatmap"}

  assert should_render_chart_as_image(chart) is False

  _write_png(image_path)
  chart["pptx_capture_path"] = str(image_path)
  assert should_render_chart_as_image(chart) is True


def test_resolve_content_slide_builder_prefers_capture_image(tmp_path: Path):
  image_path = tmp_path / "criteria_table.png"
  _write_png(image_path)
  resolver = PPTXChartResolver()
  chart = {
    "chart_id": "criteria_table",
    "chart_type": "criteria_table",
    "pptx_capture_path": str(image_path),
  }

  resolution, builder_class = resolve_content_slide_builder(chart, resolver)

  assert builder_class is PPTXImageChart
  assert resolution.registry_key == "image_capture"
  assert resolution.source == "captured_image"


def test_resolve_content_slide_builder_falls_back_to_native_without_capture():
  resolver = PPTXChartResolver()
  chart = {
    "chart_id": "brand_awareness",
    "chart_type": "horizontal_bar",
    "title": "Awareness",
  }

  resolution, builder_class = resolve_content_slide_builder(chart, resolver)

  assert builder_class is PPTXWaterfallBar
  assert resolution.registry_key == "brand_awareness"


def test_pptx_image_chart_inserts_picture(marketeers_presentation, marketeers_layout, theme, tmp_path: Path):
  image_path = tmp_path / "overall_scatter.png"
  _write_png(image_path)
  builder = PPTXImageChart(theme, marketeers_layout)
  slide = marketeers_presentation.slides.add_slide(marketeers_presentation.slide_layouts[1])

  result = builder.build(
    slide,
    {
      "chart_id": "overall_scatter",
      "chart_type": "scatter_plot",
      "pptx_capture_path": str(image_path),
    },
    chrome_owned_by_engine=True,
  )

  assert result.status.value == "rendered"
  assert any(shape.shape_type == 13 for shape in slide.shapes)


def test_engine_render_journal_records_image_capture_mode(
  marketeers_template_path,
  tmp_path: Path,
):
  from pptx import Presentation

  image_path = tmp_path / "audience_affinity.png"
  _write_png(image_path)
  engine = PPTXEngine(template_path=str(marketeers_template_path))
  prs = Presentation(str(marketeers_template_path))
  chart = {
    "chart_id": "audience_affinity",
    "chart_type": "affinity_heatmap",
    "title": "Audience Affinity",
    "pptx_capture_path": str(image_path),
  }

  engine._add_content_slide(prs, chart)

  assert engine.render_journal[-1]["chart_id"] == "audience_affinity"
  assert engine.render_journal[-1]["render_mode"] == "image_capture"
  assert engine.render_journal[-1]["registry_key"] == "image_capture"
  assert engine.render_journal[-1]["render_status"] == "rendered"