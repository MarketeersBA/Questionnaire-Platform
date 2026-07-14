import pytest
from pptx import Presentation
from backend.analytics_module.pptx_builder.pptx_strategic_narrative import PPTXStrategicNarrative
from backend.analytics_module.pptx_builder.theme import PPTXTheme
from backend.analytics_module.pptx_builder.layout import PPTXLayout
from backend.analytics_module.pptx_builder.presentation_planner import PresentationPlanner, SlideType

def test_strategic_narrative_builder_renders():
    prs = Presentation()
    theme = PPTXTheme()
    layout = PPTXLayout.from_presentation(prs)
    builder = PPTXStrategicNarrative(theme, layout)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Using index 6 as a dummy layout
    
    data = {
        "insights": {
            "executive_summary": "Test Executive Summary",
            "strategic_narrative": "Test Strategic Narrative",
            "business_question": "Test Business Question",
            "key_findings": [
                {"label": "F1", "finding": "Finding 1"},
                {"label": "F2", "finding": "Finding 2"},
                {"label": "F3", "finding": "Finding 3"},
            ]
        },
        "metadata": {
            "category": "Test Category",
            "brand": "Test Brand"
        }
    }
    
    # This should run without crashing
    builder.build(slide, data)
    
    # Verify some shapes were added
    assert len(slide.shapes) > 5

def test_presentation_planner_includes_strategic_narrative():
    report_doc = {
        "metadata": {"category": "Cats", "brand": "BrandX"},
        "insights": {"executive_summary": "Something"},
        "charts": []
    }
    intents = PresentationPlanner.define_slide_intents(report_doc)
    
    types = [i.type for i in intents]
    assert SlideType.STRATEGIC_NARRATIVE in types
    
    # Verify it's at index 2 (Slide 3)
    assert types[2] == SlideType.STRATEGIC_NARRATIVE
