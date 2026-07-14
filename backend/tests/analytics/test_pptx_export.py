import pytest
import os
import shutil
import asyncio
from pathlib import Path
from pptx import Presentation
from unittest.mock import MagicMock

from backend.analytics_module.pptx_facade import PPTXGenerator
import logging
logger = logging.getLogger(__name__)
from backend.analytics_module.pptx_builder.engine import PPTXEngine
from backend.analytics_module.pptx_builder.presentation_planner import SlideIntent, SlideType

@pytest.mark.asyncio
async def test_pptx_generation_pipeline_fidelity():
    """
    End-to-End Test: Validates the entire PPTX generation pipeline from 
    JSON payload to physical artifact validation.
    """
    # 1. Setup Mock Environment
    survey_id = "test_fidelity_001"
    from unittest.mock import AsyncMock
    mock_db = MagicMock()
    mock_db.get_collection = MagicMock(return_value=AsyncMock())
    
    # Mock Report Payload reflecting modern data structures
    report_doc = {
        "survey_id": survey_id,
        "project_name": "Phoenix Integration Test",
        "charts": [
            {
                "chart_id": "brand_awareness",
                "chart_type": "horizontal_bar",
                "title": "Brand Equity Awareness",
                "ai_headline": "Dominant market presence observed with 85% recall.",
                "data": {
                    "labels": ["Own Brand", "Competitor A", "Competitor B"],
                    "datasets": [{"label": "Awareness", "data": [0.85, 0.45, 0.12]}],
                },
            },
            {
                "chart_id": "purchase_funnel",
                "chart_type": "snake_line",
                "title": "Conversion Thresholds",
                "ai_headline": "Significant drop-off between trial and loyalty phases.",
                "data": {
                    "labels": ["Awareness", "Trial", "Loyalty"],
                    "datasets": [{"label": "Funnel", "data": [1.0, 0.65, 0.22]}],
                },
            },
        ]
    }
    
    # 2. Initialize Generator with Sandbox Paths
    generator = PPTXGenerator(mock_db, survey_id)
    test_output_dir = Path("backend/tests/analytics/tmp_pptx_sandbox")
    generator.output_dir = test_output_dir
    
    # Ensure template exists for integration test
    template_path = Path("backend/resources/analytics/marketeers_template.pptx")
    if not template_path.parent.exists():
        os.makedirs(template_path.parent, exist_ok=True)
    
    os.makedirs(test_output_dir, exist_ok=True)
    
    try:
        # 3. Execute Pipeline
        pptx_path = await generator.generate_from_report(report_doc)
        
        # 4. Perform High-Fidelity Validations
        assert os.path.exists(pptx_path), "PPTX artifact was not written to disk."
        
        # Load the generated file for internal inspection
        pres = Presentation(pptx_path)
        
        # Cover + Survey Info + section divider + chart + section divider + chart + closing
        actual_slides = len(pres.slides)
        assert actual_slides >= 6, f"Expected at least 6 slides, but found {actual_slides}."
        
        # Title Integrity Check (Engine uses .upper())
        assert pres.slides[0].shapes.title.text == "PHOENIX INTEGRATION TEST", "Cover title hydration failed."
        assert pres.slides[1].shapes.title.text == "SURVEY OVERVIEW & METHODOLOGY", "Survey info slide title mismatch."
        
        assert pres.slides[2].shapes.title.text in {"PURCHASE FUNNEL", "COMPARISONS", "DASHBOARD"}
        
        slide_titles = [slide.shapes.title.text for slide in pres.slides if slide.shapes.title]
        assert "BRAND EQUITY AWARENESS" in slide_titles
        assert "CONVERSION THRESHOLDS" in slide_titles

    finally:
        # 5. Environment Sanitization
        if test_output_dir.exists():
            shutil.rmtree(test_output_dir)

def test_pptx_engine_survey_info_injection():
    """
    Unit Test: Validates that the PPTXEngine correctly injects and 
    populates the new Survey Info slide using enriched metadata.
    """
    from datetime import datetime
    
    # 1. Initialize Engine with explicit Template Path resolution
    template_path = Path("backend/resources/analytics/marketeers_template.pptx")
    engine = PPTXEngine(template_path=str(template_path))
    
    # 2. Mock Enriched Metadata (Phase 1, 2, 5 Payload)
    report_data = {
        "metadata": {
            "title": "V2 Engine Unit Test",
            "brand": "Hero Brand",
            "target_brand": "Hero Brand",
            "company_name": "Antigravity Corp",
            "brands": ["Hero Brand", "Comp A", "Comp B"],
            "pf_brands": ["Hero Brand", "Comp A"],
            "pf_active": True,
            "survey_created_at": datetime(2026, 4, 1),
            "report_generated_at": datetime(2026, 5, 10),
            "total_responses": 412,
            "base_n": 400,
            "sample_capacity": 400,
            "research_type": "Taste Test",
            "total_attributes": 15,
            "attribute_categories": ["Sensory (8)", "Brand Affinity (4)", "Purchase Intent (3)"],
            "attributes_summary": "Taste, Aroma, Texture, Value, Loyalty, Trust, Quality, Price, Design, Packaging"
        },
        "charts": [] # No charts for focused slide-2 testing
    }
    
    from backend.analytics_module.pptx_builder.presentation_planner import PresentationPlanner, SlideIntent, SlideType

    intents = [
        SlideIntent(SlideType.COVER, data=report_data["metadata"]),
        SlideIntent(SlideType.SURVEY_OVERVIEW, data=report_data["metadata"]),
        SlideIntent(SlideType.CLOSING, data=report_data["metadata"]),
    ]

    # 3. Generate Presentation
    pptx_stream, slide_count = engine.generate_presentation(intents)
    pres = Presentation(pptx_stream)
    
    # 4. Assertions
    # Expected: Cover (1) + Survey Info (1) + Closing (1) = 3 total
    assert len(pres.slides) == 3, f"Expected 3 slides, found {len(pres.slides)}"
    
    # Verify Slide 2 Content
    cover_slide = pres.slides[0]
    cover_text = " | ".join(
        shape.text for shape in cover_slide.shapes if hasattr(shape, "text") and shape.text
    )
    assert "ANTIGRAVITY CORP" in cover_text.upper(), "Cover client company metadata missing."

    info_slide = pres.slides[1]
    all_text = ""
    for shape in info_slide.shapes:
        if hasattr(shape, "text"):
            all_text += shape.text + " | "
            
    # Check for Title
    assert "SURVEY OVERVIEW" in all_text, "Survey Info slide title missing."
    
    # Check for Enriched Metadata presence (Builder uses .upper() for some fields)
    assert "HERO BRAND" in all_text.upper(), "Target brand info missing from slide."
    assert "ANTIGRAVITY CORP" in all_text.upper(), "Company name missing from slide."
    assert "412" in all_text, "Response count missing from card."
    assert "TASTE TEST" in all_text.upper(), "Research type missing."
    
    # Phase 2 & 5 Diagnostic Assertions
    assert "DIAGNOSTIC PILLARS" in all_text.upper(), "Diagnostic pillars header missing."
    assert "15 ATTRIBUTES TESTED" in all_text.upper(), "Attribute count missing."
    assert "SENSORY (8)" in all_text.upper(), "Attribute categories missing."
    
    # Duration Check (Apr 1 -> May 10 = 39 days)
    assert "39 DAYS" in all_text.upper(), "Study duration calculation failed."

    closing_text = " | ".join(
        shape.text for shape in pres.slides[2].shapes if hasattr(shape, "text") and shape.text
    )
    assert "V2 ENGINE UNIT TEST" in closing_text.upper(), "Closing project metadata missing."
    assert "ANTIGRAVITY CORP" in closing_text.upper(), "Closing company metadata missing."
    
    logger.info("PPTXEngine Survey Info injection verified successfully.")


def test_pptx_engine_adds_analysis_followup_slide_for_long_deep_analysis():
    template_path = Path("backend/resources/analytics/marketeers_template.pptx")
    engine = PPTXEngine(template_path=str(template_path))

    deep_analysis = [
        {
            "title": "Recall Momentum",
            "body": "Hero brand recall grows strongly across key segments and outpaces the nearest competitor by double digits.",
            "sentiment": "positive",
            "recommended_action": "Reinforce memory structures with high-frequency top-funnel creative.",
        },
        {
            "title": "Conversion Risk",
            "body": "Trial-to-loyalty conversion still shows friction despite healthy awareness and initial trial pull.",
            "sentiment": "negative",
            "recommended_action": "Deploy targeted retention messaging and optimize repeat-trigger touchpoints.",
        },
    ]

    intents = [
        SlideIntent(
            SlideType.CONTENT_SLIDE,
            data={
                "chart_id": "purchase_intent",
                "chart_type": "stacked_bar",
                "title": "Purchase Intent",
                "ai_headline": "Intent is strong but loyalty conversion requires intervention.",
                "ai_deep_analysis": deep_analysis,
                "data": {
                    "labels": ["Top 2 Box", "Neutral", "Bottom 2 Box"],
                    "datasets": [{"label": "Hero Brand", "data": [0.62, 0.24, 0.14]}],
                },
            },
        )
    ]

    pptx_stream, slide_count = engine.generate_presentation(intents)
    prs = Presentation(pptx_stream)

    assert slide_count == 2
    titles = [slide.shapes.title.text for slide in prs.slides if slide.shapes.title]
    assert any("AI DEEP ANALYSIS" in title for title in titles)
