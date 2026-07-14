from pathlib import Path

import pytest
from pptx import Presentation

from backend.analytics_module.pptx_builder.layout import PPTXLayout
from backend.analytics_module.pptx_builder.pptx_horizontal_bar import PPTXHorizontalBar
from backend.analytics_module.pptx_builder.slide_chrome import ContentSlideChromeApplier
from backend.analytics_module.pptx_builder.template_adapter import TemplateAdapter
from backend.analytics_module.pptx_builder.theme import PPTXTheme


TEMPLATE_PATH = Path("backend/resources/analytics/marketeers_template.pptx")


@pytest.fixture
def theme():
    return PPTXTheme()


def _slide_text(slide) -> str:
    chunks = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            chunks.append(shape.text)
    return " | ".join(chunks)


def test_content_slide_chrome_uses_template_title_and_ai_headline():
    prs = Presentation(str(TEMPLATE_PATH))
    layout = PPTXLayout.from_presentation(prs)
    theme = PPTXTheme()
    slide = prs.slides.add_slide(TemplateAdapter().get_layout(prs, "content"))

    ContentSlideChromeApplier(theme, layout).apply(
        slide,
        {
            "title": "Brand Equity Awareness",
            "ai_headline": "Dominant market presence observed with 85% recall.",
            "insight": "Should not appear when deep analysis exists.",
            "ai_deep_analysis": [{"title": "Recall", "body": "Hero brand leads the category."}],
            "footnote": "Base: total sample.",
        },
    )

    assert slide.shapes.title.text == "BRAND EQUITY AWARENESS"
    assert "Dominant market presence observed with 85% recall." in _slide_text(slide)
    assert "INSIGHT: Should not appear when deep analysis exists." not in _slide_text(slide)
    assert "Base: total sample." in _slide_text(slide)
    assert "Recall: Hero brand leads the category." in slide.notes_slide.notes_text_frame.text


def test_content_slide_chrome_uses_template_subtitle_placeholder():
    prs = Presentation(str(TEMPLATE_PATH))
    layout = PPTXLayout.from_presentation(prs)
    theme = PPTXTheme()
    slide = prs.slides.add_slide(TemplateAdapter().get_layout(prs, "content"))

    ContentSlideChromeApplier(theme, layout).apply(
        slide,
        {
            "title": "Brand Equity Awareness",
            "subtitle": "Awareness vs consideration",
            "chart_type": "horizontal_bar",
            "data": {"labels": ["A"], "datasets": [{"label": "Awareness", "data": [0.5]}]},
        },
    )

    subtitle_texts = []
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1 and ph.has_text_frame:
            subtitle_texts.append(ph.text)
    if subtitle_texts:
        assert subtitle_texts == ["Awareness vs consideration"]
    else:
        assert "Awareness vs consideration" in _slide_text(slide)


def test_standalone_builder_reuses_engine_chrome_contract(marketeers_presentation, marketeers_layout, theme):
    slide = marketeers_presentation.slides.add_slide(marketeers_presentation.slide_layouts[1])
    chart_data = {
        "title": "Standalone Chart",
        "chart_type": "horizontal_bar",
        "data": {
            "labels": ["Own Brand", "Competitor A"],
            "datasets": [{"label": "Awareness", "data": [0.85, 0.45]}],
        },
    }

    PPTXHorizontalBar(theme, marketeers_layout).build(slide, chart_data)

    assert slide.shapes.title.text == "STANDALONE CHART"
    assert _slide_text(slide).upper().count("STANDALONE CHART") == 1


def test_engine_owned_chrome_avoids_duplicate_builder_title():
    prs = Presentation(str(TEMPLATE_PATH))
    layout = PPTXLayout.from_presentation(prs)
    theme = PPTXTheme()
    slide = prs.slides.add_slide(TemplateAdapter().get_layout(prs, "content"))

    chart_data = {
        "title": "Brand Equity Awareness",
        "chart_type": "horizontal_bar",
        "data": {
            "labels": ["Own Brand", "Competitor A"],
            "datasets": [{"label": "Awareness", "data": [0.85, 0.45]}],
        },
    }

    ContentSlideChromeApplier(theme, layout).apply(slide, chart_data)
    PPTXHorizontalBar(theme, layout).build(slide, chart_data, chrome_owned_by_engine=True)

    title_occurrences = _slide_text(slide).upper().count("BRAND EQUITY AWARENESS")
    assert title_occurrences == 1
