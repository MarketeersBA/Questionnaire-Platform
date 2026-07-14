from __future__ import annotations

from typing import Any, Dict, List, Optional

from pptx import Presentation

from backend.analytics_module.pptx_builder.narrative_expansion import estimate_extra_slides_for_intents
from backend.tests.analytics.pptx_acceptance_contract import (
    chart_titles_from_report,
    charts_with_deep_analysis,
    expected_backend_key,
)


def collect_slide_text(presentation: Presentation) -> str:
    chunks: List[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
    return "\n".join(chunks)


def collect_slide_titles(presentation: Presentation) -> List[str]:
    titles: List[str] = []
    for slide in presentation.slides:
        title_shape = slide.shapes.title
        if title_shape is not None and title_shape.text:
            titles.append(title_shape.text.strip())
    return titles


def collect_notes_text(presentation: Presentation) -> str:
    chunks: List[str] = []
    for slide in presentation.slides:
        try:
            notes = slide.notes_slide.notes_text_frame.text
        except Exception:
            notes = ""
        if notes:
            chunks.append(notes)
    return "\n".join(chunks)


def _check(
    check: str,
    *,
    expected: Any,
    actual: Any,
    passed: bool,
) -> Dict[str, Any]:
    return {
        "check": check,
        "expected": expected,
        "actual": actual,
        "passed": passed,
    }


def evaluate_resolver_dispatch(
    report_doc: Dict[str, Any],
    render_journal: List[Dict[str, Any]],
    narrative_journal: Optional[List[Dict[str, Any]]] = None,
) -> List[Dict[str, Any]]:
    checks: List[Dict[str, Any]] = []
    journal_by_chart_id = {
        entry.get("chart_id"): entry
        for entry in render_journal
        if entry.get("chart_id")
    }

    for chart in report_doc.get("charts", []):
        chart_id = chart.get("chart_id")
        chart_type = chart.get("chart_type")
        # Brand cards are intentionally grouped through BRAND_PROFILE narrative flow.
        if isinstance(chart_id, str) and chart_id.startswith("brand_card_") and chart_type == "scorecard":
            rendered_via_narrative = any(
                str(entry.get("section_id", "")).startswith("brand_profile::")
                and chart.get("title", "").lower() in str(entry.get("title", "")).lower()
                for entry in (narrative_journal or [])
            )
            checks.append(
                _check(
                    f"resolver::{chart_id}",
                    expected="brand_profile_narrative",
                    actual="present" if rendered_via_narrative else "missing",
                    passed=rendered_via_narrative,
                )
            )
            continue
        expected_key = expected_backend_key(chart)
        entry = journal_by_chart_id.get(chart_id)
        actual_key = entry.get("registry_key") if entry else "missing"
        checks.append(
            _check(
                f"resolver::{chart_id}",
                expected=expected_key,
                actual=actual_key,
                passed=entry is not None and actual_key == expected_key,
            )
        )
    return checks


def evaluate_failure_markers(
    validation_report: Dict[str, Any],
    combined_text: str,
) -> List[Dict[str, Any]]:
    export_audit = validation_report.get("export_audit", {})
    text_markers = export_audit.get("text_markers", {})
    analysis_interrupted = int(text_markers.get("analysis_interrupted", 0))
    unsupported_count = int(validation_report.get("unsupported_count", 0))
    error_placeholder_count = int(validation_report.get("error_placeholder_count", 0))

    return [
        _check(
            "analysis_interrupted_absent",
            expected="0",
            actual=str(analysis_interrupted),
            passed=analysis_interrupted == 0 and "ANALYSIS INTERRUPTED" not in combined_text,
        ),
        _check(
            "unsupported_placeholders_absent",
            expected="0",
            actual=str(unsupported_count),
            passed=unsupported_count == 0,
        ),
        _check(
            "error_placeholders_absent",
            expected="0",
            actual=str(error_placeholder_count),
            passed=error_placeholder_count == 0,
        ),
    ]


def evaluate_required_ai_sections(
    report_doc: Dict[str, Any],
    combined_text: str,
) -> List[Dict[str, Any]]:
    checks: List[Dict[str, Any]] = []
    insights = report_doc.get("insights", {})

    if insights.get("executive_summary"):
        checks.append(
            _check(
                "executive_summary_present",
                expected="EXECUTIVE SUMMARY",
                actual="present" if "EXECUTIVE SUMMARY" in combined_text else "missing",
                passed="EXECUTIVE SUMMARY" in combined_text,
            )
        )
    if insights.get("key_findings"):
        checks.append(
            _check(
                "critical_findings_present",
                expected="CRITICAL FINDINGS",
                actual="present" if "CRITICAL FINDINGS" in combined_text else "missing",
                passed="CRITICAL FINDINGS" in combined_text,
            )
        )
    if insights.get("opportunity_insights"):
        checks.append(
            _check(
                "opportunity_insights_present",
                expected="STRATEGIC INTELLIGENCE or EXECUTION PLAYBOOK",
                actual="present"
                if ("STRATEGIC INTELLIGENCE" in combined_text or "EXECUTION PLAYBOOK" in combined_text)
                else "missing",
                passed="STRATEGIC INTELLIGENCE" in combined_text or "EXECUTION PLAYBOOK" in combined_text,
            )
        )
    if insights.get("market_position_report"):
        checks.append(
            _check(
                "market_position_present",
                expected="MARKET ARCHETYPE or STRATEGIC POSITIONING",
                actual="present"
                if ("MARKET ARCHETYPE" in combined_text or "STRATEGIC POSITIONING" in combined_text)
                else "missing",
                passed="MARKET ARCHETYPE" in combined_text or "STRATEGIC POSITIONING" in combined_text,
            )
        )
    if insights.get("brand_swot"):
        checks.append(
            _check(
                "swot_present",
                expected="COMPETITIVE SWOT",
                actual="present" if "COMPETITIVE SWOT" in combined_text else "missing",
                passed="COMPETITIVE SWOT" in combined_text,
            )
        )
    if insights.get("recommendations_4p"):
        checks.append(
            _check(
                "recommendations_present",
                expected="4P RECOMMENDATIONS",
                actual="present" if "4P RECOMMENDATIONS" in combined_text else "missing",
                passed="4P RECOMMENDATIONS" in combined_text,
            )
        )
    return checks


def evaluate_chart_title_parity(
    report_doc: Dict[str, Any],
    presentation: Presentation,
    narrative_journal: Optional[List[Dict[str, Any]]] = None,
) -> List[Dict[str, Any]]:
    expected_titles = [title.upper() for title in chart_titles_from_report(report_doc)]
    combined_text = collect_slide_text(presentation).upper()
    slide_titles = [title.upper() for title in collect_slide_titles(presentation)]
    present_count = sum(1 for title in expected_titles if title and title in combined_text)

    checks = [
        _check(
            "chart_title_count_matches_source_report",
            expected=str(len(expected_titles)),
            actual=str(present_count),
            passed=present_count == len(expected_titles),
        )
    ]

    for chart in report_doc.get("charts", []):
        title = str(chart.get("title", "")).upper()
        chart_id = str(chart.get("chart_id", ""))
        chart_type = str(chart.get("chart_type", ""))
        if chart_id.startswith("brand_card_") and chart_type == "scorecard":
            narrative_hit = any(
                chart.get("title", "").lower() in str(entry.get("title", "")).lower()
                for entry in (narrative_journal or [])
                if str(entry.get("section_id", "")).startswith("brand_profile::")
            )
            checks.append(
                _check(
                    f"chart_title::{chart.get('chart_id')}",
                    expected=title,
                    actual="present" if narrative_hit else "missing",
                    passed=narrative_hit,
                )
            )
            continue

        if chart_id == "sigma_intent":
            sigma_present = any(
                entry.startswith(f"{title} (")
                for entry in slide_titles
            ) or title in slide_titles
            checks.append(
                _check(
                    f"chart_title::{chart.get('chart_id')}",
                    expected=f"{title} or chunked variant",
                    actual="present" if sigma_present else "missing",
                    passed=sigma_present,
                )
            )
            continue

        checks.append(
            _check(
                f"chart_title::{chart.get('chart_id')}",
                expected=title,
                actual="present" if title and title in slide_titles else "missing",
                passed=bool(title) and title in slide_titles,
            )
        )
    return checks


def evaluate_shape_bounds(validation_report: Dict[str, Any]) -> List[Dict[str, Any]]:
    layout_warning_count = int(validation_report.get("layout_warning_count", 0))
    return [
        _check(
            "shape_bounds_pass",
            expected="0",
            actual=str(layout_warning_count),
            passed=layout_warning_count == 0,
        )
    ]


def evaluate_deep_analysis_notes(
    report_doc: Dict[str, Any],
    presentation: Presentation,
    validation_report: Dict[str, Any],
) -> List[Dict[str, Any]]:
    deep_analysis_charts = charts_with_deep_analysis(report_doc)
    if not deep_analysis_charts:
        return []

    notes_count = int(validation_report.get("notes_count", 0))
    notes_text = collect_notes_text(presentation)
    checks: List[Dict[str, Any]] = [
        _check(
            "notes_present_for_deep_analysis",
            expected=str(len(deep_analysis_charts)),
            actual=str(notes_count),
            passed=notes_count > 0 and bool(notes_text.strip()),
        )
    ]

    for chart in deep_analysis_charts:
        chart_id = chart.get("chart_id")
        deep_analysis = chart.get("ai_deep_analysis")
        expected_snippet = ""
        if isinstance(deep_analysis, list) and deep_analysis:
            first = deep_analysis[0]
            if isinstance(first, dict):
                expected_snippet = str(first.get("body") or first.get("title") or "").strip()
        elif isinstance(deep_analysis, str):
            expected_snippet = deep_analysis.strip()

        checks.append(
            _check(
                f"deep_analysis_notes::{chart_id}",
                expected=expected_snippet or "notes content",
                actual="present" if expected_snippet and expected_snippet in notes_text else notes_text[:120],
                passed=bool(expected_snippet) and expected_snippet in notes_text,
            )
        )
    return checks


def build_acceptance_checklist(
    report_doc: Dict[str, Any],
    presentation: Presentation,
    intents: List[Any],
    validation_report: Dict[str, Any],
    *,
    render_journal: Optional[List[Dict[str, Any]]] = None,
    narrative_journal: Optional[List[Dict[str, Any]]] = None,
) -> Dict[str, Any]:
    combined_text = collect_slide_text(presentation).upper()
    slide_titles = [title.upper() for title in collect_slide_titles(presentation)]
    expected_min = len(intents)
    expected_max = len(intents) + estimate_extra_slides_for_intents(intents) + 2

    checks: List[Dict[str, Any]] = [
        _check(
            "cover_title",
            expected=report_doc.get("project_name", "").upper(),
            actual=slide_titles[0] if slide_titles else "",
            passed=bool(slide_titles) and slide_titles[0] == report_doc.get("project_name", "").upper(),
        ),
        _check(
            "survey_overview_present",
            expected="SURVEY OVERVIEW",
            actual="present" if "SURVEY OVERVIEW" in combined_text else "missing",
            passed="SURVEY OVERVIEW" in combined_text,
        ),
        _check(
            "slide_count_within_expected_range",
            expected=f"{expected_min}-{expected_max}",
            actual=str(len(presentation.slides)),
            passed=expected_min <= len(presentation.slides) <= expected_max,
        ),
    ]

    checks.extend(evaluate_failure_markers(validation_report, combined_text))
    checks.extend(evaluate_required_ai_sections(report_doc, combined_text))
    checks.extend(evaluate_chart_title_parity(report_doc, presentation, narrative_journal=narrative_journal))
    checks.extend(evaluate_shape_bounds(validation_report))
    checks.extend(
        evaluate_deep_analysis_notes(report_doc, presentation, validation_report)
    )
    if render_journal is not None:
        checks.extend(
            evaluate_resolver_dispatch(
                report_doc,
                render_journal,
                narrative_journal=narrative_journal,
            )
        )

    failed = [item for item in checks if not item["passed"]]
    return {
        "checks": checks,
        "failed_checks": failed,
        "passed": not failed,
    }


def build_smoke_checklist(
    report_doc: Dict[str, Any],
    presentation: Presentation,
    intents: List[Any],
    validation_report: Dict[str, Any],
    *,
    render_journal: Optional[List[Dict[str, Any]]] = None,
    narrative_journal: Optional[List[Dict[str, Any]]] = None,
) -> List[Dict[str, Any]]:
    """Backward-compatible smoke checklist for known-report golden export."""
    return build_acceptance_checklist(
        report_doc,
        presentation,
        intents,
        validation_report,
        render_journal=render_journal,
        narrative_journal=narrative_journal,
    )["checks"]


def build_golden_smoke_checklist(
    report_doc: Dict[str, Any],
    presentation: Presentation,
    intents: List[Any],
    validation_report: Dict[str, Any],
    *,
    render_journal: Optional[List[Dict[str, Any]]] = None,
    narrative_journal: Optional[List[Dict[str, Any]]] = None,
) -> Dict[str, Any]:
    """Single structured smoke checklist for known report vs exported PPTX."""
    return build_acceptance_checklist(
        report_doc,
        presentation,
        intents,
        validation_report,
        render_journal=render_journal,
        narrative_journal=narrative_journal,
    )
