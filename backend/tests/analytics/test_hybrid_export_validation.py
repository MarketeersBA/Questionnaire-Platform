from __future__ import annotations

from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation

from backend.analytics_module.pptx_builder.theme import PPTXTheme
from backend.analytics_module.pptx_builder.engine import PPTXEngine
from backend.analytics_module.pptx_builder.export_validation_manifest import (
  build_export_manifest,
  summarize_render_journal,
)
from backend.analytics_module.pptx_builder.hybrid_export.capture_planning import (
  build_capture_requests_from_report,
)
from backend.analytics_module.pptx_builder.hybrid_export.capture_validation import (
  evaluate_download_readiness,
)
from backend.analytics_module.pptx_builder.hybrid_export.success_criteria import (
  build_phase0_scope_manifest,
)
from backend.analytics_module.pptx_builder.layout import PPTXLayout
from backend.analytics_module.pptx_builder.pptx_export_forensics import build_export_forensics_manifest
from backend.analytics_module.pptx_builder.pptx_image_chart import PPTXImageChart
from backend.tests.analytics.hybrid_export_golden_fixture import (
  build_golden_hybrid_capture_manifest,
  build_golden_image_backed_render_journal,
)
from backend.tests.analytics.pptx_acceptance_contract import build_representative_screen_report


@pytest.fixture
def theme():
  return PPTXTheme()


def _write_png(path: Path) -> None:
  Image.new("RGB", (32, 32), color=(0, 128, 255)).save(path)


def test_planner_capture_list_matches_phase0_manifest():
  report = build_representative_screen_report()
  phase0 = build_phase0_scope_manifest(report)
  capture_requests = build_capture_requests_from_report(report)

  assert {item.chart_id for item in capture_requests} == {
    item["chart_id"] for item in phase0.capture_candidates
  }


def test_failed_capture_is_visible_in_forensics_manifest():
  report = build_representative_screen_report()
  capture_manifest = build_golden_hybrid_capture_manifest()
  render_journal = build_golden_image_backed_render_journal()

  forensics = build_export_forensics_manifest(
    report_doc=report,
    intents=[],
    preparation_snapshot={"charts": []},
    normalization_notes=[],
    render_journal=render_journal,
    narrative_journal=[],
    certification={"passes_gate": True, "export_audit": {"text_markers": {}, "slide_summaries": []}},
    capture_manifest=capture_manifest,
  )

  failed_rows = [
    row
    for row in forensics["chart_forensics"]
    if row.get("capture", {}).get("capture_status") == "failed"
  ]
  assert failed_rows
  assert failed_rows[0]["capture"]["capture_error"] == "navigation_failed"


def test_image_backed_forensics_pass_without_native_chart_shape():
  report = build_representative_screen_report()
  capture_manifest = build_golden_hybrid_capture_manifest()
  render_journal = build_golden_image_backed_render_journal()

  forensics = build_export_forensics_manifest(
    report_doc=report,
    intents=[],
    preparation_snapshot={"charts": []},
    normalization_notes=[],
    render_journal=render_journal,
    narrative_journal=[],
    certification={
      "passes_gate": True,
      "export_audit": {
        "text_markers": {},
        "slide_summaries": [
          {
            "slide_index": 8,
            "chart_count": 0,
            "picture_count": 1,
            "has_error_placeholder": False,
          }
        ],
      },
    },
    capture_manifest=capture_manifest,
  )

  affinity_row = next(
    row for row in forensics["chart_forensics"] if row["chart_id"] == "audience_affinity"
  )
  assert affinity_row["passed"] is True
  assert "missing_native_chart_shape" not in affinity_row["blocked_reasons"]
  assert affinity_row["capture"]["viewport_url"].startswith("http://frontend.test/")


def test_image_chart_uses_template_chart_body_bounds(
  marketeers_presentation,
  marketeers_layout,
  theme,
  tmp_path: Path,
):
  image_path = tmp_path / "audience_affinity.png"
  _write_png(image_path)
  builder = PPTXImageChart(theme, marketeers_layout)
  slide = marketeers_presentation.slides.add_slide(marketeers_presentation.slide_layouts[1])
  builder.build(
    slide,
    {
      "chart_id": "audience_affinity",
      "chart_type": "affinity_heatmap",
      "pptx_capture_path": str(image_path),
    },
    chrome_owned_by_engine=True,
  )

  picture_shapes = [shape for shape in slide.shapes if shape.shape_type == 13]
  assert picture_shapes
  left, top, width, height = marketeers_layout.chart_body_bounds()
  picture = picture_shapes[-1]
  assert int(picture.left) == int(left)
  assert int(picture.top) == int(top)
  assert int(picture.width) == int(width)
  assert int(picture.height) == int(height)


def test_export_manifest_includes_capture_validation_summary():
  capture_manifest = build_golden_hybrid_capture_manifest()
  manifest = build_export_manifest(
    report_id="golden-report",
    generated_at="2026-05-14T10:00:00",
    template_hash="hash",
    certification={
      "passes_gate": True,
      "render_tally": summarize_render_journal(build_golden_image_backed_render_journal()),
      "export_audit": {},
    },
    report_doc=build_representative_screen_report(),
    preparation_snapshot={"charts": []},
    chart_normalization_notes=[],
    chart_parity={"order_mismatch": False},
    narrative_render_manifest={"missing_sections": []},
    layout_geometry={"chart_frame_fits_slide": True},
    actual_slide_count=10,
    capture_manifest=capture_manifest,
  )

  assert manifest["capture_validation"]["capture_candidate_count"] == 5
  assert manifest["capture_validation"]["capture_failure_count"] == 1
  assert manifest["image_capture_count"] == 1
  assert manifest["capture_manifest"]["artifact_root"] == "/tmp/capture_artifacts"


def test_download_readiness_requires_gate_and_path():
  manifest = {
    "passes_gate": True,
    "failed_chart_count": 0,
    "capture_validation": {"capture_failure_count": 0},
  }
  ready = evaluate_download_readiness(export_manifest=manifest, pptx_path="/tmp/report.pptx")
  assert ready["ready_for_download"] is True

  blocked = evaluate_download_readiness(export_manifest={**manifest, "passes_gate": False}, pptx_path=None)
  assert blocked["ready_for_download"] is False
  assert "validation_gate_failed" in blocked["issues"]


def test_golden_acceptance_image_backed_slide_records_picture_audit(
  marketeers_template_path,
  tmp_path: Path,
):
  image_path = tmp_path / "audience_affinity.png"
  _write_png(image_path)
  engine = PPTXEngine(template_path=str(marketeers_template_path))
  prs = Presentation(str(marketeers_template_path))
  engine._add_content_slide(
    prs,
    {
      "chart_id": "audience_affinity",
      "chart_type": "affinity_heatmap",
      "title": "Audience Affinity",
      "pptx_capture_path": str(image_path),
    },
  )

  from backend.analytics_module.pptx_builder.pptx_export_audit import audit_pptx_bytes

  stream = __import__("io").BytesIO()
  prs.save(stream)
  audit = audit_pptx_bytes(stream.getvalue())
  slide_summary = audit["slide_summaries"][-1]

  assert engine.render_journal[-1]["render_mode"] == "image_capture"
  assert slide_summary["picture_count"] >= 1
