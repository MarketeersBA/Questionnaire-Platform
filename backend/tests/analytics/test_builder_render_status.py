import io

import pytest
from pptx import Presentation

from backend.analytics_module.pptx_builder.builder_render_status import ChartRenderStatus
from backend.analytics_module.pptx_builder.chart_payload_adapters import adapt_chart_data_for_builder
from backend.analytics_module.pptx_builder.pptx_nps_gauge import PPTXNPSGauge
from backend.analytics_module.pptx_builder.pptx_profile_chart import PPTXProfileChart
from backend.analytics_module.pptx_builder.pptx_scatter import PPTXScatter
from backend.analytics_module.pptx_builder.theme import PPTXTheme
from backend.analytics_module.pptx_builder.validation_gating import PPTXValidationMode
from backend.analytics_module.pptx_builder.validator import PPTXIntegrityValidator


@pytest.fixture
def theme():
    return PPTXTheme()


def test_scatter_empty_payload_reports_skipped_empty_data(marketeers_presentation, marketeers_layout, theme):
    builder = PPTXScatter(theme, marketeers_layout)
    slide = marketeers_presentation.slides.add_slide(marketeers_presentation.slide_layouts[1])

    result = builder.build(slide, {"title": "Scatter", "data": {"datasets": []}})

    assert result.status is ChartRenderStatus.SKIPPED_EMPTY_DATA


def test_profile_chart_renders_with_labels_and_datasets(marketeers_presentation, marketeers_layout, theme):
    builder = PPTXProfileChart(theme, marketeers_layout)
    slide = marketeers_presentation.slides.add_slide(marketeers_presentation.slide_layouts[1])

    result = builder.build(
        slide,
        {
            "title": "Profile",
            "data": {
                "labels": ["Taste", "Value"],
                "datasets": [{"label": "Hero Brand", "brand": "Hero Brand", "data": [4.5, 3.8]}],
            },
        },
    )

    assert result.status is ChartRenderStatus.RENDERED
    assert any(hasattr(shape, "chart") for shape in slide.shapes)


def test_nps_gauge_renders_with_transposed_payload(marketeers_presentation, marketeers_layout, theme):
    adapted, notes = adapt_chart_data_for_builder(
        "nps_recommend",
        {
            "title": "NPS",
            "data": {
                "labels": ["Promoters_Pct", "Passives_Pct", "Detractors_Pct"],
                "datasets": [{"label": "Hero Brand", "data": [0.5, 0.3, 0.2]}],
                "nps_scores": {"Hero Brand": 30},
            },
        },
    )
    assert "transposed_nps_segment_rows_to_brand_rows" in notes

    builder = PPTXNPSGauge(theme, marketeers_layout)
    slide = marketeers_presentation.slides.add_slide(marketeers_presentation.slide_layouts[1])
    result = builder.build(slide, {"title": "NPS", "data": adapted})

    assert result.status is ChartRenderStatus.RENDERED


def test_nps_gauge_renders_with_canonical_payload(marketeers_presentation, marketeers_layout, theme):
    from backend.analytics_module.pptx_builder.chart_contracts import canonical_nps_gauge_sample

    adapted, notes = adapt_chart_data_for_builder(
        "nps_recommend",
        {
            "title": "Net Promoter Score",
            "data": canonical_nps_gauge_sample(),
        },
    )
    assert "transposed_nps_segment_rows_to_brand_rows" not in notes

    builder = PPTXNPSGauge(theme, marketeers_layout)
    slide = marketeers_presentation.slides.add_slide(marketeers_presentation.slide_layouts[1])
    result = builder.build(slide, {"title": "Net Promoter Score", "data": adapted})

    assert result.status is ChartRenderStatus.RENDERED
    assert any(hasattr(shape, "chart") for shape in slide.shapes)


@pytest.mark.asyncio
async def test_validator_flags_skipped_empty_data_render_journal():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Empty Chart"
    stream = io.BytesIO()
    presentation.save(stream)

    validator = PPTXIntegrityValidator(stream)
    result = await validator.validate(
        {"charts": []},
        intents=[],
        mode=PPTXValidationMode.PRODUCTION,
        render_journal=[
            {
                "chart_id": "overall_scatter",
                "slide_index": 1,
                "render_status": "skipped_empty_data",
                "render_message": "No datasets found for scatter chart.",
            }
        ],
    )

    assert result["passes_gate"] is False
    assert any("rendered without native content" in message for message in result["validation_errors"])
