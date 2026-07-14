from __future__ import annotations

import io
import json
import re
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Mapping, Optional, Sequence, Tuple

from pptx import Presentation

from backend.analytics_module.pptx_builder.chart_contracts import (
    TASTE_TEST_CONTRACTS,
    resolve_taste_test_contract,
    validate_against_contract,
)
from backend.analytics_module.pptx_builder.chart_payload_contract import prepare_report_for_pptx
from backend.analytics_module.pptx_builder.chart_render_manifest import (
    build_chart_parity_manifest,
    collect_screen_chart_ids,
)
from backend.analytics_module.pptx_builder.chart_resolver import PPTXChartResolver
from backend.analytics_module.pptx_builder.engine import PPTXEngine
from backend.analytics_module.pptx_builder.export_validation_manifest import SHAPE_BASED_CHART_TYPES
from backend.analytics_module.pptx_builder.hybrid_export.render_mode import (
    PPTXRenderMode,
    resolve_render_mode,
)
from backend.analytics_module.pptx_builder.presentation_planner import PresentationPlanner
from backend.analytics_module.pptx_builder.pptx_export_audit import audit_pptx_bytes
from backend.analytics_module.pptx_builder.validation_gating import PPTXValidationMode
from backend.analytics_module.pptx_builder.validator import PPTXIntegrityValidator
from backend.tests.analytics.pptx_acceptance_contract import expected_backend_key

_NUMERIC_TOKEN = re.compile(r"\d")
_PERCENT_TOKEN = re.compile(r"\d+\s*%|\d+\.\d+\s*%")

# Charts that communicate via shapes/tables/text rather than embedded chart series.
_NON_SERIES_CHART_TYPES = frozenset(
    {
        *SHAPE_BASED_CHART_TYPES,
        "importance_combined",
        "sigma_intent_scatter",
        "scatter_bubble",
        "scatter_plot",
        "scatter",
    }
)


@dataclass(frozen=True)
class NativeE2EExportResult:
    report_doc: Dict[str, Any]
    preparation_snapshot: Dict[str, Any]
    normalized_charts: List[Dict[str, Any]]
    normalization_notes: List[Dict[str, Any]]
    engine: PPTXEngine
    intents: List[Any]
    presentation: Presentation
    pptx_stream: io.BytesIO
    slide_count: int
    validation: Dict[str, Any]
    render_mode: str
    chart_parity: Dict[str, Any]
    export_audit: Dict[str, Any]


def _generate_native_deck(
    report_doc: Dict[str, Any],
    *,
    template_path: Path,
) -> Tuple[PPTXEngine, List[Any], Presentation, io.BytesIO, int, Any]:
    preparation = prepare_report_for_pptx(report_doc)
    prepared_report = preparation.report_doc
    intents = PresentationPlanner.define_slide_intents(prepared_report)

    engine = PPTXEngine(template_path=str(template_path))
    pptx_stream, slide_count = engine.generate_presentation(intents)
    presentation = Presentation(pptx_stream)
    return engine, intents, presentation, pptx_stream, slide_count, preparation


async def run_native_e2e_export(
    report_doc: Dict[str, Any],
    *,
    template_path: Path,
    render_mode: str = "native",
) -> NativeE2EExportResult:
    """
    Full production-shaped native export path:
    prepare_report_for_pptx -> planner -> engine -> PRODUCTION validator.
    """
    if resolve_render_mode(render_mode).value != PPTXRenderMode.NATIVE.value:
        raise ValueError(f"Phase F E2E requires native render mode, got {render_mode!r}")

    engine, intents, presentation, pptx_stream, slide_count, preparation = _generate_native_deck(
        report_doc,
        template_path=template_path,
    )
    prepared_report = preparation.report_doc

    validator = PPTXIntegrityValidator(io.BytesIO(pptx_stream.getvalue()))
    validation = await validator.validate(
        prepared_report,
        intents,
        mode=PPTXValidationMode.PRODUCTION,
        render_journal=engine.render_journal,
        narrative_journal=engine.narrative_render_journal,
    )
    export_audit = audit_pptx_bytes(pptx_stream.getvalue())
    chart_parity = build_chart_parity_manifest(
        screen_chart_ids=collect_screen_chart_ids(prepared_report),
        normalized_charts=preparation.normalized_charts,
        render_journal=engine.render_journal,
    )

    return NativeE2EExportResult(
        report_doc=prepared_report,
        preparation_snapshot=preparation.snapshot,
        normalized_charts=preparation.normalized_charts,
        normalization_notes=preparation.normalization_notes,
        engine=engine,
        intents=intents,
        presentation=presentation,
        pptx_stream=pptx_stream,
        slide_count=slide_count,
        validation=validation,
        render_mode=render_mode,
        chart_parity=chart_parity,
        export_audit=export_audit,
    )


def run_native_e2e_export_sync(
    report_doc: Dict[str, Any],
    *,
    template_path: Path,
    render_mode: str = "native",
) -> NativeE2EExportResult:
    import asyncio

    return asyncio.run(
        run_native_e2e_export(
            report_doc,
            template_path=template_path,
            render_mode=render_mode,
        )
    )


def collect_numeric_samples(value: Any, *, limit: int = 8) -> List[str]:
    """Extract representative numeric tokens from chart payload data."""
    samples: List[str] = []

    def _walk(node: Any) -> None:
        if len(samples) >= limit:
            return
        if isinstance(node, bool):
            return
        if isinstance(node, (int, float)):
            if isinstance(node, float) and node.is_integer():
                samples.append(str(int(node)))
            else:
                samples.append(f"{node:g}")
            return
        if isinstance(node, str):
            for match in _PERCENT_TOKEN.findall(node):
                if match not in samples:
                    samples.append(match.strip())
            for match in _NUMERIC_TOKEN.findall(node):
                if match not in samples:
                    samples.append(match)
            return
        if isinstance(node, Mapping):
            for child in node.values():
                _walk(child)
            return
        if isinstance(node, Sequence) and not isinstance(node, (str, bytes, bytearray)):
            for child in node:
                _walk(child)

    _walk(value)
    return samples[:limit]


def _slide_summaries_by_index(export_audit: Dict[str, Any]) -> Dict[int, Dict[str, Any]]:
    return {
        int(summary["slide_index"]): summary
        for summary in export_audit.get("slide_summaries", [])
        if summary.get("slide_index") is not None
    }


def _slide_text_by_index(presentation: Presentation) -> Dict[int, str]:
    texts: Dict[int, str] = {}
    for slide_index, slide in enumerate(presentation.slides, start=1):
        chunks: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
            if getattr(shape, "has_table", False):
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            chunks.append(cell.text)
        texts[slide_index] = "\n".join(chunks)
    return texts


def _find_render_entry(
    chart_id: str,
    render_journal: Sequence[Mapping[str, Any]],
) -> Optional[Dict[str, Any]]:
    for entry in render_journal:
        if entry.get("chart_id") == chart_id:
            return dict(entry)
    return None


def _brand_card_rendered_via_narrative(
    chart: Mapping[str, Any],
    narrative_journal: Sequence[Mapping[str, Any]],
) -> bool:
    chart_id = str(chart.get("chart_id") or "")
    chart_type = str(chart.get("chart_type") or "")
    if not (chart_id.startswith("brand_card_") and chart_type == "scorecard"):
        return False
    title = str(chart.get("title") or "").lower()
    return any(
        str(entry.get("section_id", "")).startswith("brand_profile::")
        and title in str(entry.get("title", "")).lower()
        for entry in narrative_journal
    )


def _slide_has_render_evidence(
    chart: Mapping[str, Any],
    slide_summary: Optional[Mapping[str, Any]],
    *,
    slide_text: str = "",
    brand_narrative: bool = False,
) -> Tuple[bool, str]:
    if brand_narrative:
        return True, "brand_profile_narrative"

    chart_type = str(chart.get("chart_type") or "")
    if not slide_summary:
        return False, "missing_slide_summary"

    if chart_type in _NON_SERIES_CHART_TYPES:
        shape_count = int(slide_summary.get("shape_count", 0))
        if shape_count > 0 or bool(slide_text.strip()):
            return True, "shape_or_text_content"
        return False, "no_shape_content"

    chart_count = int(slide_summary.get("chart_count", 0))
    if chart_count > 0:
        return True, "native_chart_shape"
    if bool(slide_text.strip()):
        return True, "slide_text_content"
    if int(slide_summary.get("picture_count", 0)) > 0:
        return True, "picture_shape"
    return False, "no_chart_shape"


def _slide_has_numeric_evidence(
    chart: Mapping[str, Any],
    numeric_samples: Sequence[str],
    slide_summary: Optional[Mapping[str, Any]],
    *,
    slide_text: str = "",
    brand_narrative: bool = False,
) -> Tuple[bool, str]:
    chart_type = str(chart.get("chart_type") or "")
    if chart_type in {"wordcloud", "verbatim_analysis"}:
        return True, "qualitative_exempt"

    if brand_narrative and numeric_samples:
        return True, "brand_card_source_numbers_present"

    if not numeric_samples:
        return False, "no_numeric_source_data"

    for sample in numeric_samples:
        if sample and sample in slide_text:
            return True, "visible_in_slide_text"

    # Percent tokens may be formatted differently on slide (e.g. 65 vs 65%).
    slide_digits = re.sub(r"[^\d]", "", slide_text)
    for sample in numeric_samples:
        digits = re.sub(r"[^\d]", "", sample)
        if digits and digits in slide_digits:
            return True, "visible_digits_in_slide_text"

    if chart_type in _NON_SERIES_CHART_TYPES:
        return True, "shape_chart_source_numbers_present"

    if slide_summary and int(slide_summary.get("chart_count", 0)) > 0:
        return True, "chart_shape_with_source_numbers"

    return False, "numbers_not_visible_on_slide"


def build_per_chart_coverage_row(
    chart: Mapping[str, Any],
    *,
    render_journal: Sequence[Mapping[str, Any]],
    narrative_journal: Sequence[Mapping[str, Any]],
    slide_summaries: Mapping[int, Mapping[str, Any]],
    slide_texts: Mapping[int, str],
    validation_issues: Sequence[Mapping[str, Any]],
) -> Dict[str, Any]:
    chart_id = str(chart.get("chart_id") or "")
    chart_type = str(chart.get("chart_type") or "")
    contract = resolve_taste_test_contract(chart)
    contract_errors = list(chart.get("_contract_errors") or [])
    resolution = chart.get("_resolution") or PPTXChartResolver().describe_export_contract(dict(chart))
    registry_key = resolution.get("registry_key") or expected_backend_key(dict(chart))
    uses_fallback = bool(resolution.get("uses_fallback_table"))

    render_entry = _find_render_entry(chart_id, render_journal)
    brand_narrative = _brand_card_rendered_via_narrative(chart, narrative_journal)
    slide_index = render_entry.get("slide_index") if render_entry else None
    slide_summary = slide_summaries.get(int(slide_index)) if slide_index is not None else None

    slide_text = slide_texts.get(int(slide_index), "") if slide_index is not None else ""
    numeric_samples = collect_numeric_samples(chart.get("data"))
    render_evidence, render_evidence_reason = _slide_has_render_evidence(
        chart,
        slide_summary,
        slide_text=slide_text,
        brand_narrative=brand_narrative,
    )
    numeric_evidence, numeric_reason = _slide_has_numeric_evidence(
        chart,
        numeric_samples,
        slide_summary,
        slide_text=slide_text,
        brand_narrative=brand_narrative,
    )

    chart_issues = [
        issue
        for issue in validation_issues
        if chart_id and chart_id in str(issue.get("message", ""))
    ]

    if brand_narrative:
        render_status = "rendered_via_narrative"
        resolution_source = "brand_profile_narrative"
    elif render_entry:
        render_status = str(render_entry.get("render_status") or "unknown")
        resolution_source = str(render_entry.get("resolution_source") or resolution.get("source") or "")
    else:
        render_status = "missing"
        resolution_source = str(resolution.get("source") or "")

    passed = (
        contract is not None
        and not contract_errors
        and not uses_fallback
        and resolution_source != "fallback_table"
        and render_status in {"rendered", "rendered_via_narrative"}
        and render_evidence
        and numeric_evidence
        and not chart_issues
    )

    return {
        "chart_id": chart_id,
        "chart_type": chart_type,
        "title": chart.get("title"),
        "contract_pattern": contract.chart_id_pattern if contract else None,
        "contract_builder_key": contract.builder_registry_key if contract else None,
        "expected_registry_key": expected_backend_key(dict(chart)),
        "actual_registry_key": registry_key,
        "resolution_source": resolution_source,
        "uses_fallback_table": uses_fallback,
        "contract_errors": contract_errors,
        "render_status": render_status,
        "render_mode": (render_entry or {}).get("render_mode", "native"),
        "slide_index": slide_index,
        "render_evidence": render_evidence,
        "render_evidence_reason": render_evidence_reason,
        "numeric_samples": numeric_samples,
        "numeric_evidence": numeric_evidence,
        "numeric_evidence_reason": numeric_reason,
        "validation_issue_count": len(chart_issues),
        "validation_issues": chart_issues,
        "passed": passed,
    }


def build_per_chart_coverage_report(export_result: NativeE2EExportResult) -> Dict[str, Any]:
    slide_summaries = _slide_summaries_by_index(export_result.export_audit)
    slide_texts = _slide_text_by_index(export_result.presentation)
    validation_issues = export_result.validation.get("issues", [])

    rows = [
        build_per_chart_coverage_row(
            chart,
            render_journal=export_result.engine.render_journal,
            narrative_journal=export_result.engine.narrative_render_journal,
            slide_summaries=slide_summaries,
            slide_texts=slide_texts,
            validation_issues=validation_issues,
        )
        for chart in export_result.normalized_charts
    ]

    passed_rows = [row for row in rows if row["passed"]]
    failed_rows = [row for row in rows if not row["passed"]]
    narrative_rendered_ids = [
        row["chart_id"] for row in rows if row.get("render_status") == "rendered_via_narrative"
    ]
    effective_rendered_ids = list(export_result.chart_parity.get("rendered_chart_ids", []))
    for chart_id in narrative_rendered_ids:
        if chart_id not in effective_rendered_ids:
            effective_rendered_ids.append(chart_id)
    screen_chart_ids = export_result.chart_parity.get("screen_chart_ids", [])
    adjusted_missing_from_pptx = [
        chart_id for chart_id in screen_chart_ids if chart_id not in effective_rendered_ids
    ]
    contract_registry = [
        {
            "pattern": contract.chart_id_pattern,
            "chart_type": contract.chart_type,
            "builder_registry_key": contract.builder_registry_key,
            "active_in_compute_all": contract.active_in_compute_all,
        }
        for contract in TASTE_TEST_CONTRACTS
    ]

    return {
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "render_mode": export_result.render_mode,
        "report_name": export_result.report_doc.get("project_name"),
        "research_type": (export_result.report_doc.get("metadata") or {}).get("research_type"),
        "survey_id": (export_result.report_doc.get("metadata") or {}).get("survey_id"),
        "chart_count": len(rows),
        "passed_chart_count": len(passed_rows),
        "failed_chart_count": len(failed_rows),
        "production_passes_gate": bool(export_result.validation.get("passes_gate")),
        "validation_mode": export_result.validation.get("validation_mode"),
        "slide_count": export_result.slide_count,
        "render_tally": export_result.validation.get("render_tally", {}),
        "chart_parity": export_result.chart_parity,
        "adjusted_chart_parity": {
            **export_result.chart_parity,
            "effective_rendered_chart_ids": effective_rendered_ids,
            "missing_from_pptx": adjusted_missing_from_pptx,
        },
        "normalization_notes": export_result.normalization_notes,
        "contract_registry_count": len(contract_registry),
        "rows": rows,
        "failed_rows": failed_rows,
        "contract_registry": contract_registry,
        "summary": {
            "all_charts_passed": not failed_rows,
            "native_only": export_result.render_mode == "native",
            "unsupported_count": export_result.validation.get("unsupported_count", 0),
            "error_placeholder_count": export_result.validation.get("error_placeholder_count", 0),
            "layout_warning_count": export_result.validation.get("layout_warning_count", 0),
        },
    }


def format_coverage_report_markdown(report: Mapping[str, Any]) -> str:
    lines = [
        "# PPTX Native E2E Coverage Report",
        "",
        f"- Generated: {report.get('generated_at')}",
        f"- Report: {report.get('report_name')}",
        f"- Render mode: `{report.get('render_mode')}`",
        f"- PRODUCTION gate: **{'PASS' if report.get('production_passes_gate') else 'FAIL'}**",
        f"- Charts: {report.get('passed_chart_count')}/{report.get('chart_count')} passed",
        f"- Slides: {report.get('slide_count')}",
        "",
        "## Per-chart coverage",
        "",
        "| Chart ID | Type | Registry | Render | Numbers | Status |",
        "| --- | --- | --- | --- | --- | --- |",
    ]

    for row in report.get("rows", []):
        status = "PASS" if row.get("passed") else "FAIL"
        lines.append(
            "| {chart_id} | {chart_type} | {registry} | {render} | {numbers} | {status} |".format(
                chart_id=row.get("chart_id"),
                chart_type=row.get("chart_type"),
                registry=row.get("actual_registry_key"),
                render=row.get("render_status"),
                numbers="yes" if row.get("numeric_evidence") else "no",
                status=status,
            )
        )

    failed_rows = report.get("failed_rows") or []
    if failed_rows:
        lines.extend(["", "## Failures", ""])
        for row in failed_rows:
            lines.append(
                f"- **{row.get('chart_id')}**: render={row.get('render_status')}, "
                f"fallback={row.get('uses_fallback_table')}, "
                f"numeric={row.get('numeric_evidence_reason')}, "
                f"render_evidence={row.get('render_evidence_reason')}"
            )
            if row.get("contract_errors"):
                lines.append(f"  - contract_errors: {len(row['contract_errors'])}")

    return "\n".join(lines) + "\n"


def write_coverage_artifacts(
    report: Mapping[str, Any],
    *,
    output_dir: Path,
    deck_bytes: Optional[bytes] = None,
    deck_filename: str = "native_e2e_deck.pptx",
) -> Dict[str, str]:
    output_dir.mkdir(parents=True, exist_ok=True)
    json_path = output_dir / "native_e2e_coverage.json"
    md_path = output_dir / "native_e2e_coverage.md"
    json_path.write_text(json.dumps(report, indent=2, default=str), encoding="utf-8")
    md_path.write_text(format_coverage_report_markdown(report), encoding="utf-8")

    paths = {"json": str(json_path), "markdown": str(md_path)}
    if deck_bytes is not None:
        deck_path = output_dir / deck_filename
        deck_path.write_bytes(deck_bytes)
        paths["deck"] = str(deck_path)
    return paths
