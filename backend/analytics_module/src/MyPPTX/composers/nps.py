import logging
import pandas as pd
from pptx.util import Inches, Pt
from .base_composer import BaseChartComposer
from backend.analytics_module.schemas.export import ChartDefinition
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

class NPSGaugeComposer(BaseChartComposer):
    """
    Advanced Marketing Engine: NPS Strategic Gauge.
    Creates a composite visualization using native shapes (Detractors/Passives/Promoters).
    """

    def compose(self, slide, chart_definition: ChartDefinition, data: pd.DataFrame):
        # We expect a Series with index: ['Promoters', 'Passives', 'Detractors']
        # and a total NPS score value.
        
        nps_score = int(chart_definition.base_n) if chart_definition.base_n is not None else 0
        try:
            p_val = data.loc["Promoters"].values[0] if "Promoters" in data.index else 0
            n_val = data.loc["Passives"].values[0] if "Passives" in data.index else 0
            d_val = data.loc["Detractors"].values[0] if "Detractors" in data.index else 0
        except:
            p_val, n_val, d_val = 40, 40, 20 # Fallback for layout testing

        # 1. Base Dimensions
        left, top = Inches(2), Inches(3)
        total_w = Inches(6)
        h = Inches(0.4)

        # 2. Draw Detractors (Red)
        d_w = total_w * (d_val / 100)
        d_rect = slide.shapes.add_shape(1, left, top, d_w, h)
        d_rect.fill.solid()
        d_rect.fill.fore_color.rgb = RGBColor(220, 53, 69) # Brand Pink/Red

        # 3. Draw Passives (Gray)
        n_w = total_w * (n_val / 100)
        n_rect = slide.shapes.add_shape(1, left + d_w, top, n_w, h)
        n_rect.fill.solid()
        n_rect.fill.fore_color.rgb = RGBColor(180, 180, 180) # Gray

        # 4. Draw Promoters (Green)
        p_w = total_w * (p_val / 100)
        p_rect = slide.shapes.add_shape(1, left + d_w + n_w, top, p_w, h)
        p_rect.fill.solid()
        p_rect.fill.fore_color.rgb = RGBColor(40, 167, 69) # Emerald Green

        # 5. Insert Score Callout
        score_box = slide.shapes.add_textbox(left + (total_w/2) - Inches(0.5), top - Inches(1), Inches(1), Inches(1))
        tf = score_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(nps_score)
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.name = self.branding.font_family
        p.alignment = PP_ALIGN.CENTER

        label_box = slide.shapes.add_textbox(left + (total_w/2) - Inches(1), top - Inches(0.3), Inches(2), Inches(0.5))
        label_box.text_frame.text = "NET PROMOTER SCORE"
        label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return d_rect
