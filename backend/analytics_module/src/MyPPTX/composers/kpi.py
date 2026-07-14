import logging
import pandas as pd
from pptx.util import Inches, Pt
from .base_composer import BaseChartComposer
from backend.analytics_module.schemas.export import ChartDefinition
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

class KPICardComposer(BaseChartComposer):
    """
    Advanced Engine: Strategic KPI Cards.
    Renders high-impact metric boxes for cover slides or executive summaries.
    """

    def compose(self, slide, chart_definition: ChartDefinition, data: pd.DataFrame):
        # We expect a Series where each Row is a KPI Label and value is the Metric
        # Example: ['Total N', 'Aided Awareness', 'Top Choice']
        
        cards = data.to_dict().get(data.columns[0], {})
        num_cards = len(cards)
        if num_cards == 0: return None

        card_w = Inches(2)
        card_h = Inches(1.5)
        gap = Inches(0.2)
        
        # Center the row of cards
        total_content_w = (card_w * num_cards) + (gap * (num_cards - 1))
        start_left = (Inches(10) - total_content_w) / 2
        top = Inches(2.5)

        for i, (label, val) in enumerate(cards.items()):
            current_left = start_left + (i * (card_w + gap))
            
            # 1. Background Box (Glassmorphic look)
            rect = slide.shapes.add_shape(1, current_left, top, card_w, card_h)
            rect.fill.solid()
            rect.fill.fore_color.rgb = RGBColor(245, 245, 250)
            rect.line.color.rgb = RGBColor(200, 200, 200)

            # 2. Value Text
            val_box = slide.shapes.add_textbox(current_left, top + Inches(0.2), card_w, Inches(0.8))
            tf_v = val_box.text_frame
            p_v = tf_v.paragraphs[0]
            p_v.text = str(val)
            p_v.font.size = Pt(32)
            p_v.font.bold = True
            p_v.font.color.rgb = RGBColor(0, 0, 128) # Navy
            p_v.alignment = PP_ALIGN.CENTER

            # 3. Label Text
            label_box = slide.shapes.add_textbox(current_left, top + Inches(0.9), card_w, Inches(0.4))
            tf_l = label_box.text_frame
            p_l = tf_l.paragraphs[0]
            p_l.text = str(label).upper()
            p_l.font.size = Pt(10)
            p_l.font.color.rgb = RGBColor(100, 100, 100)
            p_l.alignment = PP_ALIGN.CENTER
            
        return rect
