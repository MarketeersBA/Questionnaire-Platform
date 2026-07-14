import logging
import math
from typing import Dict, List, Optional, Tuple, Union, Any

import pandas as pd
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Pt

from . import design_config
logger = logging.getLogger(__name__)

# =============================================================================
# CONSTANTS
# =============================================================================

LINE_CHART_TYPES = frozenset({
    XL_CHART_TYPE.LINE,
    XL_CHART_TYPE.LINE_MARKERS,
    XL_CHART_TYPE.LINE_STACKED,
    XL_CHART_TYPE.LINE_STACKED_100,
    XL_CHART_TYPE.LINE_MARKERS_STACKED,
    XL_CHART_TYPE.LINE_MARKERS_STACKED_100,
})

LEGEND_POSITION_MAP = {
    "bottom": XL_LEGEND_POSITION.BOTTOM,
    "top": XL_LEGEND_POSITION.TOP,
    "left": XL_LEGEND_POSITION.LEFT,
    "right": XL_LEGEND_POSITION.RIGHT,
    "corner": XL_LEGEND_POSITION.CORNER,
}

class ChartFormatter:
    """Encapsulates common chart formatting operations."""
    
    @staticmethod
    def apply_title_format(chart, new_title: Optional[str] = None) -> None:
        """Apply consistent title formatting to a chart.
        
        Args:
            chart: The chart object to format
            new_title: Optional new title text to set
        """
        if new_title is not None:
            chart.has_title = True
            chart.chart_title.text_frame.text = new_title
        
        if not chart.has_title:
            return
            
        title_par = chart.chart_title.text_frame.paragraphs[0]
        title_color = design_config.get_chart_title_color()
        title_size = design_config.get_chart_title_size()
        title_font = design_config.get_chart_font()
        title_bold = design_config.get_chart_title_bold()
        
        # Only apply when config provides a value; otherwise leave template as-is
        if title_size is not None:
            title_par.font.size = Pt(title_size)
        if title_font is not None:
            title_par.font.name = title_font
        if title_bold is not None:
            title_par.font.bold = title_bold
        if title_color is not None:
            title_par.font.color.rgb = title_color
            for run in title_par.runs:
                run.font.color.rgb = title_color
                
    @staticmethod
    def apply_legend_format(chart) -> None:
        """Apply consistent legend formatting to a chart. Only applies when config provides values."""
        try:
            legend_show = design_config.get_legend_show()
            if legend_show is False:
                chart.has_legend = False
                return
            if legend_show is not True:
                return
            chart.has_legend = True
            legend = chart.legend
            position = design_config.get_legend_position()
            if position is not None:
                pos_lower = position.lower()
                if pos_lower in LEGEND_POSITION_MAP:
                    legend.position = LEGEND_POSITION_MAP[pos_lower]
            font_size = design_config.get_legend_font_size()
            if font_size is not None:
                legend.font.size = Pt(font_size)
            font_name = design_config.get_legend_font_name()
            if font_name is not None:
                legend.font.name = font_name
            legend_color = design_config.get_legend_color()
            if legend_color is not None:
                legend.font.color.rgb = legend_color
        except Exception as e:
            logger.warning("Couldn't apply legend formatting: %s", e)
    
    @staticmethod
    def apply_data_labels(chart, show_percentage: bool = None, show_value: bool = None,
                          number_format: str = None) -> None:
        """Apply consistent data label formatting to all series in a chart.
        
        Args:
            chart: The chart object
            show_percentage: Override for showing percentage (uses config default if None)
            show_value: Override for showing value (uses config default if None)
            number_format: Override for number format (uses config default if None)
        """
        _show_percentage = show_percentage if show_percentage is not None else design_config.get_data_label_show_percentage()
        _show_value = show_value if show_value is not None else design_config.get_data_label_show_value()
        _number_format = number_format if number_format is not None else design_config.get_data_label_number_format()
        
        for series in chart.series:
            series.has_data_labels = True
            dls = series.data_labels
            if _show_percentage is not None:
                dls.show_percentage = _show_percentage
            if _show_value is not None:
                dls.show_value = _show_value
            if _number_format is not None:
                dls.number_format = _number_format
            try:
                font = getattr(dls, "font", None) or dls.format.font
                dl_font_size = design_config.get_data_label_font_size()
                if dl_font_size is not None:
                    font.size = Pt(dl_font_size)
                dl_color = design_config.get_data_label_color()
                if dl_color is not None:
                    font.color.rgb = dl_color
            except Exception as e:
                logger.error("Couldn't set data label font: %s", e, exc_info=True)
    
    @staticmethod
    def apply_series_colors(chart, colors: List[RGBColor] = None) -> None:
        """Apply color palette to chart series.
        
        Args:
            chart: The chart object
            colors: Optional list of RGBColor objects (uses config palette if None)
        """
        try:
            plot = chart.plots[0]
            num_series = len(plot.series)
        except Exception:
            num_series = 0
        
        if colors is None:
            colors = design_config.get_color_palette(num_series if num_series > 0 else None)
        if not colors:
            return
        try:
            plot = chart.plots[0]
            
            # Clean up stale data point format overrides from template
            try:
                for ser_el in plot._element.iterchildren(qn('c:ser')):
                    for dPt in list(ser_el.iterchildren(qn('c:dPt'))):
                        ser_el.remove(dPt)
            except Exception as e:
                logger.warning("Couldn't clean up old data point formatting: %s", e)
            
            # Disable vary_by_categories so series colors are used
            try:
                plot.vary_by_categories = False
            except AttributeError:
                pass
            
            # Apply colors to each series
            for series_idx, series in enumerate(plot.series):
                color = colors[series_idx % len(colors)]
                try:
                    series_fill = series.format.fill
                    series_fill.solid()
                    series_fill.fore_color.rgb = color
                except Exception as e:
                    logger.warning("Couldn't set color for series '%s': %s", series.name, e)
        except Exception as e:
            logger.warning("Couldn't apply color palette: %s", e)
    
    @staticmethod
    def apply_line_series_colors(chart, colors: List[RGBColor] = None) -> None:
        """Apply color palette to line chart series (lines and markers).
        
        Args:
            chart: The chart object
            colors: Optional list of RGBColor objects (uses config palette if None)
        """
        try:
            plot = chart.plots[0]
            num_series = len(plot.series)
        except Exception:
            num_series = 0
        
        if colors is None:
            colors = design_config.get_color_palette(num_series if num_series > 0 else None)
        if not colors:
            return
        try:
            plot = chart.plots[0]
            for series_idx, series in enumerate(plot.series):
                color = colors[series_idx % len(colors)]
                try:
                    # Set line color
                    line = series.format.line
                    line.color.rgb = color
                    
                    # Set marker fill color if present
                    if hasattr(series, 'marker') and series.marker:
                        marker_fill = series.marker.format.fill
                        marker_fill.solid()
                        marker_fill.fore_color.rgb = color
                except Exception as e:
                    logger.warning("Couldn't set color for line series '%s': %s", series.name, e)
        except Exception as e:
            logger.warning("Couldn't apply color palette to line chart: %s", e)
    
    @staticmethod
    def apply_axis_title_format(axis_title) -> None:
        """Apply consistent formatting to an axis title. Only applies when config provides values."""
        try:
            paragraph = axis_title.text_frame.paragraphs[0]
            size = design_config.get_axis_title_font_size()
            if size is not None:
                paragraph.font.size = Pt(size)
            name = design_config.get_axis_title_font_name()
            if name is not None:
                paragraph.font.name = name
            color = design_config.get_axis_title_color()
            if color is not None:
                paragraph.font.color.rgb = color
        except Exception as e:
            logger.warning("Couldn't apply axis title formatting: %s", e)

    @staticmethod
    def apply_axis_label_format(chart) -> None:
        """Apply consistent formatting to axis tick labels (category and value axes)."""
        try:
            # Format category axis (X-axis for column charts, Y-axis for bar charts)
            axis_label_size = design_config.get_axis_label_font_size()
            axis_label_name = design_config.get_axis_label_font_name()
            axis_label_color = design_config.get_axis_label_color()
            if hasattr(chart, 'category_axis') and chart.category_axis:
                category_axis = chart.category_axis
                if hasattr(category_axis, 'tick_labels'):
                    tick_labels = category_axis.tick_labels
                    if axis_label_size is not None:
                        tick_labels.font.size = Pt(axis_label_size)
                    if axis_label_name is not None:
                        tick_labels.font.name = axis_label_name
                    if axis_label_color is not None:
                        tick_labels.font.color.rgb = axis_label_color

            if hasattr(chart, 'value_axis') and chart.value_axis:
                value_axis = chart.value_axis
                if hasattr(value_axis, 'tick_labels'):
                    tick_labels = value_axis.tick_labels
                    if axis_label_size is not None:
                        tick_labels.font.size = Pt(axis_label_size)
                    if axis_label_name is not None:
                        tick_labels.font.name = axis_label_name
                    if axis_label_color is not None:
                        tick_labels.font.color.rgb = axis_label_color
        except Exception as e:
            logger.warning("Couldn't apply axis label formatting: %s", e)

# =============================================================================
# CONTENT HYDRATION ENGINE
# =============================================================================

class ContentHydrator:
    """
    Expert system for mapping report metadata (title, insights, sample size)
    to slide shape placeholders in real-time.
    """
    
    @staticmethod
    def hydrate_slide(slide, metadata: Dict[str, Any]) -> None:
        """
        Populate all text components on a slide based on chart metadata.
        """
        # 1. Main Slide Title
        title = metadata.get("title")
        if title and slide.shapes.title:
            slide.shapes.title.text = title
            ContentHydrator._apply_premium_style(slide.shapes.title, size=32, bold=True)

        # 2. Insight Box (Dynamic search for placeholders like 'headline' or 'insight')
        headline = metadata.get("ai_headline") or metadata.get("insight")
        if headline:
            # We look for shapes named "insight", "analysis", or "subtitle"
            success = ContentHydrator._set_placeholder_text(
                slide, 
                ["insight", "headline", "subtitle", "analysis", "box"], 
                headline,
                size=18
            )
            if not success:
                 logger.debug(f"Slide {getattr(slide, 'name', '')}: No insight placeholder found.")

        # 3. Base / Sample Size (N=...)
        sample_n = metadata.get("sample_n")
        if sample_n:
            ContentHydrator._set_placeholder_text(
                slide, 
                ["n_label", "sample", "base", "note"], 
                f"Base: N={sample_n}",
                size=12,
                italic=True
            )

    @staticmethod
    def _set_placeholder_text(slide, keys: List[str], text: str, size: int = 16, italic: bool = False) -> bool:
        """Finds a shape by keyword match and injects text with corporate styling."""
        for shape in slide.shapes:
            shape_name = (shape.name or "").lower()
            # Alt text match for accessibility-compliant templates
            alt_text = getattr(shape, "alternative_text", "").lower()
            
            if any(k in shape_name for k in keys) or any(k in alt_text for k in keys):
                if hasattr(shape, "text_frame"):
                    shape.text_frame.text = text
                    ContentHydrator._apply_premium_style(shape, size=size, italic=italic)
                    return True
        return False

    @staticmethod
    def _apply_premium_style(shape, size: Optional[int] = None, bold: bool = False, italic: bool = False):
        """Standardizes font, color, and weight for hydrated text."""
        if not hasattr(shape, "text_frame"):
            return
            
        font_name = design_config.get_chart_font() or "Pangram"
        brand_color = None
        try:
            brand_color = design_config.resolve_color("brand_navy")
        except:
             brand_color = (0, 32, 96) # Default Navy

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                if size:
                    run.font.size = Pt(size)
                run.font.bold = bold
                run.font.italic = italic
                if brand_color:
                    run.font.color.rgb = RGBColor(*brand_color)
