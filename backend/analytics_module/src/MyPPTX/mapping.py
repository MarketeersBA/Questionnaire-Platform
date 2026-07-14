"""
Mapping utilities for finding shapes by semantic keys (Alt Text / Names).
This decouples the presentation design from the code logic.
"""
import logging
from typing import Optional, List, Union, Any
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

class ShapeMapper:
    """Find shapes on a slide based on semantic keys."""

    @staticmethod
    def find_all_by_key(slide, key: str) -> List:
        """
        Find all shapes on a slide whose name or alt text matches the key.
        Matches are case-insensitive and stripped.
        """
        matches = []
        target = key.strip().lower()
        for shape in slide.shapes:
            # shape.name is the most reliable place for 'Alt Text' / 'Name' in PPT
            name = (shape.name or "").strip().lower()
            if name == target:
                matches.append(shape)
                continue
                
            # Fallback to description (Alt Text Description)
            try:
                # Some versions of python-pptx / PPT might put it here
                desc = (getattr(shape, "description", "") or "").strip().lower()
                if desc == target:
                    matches.append(shape)
            except Exception:
                pass
        return matches

    @staticmethod
    def find_first_by_key(slide, key: str) -> Optional[Any]:
        """Find the first shape matching the key."""
        matches = ShapeMapper.find_all_by_key(slide, key)
        return matches[0] if matches else None

    @staticmethod
    def find_chart_by_key(slide, key: str):
        """Find a chart shape by key."""
        shape = ShapeMapper.find_first_by_key(slide, key)
        if shape and shape.has_chart:
            return shape.chart
        return None

    @staticmethod
    def find_table_by_key(slide, key: str):
        """Find a table shape by key."""
        shape = ShapeMapper.find_first_by_key(slide, key)
        if shape and shape.has_table:
            return shape.table
        return None

    @staticmethod
    def find_text_frame_by_key(slide, key: str):
        """Find a text frame (textbox/placeholder) by key."""
        shape = ShapeMapper.find_first_by_key(slide, key)
        if shape and shape.has_text_frame:
            return shape.text_frame
        return None
