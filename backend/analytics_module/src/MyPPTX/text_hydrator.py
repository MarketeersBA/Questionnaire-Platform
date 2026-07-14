"""
Text & Insight Hydrator — Phase 1, Task 3.
Centralized engine for injecting AI narratives and research metadata into slides.
"""
import logging
from typing import Dict, Any, Optional
from pptx.util import Pt
from .mapping import ShapeMapper
from . import design_config

logger = logging.getLogger(__name__)

class TextHydrator:
    """
    Expert system for mapping AI narratives and metadata to slide text frames.
    Uses semantic key discovery to decouple code from layout variations.
    """
    
    # Strict Semantic Mappings (Placeholder names in PowerPoint 'Selection Pane')
    SEMANTIC_KEYS = {
        "TITLE": "title",
        "SUBTITLE": "subtitle",
        "INSIGHT": "insight_box",
        "BASE_N": "base_n",
        "FOOTNOTE": "footnote"
    }

    def __init__(self):
        self.font_name = design_config.get_chart_font() or "Pangram"

    def hydrate_slide(self, slide, chart_data: Dict[str, Any]):
        """
        Systematically injects text components into the slide archetype.
        """
        # 1. Extraction of Narratives & Meta-layers
        title = chart_data.get("title", "")
        # Some frontend components use 'headline' or 'subtitle'
        subtitle = chart_data.get("subtitle", "") or chart_data.get("headline", "")
        
        # Aggregate AI commentary
        insight = (
            chart_data.get("ai_headline") or 
            chart_data.get("insight") or 
            chart_data.get("ai_deep_analysis") or ""
        )
        
        base_n = chart_data.get("base_n") or chart_data.get("sample_size")
        footnote = chart_data.get("footnote", "")

        # 2. Systematic Injection Sequence
        
        # Primary Title
        if title:
            self._apply_text_layer(slide, self.SEMANTIC_KEYS["TITLE"], title, bold=True)

        # Secondary Context (Subtitle)
        if subtitle:
            self._apply_text_layer(slide, self.SEMANTIC_KEYS["SUBTITLE"], subtitle)

        # Narrative Layer (Insights)
        if insight:
            self._apply_text_layer(slide, self.SEMANTIC_KEYS["INSIGHT"], insight, italic=True)

        # Metadata Layer (Base size / N)
        if base_n:
            formatted_n = f"Base: n={base_n}"
            self._apply_text_layer(slide, self.SEMANTIC_KEYS["BASE_N"], formatted_n, font_size=10)
            
        # Footnotes / Methodology details
        if footnote:
            self._apply_text_layer(slide, self.SEMANTIC_KEYS["FOOTNOTE"], footnote, font_size=9)

    def _apply_text_layer(self, slide, key: str, text: str, 
                         bold=False, italic=False, font_size: Optional[int] = None):
        """
        Discovers a shape by semantic key and applies text with corporate styling.
        Includes fallback discovery for slides without explicit naming.
        """
        # 1. Discover target shape
        shape = ShapeMapper.find_first_by_key(slide, key)
        
        # 2. Contextual fallbacks
        if not shape:
            if key == self.SEMANTIC_KEYS["TITLE"] and slide.shapes.title:
                shape = slide.shapes.title
            elif key == self.SEMANTIC_KEYS["INSIGHT"]:
                # Broad capture for insight placeholders
                fallbacks = ["analysis", "headline", "commentary", "conclusion", "deep_analysis"]
                for f in fallbacks:
                    shape = ShapeMapper.find_first_by_key(slide, f)
                    if shape: break
        
        # 3. Apply Hydration & Styling
        if shape and hasattr(shape, "text_frame"):
            tf = shape.text_frame
            tf.text = str(text)
            
            # Ensure all text runs inherit corporate aesthetics
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    run.font.name = self.font_name
                    run.font.bold = bold
                    run.font.italic = italic
                    if font_size:
                        run.font.size = Pt(font_size)
            return True
            
        return False
