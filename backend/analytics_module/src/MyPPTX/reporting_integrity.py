import pandas as pd
from typing import List, Dict, Any
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class ReportingIntegrity:
    """
    Appends a 'Reporting Integrity' slide to the presentation.
    Acts as a quality certificate detailing the health of the data ingestion.
    """

    @staticmethod
    def append_diagnostic_slide(pres, audit_report: Dict[str, Any]):
        """
        Creates a new slide summarizing the analytical health.
        """
        # Add a blank slide
        blank_slide_layout = pres.slide_layouts[6] # Often blank
        slide = pres.slides.add_slide(blank_slide_layout)
        
        # 1. Add Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = title_box.text_frame
        tf.text = "Analytical Integrity Report"
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        
        # 2. Add Status Summary
        status = audit_report.get("status", "Unknown")
        coverage = audit_report.get("coverage_percent", 0)
        
        summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
        stf = summary_box.text_frame
        stf.word_wrap = True
        
        p1 = stf.add_paragraph()
        p1.text = f"Deployment Status: {status}"
        p1.font.bold = True
        p1.font.size = Pt(14)
        
        p2 = stf.add_paragraph()
        p2.text = f"Data Coverage: {coverage}%"
        p2.font.size = Pt(14)

        # 2.1 Quality Scrub Summary
        scrub = audit_report.get("quality_scrub", {})
        outliers = scrub.get("outlier_ids", [])
        if outliers:
            p_scrub = stf.add_paragraph()
            p_scrub.text = f"Quality Alert: {len(outliers)} responses excluded (Straight-liners/Speeders)"
            p_scrub.font.color.rgb = RGBColor(255, 102, 0) # Amber
            p_scrub.font.size = Pt(12)
            p_scrub.font.bold = True
        
        # 3. Add Missing Fields Detail
        missing = audit_report.get("missing_fields", [])
        if missing:
            detail_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(4))
            dtf = detail_box.text_frame
            dtf.word_wrap = True
            
            p3 = dtf.add_paragraph()
            p3.text = "The following fields were suppressed or missing from the source dataset:"
            p3.font.underline = True
            p3.font.size = Pt(12)
            
            for field in missing[:15]: # Cap at 15 for space
                p = dtf.add_paragraph()
                p.text = f"• {field} (Resolution: Suppressed)"
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(128, 128, 128) # Gray
            
            if len(missing) > 15:
                p = dtf.add_paragraph()
                p.text = f"... and {len(missing) - 15} more fields."
                p.font.size = Pt(10)
        else:
            p3 = stf.add_paragraph()
            p3.text = "100% Structural Match Achieved. All visual components are data-driven."
            p3.font.color.rgb = RGBColor(0, 128, 0) # Green
            p3.font.size = Pt(14)

        return slide
