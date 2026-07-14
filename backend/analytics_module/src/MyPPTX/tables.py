import numbers
import numpy as np
import pandas as pd
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from copy import deepcopy
from dataclasses import dataclass
from typing import Tuple, Optional, Dict, Any
from pptx.dml.color import RGBColor

from . import design_config
from .mapping import ShapeMapper

def create_native_table(slide, df, x, y, width, height):
    """
    Creates a brand-new native PPTX table from a DataFrame.
    """
    rows, cols = df.shape
    table_shape = slide.shapes.add_table(rows + 1, cols, x, y, width, height)
    table = table_shape.table
    
    # Header
    for c, col in enumerate(df.columns):
        write_cell_text(table.cell(0, c), str(col))
        
    # Data
    for r in range(rows):
        for c in range(cols):
            write_cell_text(table.cell(r+1, c), str(df.iloc[r, c]))
            
    return table



_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


def _resolve_table_theme(theme_name: str, table_config: Optional[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
    """Resolve table theme from table_themes.json. Checks themes[theme_name] then root-level theme_name."""
    if not table_config:
        return None
    themes = table_config.get("themes", {})
    if theme_name in themes:
        raw = themes[theme_name]
    elif theme_name in table_config and isinstance(table_config[theme_name], dict):
        raw = table_config[theme_name]
    else:
        return None
    return {k: v for k, v in raw.items() if not (isinstance(k, str) and k.startswith("comment"))}


def _resolve_column_style(
    style_name: str,
    column_config: Dict[str, Any],
    visited: Optional[set] = None,
) -> Dict[str, Any]:
    """Resolve a column theme style with 'extends' inheritance. Returns full style dict."""
    if visited is None:
        visited = set()
    if style_name in visited or not column_config:
        return {}
    themes = column_config.get("themes", {})
    if style_name not in themes:
        return {}
    visited.add(style_name)
    raw = deepcopy(themes[style_name])
    raw = {k: v for k, v in raw.items() if not (isinstance(k, str) and k.startswith("comment"))}
    if "extends" in raw:
        parent_name = raw.pop("extends")
        parent = _resolve_column_style(parent_name, column_config, visited.copy())
        merged = deepcopy(parent)
        merged.update(raw)
        return merged
    return raw


def _column_style_to_cell_format(style: Dict[str, Any]) -> Dict[str, Any]:
    """Convert resolved column theme dict to kwargs for write_cell_text / cell styling."""
    out = {}
    if "font" in style:
        out["font_name"] = style["font"]
    if "font_size" in style:
        out["font_size"] = style["font_size"]
    if "bold" in style:
        out["bold"] = style["bold"]
    if "align" in style:
        out["align"] = _ALIGN_MAP.get(style["align"], PP_ALIGN.CENTER)
    try:
        if "font_color" in style:
            rgb = design_config.resolve_color(style["font_color"])
            out["color_rgb"] = rgb
    except (ValueError, TypeError):
        pass
    try:
        if "fill_color" in style:
            out["fill_color"] = design_config.resolve_color(style["fill_color"])
    except (ValueError, TypeError):
        pass
    return out


def _get_table_cell_style(cell_type: str, column_header: Optional[str] = None) -> Optional[Dict[str, Any]]:
    """
    Get resolved style for a table cell using table_theme + column_theme (Option 2).
    cell_type: 'header' | 'index' | 'value'
    Returns None if new theme path is not available (caller should use design_config getters).
    """
    theme_name = design_config.get_table_theme()
    table_config = design_config.get_table_config()
    column_config = design_config.get_table_column_config()
    if not theme_name or not table_config or not column_config:
        return None
    table_theme = _resolve_table_theme(theme_name, table_config)
    if not table_theme:
        return None
    # For value cells, prefer per-column mapping so columns_by_header overrides the default "value"
    if cell_type == "value" and column_header is not None:
        columns_by_header = table_theme.get("columns_by_header") or {}
        style_key = columns_by_header.get(column_header)
        if style_key is None:
            style_key = table_theme.get("value")
    else:
        style_key = table_theme.get(cell_type) or table_theme.get("value")
    if not style_key:
        return None
    resolved = _resolve_column_style(style_key, column_config)
    if not resolved:
        return None
    return _column_style_to_cell_format(resolved)


def apply_table_theme(table, *, header_row=0, data_start_row=1, index_col=0, column_headers=None):
    rows = len(table.rows)
    cols = len(table.columns)
    theme_name = design_config.get_table_theme()
    table_config = design_config.get_table_config()
    if not theme_name and table_config:
        theme_name = table_config.get("default_theme") or None
    table_theme = _resolve_table_theme(theme_name, table_config) if theme_name else None
    header_style = _get_table_cell_style("header")
    use_new_theme = table_theme is not None and column_headers is not None and len(column_headers) >= cols

    if use_new_theme:
        # New path: apply only what's in the theme definition (no default theme)
        # ---- header row (only if theme defines header) ----
        if header_style:
            for c in range(cols):
                cell = table.cell(header_row, c)
                style = header_style
                tf = cell.text_frame
                for p in tf.paragraphs:
                    for r in p.runs:
                        if style.get("font_name"):
                            r.font.name = style["font_name"]
                        if style.get("font_size") is not None:
                            r.font.size = Pt(style["font_size"])
                        r.font.bold = style.get("bold", False)
                        if style.get("color_rgb"):
                            r.font.color.rgb = RGBColor(*style["color_rgb"])
                if style.get("align") is not None:
                    tf.paragraphs[0].alignment = style["align"]
                cell.fill.solid()
                if style.get("fill_color"):
                    cell.fill.fore_color.rgb = RGBColor(*style["fill_color"])

        # ---- body (only cells with a style in theme: index, value, or columns_by_header) ----
        for r in range(data_start_row, rows):
            for c in range(cols):
                cell = table.cell(r, c)
                tf = cell.text_frame
                # if not tf.paragraphs or not tf.paragraphs[0].runs:
                #     continue
                is_index = c == index_col
                col_header = column_headers[c] if c < len(column_headers) else None
                style = _get_table_cell_style("index" if is_index else "value", column_header=col_header)
                if not style:
                    continue
                for p in tf.paragraphs:
                    for run in p.runs:
                        if style.get("font_name"):
                            run.font.name = style["font_name"]
                        if style.get("font_size") is not None:
                            run.font.size = Pt(style["font_size"])
                        run.font.bold = style.get("bold", False)
                        if style.get("color_rgb"):
                            run.font.color.rgb = RGBColor(*style["color_rgb"])
                if style.get("align") is not None:
                    tf.paragraphs[0].alignment = style["align"]
                cell.fill.solid()
                if style.get("fill_color"):
                    cell.fill.fore_color.rgb = RGBColor(*style["fill_color"])
        return

    # Legacy path: apply only when config provides values (null = leave template as-is)
    for c in range(cols):
        cell = table.cell(header_row, c)
        tf = cell.text_frame
        font_name = design_config.get_table_font()
        header_size = design_config.get_table_header_font_size()
        header_bold = design_config.get_table_header_bold()
        header_color = design_config.get_table_header_font_color()
        header_align = design_config.get_table_header_align()
        header_fill = design_config.get_table_header_fill_color()
        for p in tf.paragraphs:
            for r in p.runs:
                if font_name is not None:
                    r.font.name = font_name
                if header_size is not None:
                    r.font.size = Pt(header_size)
                if header_bold is not None:
                    r.font.bold = header_bold
                if header_color is not None:
                    r.font.color.rgb = header_color
        if header_align is not None and header_align in _ALIGN_MAP:
            tf.paragraphs[0].alignment = _ALIGN_MAP[header_align]
        cell.fill.solid()
        if header_fill is not None:
            cell.fill.fore_color.rgb = header_fill

    for r in range(data_start_row, rows):
        for c in range(cols):
            cell = table.cell(r, c)
            tf = cell.text_frame
            if not tf.paragraphs or not tf.paragraphs[0].runs:
                continue
            is_index = c == index_col
            font_name = design_config.get_table_font()
            index_size = design_config.get_table_index_font_size()
            index_bold = design_config.get_table_index_bold()
            index_color = design_config.get_table_index_font_color()
            value_size = design_config.get_table_value_font_size()
            value_color = design_config.get_table_value_font_color()
            align = design_config.get_table_index_align() if is_index else design_config.get_table_value_align()
            fill = design_config.get_table_index_fill_color() if is_index else design_config.get_table_value_fill_color()
            for p in tf.paragraphs:
                for run in p.runs:
                    if font_name is not None:
                        run.font.name = font_name
                    if is_index:
                        if index_size is not None:
                            run.font.size = Pt(index_size)
                        if index_bold is not None:
                            run.font.bold = index_bold
                        if index_color is not None:
                            run.font.color.rgb = index_color
                    else:
                        if value_size is not None:
                            run.font.size = Pt(value_size)
                        if value_color is not None:
                            run.font.color.rgb = value_color
            if align is not None and align in _ALIGN_MAP:
                tf.paragraphs[0].alignment = _ALIGN_MAP[align]
            cell.fill.solid()
            if fill is not None:
                cell.fill.fore_color.rgb = fill


def apply_highlight_rules_by_value(table, highlight_rules):
    """
    Apply background color to cells based on their text value.
    """
    if not highlight_rules:
        return

    rows = len(table.rows)
    cols = len(table.columns)

    for i in range(1, rows):          # skip header
        for j in range(cols):
            cell = table.cell(i, j)
            val = cell.text.strip()

            if val in highlight_rules:
                r, g, b = highlight_rules[val]
                fill = cell.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(r, g, b)


def _apply_highlight_by_index(cell, row_idx, col_name, index_highlight_rules):
    if index_highlight_rules is None:
        return

    rules = index_highlight_rules.get(col_name)
    if rules is None:
        return

    # If user passed a single dict instead of a list, normalize it.
    if isinstance(rules, dict):
        rules = [rules]

    for rule in rules:
        indices = rule.get("indices")
        bg_color = rule.get("bg_color")

        if indices is None or bg_color is None:
            continue

        if row_idx in indices:
            fill = cell.fill
            fill.solid()
            r, g, b = bg_color
            fill.fore_color.rgb = RGBColor(r, g, b)
            return  # stop at first match


def _apply_highlight_by_col_condition(cell, value, col_name, highlight_rules):
    """
    Apply background color based on highlight_rules.
    highlight_rules example:
        {
            "ConversionRate": {"op": ">", "value": 0.5, "bg_color": (255, 255, 153)},
            "Margin":         {"op": "<", "value": 0.0, "bg_color": (255, 204, 204)},
        }
    """
    if highlight_rules is None:
        return

    if col_name not in highlight_rules:
        return

    rule = highlight_rules[col_name]
    op = rule.get("op")
    threshold = rule.get("value")
    bg_color = rule.get("bg_color")  # (R, G, B) tuple

    if not isinstance(value, (numbers.Number, np.number)):
        return

    if op is None or threshold is None or bg_color is None:
        return

    if _check_condition(float(value), op, float(threshold)):
        fill = cell.fill
        fill.solid()
        r, g, b = bg_color
        fill.fore_color.rgb = RGBColor(r, g, b)


def write_cell_text(
    cell,
    text,
    *,
    font_name=None,
    font_size=18,
    bold=False,
    color_rgb=None,
    align=None,
):
    """
    Safely write text to a PPTX table cell WITHOUT resetting runs later.

    This function is the ONLY correct way to write text if:
    - you apply themes
    - you apply highlights
    - you want font colors to persist
    """
    tf = cell.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "" if text is None else str(text)

    if font_name:
        run.font.name = font_name
    if font_size:
        run.font.size = Pt(font_size)
    run.font.bold = bold
    if color_rgb:
        run.font.color.rgb = RGBColor(*color_rgb)

    if align:
        p.alignment = align

    return run


def _apply_theme1_direct(table, *, header_row=0, data_start_row=1, index_col=0):
    """
    Apply Theme1 directly: header_dark, index_blue, value_light.
    Hardcoded to avoid config dependency issues.
    """
    rows = len(table.rows)
    cols = len(table.columns)
    
    # Resolve colors
    try:
        header_fill_rgb = RGBColor(*design_config.resolve_color("brand_navy"))
        header_font_rgb = RGBColor(*design_config.resolve_color("white"))
        index_fill_rgb = RGBColor(*design_config.resolve_color("brand_index_blue"))
        value_fill_rgb = RGBColor(*design_config.resolve_color("brand_gray"))
        value_font_rgb = RGBColor(*design_config.resolve_color("black"))
    except (ValueError, AttributeError):
        # Fallback if colors can't be resolved
        header_fill_rgb = RGBColor(0, 32, 96)  # dark blue
        header_font_rgb = RGBColor(255, 255, 255)  # white
        index_fill_rgb = RGBColor(68, 114, 196)  # blue
        value_fill_rgb = RGBColor(242, 242, 242)  # light gray
        value_font_rgb = RGBColor(0, 0, 0)  # black
    
    # Header row: header_dark (Pangram 18pt bold, white text, brand_navy fill, center)
    for c in range(cols):
        cell = table.cell(header_row, c)
        tf = cell.text_frame
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.name = "Pangram"
                r.font.size = Pt(18)
                r.font.bold = True
                r.font.color.rgb = header_font_rgb
        if tf.paragraphs:
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill_rgb
    
    # Body rows: index_blue for index col, value_light for others
    for r in range(data_start_row, rows):
        for c in range(cols):
            cell = table.cell(r, c)
            tf = cell.text_frame
            is_index = c == index_col
            
            if is_index:
                # index_blue: extends header_dark, font_size=16, fill_color=brand_index_blue
                for p in tf.paragraphs:
                    for run in p.runs:
                        run.font.name = "Pangram"
                        run.font.size = Pt(16)
                        run.font.bold = True
                        run.font.color.rgb = header_font_rgb
                if tf.paragraphs:
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.fill.solid()
                cell.fill.fore_color.rgb = index_fill_rgb
            else:
                # value_light: Pangram 18pt, not bold, black text, brand_gray fill, center
                for p in tf.paragraphs:
                    for run in p.runs:
                        run.font.name = "Pangram"
                        run.font.size = Pt(18)
                        run.font.bold = False
                        run.font.color.rgb = value_font_rgb
                if tf.paragraphs:
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.fill.solid()
                cell.fill.fore_color.rgb = value_fill_rgb

def _apply_criteria_theme(table, header_row=0, data_start_row=1, index_col=0):
    """
    Expert styling for Criteria/Reference tables:
    - Header: brand_navy + white bold text
    - Rows: brand_gray / white alternation
    - Index: bold text
    """
    rows = len(table.rows)
    cols = len(table.columns)
    
    try:
        brand_navy = RGBColor(*design_config.resolve_color("brand_navy"))
        brand_gray = RGBColor(*design_config.resolve_color("brand_gray"))
        white = RGBColor(255, 255, 255)
        black = RGBColor(0, 0, 0)
    except:
        brand_navy = RGBColor(0, 32, 96)
        brand_gray = RGBColor(217, 217, 217)
        white = RGBColor(255, 255, 255)
        black = RGBColor(0, 0, 0)
    
    pangram = "Pangram"
    
    # 1. Header
    for c in range(cols):
        cell = table.cell(header_row, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = brand_navy
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = pangram
                r.font.size = Pt(16)
                r.font.bold = True
                r.font.color.rgb = white
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 2. Body
    for r in range(data_start_row, rows):
        is_even = (r - data_start_row) % 2 == 0
        fill_color = brand_gray if is_even else white
        
        for c in range(cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color
            
            is_index = (c == index_col)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = pangram
                    r.font.size = Pt(14)
                    r.font.bold = is_index
                    r.font.color.rgb = black
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def new_table(
        prs,
        df,
        *,
        slide_index=None,
        slide_layout=5,
        title=None,
        left=0.5,
        top=1.5,
        width=20,
        height=None,
        include_index=True,
        index_header=None,
        apply_theme=False,
        table_type: Optional[str] = None,
        highlight_rules=None,
):
    if df.empty:
        raise ValueError("DataFrame is empty")

    df_display = df.copy()

    # Identify CBI rows before modifying the dataframe
    cbi_row_mask = df.index.astype(str).str.upper() == "CBI"

    if include_index:
        df_display.insert(0, index_header or "", df_display.index.astype(str))

    # Format numeric columns: rows with index "CBI" without decimals, others with 3 decimals
    # Check if columns contain numeric values (even if dtype is object)
    for col in df_display.columns:
        # Skip the index column if it was inserted
        if include_index and col == (index_header or ""):
            df_display[col] = df_display[col].astype(str)
            continue
            
        # Check if column is numeric: either numeric dtype OR all non-null values are numeric
        is_numeric = False
        if pd.api.types.is_numeric_dtype(df_display[col]):
            is_numeric = True
        else:
            # For object dtype, check if all non-null values can be converted to numeric
            try:
                numeric_series = pd.to_numeric(df_display[col], errors='coerce')
                non_null_count = df_display[col].notna().sum()
                numeric_count = numeric_series.notna().sum()
                # Only treat as numeric if most values (>80%) are numeric
                if non_null_count > 0 and numeric_count / non_null_count > 0.8:
                    is_numeric = True
            except (ValueError, TypeError):
                is_numeric = False
        
        if is_numeric:
            # Column is numeric - format it
            if pd.api.types.is_numeric_dtype(df_display[col]):
                numeric_col = df_display[col].astype(float)
            else:
                # Convert object dtype column to numeric (should be safe since we checked it's mostly numeric)
                numeric_col = pd.to_numeric(df_display[col], errors='coerce')
            
            if cbi_row_mask.any():
                # CBI rows: no decimals (format as integers)
                cbi_numeric = numeric_col.loc[cbi_row_mask]
                cbi_formatted = cbi_numeric.round(0).astype('Int64').astype(str).replace('<NA>', '')
                df_display.loc[cbi_row_mask, col] = cbi_formatted
                
                # Other rows: 3 decimals
                if (~cbi_row_mask).any():
                    other_mask = ~cbi_row_mask
                    other_numeric = numeric_col.loc[other_mask]
                    other_formatted = other_numeric.round(3).astype(str).replace('nan', '')
                    df_display.loc[other_mask, col] = other_formatted
            else:
                # No CBI row found, format all with 3 decimals
                df_display[col] = numeric_col.round(3).astype(str).replace('nan', '')
        else:
            # Non-numeric column - just convert to string, keep original values
            df_display[col] = df_display[col].astype(str)

    rows, cols = df_display.shape[0] + 1, df_display.shape[1]

    slide = prs.slides[slide_index] if slide_index is not None \
        else prs.slides.add_slide(prs.slide_layouts[slide_layout])

    if title and slide.shapes.title:
        slide.shapes.title.text = title

    # Calculate slide dimensions (convert from EMU to inches)
    # 1 inch = 914400 EMU
    slide_width_inches = prs.slide_width / 914400
    slide_height_inches = prs.slide_height / 914400
    
    # Calculate table dimensions to fit the slide
    # Use small margins (0.5 inches on each side)
    margin = 0.5
    table_width = slide_width_inches - (2 * margin)
    table_left = margin
    
    # Calculate table height - account for title if present
    title_height = 1.0 if (title and slide.shapes.title) else 0.0
    table_top = title_height + 0.3  # Small gap after title
    available_height = slide_height_inches - table_top - margin
    table_height = available_height if height is None else height
    table_height = max(0.6, table_height)  # Ensure minimum height

    table = slide.shapes.add_table(
        rows, cols, Inches(table_left), Inches(table_top), Inches(table_width), Inches(table_height)
    ).table

    # headers
    for j, col in enumerate(df_display.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        # Set alignment: center horizontal, middle vertical
        tf = cell.text_frame
        if tf.paragraphs:
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # data
    for i in range(df_display.shape[0]):
        for j in range(cols):
            cell = table.cell(i + 1, j)
            cell.text = df_display.iat[i, j]
            # Set alignment: center horizontal, middle vertical
            tf = cell.text_frame
            if tf.paragraphs:
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ✅ APPLY THEME FIRST
    if apply_theme:
        if table_type in ["criteria", "reference", "criteria_table"]:
            _apply_criteria_theme(table, header_row=0, data_start_row=1, index_col=0)
        else:
            _apply_theme1_direct(table, header_row=0, data_start_row=1, index_col=0)
        
        # Re-apply vertical alignment after theme (theme might override it)
        for i in range(rows):
            for j in range(cols):
                table.cell(i, j).vertical_anchor = MSO_ANCHOR.MIDDLE

    # ✅ APPLY HIGHLIGHT AFTER THEME
    apply_highlight_rules_by_value(table, highlight_rules)

    slide_idx = prs.slides.index(slide)
    return table, slide_idx


def template_table(
        table,
        df,
        pres=None,
        *,
        title_row=0,
        header_row=0,
        data_start_row=1,
        numeric_round_cols=None,
        percent_cols=None,
        decimals=2,
        column_override=None,
        highlight_rules=None,
        index_highlight_rules=None,
        new_name=None,
        apply_theme=False,
        slide_index_hint=None,
):
    """
    Fill an existing PPTX table using a DataFrame.
    Data is written without style; then theme is applied once; then highlights; then unused rows are removed.
    When slide_index_hint is provided, that slide index is used (avoids matching the wrong slide when
    multiple slides have tables with the same title, e.g. original template vs duplicated slide).
    
    Returns:
        int: The slide index containing the table, or None if pres is not provided.
    """
    # Find the slide index BEFORE modifying the table
    slide_index = None
    if slide_index_hint is not None and pres is not None and 0 <= slide_index_hint < len(pres.slides):
        slide_index = slide_index_hint
    elif pres is not None:
        try:
            table_title = table.cell(0, 0).text_frame.text.strip().lower()
            for idx, slide in enumerate(pres.slides):
                for shape in slide.shapes:
                    if not shape.has_table:
                        continue
                    try:
                        shape_table = shape.table
                        shape_title = shape_table.cell(0, 0).text_frame.text.strip().lower()
                        if shape_title == table_title:
                            slide_index = idx
                            break
                    except Exception:
                        continue
                if slide_index is not None:
                    break
        except Exception:
            pass
    
    max_rows = len(table.rows)
    max_cols = len(table.columns)

    df_rows, df_cols = df.shape

    # ----------------------------------
    # Header mapping (needed for theme column resolution)
    # ----------------------------------
    ppt_headers = [
        table.cell(header_row, i).text.strip()
        for i in range(len(table.columns))
    ]

    column_map = {
        col_idx: col_name
        for col_idx, col_name in enumerate(ppt_headers)
        if col_name in df.columns and col_idx < max_cols
    }

    if column_override:
        for ppt_col, df_col in column_override.items():
            if df_col in df.columns:
                column_map[ppt_col] = df_col

    # ----------------------------------
    # Numeric detection
    # ----------------------------------
    if numeric_round_cols is None:
        numeric_round_cols = [
            col for col in df.columns
            if pd.api.types.is_numeric_dtype(df[col])
        ]

    percent_cols = percent_cols or []
    if percent_cols == "all":
        percent_cols = numeric_round_cols

    # ----------------------------------
    # Optional title (plain, no style)
    # ----------------------------------
    if new_name is not None:
        write_cell_text(table.cell(title_row, 0), new_name)

    # ----------------------------------
    # Fill rows (plain text only, no style)
    # ----------------------------------
    max_rows = len(table.rows) - data_start_row
    num_rows = min(max_rows, len(df))

    for i in range(num_rows):
        row_idx = data_start_row + i
        row = df.iloc[i]

        for col_idx, df_col in column_map.items():
            if col_idx >= max_cols:
                continue
            value = row[df_col]

            # ---- format value ----
            if df_col in numeric_round_cols and pd.notna(value):
                value = float(value)
                if df_col in percent_cols:
                    text = f"{value * 100:.0f}%"
                else:
                    text = f"{value:.{decimals}f}"
            else:
                text = "" if pd.isna(value) else str(value)

            cell = table.cell(row_idx, col_idx)
            write_cell_text(cell, text)

    # ----------------------------------
    # Apply theme once (after populating)
    # ----------------------------------
    if apply_theme:
        apply_table_theme(
            table,
            header_row=header_row,
            data_start_row=data_start_row,
            index_col=0,
            column_headers=ppt_headers,
        )

    # ----------------------------------
    # Apply highlight rules (after theme)
    # ----------------------------------
    for i in range(num_rows):
        row_idx = data_start_row + i
        row = df.iloc[i]
        for col_idx, df_col in column_map.items():
            if col_idx >= max_cols:
                continue
            value = row[df_col]
            if df_col in numeric_round_cols and pd.notna(value):
                value = float(value)
            cell = table.cell(row_idx, col_idx)
            _apply_highlight_by_col_condition(cell, value, df_col, highlight_rules)
            _apply_highlight_by_index(cell, row_idx, df_col, index_highlight_rules)

    # ----------------------------------
    # Remove unused rows
    # ----------------------------------
    # Calculate rows to remove (from bottom up to avoid index shifting)
    rows_to_remove = len(table.rows) - (data_start_row + num_rows)
    if rows_to_remove > 0:
        # Get the table's XML structure
        tbl = table._tbl
        # Remove rows from bottom to top
        for _ in range(rows_to_remove):
            # Get the last row index
            last_row_idx = len(table.rows) - 1
            # Only remove if it's beyond the data rows
            if last_row_idx >= data_start_row + num_rows:
                row = table.rows[last_row_idx]
                tr = row._tr
                tbl.remove(tr)

    # ----------------------------------
    # Apply High-Fidelity Styling (Zebra & Highlights)
    # ----------------------------------
    apply_high_fidelity_styling(table, data_start_row=data_start_row)

    # Return the slide index we found at the beginning
    return slide_index



def add_df_table_slide(prs, df, title="Data Table"):
    """
    Add a new slide with a dark-blue theme and a table built from a small DataFrame.

    Parameters
    ----------
    prs : pptx.Presentation
        Existing Presentation object.
    df : pandas.DataFrame
        Small DataFrame to turn into a table.
    title : str, optional
        Title text for the slide.
    """
    # --- colors (dark blue theme) ---
    DARK_BLUE = RGBColor(0, 32, 96)  # slide background + header
    HEADER_TEXT = RGBColor(255, 255, 255)  # white
    BODY_TEXT = RGBColor(0, 0, 0)  # black

    # --- create slide (Title + Content layout is usually index 1, but can vary) ---
    # Fallback to blank layout if index 1 doesn't exist
    layout_idx = 1 if len(prs.slide_layouts) > 1 else 5
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)

    # --- slide background dark blue ---
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    # --- title ---
    if slide.shapes.title:
        title_shape = slide.shapes.title
        title_shape.text = title
        # Make title text white so it pops on dark blue background
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = HEADER_TEXT
                run.font.bold = True
                run.font.size = Pt(32)

    # --- table dimensions based on df shape ---
    rows = df.shape[0] + 1  # +1 for header
    cols = df.shape[1]

    # --- positioning and size of the table ---
    left = Inches(0.7)
    top = Inches(1.8)
    width = Inches(9)
    height = Inches(0.8 + 0.3 * rows)  # rough heuristic height

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # --- set column widths (evenly) ---
    for i in range(cols):
        table.columns[i].width = width // cols

    # --- header row ---
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)

        # styling header cell
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = DARK_BLUE

        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.color.rgb = HEADER_TEXT
                run.font.size = Pt(14)
            paragraph.alignment = PP_ALIGN.CENTER

    # --- body rows ---
    for row_idx in range(df.shape[0]):
        for col_idx in range(cols):
            value = df.iat[row_idx, col_idx]
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = "" if value is None else str(value)

            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = BODY_TEXT
                    run.font.size = Pt(12)
                paragraph.alignment = PP_ALIGN.LEFT

    return slide


def _check_condition(value, op, threshold):
    """Return True if value <op> threshold (for simple numeric ops)."""
    if op == ">":
        return value > threshold
    elif op == ">=":
        return value >= threshold
    elif op == "<":
        return value < threshold
    elif op == "<=":
        return value <= threshold
    elif op == "==":
        return value == threshold
    elif op == "!=":
        return value != threshold
    else:
        # Unsupported operator -> never highlight
        return False


def get_table_by_name(pres, name, slide=None):
    """
    Find a table by name.
    1. First tries to match via 'Alt Text' (shape.name) if a slide is provided.
    2. Then tries to match cell (0,0) title text (legacy).
    """
    name = name.strip().lower()
    
    # 1. Semantic match via Alt Text / Name (Reliable if set in template)
    if slide is not None:
        table = ShapeMapper.find_table_by_key(slide, name)
        if table:
            return table

    # 2. Match the table title in (0,0) (Legacy)
    slides_to_search = [slide] if slide is not None else pres.slides
    for s in slides_to_search:
        for shape in s.shapes:
            if not shape.has_table:
                continue

            table = shape.table

            # Match the table title in (0,0)
            try:
                title_text = table.cell(0, 0).text_frame.text.strip().lower()
                if title_text == name:
                    return table
            except (ValueError, IndexError):
                continue
    return None


def create_table(slide, rows, cols, left, top, width, height):
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return shape.table

def apply_high_fidelity_styling(table, data_start_row=1):
    """
    Applies executive zebra-striping and semantic row highlighting.
    """
    rows = len(table.rows)
    cols = len(table.columns)
    
    # 1. Fetch Palette from Theme
    try:
        zebra_rgb = design_config.resolve_color("brand_glass_blue")
    except Exception:
        zebra_rgb = (241, 245, 249) # Fallback to soft slate
    
    for r in range(data_start_row, rows):
        # Zebra Striping logic
        is_zebra = (r - data_start_row) % 2 == 1
        
        for c in range(cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            
            if is_zebra:
                cell.fill.fore_color.rgb = RGBColor(*zebra_rgb)
            else:
                cell.fill.background() # White/Transparent
                
            # Semantic Highlight: If a row contains 'Total' or 'Average'
            row_text = "".join([table.cell(r, i).text.lower() for i in range(cols)])
            if "total" in row_text or "average" in row_text:
                for i in range(cols):
                    for paragraph in table.cell(r, i).text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True

