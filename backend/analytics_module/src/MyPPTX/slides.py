import copy
import logging
import re
from typing import Any, Set, Optional

from pptx.parts.chart import ChartPart
from pptx.parts.embeddedpackage import EmbeddedXlsxPart
from pptx.parts.image import ImagePart, Image
from pptx.package import Package

from . import design_config

def get_slide_title(slide):
    if slide.shapes.title:
        return slide.shapes.title.text.strip()
    return None

def set_insight_text(slide, text: str) -> bool:
    """
    Expert Insight Injector: Dynamically seeks textboxes dedicated to insights 
    (e.g., shapes named 'insight', 'analysis', 'headline') and injects 
    AI-generated commentary with brand-consistent styling.
    """
    if not text:
        return False
        
    targets = ["insight", "analysis", "headline", "commentary", "conclusion", "deep_analysis"]
    font_name = design_config.get_chart_font() or "Pangram"
    
    found = False
    for shape in slide.shapes:
        shape_name = (shape.name or "").lower()
        # Checks if any target keyword matches the shape name
        if any(t in shape_name for t in targets):
            if hasattr(shape, "text_frame"):
                shape.text_frame.text = text
                _apply_brand_styling(shape, font_name)
                found = True
                
    # Fallback to Subtitle if no specific insight box is found
    if not found:
        for shape in slide.shapes:
            if "subtitle" in (shape.name or "").lower():
                if hasattr(shape, "text_frame"):
                    shape.text_frame.text = text
                    _apply_brand_styling(shape, font_name)
                    found = True
                    
    return found

def _apply_brand_styling(shape, font_name: str):
    """Reinforces corporate typography and professional spacing."""
    if not hasattr(shape, "text_frame"):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            # Slightly italicize insights for a premium narrative feel
            run.font.italic = True


def _slide_has_no_charts_and_no_tables(slide):
    """True if the slide has no charts and no tables (e.g. section header / divider slides)."""
    for shape in slide.shapes:
        if getattr(shape, "has_chart", False) and shape.has_chart:
            return False
        if getattr(shape, "has_table", False) and shape.has_table:
            return False
    return True


def find_slide_by_title(pres, title_text):
    """Return the first slide whose title contains title_text, skipping section-header slides (no charts, no tables)."""
    for idx, slide in enumerate(pres.slides):
        if _slide_has_no_charts_and_no_tables(slide):
            continue
        title = get_slide_title(slide)
        if title:
            if title_text.lower() in title.lower():
                return idx
    return None


def find_slide_by_title_exact(pres, title_text):
    """Return the first slide whose title exactly matches title_text (case-insensitive), skipping section-header slides."""
    for idx, slide in enumerate(pres.slides):
        if _slide_has_no_charts_and_no_tables(slide):
            continue
        title = get_slide_title(slide)
        if title:
            if title_text.strip().lower() == title.strip().lower():
                return idx
    return None


def _normalize_title_for_match(title):
    """Normalize for flexible matching: lower, collapse spaces/dashes to single hyphen."""
    if not title:
        return ""
    s = title.strip().lower()
    s = re.sub(r"[\s\u2013\u2014\-]+", "-", s)  # spaces, en-dash, em-dash, hyphen
    return re.sub(r"-+", "-", s).strip("-")  # collapse multiple hyphens


def find_slide_index_by_title_exact(pres, title_text):
    """Return the first slide index whose title exactly matches title_text (case-insensitive). Checks all slides (e.g. for recommendation slides that have no charts/tables)."""
    want = title_text.strip().lower()
    want_normalized = _normalize_title_for_match(title_text)
    for idx, slide in enumerate(pres.slides):
        title = get_slide_title(slide)
        if not title:
            continue
        t = title.strip().lower()
        if t == want:
            return idx
        if want_normalized and _normalize_title_for_match(title) == want_normalized:
            return idx
    return None


def find_section_header_by_title(pres, title_text):
    """Return the index of the first slide that has no charts and no tables and whose title contains title_text."""
    for idx, slide in enumerate(pres.slides):
        if not _slide_has_no_charts_and_no_tables(slide):
            continue
        title = get_slide_title(slide)
        if title and title_text.lower() in title.lower():
            return idx
    return None


def duplicate_charts_slide_by_number(pres, number):
    """
    Search for a slide with title f"{number}-charts",
    duplicate it, and append it to the end of the presentation.
    """
    target_title = f"{number}-charts"

    slide_index = find_slide_by_title(pres, target_title)

    if slide_index is None:
        raise ValueError(f"No slide found with title '{target_title}'")

    new_slide = duplicate_slide(pres, slide_index)
    return new_slide

def duplicate_brand_card_slide_by_number(pres):
    """
    Search for a slide with title f"{number}-charts",
    duplicate it, and append it to the end of the presentation.
    """
    target_title = f"Brand Card"

    slide_index = find_slide_by_title(pres, target_title)

    if slide_index is None:
        raise ValueError(f"No slide found with title '{target_title}'")

    new_slide = duplicate_slide(pres, slide_index)

    return new_slide


def __get_blank_slide_layout(pres):
    """
    Usage: Define the __get_blank_slide_layout function, which gets the slide layout with the fewest placeholders from the slide layout collection of a PowerPoint presentation.
    It takes a PowerPoint presentation object as a parameter.

    Parameters:
    -----------
    pres : pptx.Presentation
        DESCRIPTION: A PowerPoint presentation object.

    Returns:
    --------
    pptx.slide.SlideLayout
        The slide layout with the fewest placeholders.
    """
    # Create a list of the number of placeholders in each slide layout in the presentation.
    layout_items_count = [len(layout.placeholders)
                          for layout in pres.slide_layouts]
    # Get the minimum number of placeholders.
    min_items = min(layout_items_count)
    # Get the index of the slide layout with the minimum number of placeholders.
    blank_layout_id = layout_items_count.index(min_items)
    # Return the slide layout with the fewest placeholders.
    return pres.slide_layouts[blank_layout_id]

def duplicate_slide(pres, index):
    """
    Duplicate slide at `index` and append it to the presentation.
    Safe for charts, images, and media.
    """
    source = pres.slides[index]
    # Use the source slide's layout to preserve layout structure
    dest = pres.slides.add_slide(source.slide_layout)
    
    # ----------------------------------
    # Remove ALL shapes and placeholders from the new slide
    # ----------------------------------
    # Get list of all shapes to remove (must collect first, then remove)
    shapes_to_remove = []
    for shape in dest.shapes:
        shapes_to_remove.append(shape)
    
    # Remove all shapes from the layout
    for shape in shapes_to_remove:
        sp = shape.element
        sp.getparent().remove(sp)

    # ----------------------------------
    # Copy ALL shapes from source slide (including placeholders)
    # ----------------------------------
    for shape in source.shapes:
        newel = copy.deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

    # ----------------------------------
    # Copy relationships and map old rIds to new rIds
    # ----------------------------------
    rId_map = {}  # Maps old rId to new rId

    for rId in list(source.part.rels):
        try:
            rel = source.part.rels[rId]
        except KeyError:
            continue

        target = rel._target

        # --------- CHARTS ONLY ---------
        if isinstance(target, ChartPart):
            partname = target.package.next_partname(ChartPart.partname_template)
            new_chart_part = ChartPart(partname, target.content_type, target.package, copy.deepcopy(target._element))
            xlsx_blob = target.chart_workbook.xlsx_part.blob
            new_chart_part.chart_workbook.xlsx_part = EmbeddedXlsxPart.new(xlsx_blob, target.package)
            new_rId = dest.part.rels.get_or_add(rel.reltype, new_chart_part)
            rId_map[rel.rId] = new_rId

        # --------- IMAGES ---------
        elif isinstance(target, ImagePart):
            try:
                image = Image.from_blob(target.blob, getattr(target, "_filename", None))
                new_image_part = ImagePart.new(dest.part.package, image)
                new_rId = dest.part.rels.get_or_add(rel.reltype, new_image_part)
                rId_map[rel.rId] = new_rId
            except Exception:
                new_rId = dest.part.rels.get_or_add(rel.reltype, rel._target)
                rId_map[rel.rId] = new_rId

        # --------- EVERYTHING ELSE ---------
        else:
            new_rId = dest.part.rels.get_or_add(rel.reltype, rel._target)
            rId_map[rel.rId] = new_rId

    # ----------------------------------
    # Deep Relationship Re-Bridge (Definitive XML Sweep)
    # ----------------------------------
    from lxml import etree
    xml_str = etree.tostring(dest._element).decode('utf-8')
    
    # We must be careful to only replace exact rId="..." patterns or legacy r:id="..."
    # Sort keys by length descending to avoid partial replacement (rId10 vs rId1)
    for old_rid in sorted(rId_map.keys(), key=len, reverse=True):
        new_rid = rId_map[old_rid]
        # Cover standard attributes
        xml_str = xml_str.replace(f'"{old_rid}"', f'"{new_rid}"')
        xml_str = xml_str.replace(f'\'{old_rid}\'', f'\'{new_rid}\'')

    dest._element = etree.fromstring(xml_str.encode('utf-8'))
    return dest


def remove_unselected(selections, prs, indicies):
    """Legacy unselected removal (deprecated in favor of prune_presentation)."""
    keep = indicies.get("Intro", [])
    for section in selections:
        keep.extend(indicies.get(section, []))
    keep_set = {i - 1 for i in keep}

    for i in reversed(range(len(prs.slides))):
        if i not in keep_set:
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)  # drop relationship
            del prs.slides._sldIdLst[i]  # remove slide
    return prs

# ==============================================================================
# PHASE 2 : DYNAMIC ARCHETYPE MAPPING
# ==============================================================================

# ARCHETYPE_MAP: Maps logical chart types or specific analytical IDs 
# to their canonical template title in template.pptx.
ARCHETYPE_MAP = {
    # 1. Structural Slides
    "cover": "Cover Page",
    "divider": "Brand Awareness", # Base divider
    "closing": "Thank You",
    
    # 2. Core Analytical Templates
    "brand_awareness": "Brand Awareness",
    "purchase_funnel": "Purchase Funnel",
    "funnel_table": "Purchase Funnel Table",
    "drivers": "Drivers & Barriers",
    "imagery": "Imagery",
    "performance": "Performance",
    
    # 3. Chart-Type Fallbacks (Archetypes)
    "bar": "sc template",
    "column": "sc template",
    "stacked_bar": "mcmc template",
    "line": "2-sc-charts",
    "table": "Sensory Attribute Ratings Table",
    "radar": "Performance",
    "scatter": "Drivers & Barriers",
    
    # 4. Habit Archetypes
    "scsc": "scsc template",
    "mcmc": "mcmc template",
    "scmc": "scmc template",
    "mcsc": "mcsc template",
}

def resolve_template_index(pres, key: str) -> Optional[int]:
    """
    Dynamically finds the slide index for a given archetype key or chart_id.
    It first looks up the ARCHETYPE_MAP, then tries a title search.
    """
    # 1. Direct Archetype Lookup
    target_title = ARCHETYPE_MAP.get(key.lower())
    if target_title:
        idx = find_slide_index_by_title_exact(pres, target_title)
        if idx is not None:
            return idx

    # 2. Fuzzy Title Fallback
    idx = find_slide_index_by_title_exact(pres, key)
    if idx is not None:
        return idx
        
    return None

# LEGACY_SLIDE_REGISTRY: Kept strictly for backward compatibility with 
# the old 'prune_presentation' and legacy analytical engines.
SLIDE_REGISTRY = {
    1: {"section_id": "cover", "slide_type": "cover", "title_text": "Cover Page"},
    2: {"section_id": "brand_awareness", "slide_type": "divider", "title_text": "Brand Awareness"},
    3: {"section_id": "brand_awareness", "slide_type": "data", "title_text": "Brand Awareness"},
    # ... (Rest is implied as existing in many tools, but we focus on the new dynamic path)
}

def prune_presentation(prs: Package, store) -> Package:
    """
    Deletes all slides that do NOT have a slide_state of READY or DIVIDER.
    Operates explicitly via indexing and title matching.
    """
    slides_to_remove = []
    
    # Track the exact XML element matches to delete.
    # Python-pptx requires deleting via relationship ID tracking.
    
    # Because we delete, we figure out exact positional matching based on the baseline template mapping
    # Assuming baseline format maps index [0 ... 64] -> Slide [1 ... 65]
    for idx in range(len(prs.slides)):
        slide_n = idx + 1
        
        state = store.slide_states.get(slide_n, "SKIPPED")
        if state == "SKIPPED":
            slides_to_remove.append(idx)
            
    # Iterate backwards so we don't mess up indices
    for idx in reversed(slides_to_remove):
        try:
            rId = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx]
        except Exception as e:
            logging.getLogger(__name__).warning("Failed to prune slide at index %d: %s", idx, e)

    # -------------------------------------------------------------------------
    # PHASE 5: Handle BLOCKED slides (Wipe Content & Add Failure Title)
    # -------------------------------------------------------------------------
    # Now that we deleted skipped slides, indices shifted. 
    # Use title_text referencing to accurately locate remaining slides that are BLOCKED
    for slide_n, state in store.slide_states.items():
        if state == "BLOCKED":
            reg_info = SLIDE_REGISTRY.get(slide_n)
            if not reg_info:
                continue
            
            # Find the remaining slide by title string match
            target_title = reg_info.get("title_text", "")
            slide_idx = find_slide_index_by_title_exact(prs, target_title)
            
            if slide_idx is not None:
                slide = prs.slides[slide_idx]
                reasons = " | ".join(store.slide_blocking_reasons.get(slide_n, ["Missing critical fields"]))
                
                # Delete all shapes that are NOT the title
                shapes_to_delete = []
                for shape in slide.shapes:
                    if shape != slide.shapes.title:
                        shapes_to_delete.append(shape)
                        
                for shape in shapes_to_delete:
                    try:
                        sp = shape.element
                        sp.getparent().remove(sp)
                    except Exception:
                        pass
                        
                # Update title
                if slide.shapes.title:
                    try:
                        slide.shapes.title.text_frame.text = f"[SLIDE {slide_n} BLOCKED — {reasons}]"
                    except Exception:
                        pass

    # Check Brand Cards logic (cloning) here if brand_cards is selected
    if "brand_cards" in store.selected_sections:
        brands_count = len(store.brands)
        if brands_count > 2:
            # We need to clone the last brand card (slide 59 index, which after deletions might have shifted, 
            # so we must find by title instead in robust logic)
            pass
            
    store.validation_log.append({
        "slide_number": None,
        "event_type": "SLIDES_PRUNED",
        "message": f"Successfully deleted {len(slides_to_remove)} unselected slides.",
        "severity": "INFO"
    })
    
    return prs


def _remove_unmodified_slides(
        pres: Any,
        modified_slides: Set[int],
        section_names: Set[str],
        initial_slide_count: int,
        logger: logging.Logger
) -> None:
    """
    Expert Removal System: Performs a two-pass cleanup.
    Pass 1: Removes entirely unmodified slides.
    Pass 2: Sanitizes modified slides by pruning unpopulated placeholders and sample data.
    """
    slides_to_remove = []

    # PASS 1: Identify slides to drop
    for idx in range(len(pres.slides)):
        if idx in modified_slides:
            # PASS 2: Prune unpopulated shapes within modified slides
            prune_unpopulated_shapes(pres.slides[idx], logger)
            continue
        slides_to_remove.append(idx)

    # Execute removal in reverse order
    if slides_to_remove:
        logger.info("Advanced Pruning: Removing %d unmodified slides", len(slides_to_remove))
        for idx in reversed(slides_to_remove):
            try:
                rId = pres.slides._sldIdLst[idx].rId
                pres.part.drop_rel(rId)
                del pres.slides._sldIdLst[idx]
            except Exception as e:
                logger.warning("Failed to remove slide at index %d: %s", idx, e)


def prune_unpopulated_shapes(slide: Any, logger: logging.Logger) -> None:
    """
    Deep Sanitization: Scans a slide for leftover 'zombie' shapes (placeholders, 
    sample charts, or default text) and removes them to ensure a clean final product.
    """
    to_delete = []
    targets = ["click to add", "type text here", "content placeholder", "sample chart", "[chart title]"]
    
    for shape in slide.shapes:
        # Check for unpopulated text frames
        if shape.has_text_frame:
            text = shape.text_frame.text.strip().lower()
            if not text or any(t in text for t in targets):
                to_delete.append(shape)
                continue
                
        # Check for 'Sample' charts that weren't replaced
        if shape.has_chart:
            try:
                if shape.chart.has_title:
                    title = shape.chart.chart_title.text_frame.text.lower()
                    if "chart title" in title or "sample" in title:
                        to_delete.append(shape)
            except:
                pass

    # Batch delete
    for shape in to_delete:
        try:
            sp = shape.element
            sp.getparent().remove(sp)
        except Exception:
            pass


def _duplicate_section_header(
    pres, section_name: str, modified_slides: Set[int], title_override: Optional[str] = None, logger: Optional[logging.Logger] = None
) -> None:
    """Duplicate the section header slide for section_name, optionally set title, and add new slide index to modified_slides."""
    header_idx = find_section_header_by_title(pres, section_name)
    if header_idx is None:
        if logger:
            logger.warning("Section header not found for section: '%s'", section_name)
        return
    new_slide = duplicate_slide(pres, header_idx)
    new_slide_idx = pres.slides.index(new_slide)
    modified_slides.add(new_slide_idx)
    if title_override and new_slide.shapes.title:
        try:
            new_slide.shapes.title.text_frame.text = title_override
        except Exception:
            pass
