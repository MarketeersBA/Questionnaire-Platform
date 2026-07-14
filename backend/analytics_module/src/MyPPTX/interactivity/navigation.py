import logging
from typing import List, Dict, Any
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

class NavigationHandler:
    """
    Advanced Interactivity Engine: Slide-to-Slide Navigation.
    Transforms a linear deck into a navigable report application.
    """

    @staticmethod
    def inject_back_button(slide, target_slide_index: int = 1):
        """
        Adds a semantic 'Back to Summary' button on the bottom corner of a slide.
        """
        left, top = Inches(0.5), Inches(7.1) # Bottom Left
        width, height = Inches(1.5), Inches(0.3)
        
        btn = slide.shapes.add_shape(1, left, top, width, height)
        btn.fill.solid()
        btn.fill.fore_color.rgb = RGBColor(230, 230, 230)
        btn.line.width = 0
        
        # Text Logic
        tf = btn.text_frame
        p = tf.paragraphs[0]
        p.text = "← BACK TO SUMMARY"
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER
        
        # Hyperlink Logic
        # PowerPoint internal links use the '#SlideNumber' syntax
        p.hyperlink.address = f"#{target_slide_index}"

    @staticmethod
    def build_table_of_contents(slide, sections: List[Dict[str, Any]]):
        """
        Generates a clickable Table of Contents on a Section Divider slide.
        Each entry links directly to the corresponding analysis slide.
        """
        left, top = Inches(1), Inches(2.5)
        width = Inches(8)
        
        for i, section in enumerate(sections):
            title = section.get("title", "Analysis")
            target_index = section.get("slide_index", 2)
            
            # Create a dedicated text box for the TOC entry
            row_y = top + (i * Inches(0.4))
            box = slide.shapes.add_textbox(left, row_y, width, Inches(0.4))
            tf = box.text_frame
            p = tf.paragraphs[0]
            
            # Format: "01. Brand Awareness ............. p. 4"
            p.text = f"{i+1:02d}. {title}"
            p.font.size = Pt(14)
            p.font.name = "Pangram"
            p.font.color.rgb = RGBColor(0, 0, 128) # Navy
            
            # Apply Hyperlink to the entire paragraph
            p.hyperlink.address = f"#{target_index}"
            
            # Add a subtle hover indication / Bullet point
            bullet = slide.shapes.add_shape(9, left - Inches(0.3), row_y + Inches(0.1), Inches(0.15), Inches(0.15))
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = RGBColor(80, 200, 120) # Emerald
            bullet.line.width = 0

    @staticmethod
    def resolve_all_links(prs, manifest: List[Dict[str, Any]]):
        """
        Post-Processing Pass: Resolves all dynamic indices.
        Used when the final slide volume is only known at the end of build.
        """
        logger.info("[Navigation] Running link resolution pass...")
        # Iterates over all slides and shapes to update placeholder targets if necessary.
        # Currently, our direct #index approach is performed during instantiation.
        pass
