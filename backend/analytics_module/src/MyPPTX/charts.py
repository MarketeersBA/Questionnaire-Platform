"""
Chart generation and manipulation utilities for PowerPoint presentations.

This module provides functions to populate and format various chart types
in PowerPoint presentations using python-pptx. It supports:
- Bar/Column charts
- Line charts
- Scatter (XY) charts
- Stacked column charts

All styling is driven by the design_config module which loads themes from JSON.
"""

import logging
import math
from typing import Dict, List, Optional, Tuple, Union

import pandas as pd
from lxml import etree
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Pt

from . import design_config
from .mapping import ShapeMapper
from .formatter import ChartFormatter, LINE_CHART_TYPES, LEGEND_POSITION_MAP
from .slides import _slide_has_no_charts_and_no_tables

logger = logging.getLogger(__name__)

# Categories to exclude from all charts (bar, stacked column, line, category, etc.)
EXCLUDE_CATEGORIES = [
    "أخرى (حدد) [Specify]",
    "اخري (حدد) [Specify]",  # alternate spelling (ي vs ى)
    "اخرى (حددي) [Specify]",
    "Other (specify) [Specify]",
    "Other, specify [Specify]",
    "other",
]

# =============================================================================
# DATA HELPERS
# =============================================================================

class DataHelper:
    """Utilities for preparing chart data."""
    
    @staticmethod
    def to_category_chart_data(data: Union[pd.DataFrame, pd.Series]) -> CategoryChartData:
        """Convert pandas data to CategoryChartData.
        
        Args:
            data: DataFrame (columns=series, index=categories) or Series
            
        Returns:
            CategoryChartData object ready for chart.replace_data()
        """
        chart_data = CategoryChartData()
        
        # Categories from index
        chart_data.categories = [
            "" if pd.isna(c) else str(c) for c in data.index
        ]
        
        if isinstance(data, pd.Series):
            series_name = "" if pd.isna(data.name) else str(data.name)
            chart_data.add_series(series_name, DataHelper.sanitize_values(data.tolist(), fill_value=0))
        else:
            for col in data.columns:
                if (isinstance(col, float) and math.isnan(col)) or pd.isna(col):
                    series_name = ""
                else:
                    series_name = str(col)
                chart_data.add_series(series_name, DataHelper.sanitize_values(data[col].tolist(), fill_value=0))
        
        return chart_data
    
    @staticmethod
    def apply_orientation(chart, data: Union[pd.DataFrame, pd.Series], 
                          orientation: Optional[str] = None) -> Union[pd.DataFrame, pd.Series]:
        """Apply data orientation transformation.
        
        Args:
            chart: Chart object (used for heuristic logic if orientation=None)
            data: Input DataFrame or Series
            orientation: "column" (no transpose), "row" (transpose), or None (heuristic)
            
        Returns:
            Data with appropriate orientation applied
        """
        if orientation == "row":
            return data.T
        elif orientation == "column":
            return data
        elif orientation is not None:
            logger.warning("Unknown orientation '%s', falling back to heuristic", orientation)
        
        # Heuristic logic for auto-detection
        if chart.chart_type in LINE_CHART_TYPES:
            return data.T
        
        # Check series vs category count
        try:
            series_count = len(chart.series)
            category_count = len(chart.plots[0].categories)
            
            if series_count > category_count:
                return data.T
        except (AttributeError, IndexError):
            logger.debug("Could not determine chart structure, using data as-is")
        
        return data
    
    @staticmethod
    def sanitize_values(values: List, fill_value: float = 0) -> List:
        """Replace NaN/inf/-inf with a safe numeric value.
        
        Args:
            values: List of values to sanitize
            fill_value: Value to use for invalid entries
            
        Returns:
            Sanitized list of values
        """
        clean = []
        for v in values:
            if v is None:
                clean.append(fill_value)
            elif isinstance(v, float) and (math.isnan(v) or math.isinf(v)):
                clean.append(fill_value)
            else:
                clean.append(v)
        return clean


# =============================================================================
# CHART DISCOVERY
# =============================================================================

class ChartFinder:
    """Utilities for finding charts in presentations."""
    
    @staticmethod
    def _find_all_charts_by_title(pres, title: str, chart_type: Optional[frozenset] = None,
                                   target_slide_index: Optional[int] = None):
        """Find all charts by title text (internal helper).
        
        Args:
            pres: Presentation object
            title: Chart title to search for (case-insensitive)
            chart_type: Optional set of XL_CHART_TYPE to filter by
            target_slide_index: If set, search only this slide index; otherwise search all slides
            
        Yields:
            Tuples of (slide_index, shape, chart) for matching charts
        """
        title_lower = title.lower()
        slides_to_search = [(target_slide_index, pres.slides[target_slide_index])] if target_slide_index is not None else enumerate(pres.slides)

        for idx, slide in slides_to_search:
            # Skip section-header slides (no charts, no tables) when searching all slides
            if target_slide_index is None and _slide_has_no_charts_and_no_tables(slide):
                continue

            # 1. First try direct Alt Text / Name match (Semantic Key)
            target_shapes = ShapeMapper.find_all_by_key(slide, title)
            for shape in target_shapes:
                if shape.has_chart:
                    try:
                        chart = shape.chart
                        if not chart_type or chart.chart_type in chart_type:
                            yield idx, shape, chart
                            # If we found it by Alt Text, we probably don't need to search text content
                            # on this slide for the same key, but we'll continue for other shapes.
                    except Exception as e:
                        logger.warning(f"Skipping chart missing relationship in semantic search: {e}")
            
            # 2. Legacy check: search by Chart Title text
            for shape in slide.shapes:
                if not (shape.has_chart and shape.shape_type == MSO_SHAPE_TYPE.CHART):
                    continue
                
                try:
                    chart = shape.chart
                    
                    if chart_type and chart.chart_type not in chart_type:
                        continue
                    
                    if not chart.has_title:
                        continue
                    
                    chart_title = chart.chart_title.text_frame.text or ""
                    if chart_title.strip().lower() == title_lower:
                        yield idx, shape, chart
                except Exception as e:
                    logger.warning(f"Skipping chart missing relationship in legacy search: {e}")
    
    @staticmethod
    def get_chart_with_location(pres, name: str, chart_type: Optional[frozenset] = None,
                                target_slide_index: Optional[int] = None):
        """Get the first chart by name with its slide location and shape object.
        
        Args:
            pres: Presentation object
            name: Chart title to find (case-insensitive)
            chart_type: Optional set of XL_CHART_TYPE to filter by
            target_slide_index: If set, search only this slide
            
        Returns:
            Tuple of (slide_index, shape, chart) or None if not found
        """
        for slide_idx, shape, chart in ChartFinder._find_all_charts_by_title(
                pres, name, chart_type, target_slide_index=target_slide_index):
            return slide_idx, shape, chart
        return None
    
    @staticmethod
    def get_charts_from_slide(slide) -> List:
        """Get all charts from a slide, sorted left to right.
        
        Args:
            slide: Slide object
            
        Returns:
            List of chart objects sorted by horizontal position
        """
        chart_shapes = [
            (shape.left, shape.chart) 
            for shape in slide.shapes 
            if shape.has_chart
        ]
        chart_shapes.sort(key=lambda x: x[0])
        return [chart for _, chart in chart_shapes]


# =============================================================================
# CATEGORY CHART POPULATION
# =============================================================================

def populate(pres, data: Union[pd.DataFrame, pd.Series], chart_title_in_template: str,
             orientation: Optional[str] = None,
             order_columns: Optional[List[str]] = None,
             target_slide_index: Optional[int] = None,
             new_name: Optional[str] = None,
             ymin: Optional[float] = None,
             ymax: Optional[float] = None,
             chart_type_override: Optional[Union[str, XL_CHART_TYPE]] = None,
             color_override: Optional[List[Tuple[int, int, int]]] = None) -> List[int]:
    """Populate a category chart by title in a PowerPoint presentation.
    
    Args:
        pres: PowerPoint presentation object
        data: DataFrame or Series with chart data
        chart_title_in_template: Title of chart to find and populate
        orientation: "column" (default), "row" (transpose), or None (heuristic)
        order_columns: Optional list of series/category names in desired order; series
            will be colored by palette in this order (first = first palette color).
        target_slide_index: If set, only populate charts on this slide (e.g. after duplicating a template).
        new_name: Optional chart title to replace the template title after populating.
        ymin: Optional minimum scale for the value axis.
        ymax: Optional maximum scale for the value axis.
    
    Returns:
        List of slide indices (0-based) where charts were populated
    """
    populated_slide_indices = []
    
    for idx, shape, chart in ChartFinder._find_all_charts_by_title(
            pres, chart_title_in_template, target_slide_index=target_slide_index):
        logger.debug("Found chart '%s'", chart_title_in_template)
        
        # Handle empty or insignificant data by adding a metadata overlay instead of deleting
        base_n = 0
        if isinstance(data, (pd.DataFrame, pd.Series)):
             base_n = len(data)

        if base_n < 5: # Threshold for insignificant sample
            logger.warning("Data for '%s' is insignificant (N=%d) — applying validation overlay.", chart_title_in_template, base_n)
            
            # 1. Update Title to reflect Low Sample
            new_title = f"{new_name or chart_title_in_template} (N={base_n} [LOW SAMPLE])"
            ChartFormatter.apply_title_format(chart, new_title)

            # 2. Add Translucent "VALIDATION FAILED" Overlay
            try:
                from pptx.util import Pt, RGBColor
                from pptx.enum.text import PP_ALIGN
                from pptx.enum.shapes import MSO_SHAPE
                
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                parent_slide = shape.part.slide
                
                # Add a semi-transparent rectangle over the chart area
                overlay = parent_slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, left, top, width, height
                )
                overlay.fill.solid()
                overlay.fill.fore_color.rgb = RGBColor(255, 255, 255) # White
                # Note: True transparency requires XML manipulation in python-pptx
                # We use a dashed/hollow look for now, or just a centered label
                overlay.line.color.rgb = RGBColor(255, 0, 0) # Red border
                
                # Add text to overlay
                tf = overlay.text_frame
                tf.text = "INSIGNIFICANT SAMPLE - DATA SUPPRESSED"
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                p = tf.paragraphs[0]
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(255, 0, 0) # Red Text
                
            except Exception as pe:
                logger.warning("Failed to apply resilience overlay for chart: %s", pe)
            
            if base_n == 0:
                continue # Skip data population if truly zero

        _populate_single_chart(
            chart, data, new_title=new_name, orientation=orientation, order_columns=order_columns,
            chart_type_override=chart_type_override, color_override=color_override)
        if ymin is not None or ymax is not None:
            try:
                if ymin is not None:
                    chart.value_axis.minimum_scale = ymin
                if ymax is not None:
                    chart.value_axis.maximum_scale = ymax
            except Exception as e:
                logger.warning(
                    "Could not set axis limits for chart %s: %s", chart_title_in_template, e)
        if idx not in populated_slide_indices:
            populated_slide_indices.append(idx)
    
    return populated_slide_indices


def _populate_single_chart(
    chart,
    data: Union[pd.DataFrame, pd.Series],
    new_title: Optional[str] = None,
    orientation: Optional[str] = None,
    order_columns: Optional[List[str]] = None,
    chart_type_override: Optional[Union[str, XL_CHART_TYPE]] = None,
    color_override: Optional[List[Tuple[int, int, int]]] = None
) -> None:
    """Apply the same data + formatting as populate() to a single chart (used by populate() and brand cards)."""
    if isinstance(data, (pd.DataFrame, pd.Series)) and data.empty:
        return
    prepared_data = DataHelper.apply_orientation(chart, data, orientation)
    if isinstance(prepared_data, pd.DataFrame) and order_columns:
        _ordered = [c for c in order_columns if c in prepared_data.columns]
        if _ordered:
            _extra = [c for c in prepared_data.columns if c not in _ordered]
            prepared_data = prepared_data[_ordered + _extra]
    chart_data = DataHelper.to_category_chart_data(prepared_data)
    chart.replace_data(chart_data)

    # Apply Dynamic Overrides
    if chart_type_override:
        try:
            if isinstance(chart_type_override, str):
                chart.chart_type = getattr(XL_CHART_TYPE, chart_type_override.upper())
            else:
                chart.chart_type = chart_type_override
        except Exception as e:
            logger.warning("Failed to override chart type: %s", e)

    ChartFormatter.apply_title_format(chart, new_title if new_title else None)
    
    if color_override:
        try:
            rgb_colors = [RGBColor(*c) for c in color_override]
            ChartFormatter.apply_series_colors(chart, colors=rgb_colors)
        except Exception as e:
            logger.warning("Failed to apply color override: %s", e)
    elif chart.chart_type in LINE_CHART_TYPES:
        ChartFormatter.apply_line_series_colors(chart)
    else:
        ChartFormatter.apply_series_colors(chart)
    ChartFormatter.apply_data_labels(chart)
    ChartFormatter.apply_axis_label_format(chart)
    ChartFormatter.apply_legend_format(chart)


def populate_line_chart(pres, data: Union[pd.DataFrame, pd.Series], 
                        chart_title_in_template: str, new_name: Optional[str] = None,
                        ymax: float = 10, orientation: Optional[str] = None,
                        target_slide_index: Optional[int] = None) -> None:
    """Populate a line chart by title in a PowerPoint presentation.
    
    Args:
        pres: PowerPoint presentation object
        data: DataFrame or Series with chart data
        chart_title_in_template: Title of chart to find and populate
        new_name: Optional new title for the chart
        ymax: Maximum value for y-axis (currently unused, kept for API compatibility)
        orientation: "column", "row", or None (heuristic)
        target_slide_index: If set, only populate charts on this slide (e.g. after duplicating a template).
    """
    for idx, shape, chart in ChartFinder._find_all_charts_by_title(
        pres, chart_title_in_template, chart_type=LINE_CHART_TYPES,
        target_slide_index=target_slide_index
    ):
        # Delete chart if data is empty
        if isinstance(data, (pd.DataFrame, pd.Series)) and data.empty:
            logger.warning("Data for '%s' is empty — deleting line chart.", chart_title_in_template)
            shape.element.getparent().remove(shape.element)
            continue

        # Filter out excluded categories for all charts
        data = data[~data.index.astype(str).isin(EXCLUDE_CATEGORIES)]
        if isinstance(data, (pd.DataFrame, pd.Series)) and data.empty:
            logger.warning("Data for '%s' is empty after excluding categories — deleting line chart.", chart_title_in_template)
            shape.element.getparent().remove(shape.element)
            continue

        # Prepare and replace data
        prepared_data = DataHelper.apply_orientation(chart, data, orientation)
        chart_data = DataHelper.to_category_chart_data(prepared_data)
        chart.replace_data(chart_data)
        
        # Apply formatting
        ChartFormatter.apply_title_format(chart, new_name)
        ChartFormatter.apply_line_series_colors(chart)
        ChartFormatter.apply_data_labels(chart)
        ChartFormatter.apply_legend_format(chart)


def populate_category_chart(chartHandler, newTitle: str, df: pd.DataFrame,
                            x_col: str, y_cols: List[str],
                            ymin: Optional[float] = None, 
                            ymax: Optional[float] = None) -> None:
    """Populate a category chart with DataFrame data.
    
    Args:
        chartHandler: Chart object to populate
        newTitle: New title for the chart
        df: DataFrame containing data
        x_col: Column name for x-axis categories
        y_cols: List of column names for series data
        ymin: Optional minimum value for y-axis
        ymax: Optional maximum value for y-axis
    """
    # Filter out excluded categories for all charts
    df = df[~df[x_col].astype(str).isin(EXCLUDE_CATEGORIES)]
    if df.empty:
        return
    chart_data = CategoryChartData()
    chart_data.categories = df[x_col].tolist()
    
    for col in y_cols:
        chart_data.add_series(col, df[col].tolist())
    
    chartHandler.replace_data(chart_data)
    ChartFormatter.apply_title_format(chartHandler, newTitle)
    ChartFormatter.apply_axis_label_format(chartHandler)  # Add this line

    
    if ymin is not None:
        chartHandler.value_axis.minimum_scale = ymin
    if ymax is not None:
        chartHandler.value_axis.maximum_scale = ymax


# =============================================================================
# XY SCATTER CHART POPULATION
# =============================================================================

def update_scatter_chart_by_title_top_n(prs, chart_title: str, importance: Dict[str, float],
                                        brand_feature_scores: pd.DataFrame,
                                        new_name: Optional[str] = None, top_n: int = 4,
                                        show_labels: bool = True,
                                        label_category: Optional[str] = None,
                                        target_slide_index: Optional[int] = None) -> List[int]:
    """Update scatter chart by title with top N features.
    
    Args:
        prs: Presentation object
        chart_title: Title of the chart to update
        importance: Dict mapping feature names to importance values
        brand_feature_scores: DataFrame with brand feature scores
        new_name: Optional new name for the chart title
        top_n: Number of top features to display
        show_labels: Whether to show data labels on points
        label_category: If provided, only show labels for brands starting with this string
        target_slide_index: If set, only update charts on this slide (e.g. after duplicating a template).
    
    Returns:
        List of slide indices where charts were updated
    """
    # Extract top-N by importance
    sorted_features = sorted(importance.items(), key=lambda x: x[1], reverse=True)
    top_features = sorted_features[:top_n]
    top_feature_names = [f[0] for f in top_features]
    top_x_values = [f[1] for f in top_features]
    
    # Build filtered scores per brand
    brand_scores_t = brand_feature_scores.T
    brands = list(brand_scores_t.index)
    
    filtered_scores = {}
    series_feature_names = {}
    
    for brand in brands:
        y_vals = []
        feature_names_for_brand = []
        for fname in top_feature_names:
            feature_series = brand_scores_t.get(fname)
            if feature_series is not None and brand in feature_series.index:
                y_vals.append(float(feature_series.loc[brand]))
                feature_names_for_brand.append(fname)
        
        if y_vals:
            filtered_scores[brand] = y_vals
            series_feature_names[brand] = feature_names_for_brand
    
    if not filtered_scores:
        logger.warning("No matching scores found for any brand.")
        return []
    
    # Find and update charts
    updated_slide_indices = []
    slides_to_search = [(target_slide_index, prs.slides[target_slide_index])] if target_slide_index is not None else enumerate(prs.slides)

    for idx, slide in slides_to_search:
        if target_slide_index is None and _slide_has_no_charts_and_no_tables(slide):
            continue
        for shape in slide.shapes:
            if not hasattr(shape, "has_chart") or not shape.has_chart:
                continue
            
            chart = shape.chart
            
            if not chart.has_title or chart.chart_title.text_frame.text.strip() != chart_title:
                continue
            
            if chart.chart_type != XL_CHART_TYPE.XY_SCATTER:
                raise ValueError(
                    f"Chart '{chart_title}' is not an XY Scatter. Current type: {chart.chart_type}"
                )
            
            # Build scatter chart data
            chart_data = XyChartData()
            for brand, y_values in filtered_scores.items():
                series = chart_data.add_series(brand)
                for x, y in zip(top_x_values, y_values):
                    series.add_data_point(x, y)
            
            chart.replace_data(chart_data)
            ChartFormatter.apply_title_format(chart, new_name)
            
            # Add feature names as data labels
            if show_labels:
                _apply_scatter_labels(chart, series_feature_names, label_category)
            
            # Set axis bounds
            chart.category_axis.minimum_scale = min(top_x_values) * 0.7
            chart.category_axis.maximum_scale = max(top_x_values) * 1.3
            
            # Set axis titles
            chart.category_axis.has_title = True
            chart.category_axis.axis_title.text_frame.text = "Importance"
            ChartFormatter.apply_axis_title_format(chart.category_axis.axis_title)
            
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.text = "Brand Score"
            ChartFormatter.apply_axis_title_format(chart.value_axis.axis_title)
            ChartFormatter.apply_axis_label_format(chart) 
            ChartFormatter.apply_legend_format(chart)
            
            if idx not in updated_slide_indices:
                updated_slide_indices.append(idx)
            break
    
    return updated_slide_indices


def _apply_scatter_labels(chart, series_feature_names: Dict[str, List[str]],
                          label_category: Optional[str]) -> None:
    """Apply feature name labels to scatter chart points."""
    plot = chart.plots[0]
    
    for series in plot.series:
        brand_name = series.name
        
        if label_category is not None and not brand_name.startswith(label_category):
            continue
        
        feature_names = series_feature_names.get(brand_name, [])
        
        for point, fname in zip(series.points, feature_names):
            try:
                dl = point.data_label
                dl.show_value = True
                dl.show_series_name = False
                dl.show_category_name = False
                dl.show_legend_key = False
                
                if dl.text_frame is not None:
                    tf = dl.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = fname
                    pl_font_size = design_config.get_point_label_font_size()
                    if pl_font_size is not None:
                        run.font.size = Pt(pl_font_size)
                    pl_font_name = design_config.get_point_label_font_name()
                    if pl_font_name is not None:
                        run.font.name = pl_font_name
                    pl_color = design_config.get_point_label_color()
                    if pl_color is not None:
                        run.font.color.rgb = pl_color
            except AttributeError:
                continue


def _set_category_axis_label_rotation(chart, degrees: float) -> None:
    """Rotate category-axis tick labels by *degrees* (negative = counter-clockwise).

    python-pptx exposes no rotation property for tick labels, so we write directly
    to the ``<c:txPr>/<a:bodyPr rot="…">`` element.  The OOXML ``rot`` attribute is
    in 1/60 000ths of a degree (positive = clockwise).

    Args:
        chart: Chart object whose category axis labels should be rotated.
        degrees: Rotation in degrees; -45 produces the typical slanted label.
    """
    _NSMAP = {
        "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }
    rot_val = str(int(degrees * 60_000))
    try:
        cat_ax = chart.category_axis._element
        txPr = cat_ax.find(qn("c:txPr"))
        if txPr is None:
            txPr = etree.fromstring(
                f'<c:txPr xmlns:c="{_NSMAP["c"]}" xmlns:a="{_NSMAP["a"]}">'
                f'<a:bodyPr rot="{rot_val}"/><a:lstStyle/>'
                f"</c:txPr>"
            )
            cat_ax.append(txPr)
        else:
            bodyPr = txPr.find(qn("a:bodyPr"))
            if bodyPr is None:
                bodyPr = etree.fromstring(
                    f'<a:bodyPr xmlns:a="{_NSMAP["a"]}" rot="{rot_val}"/>'
                )
                txPr.insert(0, bodyPr)
            else:
                bodyPr.set("rot", rot_val)
    except Exception as exc:
        logger.warning("Could not set category axis label rotation: %s", exc)


_NSMAP_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _get_plot_area_bounds(chart, chart_shape) -> Tuple[int, int, int, int]:
    """Get plot area bounds in EMU (left, top, width, height) relative to chart shape.

    Reads c:plotArea/c:layout/c:manualLayout from chart XML when present.
    Otherwise uses defaults: left=12%, top=12%, width=75%, height=70% of chart shape.
    """
    cx, cy = chart_shape.left, chart_shape.top
    cw, ch = chart_shape.width, chart_shape.height

    try:
        plot_area = chart._chartSpace.chart.plotArea
        layout = plot_area.find(qn("c:layout"))
        if layout is None:
            raise ValueError("no layout")
        ml = layout.find(qn("c:manualLayout"))
        if ml is None:
            raise ValueError("no manualLayout")

        # Use explicit None check — lxml elements with no children are falsy,
        # so `el or fallback` would wrongly discard a valid empty element.
        def _find(ml, *tag_names):
            for tag in tag_names:
                el = ml.find(qn(tag))
                if el is not None:
                    return el
            return None

        def _val(el, default: float) -> float:
            if el is None:
                return default
            # el.get reads the XML attribute; getattr would only work on
            # python-pptx custom element classes, not raw lxml elements.
            v = el.get("val")
            if v is None:
                v = getattr(el, "val", None)
            return float(v) if v is not None else default

        def _mode(ml, mode_tag: str) -> str:
            el = _find(ml, mode_tag)
            if el is None:
                return "factor"
            v = el.get("val")
            if v is None:
                v = getattr(el, "val", None)
            return str(v).lower() if v else "factor"

        x_el = _find(ml, "c:x", "c:left")
        y_el = _find(ml, "c:y", "c:top")
        w_el = _find(ml, "c:w", "c:width")
        h_el = _find(ml, "c:h", "c:height")

        if x_el is None and y_el is None and w_el is None and h_el is None:
            raise ValueError("manualLayout empty")

        x_val = _val(x_el, 0.12)
        y_val = _val(y_el, 0.12)
        w_val = _val(w_el, 0.75)
        h_val = _val(h_el, 0.70)

        x_mode = _mode(ml, "c:xMode")
        y_mode = _mode(ml, "c:yMode")
        w_mode = _mode(ml, "c:wMode")
        h_mode = _mode(ml, "c:hMode")

        # edge mode: value is absolute position (fraction of chart)
        # factor mode: value is size/offset
        plot_left = cx + int(x_val * cw)  # x is always the left edge
        plot_top  = cy + int(y_val * ch)  # y is always the top edge

        if w_mode == "edge":
            # w = right edge fraction → width = (right - left) * cw
            plot_width = max(1, int((w_val - x_val) * cw))
        else:
            plot_width = int(w_val * cw)

        if h_mode == "edge":
            # h = bottom edge fraction → height = (bottom - top) * ch
            plot_height = max(1, int((h_val - y_val) * ch))
        else:
            plot_height = int(h_val * ch)

        return plot_left, plot_top, plot_width, plot_height
    except Exception:
        plot_left = cx + int(0.12 * cw)
        plot_top = cy + int(0.12 * ch)
        plot_width = int(0.75 * cw)
        plot_height = int(0.70 * ch)
        return plot_left, plot_top, plot_width, plot_height


def _get_plot_area_height_fraction(chart) -> float:
    """
    Read the plot-area height (c:plotArea/c:layout/c:manualLayout/c:h) as a fraction of chart height.

    Defaults to 0.70 when not present or not readable.
    """
    try:
        plot_area = chart._chartSpace.chart.plotArea
        layout = plot_area.find(qn("c:layout"))
        if layout is None:
            return 0.70
        ml = layout.find(qn("c:manualLayout"))
        if ml is None:
            return 0.70
        h_el = ml.find(qn("c:h")) or ml.find(qn("c:height"))
        if h_el is None:
            return 0.70
        v = h_el.get("val")
        if v is None:
            v = getattr(h_el, "val", None)
        return float(v) if v is not None else 0.70
    except Exception:
        return 0.70


def _set_plot_area_height_fraction(chart, height_frac: float) -> None:
    """
    Set plot-area manualLayout height fraction (factor mode) via OOXML.

    This adjusts the plot area only, not the outer chart shape.
    """
    # Keep within sane bounds; 1.0 can overlap legend/title depending on template.
    h = max(0.05, min(float(height_frac), 0.95))

    plot_area = chart._chartSpace.chart.plotArea

    layout = plot_area.find(qn("c:layout"))
    if layout is None:
        layout = etree.Element(qn("c:layout"))
        plot_area.insert(0, layout)

    ml = layout.find(qn("c:manualLayout"))
    if ml is None:
        ml = etree.Element(qn("c:manualLayout"))
        layout.append(ml)

    # Ensure factor mode (so c:h is interpreted as a fraction of chart height).
    h_mode = ml.find(qn("c:hMode"))
    if h_mode is None:
        h_mode = etree.Element(qn("c:hMode"))
        ml.append(h_mode)
    h_mode.set("val", "factor")

    h_el = ml.find(qn("c:h"))
    if h_el is None:
        h_el = etree.Element(qn("c:h"))
        ml.append(h_el)
    h_el.set("val", f"{h:.4f}")


def _get_plot_area_layout_yh_fractions(chart) -> Tuple[float, float]:
    """
    Read plot-area manualLayout y/h as fractions of chart height (factor mode).

    Returns (y, h). Defaults to (0.12, 0.70) when not present/readable.
    """
    try:
        plot_area = chart._chartSpace.chart.plotArea
        layout = plot_area.find(qn("c:layout"))
        if layout is None:
            return 0.12, 0.70
        ml = layout.find(qn("c:manualLayout"))
        if ml is None:
            return 0.12, 0.70
        y_el = ml.find(qn("c:y")) or ml.find(qn("c:top"))
        h_el = ml.find(qn("c:h")) or ml.find(qn("c:height"))

        def _val(el, default: float) -> float:
            if el is None:
                return default
            v = el.get("val")
            if v is None:
                v = getattr(el, "val", None)
            return float(v) if v is not None else default

        return _val(y_el, 0.12), _val(h_el, 0.70)
    except Exception:
        return 0.12, 0.70


def _set_plot_area_layout_yh_fractions(chart, *, y: Optional[float] = None, h: Optional[float] = None) -> None:
    """
    Set plot-area manualLayout y/h (factor mode) via OOXML.

    This adjusts the plot area only, not the outer chart shape.
    """
    plot_area = chart._chartSpace.chart.plotArea

    layout = plot_area.find(qn("c:layout"))
    if layout is None:
        layout = etree.Element(qn("c:layout"))
        plot_area.insert(0, layout)

    ml = layout.find(qn("c:manualLayout"))
    if ml is None:
        ml = etree.Element(qn("c:manualLayout"))
        layout.append(ml)

    # Ensure factor mode so y/h are interpreted as fractions of chart height.
    for mode_tag in ("c:yMode", "c:hMode"):
        mode_el = ml.find(qn(mode_tag))
        if mode_el is None:
            mode_el = etree.Element(qn(mode_tag))
            ml.append(mode_el)
        mode_el.set("val", "factor")

    def _clamp(v: float) -> float:
        # Keep within sane bounds; values near 1.0 can overlap legend/title.
        return max(0.01, min(float(v), 0.95))

    if y is not None:
        y_el = ml.find(qn("c:y"))
        if y_el is None:
            y_el = etree.Element(qn("c:y"))
            ml.append(y_el)
        y_el.set("val", f"{_clamp(y):.4f}")

    if h is not None:
        h_el = ml.find(qn("c:h"))
        if h_el is None:
            h_el = etree.Element(qn("c:h"))
            ml.append(h_el)
        h_el.set("val", f"{_clamp(h):.4f}")


def _add_importance_highlight_rect(
        slide, chart_shape, chart,
        categories: List[str],
        feature_positions: List[int],
        brand_y_values: Dict[str, List],
        highlight_top_n: int = 1,
        line_color: Tuple[int, int, int] = (91, 155, 213),
        line_width_pt: float = 1.5,
        x_padding_slots: float = 1.5,
        y_padding_frac: float = 0.15,
        rect_overrides: Optional[Dict] = None) -> None:
    """Draw a border rectangle on *slide* around the highest-importance feature points.

    Frames the rightmost (highest-importance) category's data points, with height
    matching the Y range of those points (not full plot height). Uses the chart's
    plot area bounds and value axis scale for correct positioning.

    Args:
        slide: Slide to add the rectangle to.
        chart_shape: The chart shape (provides position / size in EMU).
        chart: The chart object (for reading plot area layout from XML).
        categories: Full category list (features + dummy empty strings).
        feature_positions: Indices of actual feature categories within *categories*.
        brand_y_values: Dict of brand name -> list of Y values at each category slot.
        highlight_top_n: Number of rightmost (highest-importance) features to surround.
        line_color: RGB tuple for the rectangle border.
        line_width_pt: Border line width in points.
        x_padding_slots: Horizontal padding in category slots (wider = bigger box).
        y_padding_frac: Vertical padding as fraction of Y span (e.g. 0.15 = 15%).
        rect_overrides: Optional dict from config: x_padding_slots, y_padding_frac,
            x_offset_frac (nudge right +), y_offset_frac (nudge down +),
            category_index (-1=rightmost, 0=leftmost).
    """
    ro = rect_overrides or {}
    x_padding_slots = ro.get("x_padding_slots", x_padding_slots)
    y_padding_frac = ro.get("y_padding_frac", y_padding_frac)
    x_offset_frac = ro.get("x_offset_frac", 0.0)
    y_offset_frac = ro.get("y_offset_frac", 0.0)
    category_index = ro.get("category_index")

    if not feature_positions or highlight_top_n <= 0:
        logger.debug("highlight_rect: skipped (feature_positions=%s, top_n=%s)",
                     feature_positions, highlight_top_n)
        return

    n_cats = len(categories)
    if n_cats <= 1:
        return

    # Which category to highlight: -1=rightmost, 0=leftmost, 1=second from right, etc.
    if category_index is not None:
        idx = category_index if category_index >= 0 else len(feature_positions) + category_index
        idx = max(0, min(idx, len(feature_positions) - 1))
        highlight_pos = feature_positions[idx:idx + highlight_top_n]
    else:
        highlight_pos = feature_positions[-highlight_top_n:]
    left_cat = min(highlight_pos)
    right_cat = max(highlight_pos)

    # X extent: category span plus padding (wider box around the points).
    # Allow up to 1 slot overflow past the plot edge so the rect stays
    # centered on the data point even when it's the first or last category.
    slot_width = 1.0 / (n_cats - 1) if n_cats > 1 else 1.0
    x_left_frac = left_cat / (n_cats - 1) - x_padding_slots * slot_width
    x_right_frac = right_cat / (n_cats - 1) + x_padding_slots * slot_width
    x_left_frac = max(-slot_width, x_left_frac)
    x_right_frac = min(1.0 + slot_width, x_right_frac)

    plot_left, plot_top, plot_width, plot_height = _get_plot_area_bounds(chart, chart_shape)
    plot_bot = plot_top + plot_height
    cx, cy = chart_shape.left, chart_shape.top
    cw, ch = chart_shape.width, chart_shape.height

    # Collect Y values at highlighted positions
    y_vals: List[float] = []
    for pos in highlight_pos:
        for vals in brand_y_values.values():
            if pos < len(vals) and vals[pos] is not None:
                y_vals.append(float(vals[pos]))

    # Value axis scale (chart Y: 0 at bottom, max at top)
    try:
        y_axis_min = getattr(chart.value_axis, "minimum_scale", None) or 0.0
        y_axis_max = getattr(chart.value_axis, "maximum_scale", None) or 10.0
    except Exception:
        y_axis_min, y_axis_max = 0.0, 10.0
    y_range = y_axis_max - y_axis_min if y_axis_max > y_axis_min else 10.0

    if not y_vals:
        y_min_data, y_max_data = y_axis_min, y_axis_max
    else:
        y_min_data = min(y_vals)
        y_max_data = max(y_vals)
        pad = (y_max_data - y_min_data) * y_padding_frac if y_max_data > y_min_data else 0.5
        y_min_data = max(y_axis_min, y_min_data - pad)
        y_max_data = min(y_axis_max, y_max_data + pad)

    # Map data Y to plot coords (Y increases upward in plot)
    def _y_to_plot(y_data: float) -> int:
        frac = (y_data - y_axis_min) / y_range
        frac = max(0.0, min(1.0, frac))
        return plot_top + int((1.0 - frac) * plot_height)

    rect_top_plot = _y_to_plot(y_max_data)
    rect_bot_plot = _y_to_plot(y_min_data)
    rect_height = max(1, rect_bot_plot - rect_top_plot)

    rect_left   = plot_left + int(x_left_frac * plot_width) + int(x_offset_frac * plot_width)
    rect_width  = max(int(0.05 * plot_width), int((x_right_frac - x_left_frac) * plot_width))
    rect_top    = rect_top_plot + int(y_offset_frac * plot_height)

    logger.info(
        "highlight_rect: chart at (%s, %s) %sx%s  →  rect (%s, %s) %sx%s  "
        "[cats=%s, highlight_pos=%s]",
        cx, cy, cw, ch, rect_left, rect_top, rect_width, rect_height,
        n_cats, highlight_pos,
    )

    # Add rectangle shape
    rect = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        rect_left, rect_top, rect_width, rect_height,
    )

    # Transparent fill via XML <a:noFill/> – more reliable than fill.background()
    sp_pr = rect._element.spPr
    _FILL_TAGS = {qn("a:solidFill"), qn("a:gradFill"), qn("a:pattFill"),
                  qn("a:blipFill"), qn("a:grpFill"), qn("a:noFill")}
    for child in list(sp_pr):
        if child.tag in _FILL_TAGS:
            sp_pr.remove(child)
    sp_pr.append(etree.fromstring(f'<a:noFill xmlns:a="{_NSMAP_A}"/>'))

    # Coloured border
    rect.line.color.rgb = RGBColor(*line_color)
    rect.line.width = Pt(line_width_pt)


def _build_importance_spaced_categories(
        sorted_features: List[Tuple[str, float]],
        max_slots: int = 20) -> Tuple[List[str], List[int]]:
    """Build category list with dummy filler slots proportional to importance gaps.

    Features are placed at positions proportional to their importance values.
    Empty-string dummies fill the gaps so horizontal spacing reflects relative
    importance on a categorical axis.

    Args:
        sorted_features: (feature_name, importance) pairs sorted ascending by importance.
        max_slots: Number of slots the full importance range spans.

    Returns:
        (categories, positions) where *categories* contains feature names at their
        importance-proportional indices and ``""`` elsewhere.
    """
    if not sorted_features:
        return [], []

    importances = [f[1] for f in sorted_features]
    min_imp, max_imp = importances[0], importances[-1]

    if max_imp == min_imp or len(sorted_features) == 1:
        step = max_slots // max(len(sorted_features) - 1, 1)
        positions = [i * step for i in range(len(sorted_features))]
    else:
        positions = [
            round((imp - min_imp) / (max_imp - min_imp) * max_slots)
            for imp in importances
        ]

    # Resolve collisions: nudge duplicates one slot forward
    for i in range(1, len(positions)):
        if positions[i] <= positions[i - 1]:
            positions[i] = positions[i - 1] + 1

    categories: List[str] = [""] * (positions[-1] + 1)
    for (fname, _), pos in zip(sorted_features, positions):
        categories[pos] = fname

    return categories, positions


def update_importance_line_chart_by_title_top_n(
        prs, chart_title: str,
        importance,
        brand_feature_scores: pd.DataFrame,
        new_name: Optional[str] = None,
        top_n: int = 4,
        show_labels: bool = True,
        label_category: Optional[str] = None,
        target_slide_index: Optional[int] = None,
        max_slots: int = 20,
        highlight_top_n: int = 0,
        rect_overrides: Optional[Dict] = None) -> List[int]:
    """Update a line-with-markers chart, spacing features proportionally by importance.

    Y axis = average score per brand/comparator; X axis = feature names placed at
    positions proportional to their importance so the gaps between marks reflect
    relative importance distances (filled with dummy empty categories).

    Top-N features are selected by importance and sorted ascending left-to-right.
    When *highlight_top_n* > 0 a border rectangle is drawn on the slide around
    the *highlight_top_n* rightmost (highest-importance) feature points.

    Args:
        prs: Presentation object.
        chart_title: Title of the chart to update.
        importance: Mapping of feature name → importance value (dict or Series).
        brand_feature_scores: DataFrame with features as rows and brand-score columns.
        new_name: Optional new chart title.
        top_n: Number of top features to display.
        show_labels: Kept for API compatibility (not used for line charts).
        label_category: Kept for API compatibility.
        target_slide_index: If set, only update charts on this slide.
        max_slots: Number of slots the full importance range is spread across.
        highlight_top_n: Number of highest-importance features to surround with a
            rectangle (0 = no rectangle).

    Returns:
        List of slide indices where charts were updated.
    """
    # Select top-N by importance, then sort ascending for left-to-right display
    sorted_all = sorted(importance.items(), key=lambda x: x[1], reverse=True)
    top_features = sorted_all[:top_n]
    sorted_features = sorted(top_features, key=lambda x: x[1])
    top_feature_names = [f[0] for f in sorted_features]

    # Build importance-proportional category list
    categories, positions = _build_importance_spaced_categories(sorted_features, max_slots)
    if not categories:
        logger.warning("No categories built for importance line chart '%s'.", chart_title)
        return []

    # brand_feature_scores: index=features, columns=brand avg-score labels
    brand_scores_t = brand_feature_scores.T  # rows=brands, cols=features
    brands = list(brand_scores_t.index)

    # Build Y value arrays: score at feature slot, None at dummy slots
    brand_y_values: Dict[str, List] = {}
    for brand in brands:
        y_vals: List = [None] * len(categories)
        for fname, pos in zip(top_feature_names, positions):
            feature_series = brand_scores_t.get(fname)
            if feature_series is not None and brand in feature_series.index:
                y_vals[pos] = float(feature_series.loc[brand])
        brand_y_values[brand] = y_vals

    if not any(v is not None for vals in brand_y_values.values() for v in vals):
        logger.warning("No matching scores found for any brand in chart '%s'.", chart_title)
        return []

    updated_slide_indices: List[int] = []
    for idx, chart_shape, chart in ChartFinder._find_all_charts_by_title(
            prs, chart_title, chart_type=LINE_CHART_TYPES,
            target_slide_index=target_slide_index):
        chart_data = CategoryChartData()
        chart_data.categories = categories
        for brand, y_vals in brand_y_values.items():
            chart_data.add_series(brand, y_vals)

        chart.replace_data(chart_data)
        ChartFormatter.apply_title_format(chart, new_name)
        ChartFormatter.apply_line_series_colors(chart)
        ChartFormatter.apply_legend_format(chart)
        ChartFormatter.apply_axis_label_format(chart)

        chart.value_axis.has_title = True
        chart.value_axis.axis_title.text_frame.text = "Avg Score"
        ChartFormatter.apply_axis_title_format(chart.value_axis.axis_title)

        _set_category_axis_label_rotation(chart, -45)

        if highlight_top_n > 0:
            _add_importance_highlight_rect(
                prs.slides[idx], chart_shape, chart,
                categories, positions, brand_y_values,
                highlight_top_n=highlight_top_n,
                rect_overrides=rect_overrides,
            )

        if idx not in updated_slide_indices:
            updated_slide_indices.append(idx)

    return updated_slide_indices


def populate_xy_chart(rawDataFrame: pd.DataFrame, chartHandler, xAxis,
                      values: List[str], chartTitle: str,
                      label_column: Optional[str] = None,
                      xmax: float = 10, ymax: float = 26.5, ymin: float = 1,
                      show_label: bool = False,
                      chart_shape = None) -> None:
    """Populate an XY chart with DataFrame data.
    
    Args:
        rawDataFrame: DataFrame containing chart data
        chartHandler: Chart object to populate
        xAxis: Column name(s) for x-axis data
        values: List of column names for y-axis data
        chartTitle: Title for the chart
        label_column: Optional column for data labels
        xmax: Maximum x-axis scale
        ymax: Maximum y-axis scale
        ymin: Minimum y-axis scale
        show_label: Whether to display data labels
        chart_shape: Optional chart shape for plot area adjustment
    """
    chart_data = XyChartData()
    
    # Prepare labels
    labels_dict = {i: [] for i in range(len(values))}
    if label_column and label_column in rawDataFrame.columns:
        base_labels = list(rawDataFrame[label_column])
        for i in range(len(values)):
            labels_dict[i] = base_labels
    
    # Build chart data
    if isinstance(xAxis, str):
        x_values_list = list(rawDataFrame[xAxis])
        for i, value in enumerate(values):
            y_values_list = list(rawDataFrame[value])
            cd = chart_data.add_series(str(value), number_format=None)
            for x, y in zip(x_values_list, y_values_list):
                cd.add_data_point(x, y, number_format=None)
    else:
        for i, (x_col, y_col) in enumerate(zip(xAxis, values)):
            x_values_list = list(rawDataFrame[x_col])
            y_values_list = list(rawDataFrame[y_col])
            cd = chart_data.add_series(str(y_col), number_format=None)
            for x, y in zip(x_values_list, y_values_list):
                cd.add_data_point(x, y, number_format=None)
    
    chartHandler.replace_data(chart_data)
    ChartFormatter.apply_title_format(chartHandler, chartTitle)
    ChartFormatter.apply_line_series_colors(chartHandler)
    
    # Apply data labels
    if show_label:
        _apply_xy_data_labels(chartHandler, labels_dict)
    else:
        _hide_xy_data_labels(chartHandler)
    
    # Set axis scales
    chartHandler.category_axis.minimum_scale = 1
    chartHandler.category_axis.maximum_scale = xmax
    chartHandler.value_axis.minimum_scale = ymin
    chartHandler.value_axis.maximum_scale = ymax
    
    # Adjust plot area height based on data points
    _adjust_plot_area_height(chartHandler, len(rawDataFrame), chart_shape)
    
    ChartFormatter.apply_legend_format(chartHandler)
    ChartFormatter.apply_axis_label_format(chartHandler)  # Add this line


def _apply_xy_data_labels(chartHandler, labels_dict: Dict[int, List]) -> None:
    """Apply data labels to XY chart points."""
    for i, labels in labels_dict.items():
        if not labels:
            continue
        
        for point, label in zip(chartHandler.series[i].points, labels):
            data_label = point.data_label
            text_frame = data_label.text_frame
            text_frame.text = str(label)
            
            axis_label_size = design_config.get_axis_label_font_size()
            axis_label_name = design_config.get_axis_label_font_name()
            axis_label_color = design_config.get_axis_label_color()
            for paragraph in text_frame.paragraphs:
                if axis_label_size is not None:
                    paragraph.font.size = Pt(axis_label_size)
                if axis_label_name is not None:
                    paragraph.font.name = axis_label_name
                if axis_label_color is not None:
                    paragraph.font.color.rgb = axis_label_color
            point_label_color = design_config.get_point_label_color()
            if point_label_color is not None:
                data_label.font.color.rgb = point_label_color


def _hide_xy_data_labels(chartHandler) -> None:
    """Hide all data labels on XY chart."""
    for series in chartHandler.series:
        for point in series.points:
            point.data_label.text_frame.text = ""


def _adjust_plot_area_height(chartHandler, num_points: int, chart_shape) -> None:
    """Adjust chart height based on number of data points."""
    if chart_shape is None:
        chart_shape = _find_chart_shape(chartHandler)
    
    if chart_shape is None or num_points <= 0:
        return
    
    try:
        # Prefer adjusting plot area via manualLayout (OOXML) so the outer chart
        # shape (title/legend/margins) stays the same.
        #
        # For "rows-style" XY charts (e.g., Taste Test averages), the visual looks best
        # when the plot area height scales with the number of points/rows.
        base_y, base_h = _get_plot_area_layout_yh_fractions(chartHandler)

        # Normalize to typical maxima: Taste Test uses <= 9 rows, others may go higher.
        fit_max = 9 if num_points <= 9 else 20
        # Add a bit of padding so small counts (e.g., 3) aren't overly cramped.
        effective_points = num_points + 1
        effective_max = fit_max + 1

        target_h = base_h * (effective_points / effective_max)
        # Clamp: never exceed base_h, and don't shrink below a small fraction.
        target_h = max(base_h * 0.25, min(target_h, base_h))

        # Anchor to the top edge (high edge): keep y fixed, shrink only height.
        # This makes the bottom edge move up as rows decrease.
        target_y = base_y
        _set_plot_area_layout_yh_fractions(chartHandler, y=target_y, h=target_h)

        logger.debug(
            "Adjusted plot area: %d points, y=%.3f h=%.3f (base y=%.3f h=%.3f)",
            num_points,
            target_y,
            target_h,
            base_y,
            base_h,
        )
    except Exception as e:
        logger.warning("Could not adjust plot area height: %s", e)


def _find_chart_shape(chartHandler):
    """Attempt to find the shape containing a chart."""
    try:
        try:
            presentation = chartHandler._chart_part.package.presentation_part.presentation
        except AttributeError:
            try:
                presentation = chartHandler.part.package.presentation_part.presentation
            except AttributeError:
                return None
        
        if presentation is None:
            return None
        
        chart_id = id(chartHandler)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'chart'):
                    try:
                        if id(shape.chart) == chart_id:
                            return shape
                    except Exception:
                        continue
    except Exception as e:
        logger.debug("Could not find chart shape: %s", e)
    
    return None


# =============================================================================
# STACKED COLUMN CHART POPULATION
# =============================================================================

def populate_charts_from_columns(slide, df: pd.DataFrame, metric: str,
                                  same_color: bool = True,
                                  sync_horizontal_axis: bool = True) -> None:
    """Populate each chart in slide from each column in DataFrame.
    
    Args:
        slide: Slide object
        df: DataFrame where columns are chart data sources
        metric: Metric name for chart titles
        same_color: If True (default), all charts use the first theme color; if False, each chart uses a different palette color.
        sync_horizontal_axis: If True (default), all charts use the same horizontal axis scale as the first chart.
    """
    # Get charts and their corresponding shapes, sorted left to right
    chart_shapes_pairs = [
        (shape.left, shape, shape.chart) 
        for shape in slide.shapes 
        if shape.has_chart
    ]
    chart_shapes_pairs.sort(key=lambda x: x[0])
    charts = [chart for _, _, chart in chart_shapes_pairs]
    chart_shapes = [shape for _, shape, _ in chart_shapes_pairs]
    
    if len(charts) != len(df.columns):
        raise ValueError(f"Charts count ({len(charts)}) != columns count ({len(df.columns)})")
    
    # Filter out excluded categories for all charts
    df = df[~df.index.astype(str).isin(EXCLUDE_CATEGORIES)]
    if df.empty:
        logger.warning("No categories left after excluding for '%s'.", metric)
        return

    categories = df.index.tolist()
    # Get colors based on number of charts (if not same_color) or just 1 color
    num_colors_needed = len(df.columns) if not same_color else 1
    colors = design_config.get_color_palette(num_colors_needed)
    
    # Step 1: Initialize variables for storing chart properties
    first_chart_value_axis_min = None
    first_chart_value_axis_max = None
    first_chart_category_tick_offset = None  # Store tick label offset from first chart
    
    # Step 2: Truncate long category labels to prevent plot area expansion
    # Limit category labels to a reasonable length to prevent plot area width changes
    MAX_CATEGORY_LABEL_LENGTH = 30  # Maximum characters per label
    truncated_categories = []
    for cat in categories:
        if isinstance(cat, str) and len(cat) > MAX_CATEGORY_LABEL_LENGTH:
            # Truncate and add ellipsis
            truncated_categories.append(cat[:MAX_CATEGORY_LABEL_LENGTH - 3] + "...")
        else:
            truncated_categories.append(str(cat) if cat is not None else "")
    
    # Step 3: Populate all charts with data first and collect all data values
    all_data_max = None
    all_chart_values = []  # Store values for each chart to calculate max across all
    
    for chart_idx, (chart, col) in enumerate(zip(charts, df.columns)):
        values = DataHelper.sanitize_values(df[col].tolist(), fill_value=0)
        all_chart_values.append(values)
        
        chart_data = CategoryChartData()
        chart_data.categories = truncated_categories  # Use truncated categories
        chart_data.add_series("Series 1", values)
        
        chart.replace_data(chart_data)
    
    # Step 3: Calculate scale from maximum value across ALL charts
    if sync_horizontal_axis:
        try:
            # Find maximum value across all charts
            all_max_values = [max(values) for values in all_chart_values if values]
            if all_max_values:
                all_data_max = max(all_max_values)
                
                # Calculate scale: min at 0, max rounded UP to next nice number with headroom
                if all_data_max > 0:
                    # Add 10% headroom to ensure bars aren't cut off
                    max_with_headroom = all_data_max * 1.1
                    
                    # Round up to next nice increment based on magnitude
                    if max_with_headroom <= 0.1:
                        first_chart_value_axis_max = 0.1
                    elif max_with_headroom <= 0.2:
                        first_chart_value_axis_max = 0.2
                    elif max_with_headroom <= 0.3:
                        first_chart_value_axis_max = 0.3
                    elif max_with_headroom <= 0.4:
                        first_chart_value_axis_max = 0.4
                    elif max_with_headroom <= 0.5:
                        first_chart_value_axis_max = 0.5
                    elif max_with_headroom <= 0.6:
                        first_chart_value_axis_max = 0.6
                    elif max_with_headroom <= 0.7:
                        first_chart_value_axis_max = 0.7
                    elif max_with_headroom <= 0.8:
                        first_chart_value_axis_max = 0.8
                    elif max_with_headroom <= 0.9:
                        first_chart_value_axis_max = 0.9
                    elif max_with_headroom <= 1.0:
                        first_chart_value_axis_max = 1.0
                    else:
                        # For values > 1.0, round up to next 0.1 increment
                        first_chart_value_axis_max = math.ceil(max_with_headroom * 10) / 10
                else:
                    first_chart_value_axis_max = 0.1
                first_chart_value_axis_min = 0.0
                
                logger.debug("Calculated value axis scale from all charts: min=%s, max=%s (max_data_value=%s, max_with_headroom=%s)", 
                           first_chart_value_axis_min, first_chart_value_axis_max, all_data_max, all_data_max * 1.1 if all_data_max > 0 else 0)
                
                # Apply the calculated scale to ALL charts immediately after data replacement
                for chart_idx, chart in enumerate(charts):
                    try:
                        if hasattr(chart, 'value_axis') and chart.value_axis:
                            chart.value_axis.minimum_scale = first_chart_value_axis_min
                            chart.value_axis.maximum_scale = first_chart_value_axis_max
                            logger.debug("Set chart %d value axis scale after data: min=%s, max=%s", 
                                       chart_idx, first_chart_value_axis_min, first_chart_value_axis_max)
                    except Exception as e:
                        logger.debug("Couldn't set scale after data replacement for chart %d: %s", chart_idx, e)
        except Exception as e:
            logger.warning("Couldn't calculate value axis scale from all charts: %s", e)
    
    # Step 3: Apply basic formatting to all charts (title, colors, data labels)
    for chart_idx, (chart, col) in enumerate(zip(charts, df.columns)):
        title = f"{metric} by {col}" if str(col) == "Total" else str(col)
        ChartFormatter.apply_title_format(chart, title)
        if colors:
            color_idx = 0 if same_color else (chart_idx % len(colors))
            series_color = [colors[color_idx]]
            if chart.chart_type in LINE_CHART_TYPES:
                ChartFormatter.apply_line_series_colors(chart, colors=series_color)
            else:
                ChartFormatter.apply_series_colors(chart, colors=series_color)
        ChartFormatter.apply_data_labels(chart)
    
    # Step 4: Apply axis formatting to ALL charts FIRST to ensure consistent label space allocation
    for chart_idx, chart in enumerate(charts):
        ChartFormatter.apply_axis_label_format(chart)
    
    # Step 5: Read first chart's category tick offset to sync to other charts
    if len(charts) > 0:
        first_chart = charts[0]
        try:
            # Read category axis tick label offset to sync to other charts
            if hasattr(first_chart, 'category_axis') and first_chart.category_axis:
                if hasattr(first_chart.category_axis, 'tick_labels'):
                    try:
                        first_chart_category_tick_offset = first_chart.category_axis.tick_labels.offset
                    except AttributeError:
                        first_chart_category_tick_offset = None
                else:
                    first_chart_category_tick_offset = None
            else:
                first_chart_category_tick_offset = None
        except Exception as e:
            logger.debug("Couldn't read first chart category tick offset: %s", e)
    
    # Step 6: Now hide category axis labels on subsequent charts and sync properties
    for chart_idx, chart in enumerate(charts):
        # Sync scale after formatting (in case formatting reset it)
        if sync_horizontal_axis and first_chart_value_axis_min is not None and first_chart_value_axis_max is not None:
            try:
                if hasattr(chart, 'value_axis') and chart.value_axis:
                    chart.value_axis.minimum_scale = first_chart_value_axis_min
                    chart.value_axis.maximum_scale = first_chart_value_axis_max
                    logger.debug("Synced chart %d value axis scale after formatting: min=%s, max=%s", 
                               chart_idx, first_chart_value_axis_min, first_chart_value_axis_max)
            except Exception as e:
                logger.warning("Couldn't sync value axis scale for chart %d after formatting: %s", chart_idx, e)
        
        # For subsequent charts: hide category axis labels but keep axis line visible
        # IMPORTANT: Sync tick offset BEFORE hiding labels to ensure same label space
        if chart_idx > 0:
            try:
                if hasattr(chart, 'category_axis') and chart.category_axis:
                    category_axis = chart.category_axis
                    
                    # First, ensure category axis has the same tick offset as first chart
                    # This ensures the same label space is reserved even when labels are hidden
                    if hasattr(category_axis, 'tick_labels') and first_chart_category_tick_offset is not None:
                        tick_labels = category_axis.tick_labels
                        try:
                            if hasattr(tick_labels, 'offset'):
                                tick_labels.offset = first_chart_category_tick_offset
                                logger.debug("Set chart %d category tick offset to %s", chart_idx, first_chart_category_tick_offset)
                        except Exception as e:
                            logger.debug("Couldn't set category tick offset for chart %d: %s", chart_idx, e)
                    
                    # Now hide tick labels (but space is already reserved via offset)
                    try:
                        category_axis.tick_label_position = XL_TICK_LABEL_POSITION.NONE
                    except AttributeError:
                        # Fallback: try XML manipulation if property doesn't exist
                        try:
                            tick_labels_element = category_axis._element.find(qn('c:tickLblPos'))
                            if tick_labels_element is not None:
                                tick_labels_element.set(qn('val'), 'none')
                        except Exception:
                            pass
                    
                    # Ensure axis line is visible (for the vertical line at x=0)
                    try:
                        if hasattr(category_axis, 'format'):
                            axis_format = category_axis.format
                            if hasattr(axis_format, 'line'):
                                line = axis_format.line
                                axis_label_color = design_config.get_axis_label_color()
                                if axis_label_color:
                                    line.color.rgb = axis_label_color
                                line.width = Pt(1)
                    except Exception:
                        pass
            except Exception as e:
                logger.warning("Couldn't hide category axis labels for chart %d: %s", chart_idx, e)
        
        ChartFormatter.apply_legend_format(chart)
        
        # Sync scale again after legend formatting (in case it triggered recalculation)
        if sync_horizontal_axis and first_chart_value_axis_min is not None and first_chart_value_axis_max is not None:
            try:
                if hasattr(chart, 'value_axis') and chart.value_axis:
                    chart.value_axis.minimum_scale = first_chart_value_axis_min
                    chart.value_axis.maximum_scale = first_chart_value_axis_max
            except Exception as e:
                logger.warning("Couldn't sync value axis scale for chart %d after legend: %s", chart_idx, e)
    
    # Step 7: Final pass to ensure scales are synced for ALL charts (including first)
    # This is a safety net to ensure all charts have exactly the same scale
    if sync_horizontal_axis and first_chart_value_axis_min is not None and first_chart_value_axis_max is not None:
        for chart_idx, chart in enumerate(charts):
            try:
                if hasattr(chart, 'value_axis') and chart.value_axis:
                    # Force set scale explicitly for ALL charts to ensure consistency
                    chart.value_axis.minimum_scale = first_chart_value_axis_min
                    chart.value_axis.maximum_scale = first_chart_value_axis_max
                    logger.debug("Final sync chart %d value axis scale: min=%s, max=%s", 
                               chart_idx, first_chart_value_axis_min, first_chart_value_axis_max)
                    
                    # Verify it was set correctly
                    actual_min = chart.value_axis.minimum_scale
                    actual_max = chart.value_axis.maximum_scale
                    if actual_min != first_chart_value_axis_min or actual_max != first_chart_value_axis_max:
                        logger.warning("Chart %d axis scale not set correctly. Expected: min=%s, max=%s. Got: min=%s, max=%s",
                                     chart_idx, first_chart_value_axis_min, first_chart_value_axis_max, actual_min, actual_max)
            except Exception as e:
                logger.warning("Couldn't set value axis scale for chart %d in final pass: %s", chart_idx, e)


def populate_charts_from_columns_sc(slide, df: pd.DataFrame, metric: str,
                                    aggregation_method: str = "none") -> None:
    """Populate stacked column charts for Single Choice questions.
    
    Each chart represents one segment, with options as stacked series.
    
    Args:
        slide: Slide object
        df: DataFrame where rows are segments and columns are options (transposed)
        metric: Metric name for chart titles
        aggregation_method: "top6", "80_20", or "none"
    """
    charts = ChartFinder.get_charts_from_slide(slide)

    # Drop excluded option columns (e.g. "Other (Specify)") for all charts
    cols_to_keep = [c for c in df.columns if str(c) not in EXCLUDE_CATEGORIES]
    df = df[cols_to_keep]
    if df.empty:
        return

    if len(charts) != len(df.index):
        raise ValueError(f"Charts count ({len(charts)}) != segments count ({len(df.index)})")

    # Calculate number of unique options to determine palette size
    all_option_names = set()
    # print(df)
    # input("Press Enter to continue...df.index")
    # Sort columns based on the value in the "Total" index (row) if present
    
    if "Total" in df.index:
        # Get index of the "Total" row
        total_row = df.loc["Total"]
        # Sort columns descending by Total value
        sorted_columns = total_row.sort_values(ascending=False).index
        df = df[sorted_columns]
        # print(df)
        # input("Press Enter to continue...df.index after sorting")


    for segment in df.index:
        series_data = _aggregate_segment_data(df, segment, aggregation_method)
        for option_name, _ in series_data:
            all_option_names.add(str(option_name))
    num_colors_needed = len(all_option_names)
    
    colors = design_config.get_color_palette(num_colors_needed)
    # Build color map from Total row order (biggest = first color) so all charts share same mapping
    if colors and "Total" in df.index:
        total_series_data = _aggregate_segment_data(df, "Total", aggregation_method)
        ordered_options = [str(opt) for opt, _ in total_series_data]
        option_color_map = {
            opt: colors[idx % len(colors)] for idx, opt in enumerate(ordered_options)
        }
    else:
        option_color_map = _build_option_color_map(df, aggregation_method, colors) if colors else None
    
    for chart, segment in zip(charts, df.index):
        series_data = _aggregate_segment_data(df, segment, aggregation_method)
        
        chart_data = CategoryChartData()
        chart_data.categories = [str(segment)]
        
        for option_name, value in series_data:
            chart_data.add_series(str(option_name), [value])
        
        chart.replace_data(chart_data)
        title = f"{metric} by {segment}" if str(segment) == "Total" else str(segment)
        ChartFormatter.apply_title_format(chart, title)
        
        try:
            chart.value_axis.minimum_scale = 0
            chart.value_axis.reverse_order = False
        except Exception as e:
            logger.warning("Couldn't set value axis minimum scale: %s", e)
        
        if option_color_map and colors:
            _apply_series_colors_by_name(chart, option_color_map, colors[0])

        ChartFormatter.apply_data_labels(chart)
        ChartFormatter.apply_legend_format(chart)
        ChartFormatter.apply_axis_label_format(chart)

def _build_option_color_map(df: pd.DataFrame, aggregation_method: str,
                            colors: List[RGBColor]) -> Dict[str, RGBColor]:
    """Build consistent color mapping for all options across charts."""
    all_option_names = set()
    
    for segment in df.index:
        series_data = _aggregate_segment_data(df, segment, aggregation_method)
        for option_name, _ in series_data:
            all_option_names.add(str(option_name))
    
    # Sort options (Others always last)
    sorted_options = sorted([opt for opt in all_option_names if opt != "Others"])
    if "Others" in all_option_names:
        sorted_options.append("Others")
    
    return {
        option: colors[idx % len(colors)]
        for idx, option in enumerate(sorted_options)
    }


def _aggregate_segment_data(df: pd.DataFrame, segment, 
                            aggregation_method: str) -> List[Tuple[str, float]]:
    """Aggregate segment data based on method."""
    segment_data = {}
    for option in df.columns:
        value = df.loc[segment, option]
        if isinstance(value, pd.Series):
            value = value.iloc[0] if len(value) > 0 else 0
        elif pd.isna(value):
            value = 0
        else:
            value = float(value)
        segment_data[option] = value
    
    # Sort by value descending so largest-first order matches plot (bottom-up) and legend
    sorted_options = sorted(segment_data.items(), key=lambda x: x[1], reverse=True)
    total = sum(v for _, v in sorted_options)
    
    if aggregation_method == "none":
        return sorted_options  # largest first = bottom of stack, legend matches plot
        # return sorted_options[::-1]  # smallest first → stack bottom-up; legend shows largest first



def _apply_series_colors_by_name(chart, color_map: Dict[str, RGBColor],
                                  default_color: RGBColor) -> None:
    """Apply colors to series based on name mapping."""
    plot = chart.plots[0]
    
    # Clean up stale data point formatting
    try:
        for ser_el in plot._element.iterchildren(qn('c:ser')):
            for dPt in list(ser_el.iterchildren(qn('c:dPt'))):
                ser_el.remove(dPt)
    except Exception as e:
        logger.warning("Couldn't clean up old data point formatting: %s", e)
    
    try:
        plot.vary_by_categories = False
    except AttributeError:
        pass
    
    for series in plot.series:
        option_name = str(series.name)
        color = color_map.get(option_name, default_color)
        try:
            series_fill = series.format.fill
            series_fill.solid()
            series_fill.fore_color.rgb = color
        except Exception as e:
            logger.warning("Couldn't set color for series '%s': %s", option_name, e)


def populate_habits_chart(chart, df: pd.DataFrame, metric: str, is_sc: bool) -> None:
    """Populate a single chart for habits/opinions: SC = stacked column (1 cat, N series), MC = bar (N cats, 1 series).
    
    Args:
        chart: Chart object to populate
        df: DataFrame from value_percentages (index=options, one column)
        metric: Question name for chart title
        is_sc: True for Single Choice (stacked column), False for Multiple Choice (bar)
    """
    if df is None or df.empty:
        return
    col = df.columns[0]
    if is_sc:
        # Drop excluded categories (e.g. "Other (Specify)") from stacked column
        df = df[~df.index.astype(str).isin(EXCLUDE_CATEGORIES)]
        if df.empty:
            return
        # Sort by value descending so biggest gets first color and appears first in legend
        df_sorted = df.sort_values(by=col, ascending=False)
        chart_data = CategoryChartData()
        chart_data.categories = ["Total"]
        for opt in df_sorted.index:
            val = df_sorted.loc[opt, col]
            val = 0 if pd.isna(val) else float(val)
            chart_data.add_series(str(opt), [val])
        chart.replace_data(chart_data)
        try:
            chart.value_axis.minimum_scale = 0
            chart.value_axis.reverse_order = False
        except Exception:
            pass
        num_colors = len(df_sorted.index)
        colors = design_config.get_color_palette(num_colors)
        if colors:
            # Assign colors by value order: first (biggest) → first theme color
            option_color_map = {
                str(opt): colors[idx % len(colors)]
                for idx, opt in enumerate(df_sorted.index)
            }
            _apply_series_colors_by_name(chart, option_color_map, colors[0])
    else:
        # Drop excluded categories from bar chart
        df_bar = df[~df.index.astype(str).isin(EXCLUDE_CATEGORIES)]
        if df_bar.empty:
            return
        chart_data = CategoryChartData()
        cats = [str(x) for x in df_bar.index]
        vals = DataHelper.sanitize_values(df_bar[col].tolist(), fill_value=0)
        chart_data.categories = list(reversed(cats))
        chart_data.add_series(metric, list(reversed(vals)))
        chart.replace_data(chart_data)
    ChartFormatter.apply_title_format(chart, metric)
    ChartFormatter.apply_data_labels(chart)
    ChartFormatter.apply_legend_format(chart)
    ChartFormatter.apply_axis_label_format(chart)


# =============================================================================
# BRAND CARD POPULATION
# =============================================================================

def populate_brand_card(slide, pf: pd.DataFrame, why_mou: pd.DataFrame, 
                        brand_name: str, why_mou_n: int = None) -> None:
    """Populate brand card slide with purchase funnel and why MOU charts.
    
    Args:
        slide: Slide object
        pf: Purchase funnel DataFrame
        why_mou: Why MOU DataFrame
        brand_name: Brand name for titles
        why_mou_n: Optional number of respondents (appearances) for Why MOU title
    """
    why_mou_title = f"{brand_name} Why MOU"
    if why_mou_n is not None:
        why_mou_title = f"{brand_name} Why MOU ({why_mou_n})"
    for shape in slide.shapes:
        if not shape.has_chart:
            continue
        
        try:
            chart = shape.chart
            # Safe access to chart title
            has_title = hasattr(chart, "chart_title") and chart.chart_title is not None and chart.chart_title.has_text_frame
            chart_title = chart.chart_title.text_frame.text.strip().lower() if has_title else ""
            
            # --- SEMANTIC MAPPING ---
            # 1. Direct match by title
            is_funnel = "purchase funnel" in chart_title
            is_mou = "why mou" in chart_title
            
            # 2. Fallback by series category analysis (Phase 2 Advanced)
            if not is_funnel and not is_mou:
                cats = [str(c.label).lower() for c in chart.plots[0].categories]
                if any(k in cats for k in ["total awareness", "repurchase", "trial"]):
                    is_funnel = True
                elif any(k in cats for k in ["quality", "taste", "price", "price/value"]):
                    is_mou = True
            
            if is_funnel:
                _populate_brand_chart(chart, pf.T, f"{brand_name} Purchase Funnel")
            elif is_mou:
                # Values to exclude from charts
                exclude_values = ["other", "اخرى (حددي) [Specify]"]
                # Filter out excluded values from DataFrame
                why_mou_local = why_mou[~why_mou.index.isin(exclude_values)]
                _populate_single_chart(chart, why_mou_local, new_title=why_mou_title)
        except Exception as e:
                logger.warning("Skipping chart shape on brand card due to error: %s", e)
    
    # Update ratio text boxes
    _update_brand_ratios(slide, pf)


def _populate_brand_chart(chart, data: Union[pd.DataFrame, pd.Series], title: str) -> None:
    """Populate a brand card chart with data."""
    chart.chart_title.text_frame.text = title
    ChartFormatter.apply_title_format(chart)
    
    chart_data = DataHelper.to_category_chart_data(data)
    chart.replace_data(chart_data)
    
    ChartFormatter.apply_data_labels(chart)
    ChartFormatter.apply_legend_format(chart)
    ChartFormatter.apply_axis_label_format(chart)


def _update_brand_ratios(slide, pf: pd.DataFrame) -> None:
    """Update ratio text boxes on brand card slide."""
    pf_row = pf.iloc[0] if len(pf) > 0 else pd.Series()
    
    # Calculate ratios
    ratios = {
        "Attractiveness": _safe_ratio(pf_row, 'Trial', 'Total Awareness'),
        "Conversion": _safe_ratio(pf_row, 'Repurchase', 'Trial'),
        "Loyalty": _safe_ratio(pf_row, 'MOU', 'Repurchase'),
    }
    
    # Update text boxes
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        import re
        text = shape.text_frame.text.strip().lower()
        
        for ratio_name, ratio_value in ratios.items():
            # Robust case-insensitive regex pattern focusing on 'Ratio' and the specific metric name
            pattern = fr"(?i){re.escape(ratio_name)}.*ratio"
            if re.search(pattern, text):
                new_text = f"{ratio_name} Ratio: {ratio_value:.0f}%"
                _set_textbox_text(shape.text_frame, new_text)
                break


def _safe_ratio(row: pd.Series, numerator_key: str, denominator_key: str) -> float:
    """Calculate ratio with division-by-zero protection."""
    if numerator_key not in row.index or denominator_key not in row.index:
        return 0.0
    
    denominator = row.get(denominator_key, 0)
    if denominator == 0 or pd.isna(denominator):
        return 0.0
    
    numerator = row.get(numerator_key, 0)
    return (numerator / denominator) * 100


def _set_textbox_text(text_frame, text: str) -> None:
    """Set text box text with consistent formatting."""
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Pangram"
    run.font.size = Pt(18)


# =============================================================================
# SIMPLE UTILITY FUNCTIONS
# =============================================================================

def get_chart(pres, name: str, target_slide_index: Optional[int] = None):
    """Get a chart by name from presentation.
    
    Args:
        pres: Presentation object
        name: Chart title to find
        target_slide_index: If set, search only this slide
        
    Returns:
        Chart object or None
    """
    result = ChartFinder.get_chart_with_location(pres, name, target_slide_index=target_slide_index)
    if result:
        _, _, chart = result
        return chart
    return None


def get_charts_from_slide(slide) -> List:
    """Get all charts from a slide, sorted left to right.
    
    Args:
        slide: Slide object
        
    Returns:
        List of chart objects
    """
    return ChartFinder.get_charts_from_slide(slide)


def replace_chart_data(chart, categories: List, values: List) -> None:
    """Replace chart data completely while keeping formatting.
    
    Args:
        chart: Chart object
        categories: List of category labels
        values: List of values
    """
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Series 1", values)
    chart.replace_data(chart_data)


def set_chart_title(chart, title_text: str) -> None:
    """Set chart title with consistent formatting.
    
    Args:
        chart: Chart object
        title_text: New title text
    """
    ChartFormatter.apply_title_format(chart, str(title_text))


def sanitize_chart_values(values: List, fill_value: float = 0) -> List:
    """Replace NaN/inf/-inf with a safe numeric value.
    
    Args:
        values: List of values to sanitize
        fill_value: Value to use for invalid entries
        
    Returns:
        Sanitized list
    """
    return DataHelper.sanitize_values(values, fill_value)


# Backward compatibility alias
_apply_orientation_if_needed = DataHelper.apply_orientation
