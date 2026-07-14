import re

from pptx.dml.color import RGBColor

from pptx.util import Pt

from . import design_config

def replace_exact_text_all_slides(presentation, find_text, replace_text):

    """
    Replaces text in text boxes ONLY when the entire text box
    exactly matches find_text, across all slides,
    and applies formatting (bold, Pangram, size 14).
    """

    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            # Exact match only
            if text_frame.text == find_text:
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = replace_text
                # Formatting only when config provides values
                font = run.font
                bold = design_config.get_textbox_bold()
                if bold is not None:
                    font.bold = bold
                name = design_config.get_textbox_font()
                if name is not None:
                    font.name = name
                size = design_config.get_textbox_font_size()
                if size is not None:
                    font.size = Pt(size)
                color = design_config.get_textbox_font_color()
                if color is not None:
                    font.color.rgb = color


def populate_text_box(pres, original_text, new_text):
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                if original_text in text_frame.text:
                    text_frame.text = text_frame.text.replace(original_text, new_text)

def get_ordered_base_textboxes(slide):
    bases = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.lower()
            if text.startswith("base"):
                bases.append((text, shape))

    bases.sort(key=lambda x: int(re.findall(r"\d+", x[0])[0]))
    return [s for _, s in bases]

def populate_base_textboxes(slide, segments, base_map):
    """
    segments: list in the SAME order as charts
    base_map: {segment: count}
    """
    base_boxes = get_ordered_base_textboxes(slide)

    if len(base_boxes) != len(segments):
        raise ValueError(
            f"Base boxes ({len(base_boxes)}) != segments ({len(segments)})"
        )

    for box, segment in zip(base_boxes, segments):
        base_value = base_map.get(segment, 0)
        box.text_frame.text = f"Base: {base_value}"


def populate_subtitle_textbox(slide, text: str) -> bool:
    """
    Find a textbox on the slide that contains "subtitle" and replace it with the given text.
    Uses textbox formatting from design_config. Returns True if a subtitle shape was found and set.
    """
    if not text or not slide:
        return False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if "subtitle" not in shape.text_frame.text.strip().lower():
            continue
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        font = run.font
        bold = design_config.get_textbox_bold()
        if bold is not None:
            font.bold = bold
        name = design_config.get_textbox_font()
        if name is not None:
            font.name = name
        size = design_config.get_textbox_font_size()
        if size is not None:
            font.size = Pt(size)
        font.color.rgb = RGBColor(255, 255, 255)
        return True
    return False


def set_slide_title(slide, text: str) -> bool:
    """
    Set the slide's main title placeholder to the given text.
    Returns True if the title shape existed and was set.
    """
    if not text or not slide:
        return False
    title_shape = slide.shapes.title
    if not title_shape or not title_shape.has_text_frame:
        return False
    title_shape.text_frame.clear()
    p = title_shape.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    return True


def set_recommendations_bullets(slide, bullet_texts: list) -> bool:
    """
    Find the textbox containing "Write Recommendations Here" on the slide,
    replace it with the given bullet strings, and apply textbox formatting.
    Same first-match rule as generator._populate_recommendation_slides (run_pipeline_old).
    Each item in bullet_texts becomes one bullet paragraph. Returns True if the
    shape was found (and populated, or cleared when bullet_texts is empty).
    """
    if not slide:
        return False
    bullet_texts = bullet_texts or []
    placeholder_text = "Write Recommendations Here"
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        if placeholder_text.strip().lower() not in text_frame.text.strip().lower():
            continue
        text_frame.clear()
        for i, line in enumerate(bullet_texts):
            line = (line or "").strip()
            if not line:
                continue
            # Use add_paragraph() for every bullet so the first one also gets bullet formatting
            p = text_frame.add_paragraph()
            p.level = 0
            run = p.add_run()
            run.text = line
            font = run.font
            bold = design_config.get_textbox_bold()
            if bold is not None:
                font.bold = bold
            name = design_config.get_textbox_font()
            if name is not None:
                font.name = name
            size = design_config.get_textbox_font_size()
            if size is not None:
                font.size = Pt(size)
            color = design_config.get_textbox_font_color()
            if color is not None:
                font.color.rgb = color
            # Empty line between recommendations (except after the last)
            if i < len(bullet_texts) - 1:
                text_frame.add_paragraph()
        # Found and populated (or cleared when bullet_texts empty)
        return True
    return False


def set_insight_text(slide, text: str) -> bool:
    """
    Set the slide's title placeholder to the insight text.
    Uses the same textbox formatting from design_config. Returns True if the title shape existed and was set.
    """
    if not text or not slide:
        return False
    title_shape = slide.shapes.title
    if not title_shape or not title_shape.has_text_frame:
        return False
    title_shape.text_frame.clear()
    p = title_shape.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    bold = design_config.get_insight_textbox_bold()
    if bold is not None:
        font.bold = bold
    name = design_config.get_insight_textbox_font()
    if name is not None:
        font.name = name
    size = design_config.get_insight_textbox_font_size()
    if size is not None:
        font.size = Pt(size)
    color = design_config.get_insight_textbox_font_color()
    if color is not None:
        font.color.rgb = color
    return True


def set_footer_metadata(slide, base_n=None, footnote=None) -> None:
    """
    Expert Metadata Injector: Identifies and hydrates absolute-positioned footer 
    placeholders with sample sizes and research footnotes.
    """
    font_name = design_config.get_chart_font() or "Pangram"
    
    # 1. Handle Base Size (N=xxx)
    if base_n is not None:
        base_text = f"Base: N = {base_n}"
        found_base = False
        for shape in slide.shapes:
            name = (shape.name or "").lower()
            if "base" in name or "sample" in name:
                if hasattr(shape, "text_frame"):
                    shape.text_frame.text = base_text
                    _apply_footer_style(shape, font_name)
                    found_base = True
                    
    # 2. Handle Qualitative Footnotes
    if footnote:
        found_fn = False
        for shape in slide.shapes:
            name = (shape.name or "").lower()
            if "footnote" in name or "note" in name or "source" in name:
                if hasattr(shape, "text_frame"):
                    shape.text_frame.text = footnote
                    _apply_footer_style(shape, font_name)
                    found_fn = True


def _apply_footer_style(shape, font_name: str):
    """Applies elite research-grade styling to footer metadata."""
    if not hasattr(shape, "text_frame"):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(10)
            # Use a professional muted gray for metadata
            run.font.color.rgb = RGBColor(120, 120, 120)
