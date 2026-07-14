import logging
from typing import Dict, Any, List
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from .base_builder import BaseChartBuilder
from .builder_render_status import BuilderEmptyDataError

logger = logging.getLogger(__name__)

class PPTXScorecard(BaseChartBuilder):
    """
    Builder for Executive Brand Scorecards.
    Features:
    - Multi-shape 'KPI Box' design
    - Directional trend arrows (Up/Down) color-coded
    - Secondary metadata (N-sizes, benchmarks)
    - Clean dashboard aesthetic using glassmorphism principles
    """

    def render(self, slide: Slide, chart_data: Dict[str, Any]) -> None:
        data_payload = chart_data.get("data", {})
        brand_name = data_payload.get("brand", "OVERALL")
        n_size = data_payload.get("n_size", 0)
        metrics = data_payload.get("metrics", [])
        
        if not metrics:
            raise BuilderEmptyDataError("No scorecard metrics to render.")

        # 1. Brand Header Information
        self._add_brand_header(slide, brand_name, n_size)

        # 2. Scorecard Grid Logic
        # Columns: 3, Rows: up to 2 (Focus on top 6 KPIs)
        col_count = 3
        spacing = Inches(0.3)
        box_w = (self.layout.CHART_WIDTH - (spacing * (col_count - 1))) / col_count
        box_h = Inches(2.7)
        
        start_left = self.layout.CHART_LEFT
        start_top = self.layout.CHART_TOP + Inches(0.6)
        
        for i, m in enumerate(metrics[:6]):
            row = i // col_count
            col = i % col_count
            
            left = start_left + (col * (box_w + spacing))
            top = start_top + (row * (box_h + spacing))
            
            self._draw_kpi_tile(slide, left, top, box_w, box_h, m)

    def _add_brand_header(self, slide: Slide, name: str, n: int):
        """Adds a small brand identity badge at the top of the chart area."""
        box = slide.shapes.add_textbox(self.layout.CHART_LEFT, self.layout.CHART_TOP - Inches(0.45), Inches(6), Inches(0.6))
        p = box.text_frame.paragraphs[0]
        p.text = f"{name.upper()} | N={n}"
        p.font.name = self.theme.FONT_BOLD
        p.font.size = Pt(16)
        p.font.color.rgb = self.theme.PRIMARY_BRAND

    def _draw_kpi_tile(self, slide: Slide, left: float, top: float, width: float, height: float, data: dict):
        """Draws a single high-impact KPI tile with value and trend."""
        # 1. Background
        tile = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        tile.fill.solid()
        tile.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tile.line.color.rgb = self.theme.get_rgb_by_name("brand_light_gray")
        tile.line.width = Pt(1.5)
        
        # 2. Metric Name
        label_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.45))
        lp = label_box.text_frame.paragraphs[0]
        lp.text = data.get("label", "").upper()
        lp.font.name = self.theme.FONT_MEDIUM
        lp.font.size = Pt(13)
        lp.font.color.rgb = self.theme.SUBTITLE_COLOR
        lp.alignment = PP_ALIGN.CENTER
        
        # 3. Main Value
        val_box = slide.shapes.add_textbox(left, top + (height/2) - Inches(0.6), width, Inches(1.2))
        vp = val_box.text_frame.paragraphs[0]
        val = data.get("value", 0)
        vp.text = str(val)
        vp.font.name = self.theme.FONT_BOLD
        vp.font.size = Pt(54)
        vp.font.color.rgb = self.theme.TEXT_COLOR
        vp.alignment = PP_ALIGN.CENTER
        
        # 4. Trend Indicator
        trend = data.get("trend", "neutral") # up, down, neutral
        change = data.get("change", 0)
        
        if trend != "neutral":
            arrow_size = Inches(0.375)
            arrow_x = left + (width/2) - (arrow_size/2)
            arrow_y = top + height - Inches(0.675)
            
            shape_type = MSO_SHAPE.UP_ARROW if trend == "up" else MSO_SHAPE.DOWN_ARROW
            color = RGBColor(0x16, 0xA3, 0x4A) if trend == "up" else RGBColor(0xDC, 0x26, 0x26)
            
            arrow = slide.shapes.add_shape(shape_type, arrow_x, arrow_y, arrow_size, arrow_size)
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = color
            arrow.line.visible = False
            
            # Trend Value Text
            t_box = slide.shapes.add_textbox(arrow_x + arrow_size + Inches(0.075), arrow_y, Inches(0.75), arrow_size)
            tp = t_box.text_frame.paragraphs[0]
            tp.text = f"{change:+.1f}"
            tp.font.size = Pt(14)
            tp.font.name = self.theme.FONT_BOLD
            tp.font.color.rgb = color
            tp.alignment = PP_ALIGN.LEFT
