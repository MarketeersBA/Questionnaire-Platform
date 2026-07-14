from __future__ import annotations

import io
import zipfile
from collections import Counter
from pathlib import Path
from typing import Any, BinaryIO, Dict, List, Union

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .pptx_geometry_audit import audit_presentation_geometry

UNSUPPORTED_MARKERS = (
    "DEVELOPMENT IN PROGRESS: NATIVE",
    "PREVIEW UNAVAILABLE",
)
ERROR_MARKERS = (
    "ANALYSIS INTERRUPTED",
)
NARRATIVE_MARKERS = (
    "EXECUTIVE SUMMARY",
    "STRATEGIC INTELLIGENCE",
    "EXECUTION PLAYBOOK",
    "INSIGHT:",
    "COMPETITIVE SWOT",
    "4P RECOMMENDATIONS",
    "MARKET ARCHETYPE",
    "STRATEGIC POSITIONING",
)


def audit_pptx_path(path: Union[str, Path]) -> Dict[str, Any]:
    return audit_pptx_bytes(Path(path).read_bytes())


def audit_pptx_stream(stream: BinaryIO) -> Dict[str, Any]:
    position = stream.tell()
    data = stream.read()
    stream.seek(position)
    return audit_pptx_bytes(data)


def audit_pptx_bytes(pptx_bytes: bytes) -> Dict[str, Any]:
    """Read-only structural audit suitable for manifests and diagnostics."""
    with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as archive:
        package_audit = _audit_package_archive(archive)

    presentation_audit: Dict[str, Any] = {
        "presentation_readable": False,
        "slide_titles": [],
        "chart_shape_count": 0,
        "layout_warning_count": 0,
        "duplicate_title_count": 0,
        "duplicate_titles": [],
        "out_of_bounds_shape_count": 0,
    }

    try:
        presentation = Presentation(io.BytesIO(pptx_bytes))
        presentation_audit = _audit_presentation(presentation)
    except Exception as exc:
        presentation_audit["presentation_error"] = str(exc)

    return _build_compact_audit(package_audit, presentation_audit)


def _audit_package_archive(archive: zipfile.ZipFile) -> Dict[str, Any]:
    names = archive.namelist()
    slide_files = [
        name for name in names if name.startswith("ppt/slides/slide") and name.endswith(".xml")
    ]
    notes_files = [
        name for name in names if name.startswith("ppt/notesSlides/notesSlide") and name.endswith(".xml")
    ]
    chart_files = [name for name in names if name.startswith("ppt/charts/chart") and name.endswith(".xml")]
    excel_files = [name for name in names if "embeddings/Microsoft_Excel" in name]

    combined_xml = "\n".join(
        archive.read(name).decode("utf-8", errors="ignore") for name in slide_files
    )
    upper_xml = combined_xml.upper()

    text_markers = {
        "analysis_interrupted": upper_xml.count("ANALYSIS INTERRUPTED"),
        "development_in_progress": upper_xml.count("DEVELOPMENT IN PROGRESS: NATIVE"),
        "preview_unavailable": upper_xml.count("PREVIEW UNAVAILABLE"),
        "executive_summary": upper_xml.count("EXECUTIVE SUMMARY"),
        "critical_findings": upper_xml.count("CRITICAL FINDINGS"),
        "strategic_intelligence": upper_xml.count("STRATEGIC INTELLIGENCE"),
        "execution_playbook": upper_xml.count("EXECUTION PLAYBOOK"),
        "insight": upper_xml.count("INSIGHT:"),
        "competitive_swot": upper_xml.count("COMPETITIVE SWOT"),
        "recommendations_4p": upper_xml.count("4P RECOMMENDATIONS"),
        "market_archetype": upper_xml.count("MARKET ARCHETYPE"),
        "strategic_positioning": upper_xml.count("STRATEGIC POSITIONING"),
    }

    slide_summaries: List[Dict[str, Any]] = []
    unsupported_placeholder_count = 0
    error_placeholder_count = 0
    overview_slide_count = 0

    for slide_index, slide_file in enumerate(slide_files, start=1):
        slide_xml = archive.read(slide_file).decode("utf-8", errors="ignore")
        upper_slide_xml = slide_xml.upper()
        has_unsupported = any(marker in upper_slide_xml for marker in UNSUPPORTED_MARKERS)
        has_error = any(marker in upper_slide_xml for marker in ERROR_MARKERS)
        if has_unsupported:
            unsupported_placeholder_count += 1
        if has_error:
            error_placeholder_count += 1
        if "SURVEY OVERVIEW & METHODOLOGY" in upper_slide_xml:
            overview_slide_count += 1

        slide_summaries.append(
            {
                "slide_index": slide_index,
                "has_unsupported_placeholder": has_unsupported,
                "has_error_placeholder": has_error,
            }
        )

    return {
        "slide_count": len(slide_files),
        "chart_part_count": len(chart_files),
        "embedded_excel_count": len(excel_files),
        "notes_count": len(notes_files),
        "text_markers": text_markers,
        "overview_slide_count": overview_slide_count,
        "unsupported_placeholder_count": unsupported_placeholder_count,
        "error_placeholder_count": error_placeholder_count,
        "slide_summaries": slide_summaries,
    }


def _audit_presentation(presentation: Presentation) -> Dict[str, Any]:
    slide_titles: List[str] = []
    chart_shape_count = 0
    picture_shape_count = 0
    slide_summaries: List[Dict[str, Any]] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title = title_shape.text.strip() if title_shape is not None and getattr(title_shape, "text", "") else ""
        if title:
            slide_titles.append(title)

        slide_chart_count = 0
        slide_picture_count = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_shape_count += 1
                slide_picture_count += 1
                continue
            try:
                _ = shape.chart
                chart_shape_count += 1
                slide_chart_count += 1
            except Exception:
                pass

        slide_summaries.append(
            {
                "slide_index": slide_index,
                "title": title,
                "chart_count": slide_chart_count,
                "picture_count": slide_picture_count,
                "shape_count": len(slide.shapes),
            }
        )

    geometry = audit_presentation_geometry(presentation)
    duplicate_titles = [
        {"title": title, "count": count}
        for title, count in Counter(title.strip().upper() for title in slide_titles if title.strip()).items()
        if count > 1
    ]

    return {
        "presentation_readable": True,
        "canvas_width_emu": int(presentation.slide_width),
        "canvas_height_emu": int(presentation.slide_height),
        "slide_titles": slide_titles,
        "chart_shape_count": chart_shape_count,
        "picture_shape_count": picture_shape_count,
        "layout_warning_count": geometry["layout_warning_count"],
        "layout_warnings": geometry["layout_warnings"],
        "geometry_issues": geometry["issues"],
        "duplicate_title_count": geometry["duplicate_title_count"],
        "duplicate_titles": duplicate_titles,
        "out_of_bounds_shape_count": geometry["layout_warning_count"],
        "slide_summaries": slide_summaries,
    }


def _build_compact_audit(
    package_audit: Dict[str, Any],
    presentation_audit: Dict[str, Any],
) -> Dict[str, Any]:
    slide_summaries = _merge_slide_summaries(
        package_audit.get("slide_summaries", []),
        presentation_audit.get("slide_summaries", []),
    )

    return {
        "slide_count": package_audit["slide_count"],
        "slide_titles": presentation_audit.get("slide_titles", []),
        "chart_part_count": package_audit["chart_part_count"],
        "chart_shape_count": presentation_audit.get("chart_shape_count", 0),
        "picture_shape_count": presentation_audit.get("picture_shape_count", 0),
        "embedded_excel_count": package_audit["embedded_excel_count"],
        "notes_count": package_audit["notes_count"],
        "overview_slide_count": package_audit["overview_slide_count"],
        "text_markers": package_audit["text_markers"],
        "unsupported_placeholder_count": package_audit["unsupported_placeholder_count"],
        "error_placeholder_count": package_audit["error_placeholder_count"],
        "duplicate_title_count": presentation_audit.get("duplicate_title_count", 0),
        "duplicate_titles": presentation_audit.get("duplicate_titles", []),
        "layout_warning_count": presentation_audit.get("layout_warning_count", 0),
        "layout_warnings": presentation_audit.get("layout_warnings", []),
        "geometry_issues": presentation_audit.get("geometry_issues", []),
        "out_of_bounds_shape_count": presentation_audit.get("out_of_bounds_shape_count", 0),
        "presentation_readable": presentation_audit.get("presentation_readable", False),
        "canvas_width_emu": presentation_audit.get("canvas_width_emu"),
        "canvas_height_emu": presentation_audit.get("canvas_height_emu"),
        "slide_summaries": slide_summaries,
    }


def compact_export_audit(audit: Dict[str, Any]) -> Dict[str, Any]:
    """Manifest-friendly audit payload without full geometry issue objects."""
    return {
        "slide_count": audit.get("slide_count", 0),
        "chart_part_count": audit.get("chart_part_count", 0),
        "chart_shape_count": audit.get("chart_shape_count", 0),
        "picture_shape_count": audit.get("picture_shape_count", 0),
        "embedded_excel_count": audit.get("embedded_excel_count", 0),
        "notes_count": audit.get("notes_count", 0),
        "overview_slide_count": audit.get("overview_slide_count", 0),
        "text_markers": audit.get("text_markers", {}),
        "unsupported_placeholder_count": audit.get("unsupported_placeholder_count", 0),
        "error_placeholder_count": audit.get("error_placeholder_count", 0),
        "duplicate_title_count": audit.get("duplicate_title_count", 0),
        "duplicate_titles": audit.get("duplicate_titles", []),
        "layout_warning_count": audit.get("layout_warning_count", 0),
        "out_of_bounds_shape_count": audit.get("out_of_bounds_shape_count", 0),
        "presentation_readable": audit.get("presentation_readable", False),
        "canvas_width_emu": audit.get("canvas_width_emu"),
        "canvas_height_emu": audit.get("canvas_height_emu"),
        "slide_titles": audit.get("slide_titles", []),
        "slide_summaries": [
            {
                "slide_index": summary.get("slide_index"),
                "title": summary.get("title", ""),
                "chart_count": summary.get("chart_count", 0),
                "picture_count": summary.get("picture_count", 0),
                "has_unsupported_placeholder": summary.get("has_unsupported_placeholder", False),
                "has_error_placeholder": summary.get("has_error_placeholder", False),
            }
            for summary in audit.get("slide_summaries", [])
        ],
    }


def _merge_slide_summaries(
    package_summaries: List[Dict[str, Any]],
    presentation_summaries: List[Dict[str, Any]],
) -> List[Dict[str, Any]]:
    presentation_by_index = {
        summary["slide_index"]: summary for summary in presentation_summaries
    }
    merged: List[Dict[str, Any]] = []

    for package_summary in package_summaries:
        slide_index = package_summary["slide_index"]
        merged_summary = {**package_summary, **presentation_by_index.get(slide_index, {})}
        merged.append(merged_summary)

    return merged
