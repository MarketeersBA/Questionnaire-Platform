import logging
import random
from typing import Dict, Any, List
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from .base_builder import BaseChartBuilder
from .builder_render_status import BuilderEmptyDataError

logger = logging.getLogger(__name__)

class PPTXWordcloud(BaseChartBuilder):
    """
    Builder for Word Clouds / Semantic Maps.
    Features:
    - Simulated 'Tag Cloud' layout
    - Frequency-based font scaling (10pt to 44pt)
    - Randomized professional color selection from brand palette
    - Summary Table fallback for high-density word lists
    """

    def render(self, slide: Slide, chart_data: Dict[str, Any]) -> None:
        data_payload = chart_data.get("data", {})
        words = data_payload.get("words", []) # List of {text, value}
        
        if not words:
            raise BuilderEmptyDataError("Wordcloud requires non-empty words payload.")

        normalized_words = []
        for item in words:
            if not isinstance(item, dict):
                continue
            normalized_words.append(
                {
                    "text": item.get("text") or item.get("term") or item.get("label") or "",
                    "value": float(item.get("value", item.get("weight", item.get("count", 0))) or 0),
                }
            )
        sorted_words = sorted(normalized_words, key=lambda x: x.get("value", 0), reverse=True)
        if not sorted_words:
            raise BuilderEmptyDataError("Wordcloud payload does not contain usable word entries.")
        
        # 2. Choice of Layout: Table (High Density) or Cloud (Low/Mid Density)
        if len(sorted_words) > 40:
            self._render_summary_table(slide, sorted_words)
        else:
            self._render_simulated_cloud(slide, sorted_words)

    def _render_simulated_cloud(self, slide: Slide, words: List[dict]):
        """Creates a visual tag cloud by distributing text boxes across the chart area."""
        max_val = words[0].get("value", 1)
        min_val = words[-1].get("value", 1)
        val_range = max_val - min_val if max_val != min_val else 1
        
        # Grid boundaries
        center_x = self.layout.CHART_LEFT + (self.layout.CHART_WIDTH / 2)
        center_y = self.layout.CHART_TOP + (self.layout.CHART_HEIGHT / 2)
        
        # Simple spiral/scatter logic for non-overlapping approximation
        # (In a real implementation we'd need collision detection, here we use a deterministic jittered grid)
        cols = 5
        rows = 4
        cell_w = self.layout.CHART_WIDTH / cols
        cell_h = self.layout.CHART_HEIGHT / rows
        
        for i, word_data in enumerate(words[:20]): # Limit cloud to top 20 for legibility
            r = i // cols
            c = i % cols
            
            # Base position
            left = self.layout.CHART_LEFT + (c * cell_w) + (cell_w / 4)
            top = self.layout.CHART_TOP + (r * cell_h) + (cell_h / 4)
            
            # Map frequency to font size (recalibrated for 20" scale)
            rel_score = (word_data.get("value", 0) - min_val) / val_range
            font_size = max(10.0, min(28.0, 10 + (rel_score * 18)))
            
            # Random professional color from theme
            color = self.theme.get_color(i % 6)
            
            self._add_word_box(slide, left, top, word_data['text'], font_size, color)

    def _add_word_box(self, slide: Slide, x: float, y: float, text: str, size: float, color: RGBColor):
        # Auto-size approximation
        box_w = Inches(len(text) * (size / 50)) 
        box_h = Pt(size * 1.5)
        
        box = slide.shapes.add_textbox(x, y, box_w, box_h)
        p = box.text_frame.paragraphs[0]
        p.text = text.upper()
        p.font.size = Pt(size)
        p.font.name = self.theme.FONT_BOLD if size > 24 else self.theme.FONT_MEDIUM
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

    def _render_summary_table(self, slide: Slide, words: List[dict]):
        """Standard high-density fallback."""
        num_rows = 11
        num_cols = 4 # 2 sets of [Word, Freq]
        
        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            self.layout.CHART_LEFT, self.layout.CHART_TOP,
            self.layout.CHART_WIDTH, self.layout.CHART_HEIGHT
        )
        table = table_shape.table
        
        # Headers
        for i in [0, 2]:
            for j, h in enumerate(["KEYWORD", "FREQ"]):
                cell = table.cell(0, i + j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.theme.PRIMARY_BRAND
                p = cell.text_frame.paragraphs[0]
                p.text = h
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.size = Pt(10)
                p.font.name = self.theme.FONT_BOLD
        
        # Populate
        for i in range(num_rows - 1):
            # Left half
            if i < len(words):
                table.cell(i+1, 0).text = words[i]['text'].upper()
                table.cell(i+1, 1).text = f"{words[i]['value']}"
            # Right half
            idx_right = i + (num_rows - 1)
            if idx_right < len(words):
                table.cell(i+1, 2).text = words[idx_right]['text'].upper()
                table.cell(i+1, 3).text = f"{words[idx_right]['value']}"
