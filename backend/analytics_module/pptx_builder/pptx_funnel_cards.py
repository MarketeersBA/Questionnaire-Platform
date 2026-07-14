import logging
from typing import Dict, Any, List
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from .base_builder import BaseChartBuilder
from .builder_render_status import BuilderEmptyDataError

logger = logging.getLogger(__name__)

class PPTXFunnelCards(BaseChartBuilder):
    """
    Builder for KPI Card grids (Funnel Ratios / Brand Health).
    Features:
    - Multi-shape composition (no native chart)
    - Grid of 4 rounded KPI cards
    - Comparative brand listing within cards
    - Color-coded value highlighting
    """

    def render(self, slide: Slide, chart_data: Dict[str, Any]) -> None:
        """
        Consumes the 'brand_cards' payload structure from ReportAggregator.
        Expected shape: data['brand_cards'] -> List of dicts with 'brand', 'stage_bars', 'ratio_labels'.
        Note: The planner chunks this into 2 cards per slide.
        """
        data_payload = chart_data.get("data", {})
        brand_cards = data_payload.get("brand_cards", [])
        
        if not brand_cards:
            brand_cards = data_payload.get("ratios", [])
        if not brand_cards:
            raise BuilderEmptyDataError("No funnel ratio cards to render.")

        # 1. Layout Configuration (Side-by-Side 1x2)
        spacing = Inches(0.4)
        card_w = (self.layout.CHART_WIDTH - spacing) / 2
        card_h = self.layout.CHART_HEIGHT * 0.8  # Executive premium height
        
        # Vertically center in the chart area
        top = self.layout.CHART_TOP + (self.layout.CHART_HEIGHT - card_h) / 2
        
        # Render exactly up to 2 cards
        for i, card_data in enumerate(brand_cards[:2]):
            left = self.layout.CHART_LEFT + (i * (card_w + spacing))
            self._draw_kpi_card(slide, left, top, card_w, card_h, card_data)

    def _draw_kpi_card(self, slide: Slide, left: float, top: float, width: float, height: float, card_data: dict):
        """Draws a single rounded card with brand conversion ratios."""
        brand_name = card_data.get("brand", "Brand")
        ratios = card_data.get("ratio_labels", [])
        
        # 1. Card Background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = self.theme.get_rgb_by_name("brand_light_gray")
        card.line.width = Pt(1)
        
        # 2. Card Title (Brand Name)
        title_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.4))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = brand_name.upper()
        tp.font.name = self.theme.FONT_BOLD
        tp.font.size = Pt(11)
        tp.font.color.rgb = self.theme.PRIMARY_BRAND
        tp.alignment = PP_ALIGN.CENTER
        
        # 3. Ratio List
        start_y = top + Inches(0.5)
        count = len(ratios)
        if count == 0: return
        
        line_h = (height - Inches(0.7)) / count
        
        for j, ratio in enumerate(ratios):
            y_pos = start_y + (j * line_h)
            label = ratio.get("label", "")
            value_text = ratio.get("text", "0%")
            
            # Label
            b_box = slide.shapes.add_textbox(left + Inches(0.2), y_pos, width * 0.6, line_h)
            bp = b_box.text_frame.paragraphs[0]
            bp.text = label
            bp.font.name = self.theme.FONT_MEDIUM
            bp.font.size = Pt(9)
            bp.font.color.rgb = self.theme.SUBTITLE_COLOR
            
            # Value
            v_box = slide.shapes.add_textbox(left + (width * 0.6), y_pos, width * 0.3, line_h)
            vp = v_box.text_frame.paragraphs[0]
            vp.text = value_text
            vp.font.name = self.theme.FONT_BOLD
            vp.font.size = Pt(10)
            vp.font.color.rgb = self.theme.PRIMARY_BRAND
            vp.alignment = PP_ALIGN.RIGHT
            
            # Subtle Underline
            if j < count - 1:
                line = slide.shapes.add_connector(
                    1, 
                    left + Inches(0.2), y_pos + line_h - Pt(2),
                    left + width - Inches(0.2), y_pos + line_h - Pt(2)
                )
                line.line.color.rgb = self.theme.get_rgb_by_name("brand_light_gray")
                line.line.width = Pt(0.25)
