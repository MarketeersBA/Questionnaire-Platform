import logging
from datetime import datetime
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from .theme import PPTXTheme
from .layout import PPTXLayout

logger = logging.getLogger(__name__)

class PPTXSurveyInfo:
    """
    Premium Slide Builder for the Survey Overview.
    Uses a 5-card grid layout to visualize research setup, brands, and data health.
    """

    def __init__(self, theme: PPTXTheme, layout: PPTXLayout):
        self.theme = theme
        self.layout = layout

    def build(self, slide, metadata: dict):
        """Orchestrates the creation of the Survey Info slide."""
        # 1. Background & Header
        self._add_header_and_bg(slide, metadata)

        # 2. Card Grid Geometry (Recalibrated for 20" Canvas)
        # Top Row (3 cards)
        row1_y = Inches(2.7)
        col_w = Inches(5.7)
        spacing = Inches(0.6)
        
        c1_x = self.layout.CHART_LEFT
        c2_x = c1_x + col_w + spacing
        c3_x = c2_x + col_w + spacing
        
        # Bottom Row (2 cards)
        row2_y = Inches(6.3)
        row2_w = (self.layout.CHART_WIDTH - spacing) / 2
        
        c4_x = self.layout.CHART_LEFT
        c5_x = c4_x + row2_w + spacing

        # 3. Render Cards (Heights scaled 1.5x)
        self._render_brands_card(slide, c1_x, row1_y, col_w, Inches(3.3), metadata)
        self._render_pf_card(slide, c2_x, row1_y, col_w, Inches(3.3), metadata)
        self._render_target_card(slide, c3_x, row1_y, col_w, Inches(3.3), metadata)
        
        self._render_timing_card(slide, c4_x, row2_y, row2_w, Inches(3.9), metadata)
        self._render_data_card(slide, c5_x, row2_y, row2_w, Inches(3.9), metadata)

    def _add_header_and_bg(self, slide, metadata: dict):
        """Adds slide title and decorative branding."""
        # Title
        title_box = slide.shapes.add_textbox(
            self.layout.TITLE_LEFT, self.layout.TITLE_TOP, 
            self.layout.TITLE_WIDTH, self.layout.TITLE_HEIGHT
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = "SURVEY OVERVIEW & METHODOLOGY"
        p.font.name = self.theme.FONT_BOLD
        p.font.size = self.theme.TITLE_SIZE
        p.font.color.rgb = self.theme.PRIMARY_BRAND
        p.font.letter_spacing = 1.0

        # Subtitle (Client Name)
        sub_box = slide.shapes.add_textbox(
            self.layout.SUBTITLE_LEFT, self.layout.SUBTITLE_TOP,
            self.layout.SUBTITLE_WIDTH, self.layout.SUBTITLE_HEIGHT
        )
        p = sub_box.text_frame.paragraphs[0]
        p.text = f"STUDY CONDUCTED FOR: {metadata.get('company_name', 'N/A').upper()}"
        p.font.name = self.theme.FONT_MEDIUM
        p.font.size = self.theme.SUBTITLE_SIZE
        p.font.color.rgb = self.theme.get_rgb_by_name("brand_slate")

    def _render_card_base(self, slide, x, y, w, h, title):
        """Draws a premium 3D-simulated card with shadows and high-contrast headers."""
        # 1. Shadow Simulation (Offset subtle rectangle)
        shadow_offset = Inches(0.04)
        shadow = slide.shapes.add_shape(1, x + shadow_offset, y + shadow_offset, w, h)
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = self.theme.get_rgb_by_name("brand_light_gray") if hasattr(self.theme, "get_rgb_by_name") else RGBColor(230, 230, 230)
        shadow.line.visible = False
        
        # 2. Main Card Body
        shape = slide.shapes.add_shape(1, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.width = self.theme.BORDER_WIDTH
        shape.line.color.rgb = self.theme.get_rgb_by_name("brand_light_gray") if hasattr(self.theme, "get_rgb_by_name") else RGBColor(200, 200, 200)
        
        # 3. High-Contrast Header bar (Navy)
        header_h = Inches(0.5)
        header = slide.shapes.add_shape(1, x, y, w, header_h)
        header.fill.solid()
        header.fill.fore_color.rgb = self.theme.get_rgb_by_name("brand_navy")
        header.line.visible = False
        
        # 4. Card Title (Big & Clear with shadow-ready contrast)
        title_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.08), w - Inches(0.3), header_h)
        p = title_box.text_frame.paragraphs[0]
        p.text = title.upper()
        p.font.name = self.theme.FONT_BOLD
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.letter_spacing = 1.0

    def _render_brands_card(self, slide, x, y, w, h, metadata):
        self._render_card_base(slide, x, y, w, h, "Participating Brands")
        
        brands = metadata.get("brands", [])
        if not brands: brands = ["N/A"]
            
        target = metadata.get("target_brand")
        
        # Smart Column & Multi-Font Adjustment (Phase 4 Final - Total Visibility)
        use_columns = len(brands) > 7
        col_w = (w - Inches(0.4)) / 2 if use_columns else (w - Inches(0.4))
        
        # Column 1
        tf1 = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.6), col_w, h - Inches(0.7)).text_frame
        tf1.word_wrap = True
        
        # Column 2 (Active only if high density)
        tf2 = None
        if use_columns:
            tf2 = slide.shapes.add_textbox(x + Inches(0.2) + col_w, y + Inches(0.6), col_w, h - Inches(0.7)).text_frame
            tf2.word_wrap = True

        half = (len(brands) + 1) // 2 if use_columns else len(brands)
        font_size = Pt(13)
        if len(brands) > 10: font_size = Pt(10)
        if len(brands) > 20: font_size = Pt(8) # Dynamic shrinkage for extreme lists
        
        for i, brand in enumerate(brands):
            cur_tf = tf1 if i < half else tf2
            p = cur_tf.paragraphs[0] if (i == 0 or i == half) else cur_tf.add_paragraph()
            
            is_target = str(brand).lower() == str(target).lower()
            p.text = f"• {brand}" + (" ★" if is_target else "")
            p.font.name = self.theme.FONT_BOLD if is_target else self.theme.FONT_MEDIUM
            p.font.size = font_size
            p.font.color.rgb = self.theme.PRIMARY_BRAND if is_target else self.theme.TEXT_COLOR
            p.space_after = Pt(2)

    def _render_pf_card(self, slide, x, y, w, h, metadata):
        self._render_card_base(slide, x, y, w, h, "Purchase Funnel Module")
        
        pf_brands = metadata.get("pf_brands", [])
        is_active = metadata.get("pf_active", False)
        
        if not is_active or not pf_brands:
            tf = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.6), w - Inches(0.4), h - Inches(0.7)).text_frame
            p = tf.paragraphs[0]
            p.text = "Module not enabled for this research."
            p.font.size = Pt(14)
            p.font.name = self.theme.FONT_LIGHT
            p.font.color.rgb = self.theme.get_rgb_by_name("brand_slate")
            return

        # Smart Column Logic for PF Module
        use_columns = len(pf_brands) > 7
        col_w = (w - Inches(0.4)) / 2 if use_columns else (w - Inches(0.4))
        
        tf1 = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.6), col_w, h - Inches(0.7)).text_frame
        tf2 = None
        if use_columns:
            tf2 = slide.shapes.add_textbox(x + Inches(0.2) + col_w, y + Inches(0.6), col_w, h - Inches(0.7)).text_frame
            
        half = (len(pf_brands) + 1) // 2 if use_columns else len(pf_brands)
        font_size = Pt(13)
        if len(pf_brands) > 10: font_size = Pt(10)

        for i, brand in enumerate(pf_brands):
            cur_tf = tf1 if i < half else tf2
            p = cur_tf.paragraphs[0] if (i == 0 or i == half) else cur_tf.add_paragraph()
            
            p.text = f"✓ {brand}"
            p.font.name = self.theme.FONT_BOLD
            p.font.size = font_size
            p.font.color.rgb = self.theme.get_rgb_by_name("accent_cyan") if hasattr(self.theme, "get_rgb_by_name") else RGBColor(0, 180, 210)
            p.space_after = Pt(2)

    def _render_target_card(self, slide, x, y, w, h, metadata):
        self._render_card_base(slide, x, y, w, h, "Focal Research Target")
        
        target = metadata.get("target_brand", "N/A")
        
        # Premium "Glow" Hero Box logic (Multi-layered for WOW factor)
        box_w = w * 0.85
        box_h = Inches(1.1)
        
        # Outer Glow Layer
        glow = slide.shapes.add_shape(
            1, x + (w - box_w)/2 + Inches(0.05), y + Inches(0.75) + Inches(0.05), box_w, box_h
        )
        glow.fill.solid()
        glow.fill.fore_color.rgb = self.theme.get_rgb_by_name("brand_3_blue")
        glow.fill.transparency = 0.5
        glow.line.visible = False
        
        # Main Branding Block
        hero_box = slide.shapes.add_shape(
            1, x + (w - box_w)/2, y + Inches(0.75), box_w, box_h
        )
        hero_box.fill.solid()
        hero_box.fill.fore_color.rgb = self.theme.get_rgb_by_name("brand_3_blue")
        hero_box.line.width = Pt(3)
        hero_box.line.color.rgb = self.theme.PRIMARY_BRAND
        
        p = hero_box.text_frame.paragraphs[0]
        p.text = str(target).upper()
        p.font.name = self.theme.FONT_BOLD
        p.font.size = Pt(24) # Increased for impact
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    def _render_timing_card(self, slide, x, y, w, h, metadata):
        self._render_card_base(slide, x, y, w, h, "Methodology & Lifecycle")
        
        start = metadata.get("survey_created_at")
        end = metadata.get("report_generated_at")
        duration = f"{(end - start).days} Days" if (isinstance(start, datetime) and isinstance(end, datetime)) else "N/A"
        
        content_box = slide.shapes.add_textbox(x + Inches(0.4), y + Inches(0.7), w - Inches(0.8), h - Inches(1.0))
        tf = content_box.text_frame
        
        # Display Duration as Hero Metric
        p = tf.paragraphs[0]
        p.text = "STUDY DURATION"
        p.font.name = self.theme.FONT_BOLD
        p.font.size = Pt(14)
        p.font.color.rgb = self.theme.get_rgb_by_name("brand_slate")
        
        p_val = tf.add_paragraph()
        p_val.text = duration.upper()
        p_val.font.name = self.theme.FONT_BOLD
        p_val.font.size = Pt(32) # Hero Size
        p_val.font.color.rgb = self.theme.PRIMARY_BRAND
        p_val.space_after = Pt(12)

        # Dates in sub-text
        p_dates = tf.add_paragraph()
        start_t = start.strftime("%d %b %Y") if isinstance(start, datetime) else "N/A"
        end_t = end.strftime("%d %b %Y") if isinstance(end, datetime) else "N/A"
        p_dates.text = f"Launch: {start_t}  |  Ready: {end_t}"
        p_dates.font.name = self.theme.FONT_MEDIUM
        p_dates.font.size = Pt(12)
        p_dates.font.color.rgb = self.theme.TEXT_COLOR

    def _render_data_card(self, slide, x, y, w, h, metadata):
        self._render_card_base(slide, x, y, w, h, "Data Health & Diagnostic Scope")
        
        content_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.65), w - Inches(0.6), h - Inches(0.8))
        tf = content_box.text_frame
        
        # 1. Respondent Stats
        p1 = tf.paragraphs[0]
        base_n = metadata.get('base_n', 0)
        total_r = metadata.get('total_responses', 0)
        p1.text = f"ANALYSIS N: {base_n}  |  TOTAL RESPONSES: {total_r}  |  TYPE: {str(metadata.get('research_type')).upper()}"
        p1.font.name = self.theme.FONT_BOLD
        p1.font.size = Pt(14)
        p1.font.color.rgb = self.theme.PRIMARY_BRAND
        p1.space_after = Pt(15)
        
        # 2. Attributes Summary (Auto-Wrapping High-Capacity Block)
        p2 = tf.add_paragraph()
        p2.text = f"DIAGNOSTIC PILLARS: {metadata.get('total_attributes', 0)} ATTRIBUTES TESTED"
        p2.font.name = self.theme.FONT_BOLD
        p2.font.size = Pt(14)
        p2.font.color.rgb = self.theme.get_rgb_by_name("brand_navy")
        p2.space_before = Pt(8)
        
        # Research Pillars Area
        p3 = tf.add_paragraph()
        p3.text = "RESEARCH PILLARS:"
        p3.font.name = self.theme.FONT_BOLD
        p3.font.size = Pt(12)
        p3.font.color.rgb = self.theme.get_rgb_by_name("brand_slate")
        p3.space_after = Pt(2)
        
        # Multi-line Wrapping Bullet Cluster
        cats = metadata.get("attribute_categories", [])
        if cats:
            # We use a single wrapping paragraph for all segments to ensure flow
            p_cat = tf.add_paragraph()
            p_cat.text = " • ".join(cats) 
            p_cat.font.name = self.theme.FONT_MEDIUM
            p_cat.font.size = Pt(11)
            p_cat.font.color.rgb = self.theme.TEXT_COLOR
            p_cat.space_after = Pt(0)
            # The TextFrame already has word_wrap = True at the engine level
