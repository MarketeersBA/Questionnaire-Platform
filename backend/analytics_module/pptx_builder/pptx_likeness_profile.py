import logging
from typing import Dict, Any, List
from pptx.slide import Slide
from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK, XL_MARKER_STYLE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from .base_builder import BaseChartBuilder
from .builder_render_status import BuilderEmptyDataError
from .chart_builder_runtime import ensure_chart_legend, set_marker_style

logger = logging.getLogger(__name__)

class PPTXLikenessProfile(BaseChartBuilder):
    """
    Builder for the 'Likeness Profile' (Semantic Differential) chart.
    Features:
    - Vertical Snake Line using XY Scatter
    - Dual-sided descriptor labels (Left/Right)
    - Customized 1-5 rating scale
    - Precision alignment between chart points and text descriptions
    """

    def render(self, slide: Slide, chart_data: Dict[str, Any]) -> None:
        data_payload = chart_data.get("data", {})
        metrics = data_payload.get("metrics", [])
        datasets = data_payload.get("datasets", [])
        labels_left = data_payload.get("labels_left", [])
        labels_right = data_payload.get("labels_right", [])
        
        if not metrics or not datasets:
            raise BuilderEmptyDataError("Missing data for likeness profile chart.")

        # 1. Prepare XY Chart Data
        # We map metrics to Y-axis integers (highest index at top)
        num_metrics = len(metrics)
        chart_data_obj = XyChartData()
        
        for ds in datasets:
            series_label = ds.get("label", "")
            data_points = ds.get("data", [])
            
            series = chart_data_obj.add_series(series_label)
            for i, score in enumerate(data_points):
                # X = score (1-5), Y = reversed index for top-down rendering
                y_val = num_metrics - i
                series.add_data_point(score, y_val)

        # 2. Define Chart Frame (Centered between labels)
        chart_width = self.layout.CHART_WIDTH * 0.6
        chart_left = self.layout.CHART_LEFT + (self.layout.CHART_WIDTH - chart_width) / 2
        
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER_LINES,
            chart_left, self.layout.CHART_TOP,
            chart_width, self.layout.CHART_HEIGHT,
            chart_data_obj
        )
        chart = chart_shape.chart
        
        # 3. Apply Professional Scaling
        self._style_axes(chart, num_metrics)
        self._style_series(chart, datasets)
        self._style_legend(chart)
        
        # 4. Add Behavioral Descriptors (Semantic Differential Labels)
        self._add_semantic_labels(slide, chart_left, chart_width, metrics, labels_left, labels_right)

    def _style_axes(self, chart: Any, num_metrics: int):
        """Styles the XY axes for a clean semantic looks."""
        # X-Axis (The 1-5 Score)
        x_axis = chart.value_axis # In XY Scatter, both are value axes
        x_axis.minimum_scale = 1
        x_axis.maximum_scale = 5
        x_axis.major_unit = 1
        x_axis.has_major_gridlines = True
        x_axis.major_gridlines.format.line.color.rgb = self.theme.get_rgb_by_name("brand_light_gray")
        
        x_axis.tick_labels.font.size = Pt(10)
        x_axis.tick_labels.font.name = self.theme.FONT_MEDIUM
        
        # Y-Axis (Hidden - used for alignment only)
        y_axis = chart.category_axis
        y_axis.minimum_scale = 0.5
        y_axis.maximum_scale = num_metrics + 0.5
        y_axis.visible = False # We don't want to see the indices

    def _style_series(self, chart: Any, datasets: list):
        """Applies snake line formatting with markers."""
        from pptx.enum.dml import MSO_LINE
        
        for i, series in enumerate(chart.series):
            ds_meta = datasets[i]
            is_benchmark = ds_meta.get("brand") == "Overall" or ds_meta.get("is_benchmark", False)
            
            line = series.format.line
            if is_benchmark:
                line.color.rgb = self.theme.SUBTITLE_COLOR
                line.dash_style = MSO_LINE.DASH
                line.width = Pt(1.5)
                set_marker_style(series.marker, XL_MARKER_STYLE.NONE)
            else:
                color = self._apply_brand_colors(ds_meta.get("brand", "Brand"))
                line.color.rgb = color
                line.width = Pt(2.0)
                set_marker_style(series.marker, XL_MARKER_STYLE.CIRCLE)
                series.marker.size = 6
                series.marker.format.fill.solid()
                series.marker.format.fill.fore_color.rgb = color

    def _style_legend(self, chart: Any):
        legend = ensure_chart_legend(chart, XL_LEGEND_POSITION.BOTTOM)
        legend.font.size = Pt(9)
        legend.font.name = self.theme.FONT_MEDIUM

    def _add_semantic_labels(self, slide: Slide, chart_left: float, chart_width: float, metrics: list, lefts: list, rights: list):
        """Manually places text boxes on both sides of the chart area."""
        num_items = len(metrics)
        step_h = self.layout.CHART_HEIGHT / num_items
        
        # We need to account for chart padding/margins internal to the object
        # but for a first pass, simple linear distribution works.
        start_y = self.layout.CHART_TOP + (step_h / 2) - Inches(0.15)
        
        for i in range(num_items):
            curr_y = start_y + (i * step_h)
            
            # Left Label
            l_box = slide.shapes.add_textbox(self.layout.CHART_LEFT, curr_y, chart_left - self.layout.CHART_LEFT - Inches(0.1), Inches(0.3))
            self._format_label(l_box, lefts[i] if i < len(lefts) else "", PP_ALIGN.RIGHT)
            
            # Right Label
            r_start_x = chart_left + chart_width + Inches(0.1)
            r_box = slide.shapes.add_textbox(r_start_x, curr_y, self.layout.WIDTH - r_start_x - Inches(0.5), Inches(0.3))
            self._format_label(r_box, rights[i] if i < len(rights) else "", PP_ALIGN.LEFT)

    def _format_label(self, shape: Any, text: str, align: Any):
        p = shape.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(9)
        p.font.name = self.theme.FONT_MEDIUM
        p.font.color.rgb = self.theme.TEXT_COLOR
        p.alignment = align
