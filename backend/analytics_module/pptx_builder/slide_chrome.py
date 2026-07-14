from __future__ import annotations

import logging
from dataclasses import dataclass
from typing import Any, Dict, List, Optional

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from .layout import PPTXLayout
from .theme import PPTXTheme

logger = logging.getLogger(__name__)

TITLE_PLACEHOLDER_IDX = 0
SUBTITLE_PLACEHOLDER_IDX = 1
AI_HEADLINE_PLACEHOLDER_IDX = 13
MANAGED_PLACEHOLDER_IDXS = {
    TITLE_PLACEHOLDER_IDX,
    SUBTITLE_PLACEHOLDER_IDX,
    AI_HEADLINE_PLACEHOLDER_IDX,
}


@dataclass(frozen=True)
class ContentSlideChromeSpec:
    title: str
    subtitle: Optional[str] = None
    ai_headline: Optional[str] = None
    footnote: Optional[str] = None
    insight: Optional[str] = None
    ai_deep_analysis: Any = None


def format_ai_deep_analysis(ai_deep_analysis: Any) -> str:
    """Normalize AI deep analysis payloads into presenter notes."""
    if isinstance(ai_deep_analysis, str):
        return ai_deep_analysis.strip()

    if not isinstance(ai_deep_analysis, list):
        return str(ai_deep_analysis).strip()

    points: List[str] = []
    for item in ai_deep_analysis:
        if isinstance(item, dict):
            title = item.get("title") or item.get("label")
            body = item.get("body") or item.get("point") or item.get("text") or item.get("insight")
            sentiment = item.get("sentiment")

            if title and body:
                line = f"• {title}: {body}"
            elif body:
                line = f"• {body}"
            elif title:
                line = f"• {title}"
            else:
                line = f"• {item}"

            if sentiment:
                line = f"{line} ({sentiment})"
            points.append(line)
        elif item:
            points.append(f"• {item}")

    return "\n".join(points)


class ChartFrameRenderer:
    """Shared chart-area card frame used by engine chrome and standalone builders."""

    @staticmethod
    def add(slide: Slide, theme: PPTXTheme, layout: PPTXLayout) -> None:
        left, top, width, height = layout.chart_frame_bounds()
        
        # 1. Safe Shadow Simulation (Backend Render Safe)
        shadow_offset = Inches(0.04)
        shadow_frame = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left + shadow_offset,
            top + shadow_offset,
            width,
            height,
        )
        shadow_frame.fill.solid()
        
        # Attempt to use theme subtles, fallback to light gray
        try:
            shadow_color = theme.get_rgb_by_name("brand_light_gray")
        except Exception:
            shadow_color = RGBColor(230, 230, 230)
            
        shadow_frame.fill.fore_color.rgb = shadow_color
        shadow_frame.line.visible = False

        # 2. Main Frame Shape
        frame = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left,
            top,
            width,
            height,
        )

        frame.fill.solid()
        frame.fill.fore_color.rgb = theme.get_rgb_by_name("white")
        frame.line.color.rgb = theme.BORDER_COLOR
        frame.line.width = theme.BORDER_WIDTH

        # 3. Z-Index Ordering
        # Native shape addition order (Chrome first, Builder next) naturally handles Z-index.
        # Manual _spTree manipulation is removed to ensure OXML schema validity across all desktop platforms.
        pass


class ContentSlideChromeApplier:
    """
    Owns template-facing chrome for analytical content slides.
    Builders should render chart content only after this runs.
    """

    def __init__(self, theme: PPTXTheme, layout: PPTXLayout):
        self.theme = theme
        self.layout = layout

    def specification_from_chart_data(self, chart_data: Dict[str, Any]) -> ContentSlideChromeSpec:
        data = chart_data.get("data") or {}
        footnote = chart_data.get("footnote")
        if not footnote and isinstance(data, dict):
            footnote = data.get("footnote")

        return ContentSlideChromeSpec(
            title=(chart_data.get("title") or "Insight Analysis").upper(),
            subtitle=chart_data.get("subtitle"),
            ai_headline=chart_data.get("ai_headline") or chart_data.get("insight_headline"),
            footnote=footnote,
            insight=chart_data.get("insight"),
            ai_deep_analysis=chart_data.get("ai_deep_analysis"),
        )

    def apply(self, slide: Slide, chart_data: Dict[str, Any]) -> ContentSlideChromeSpec:
        spec = self.specification_from_chart_data(chart_data)
        subtitle_applied = self._populate_placeholders(slide, spec)
        if spec.subtitle and not subtitle_applied:
            self._apply_subtitle_zone(slide, spec.subtitle)
        self._apply_on_slide_narrative(slide, spec)
        self._apply_presenter_notes(slide, spec.ai_deep_analysis)
        ChartFrameRenderer.add(slide, self.theme, self.layout)
        return spec

    def _populate_placeholders(self, slide: Slide, spec: ContentSlideChromeSpec) -> bool:
        subtitle_applied = False
        for ph in slide.placeholders:
            try:
                idx = ph.placeholder_format.idx
                if idx == TITLE_PLACEHOLDER_IDX:
                    ph.text = spec.title
                elif idx == SUBTITLE_PLACEHOLDER_IDX and spec.subtitle:
                    ph.text = spec.subtitle
                    subtitle_applied = True
                elif idx == AI_HEADLINE_PLACEHOLDER_IDX and spec.ai_headline:
                    ph.text = spec.ai_headline
                elif idx not in MANAGED_PLACEHOLDER_IDXS and ph.has_text_frame:
                    ph.text_frame.text = ""
            except Exception:
                logger.debug(
                    "Placeholder population skipped for idx=%s",
                    getattr(ph.placeholder_format, "idx", "?"),
                    exc_info=True,
                )
        return subtitle_applied

    def _apply_subtitle_zone(self, slide: Slide, subtitle: Optional[str]) -> None:
        if not subtitle:
            return

        shape = slide.shapes.add_textbox(
            self.layout.SUBTITLE_LEFT,
            self.layout.SUBTITLE_TOP,
            self.layout.SUBTITLE_WIDTH,
            self.layout.SUBTITLE_HEIGHT,
        )
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.text = subtitle
        paragraph.font.name = self.theme.FONT_LIGHT
        paragraph.font.size = self.theme.SUBTITLE_SIZE
        paragraph.font.color.rgb = self.theme.SUBTITLE_COLOR

    def _apply_on_slide_narrative(self, slide: Slide, spec: ContentSlideChromeSpec) -> None:
        cursor_top = self.layout.FOOTNOTE_TOP
        analysis_points = self.normalized_analysis_points(spec.ai_deep_analysis)
        has_deep_analysis = len(analysis_points) > 0

        if spec.insight and not has_deep_analysis:
            cursor_top = self._add_text_band(
                slide,
                f"INSIGHT: {spec.insight}",
                top=cursor_top,
                height=self.layout.FOOTNOTE_HEIGHT,
                italic=True,
            )
        elif has_deep_analysis:
            excerpt = self._analysis_excerpt(analysis_points, max_points=1)
            if excerpt:
                cursor_top = self._add_text_band(
                    slide,
                    excerpt,
                    top=cursor_top,
                    height=self.layout.FOOTNOTE_HEIGHT,
                    italic=False,
                )

        if spec.footnote:
            self._add_text_band(
                slide,
                spec.footnote,
                top=cursor_top,
                height=self.layout.FOOTNOTE_HEIGHT,
                italic=False,
            )

    def _apply_presenter_notes(self, slide: Slide, ai_deep_analysis: Any) -> None:
        if not self._has_deep_analysis(ai_deep_analysis):
            return

        formatted = format_ai_deep_analysis(ai_deep_analysis)
        if not formatted:
            return

        try:
            slide.notes_slide.notes_text_frame.text = formatted
        except Exception:
            logger.warning("Failed to add presenter notes for AI deep analysis.", exc_info=True)

    def _add_text_band(
        self,
        slide: Slide,
        text: str,
        *,
        top,
        height,
        italic: bool,
    ):
        shape = slide.shapes.add_textbox(
            self.layout.FOOTNOTE_LEFT,
            top,
            self.layout.FOOTNOTE_WIDTH,
            height,
        )
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.size = Pt(9)
        paragraph.font.name = self.theme.FONT_LIGHT
        paragraph.font.color.rgb = self.theme.get_rgb_by_name("brand_slate")
        paragraph.font.italic = italic
        return top + height

    def normalized_analysis_points(self, ai_deep_analysis: Any) -> List[Dict[str, str]]:
        if not isinstance(ai_deep_analysis, list):
            return []
        normalized: List[Dict[str, str]] = []
        for item in ai_deep_analysis:
            if not isinstance(item, dict):
                continue
            normalized.append(
                {
                    "title": str(item.get("title") or item.get("label") or "Insight").strip(),
                    "body": str(
                        item.get("body")
                        or item.get("point")
                        or item.get("text")
                        or item.get("insight")
                        or ""
                    ).strip(),
                    "sentiment": str(item.get("sentiment") or "neutral").strip().lower(),
                    "recommended_action": str(item.get("recommended_action") or "").strip(),
                }
            )
        return [point for point in normalized if point["title"] or point["body"]]

    def analysis_requires_followup_slide(self, ai_deep_analysis: Any) -> bool:
        points = self.normalized_analysis_points(ai_deep_analysis)
        if len(points) > 1:
            return True
        if not points:
            return False
        total_chars = sum(len(p["title"]) + len(p["body"]) + len(p["recommended_action"]) for p in points)
        return total_chars > 260

    @staticmethod
    def _analysis_excerpt(points: List[Dict[str, str]], *, max_points: int = 1) -> str:
        rows: List[str] = []
        for point in points[:max_points]:
            sentiment = point.get("sentiment", "neutral").upper()
            title = point.get("title", "Insight")
            body = point.get("body", "")
            action = point.get("recommended_action", "")
            line = f"AI DEEP ANALYSIS [{sentiment}] {title}: {body}".strip()
            if action:
                line = f"{line} | Action: {action}"
            rows.append(line)
        return "\n".join(rows).strip()

    @staticmethod
    def _has_deep_analysis(ai_deep_analysis: Any) -> bool:
        if isinstance(ai_deep_analysis, str):
            return bool(ai_deep_analysis.strip())
        if isinstance(ai_deep_analysis, list):
            return len(ai_deep_analysis) > 0
        return bool(ai_deep_analysis)
