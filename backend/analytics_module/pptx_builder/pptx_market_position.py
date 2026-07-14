from __future__ import annotations

import logging
from typing import Any, Dict, List

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .layout import PPTXLayout
from .market_position_sections import MarketPositionSection, MarketPositionSectionBuilder
from .narrative_pagination import chunk_sequence, split_text_blocks
from .template_adapter import TemplateAdapter
from .theme import PPTXTheme

logger = logging.getLogger(__name__)


class PPTXMarketPosition:
    """Render market_position_report as explicit, paginated narrative slides."""

    def __init__(self, theme: PPTXTheme, layout: PPTXLayout):
        self.theme = theme
        self.layout = layout

    def build_deck(self, presentation, template_adapter: TemplateAdapter, payload: Any) -> int:
        sections = MarketPositionSectionBuilder.from_payload(payload)
        if not sections:
            return 0

        rendered = 0
        for section in sections:
            pages = self._paginate_section(section)
            for page in pages:
                slide = presentation.slides.add_slide(
                    template_adapter.get_layout(presentation, "ai_narrative")
                )
                self._render_page(slide, page)
                rendered += 1
        return rendered

    def _paginate_section(self, section: MarketPositionSection) -> List[MarketPositionSection]:
        if section.kind == "archetype":
            return [section]

        pages: List[MarketPositionSection] = []
        if section.body_lines:
            for block in split_text_blocks("\n\n".join(section.body_lines)):
                pages.append(
                    MarketPositionSection(
                        kind=section.kind,
                        title=section.title,
                        body_lines=[block],
                        metadata=section.metadata,
                    )
                )

        if section.bullets:
            for index, chunk in enumerate(chunk_sequence(section.bullets, 3), start=1):
                title = section.title if index == 1 and not pages else f"{section.title} ({index})"
                pages.append(
                    MarketPositionSection(
                        kind=section.kind,
                        title=title,
                        bullets=chunk,
                        metadata=section.metadata,
                    )
                )

        return pages or [section]

    def _render_page(self, slide, section: MarketPositionSection) -> None:
        if slide.placeholders:
            slide.placeholders[0].text = section.title.upper()

        top = Inches(1.5)
        left = self.layout.CHART_LEFT
        width = self.layout.WIDTH - (left * 2)

        if section.kind == "archetype":
            position = str(section.metadata.get("market_position", "Follower"))
            confidence = str(section.metadata.get("position_confidence", "Medium"))
            headline = slide.shapes.add_textbox(left, top, width, Inches(1.0))
            headline.text_frame.paragraphs[0].text = position.upper()
            headline.text_frame.paragraphs[0].font.name = self.theme.FONT_BOLD
            headline.text_frame.paragraphs[0].font.size = Pt(36)
            headline.text_frame.paragraphs[0].font.color.rgb = self.theme.PRIMARY_BRAND

            confidence_box = slide.shapes.add_textbox(left, top + Inches(1.1), width, Inches(0.4))
            confidence_box.text_frame.paragraphs[0].text = f"{confidence.upper()} STABILITY"
            confidence_box.text_frame.paragraphs[0].font.size = Pt(12)
            confidence_box.text_frame.paragraphs[0].font.color.rgb = self.theme.SUBTITLE_COLOR

            body_top = top + Inches(1.7)
        else:
            body_top = top

        for line in section.body_lines:
            box = slide.shapes.add_textbox(left, body_top, width, Inches(1.2))
            tf = box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = line
            tf.paragraphs[0].font.name = self.theme.FONT_MEDIUM
            tf.paragraphs[0].font.size = Pt(16 if section.kind == "audience_profile" else 14)
            tf.paragraphs[0].font.color.rgb = self.theme.TEXT_COLOR
            body_top += Inches(1.0)

        if section.bullets:
            bullet_box = slide.shapes.add_textbox(left, body_top, width, Inches(3.5))
            tf = bullet_box.text_frame
            tf.word_wrap = True
            for index, bullet in enumerate(section.bullets):
                paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
                paragraph.text = f"• {bullet}"
                paragraph.font.name = self.theme.FONT_MEDIUM
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = self.theme.TEXT_COLOR
                paragraph.space_after = Pt(6)

        if section.kind == "audience_profile":
            profile = slide.shapes.add_textbox(left, top + Inches(0.4), width, Inches(0.4))
            profile.text_frame.paragraphs[0].text = "PRIMARY DEMOGRAPHIC & GEOGRAPHIC PROFILE"
            profile.text_frame.paragraphs[0].font.size = Pt(10)
            profile.text_frame.paragraphs[0].font.color.rgb = self.theme.SUBTITLE_COLOR

        if section.kind == "competitive_stance":
            label = slide.shapes.add_textbox(left, top + Inches(0.2), width, Inches(0.3))
            label.text_frame.paragraphs[0].text = "RELATIVE MARKET POSITIONING"
            label.text_frame.paragraphs[0].font.size = Pt(10)
            label.text_frame.paragraphs[0].font.color.rgb = self.theme.SUBTITLE_COLOR

        if section.kind == "strategic_implications":
            label = slide.shapes.add_textbox(left, top + Inches(0.2), width, Inches(0.3))
            label.text_frame.paragraphs[0].text = "ACTIONABLE POSITIONING INSIGHTS"
            label.text_frame.paragraphs[0].font.size = Pt(10)
            label.text_frame.paragraphs[0].font.color.rgb = self.theme.SUBTITLE_COLOR
