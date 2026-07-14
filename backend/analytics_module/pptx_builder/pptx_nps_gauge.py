import logging
from typing import Dict, Any
from pptx.slide import Slide
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from .base_builder import BaseChartBuilder
from .builder_render_status import BuilderEmptyDataError
from .chart_builder_runtime import ensure_chart_legend

logger = logging.getLogger(__name__)

class PPTXNPSGauge(BaseChartBuilder):
    """
    Builder for Net Promoter Score (NPS) gauges.
    Features:
    - Triple-segmented 100% stacked bar (Detractor-Passive-Promoter)
    - Industry-standard NPS color grading (Red-Grey-Green)
    - Dynamic NPS score callouts (text boxes) positioned on the high-end of bars
    - Brand-ranked vertical list
    """

    def render(self, slide: Slide, chart_data: Dict[str, Any]) -> None:
        data_payload = chart_data.get("data", {})
        labels = data_payload.get("labels", []) # Brand Names
        datasets = data_payload.get("datasets", []) # Detractors, Passives, Promoters
        nps_scores = data_payload.get("nps_scores", {}) # {brand: nps_val}
        
        if not labels or not datasets:
            raise BuilderEmptyDataError("Missing NPS gauge data.")

        # 1. Populate 100% Stacked Bar Data
        chart_data_obj = ChartData()
        chart_data_obj.categories = labels
        for ds in datasets:
            chart_data_obj.add_series(ds.get("label", ""), ds.get("data", []))

        # 2. Add Chart Shape (Restricted width to allow callouts)
        chart_width = self.layout.CHART_WIDTH * 0.8
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED_100,
            self.layout.CHART_LEFT, self.layout.CHART_TOP,
            chart_width, self.layout.CHART_HEIGHT,
            chart_data_obj
        )
        chart = chart_shape.chart
        
        # 3. Apply Professional NPS Styling
        self._style_axes(chart)
        self._style_nps_segments(chart)
        self._style_legend(chart)
        
        # 4. Add NPS Score Callouts
        self._add_score_callouts(slide, chart_width, labels, nps_scores)

    def _style_axes(self, chart: Any):
        """Minimal axes for cleaner gauge look."""
        c_axis = chart.category_axis
        c_axis.reverse_order = True
        c_axis.tick_labels.font.name = self.theme.FONT_BOLD
        c_axis.tick_labels.font.size = Pt(11)
        
        v_axis = chart.value_axis
        v_axis.visible = False # The bar itself shows percentage distribution
        v_axis.has_major_gridlines = False

    def _style_nps_segments(self, chart: Any):
        """Applies Detractor (Red), Passive (Grey), Promoter (Green) colors."""
        palette = self.theme.get_nps_palette()
        
        for i, series in enumerate(chart.series):
            color = palette[i] if i < len(palette) else self.theme.PRIMARY_BRAND
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = color
            
            # Internal Data Labels
            self.apply_series_data_labels(
                series,
                position=XL_DATA_LABEL_POSITION.CENTER,
                number_format="0%",
                font_size_pt=9,
                font_color=RGBColor(255, 255, 255),
            )

    def _style_legend(self, chart: Any):
        legend = ensure_chart_legend(chart, XL_LEGEND_POSITION.BOTTOM)
        legend.font.name = self.theme.FONT_MEDIUM
        legend.font.size = Pt(10)

    def _add_score_callouts(self, slide: Slide, chart_width: float, brands: list, nps_scores: dict):
        """Adds large text boxes to the right of each gauge showing the final NPS."""
        num_brands = len(brands)
        step_y = self.layout.CHART_HEIGHT / num_brands
        
        # Adjust vertical alignment to match chart categories
        start_y = self.layout.CHART_TOP + (step_y / 2) - Inches(0.3)
        callout_x = self.layout.CHART_LEFT + chart_width + Inches(0.3)
        
        for i, brand in enumerate(brands):
            score = nps_scores.get(brand, 0)
            curr_y = start_y + (i * step_y)
            
            # Simple score callout box
            box = slide.shapes.add_textbox(callout_x, curr_y, Inches(1.0), Inches(0.4))
            p = box.text_frame.paragraphs[0]
            
            # Styling: Circular highlight or simple bold text
            p.text = f"{score:+.1f}"
            p.font.size = Pt(14)
            p.font.name = self.theme.FONT_BOLD
            # Color coding the final score
            if score >= 50: p.font.color.rgb = self.theme.get_branding_secondary()
            elif score >= 0: p.font.color.rgb = self.theme.TEXT_COLOR # Neutral
            else: p.font.color.rgb = self.theme.get_rgb_by_name("brand_crimson")
            
            p.alignment = PP_ALIGN.LEFT
            
            # Label "NPS" below
            sub_p = box.text_frame.add_paragraph()
            sub_p.text = "NPS"
            sub_p.font.size = Pt(8)
            sub_p.font.name = self.theme.FONT_LIGHT
            sub_p.font.color.rgb = self.theme.SUBTITLE_COLOR
