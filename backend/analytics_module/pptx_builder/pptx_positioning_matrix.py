import logging
from typing import Dict, Any, List
from pptx.slide import Slide
from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_MARKER_STYLE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from .base_builder import BaseChartBuilder
from .builder_render_status import BuilderEmptyDataError
from .chart_builder_runtime import ensure_chart_legend, set_marker_style

logger = logging.getLogger(__name__)

class PPTXPositioningMatrix(BaseChartBuilder):
    """
    Builder for Competitive Positioning Matrices.
    Features:
    - XY Scatter with variable marker sizes (Bubble Simulation)
    - Dynamic marker scaling (maps raw size to Pt dimensions)
    - Centered crosshairs for quadrant analysis
    - Precision point-level labeling
    """

    def render(self, slide: Slide, chart_data: Dict[str, Any]) -> None:
        data_payload = chart_data.get("data", {})
        datasets = data_payload.get("datasets", [])
        
        if not datasets:
            raise BuilderEmptyDataError("No datasets found for positioning matrix.")

        # 1. Prepare XY Data
        chart_data_obj = XyChartData()
        all_points = []
        
        # We group all points into a single series for easier cross-point formatting
        # or use multiple series if we want different colors per brand automatically.
        # Here we use brand-specific series for better legend/coloring.
        for ds in datasets:
            label = ds.get("label", "Point")
            p_data = ds.get("data", [])
            series = chart_data_obj.add_series(label)
            for point in p_data:
                x, y = point.get("x", 0), point.get("y", 0)
                series.add_data_point(x, y)
                all_points.append(point)

        if not all_points:
            raise BuilderEmptyDataError("Positioning matrix has no plotted points.")
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER,
            self.layout.CHART_LEFT, self.layout.CHART_TOP,
            self.layout.CHART_WIDTH, self.layout.CHART_HEIGHT,
            chart_data_obj
        )
        chart = chart_shape.chart
        
        # 3. Apply Premium Styling & Variable Sizing
        self._style_axes(chart)
        self._apply_bubble_sizing(chart, datasets)
        self._style_legend(chart)
        
        # 4. Add Crosshairs (Midpoint 5.0 assuming 0-10 scale)
        self._add_crosshairs(slide)

    def _style_axes(self, chart: Any):
        """Standardizes positioning axes (0-10 research scale)."""
        for axis in [chart.value_axis, chart.category_axis]:
            axis.minimum_scale = 0
            axis.maximum_scale = 10
            axis.has_major_gridlines = True
            axis.major_gridlines.format.line.color.rgb = self.theme.get_rgb_by_name("brand_light_gray")
            axis.tick_labels.font.size = Pt(9)
            axis.tick_labels.font.name = self.theme.FONT_LIGHT

    def _apply_bubble_sizing(self, chart: Any, datasets: list):
        """
        Simulates bubbles by modifying marker size per point.
        Maps raw data size (e.g. 5-20) to Pt size (10-50).
        """
        for i, series in enumerate(chart.series):
            brand_name = datasets[i].get("label", "Brand")
            color = self._apply_brand_colors(brand_name)
            
            set_marker_style(series.marker, XL_MARKER_STYLE.CIRCLE)
            # Marker-level fill for OXML compliance (not series.format.fill)
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = color
            series.marker.format.line.visible = False
            series.format.line.visible = False

            points_meta = datasets[i].get("data", [])
            for point_shape, point_meta in zip(series.points, points_meta):
                raw_size = point_meta.get("size", point_meta.get("n", 10))
                marker_size = 10 + (raw_size * 0.4)
                point_shape.marker.size = int(min(max(marker_size, 8), 70))

                label = point_meta.get("label") or point_meta.get("brand") or brand_name
                # Safe data label activation (Desktop PowerPoint compatible)
                try:
                    point_shape.has_data_label = True
                    dl = point_shape.data_label
                    if hasattr(dl, "text_frame"):
                        dl.text_frame.text = str(label)
                except Exception:
                    pass

    def _add_crosshairs(self, slide: Slide):
        """Draws the mid-point axes for quadrant classification."""
        mid_x = self.layout.CHART_LEFT + (self.layout.CHART_WIDTH / 2)
        mid_y = self.layout.CHART_TOP + (self.layout.CHART_HEIGHT / 2)
        
        for coords in [
            (int(mid_x), int(self.layout.CHART_TOP), int(mid_x), int(self.layout.CHART_TOP + self.layout.CHART_HEIGHT)),
            (int(self.layout.CHART_LEFT), int(mid_y), int(self.layout.CHART_LEFT + self.layout.CHART_WIDTH), int(mid_y))
        ]:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, *coords)
            line.line.color.rgb = self.theme.get_rgb_by_name("brand_slate")
            line.line.width = Pt(1.5)

    def _style_legend(self, chart: Any):
        legend = ensure_chart_legend(chart, XL_LEGEND_POSITION.BOTTOM)
        legend.font.size = Pt(9)
        legend.font.name = self.theme.FONT_MEDIUM
