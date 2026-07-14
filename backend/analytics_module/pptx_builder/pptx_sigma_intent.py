import logging
from typing import Any, Dict, List

from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_MARKER_STYLE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from .base_builder import BaseChartBuilder
from .builder_render_status import BuilderEmptyDataError
from .chart_builder_runtime import apply_xy_point_labels, set_marker_style

logger = logging.getLogger(__name__)


class PPTXSigmaIntent(BaseChartBuilder):
    """Builder for the Sigma/Intent strategy matrix."""

    def render(self, slide: Slide, chart_data: Dict[str, Any]) -> None:
        data_payload = chart_data.get("data", {})
        datasets = data_payload.get("datasets", [])

        if not datasets:
            raise BuilderEmptyDataError("No datasets found for sigma intent chart.")

        chart_data_obj = XyChartData()
        plotted_points = 0

        for index, dataset in enumerate(datasets):
            if isinstance(dataset, str):
                label = dataset
                points: List[Any] = []
            else:
                label = dataset.get("label", f"P{index}")
                points = dataset.get("data", [])

            series = chart_data_obj.add_series(label)
            for point in points:
                if isinstance(point, dict):
                    x, y = point.get("x", 0), point.get("y", 0)
                else:
                    x, y = 0, 0
                series.add_data_point(x, y)
                plotted_points += 1

        if plotted_points == 0:
            raise BuilderEmptyDataError("Sigma intent chart has no plotted points.")

        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER,
            self.layout.CHART_LEFT,
            self.layout.CHART_TOP,
            self.layout.CHART_WIDTH,
            self.layout.CHART_HEIGHT,
            chart_data_obj,
        )
        chart = chart_shape.chart

        self._style_axes(chart)
        self._style_series(chart)
        apply_xy_point_labels(chart, datasets)
        self._add_matrix_zones(slide)

    def _style_axes(self, chart: Any):
        for axis in [chart.value_axis, chart.category_axis]:
            axis.minimum_scale = 0
            axis.maximum_scale = 100
            axis.has_major_gridlines = False
            axis.tick_labels.font.size = Pt(9)
            axis.tick_labels.font.name = self.theme.FONT_LIGHT

    def _style_series(self, chart: Any):
        for i, series in enumerate(chart.series):
            set_marker_style(series.marker, XL_MARKER_STYLE.CIRCLE)
            series.marker.size = 10
            # Marker-level fill for OXML compliance
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = self.theme.get_color(i)
            series.marker.format.line.visible = False
            series.format.line.visible = False

    def _add_matrix_zones(self, slide: Slide):
        mid_x = self.layout.CHART_LEFT + (self.layout.CHART_WIDTH / 2)
        mid_y = self.layout.CHART_TOP + (self.layout.CHART_HEIGHT / 2)

        for line_coords in [
            (int(mid_x), int(self.layout.CHART_TOP), int(mid_x), int(self.layout.CHART_TOP + self.layout.CHART_HEIGHT)),
            (int(self.layout.CHART_LEFT), int(mid_y), int(self.layout.CHART_LEFT + self.layout.CHART_WIDTH), int(mid_y)),
        ]:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, *line_coords)
            line.line.color.rgb = self.theme.get_rgb_by_name("brand_slate")
            line.line.width = Pt(2)

        self._add_zone_tag(slide, mid_x, self.layout.CHART_TOP, "STARS", self.theme.get_rgb_by_name("brand_emerald"))
        self._add_zone_tag(slide, mid_x, mid_y, "CASH COWS", self.theme.get_rgb_by_name("brand_navy"))
        self._add_zone_tag(slide, self.layout.CHART_LEFT, self.layout.CHART_TOP, "QUESTIONS", self.theme.get_rgb_by_name("brand_gold"))
        self._add_zone_tag(slide, self.layout.CHART_LEFT, mid_y, "DOGS", self.theme.get_rgb_by_name("brand_crimson"))

    def _add_zone_tag(self, slide: Slide, x: float, y: float, text: str, color: RGBColor):
        box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.1), Inches(1.5), Inches(0.4))
        paragraph = box.text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.size = Pt(12)
        paragraph.font.name = self.theme.FONT_BOLD
        paragraph.font.color.rgb = color
        paragraph.alignment = PP_ALIGN.LEFT
