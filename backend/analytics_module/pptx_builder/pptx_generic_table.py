import logging
from typing import Dict, Any, List
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from .base_builder import BaseChartBuilder
from .builder_render_status import BuilderEmptyDataError

logger = logging.getLogger(__name__)

class PPTXGenericTable(BaseChartBuilder):
    """
    Fallback builder for data that doesn't match a specialized chart type.
    Renders any valid labels/datasets structure into a professional banded table.
    """

    def render(self, slide: Slide, chart_data: Dict[str, Any]) -> None:
        data_payload = chart_data.get("data", {})
        
        # Robust Mapping: Handle List-of-Dicts (Legacy/Simple) or Dict-of-Matrix (Advanced)
        if isinstance(data_payload, list):
            # Convert list [{"category": "X", "value": 1}] to labels/datasets
            labels = ["VALUE"]
            datasets = []
            for item in data_payload:
                datasets.append({
                    "label": str(item.get("category", "Item")),
                    "data": [item.get("value", 0)]
                })
        else:
            labels = data_payload.get("labels", [])
            datasets = data_payload.get("datasets", [])
        
        if not datasets:
            raise BuilderEmptyDataError("Generic fallback table requires at least one dataset.")

        # 1. Determine Table Structure
        # Is it a list of items or a matrix?
        # Usually datasets[i].data contains the values for each label
        num_cols = len(labels) + 1
        num_rows = len(datasets) + 1
        
        # 2. Add Table
        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            self.layout.CHART_LEFT, self.layout.CHART_TOP,
            self.layout.CHART_WIDTH, self.layout.CHART_HEIGHT
        )
        table = table_shape.table
        
        # 3. Headers
        headers = ["DESCRIPTION"] + [str(l).upper() for l in labels]
        self._set_header(table, headers)
        
        # 4. Data Rows
        for i, ds in enumerate(datasets):
            row_idx = i + 1
            row_label = ds.get("label", f"Item {i+1}")
            row_values = ds.get("data", [])
            
            # Label Cell
            self._fill_cell(table.cell(row_idx, 0), row_label, align=PP_ALIGN.LEFT, is_bold=True)
            
            # Value Cells
            for j, val in enumerate(row_values):
                if j + 1 >= num_cols: break
                cell = table.cell(row_idx, j + 1)
                text = f"{val:.1f}" if isinstance(val, (int, float)) else str(val)
                self._fill_cell(cell, text)
            
            # Banding
            if i % 2 == 1:
                self._apply_banding(table, row_idx)

    def _set_header(self, table: Any, headers: List[str]):
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme.PRIMARY_BRAND
            p = cell.text_frame.paragraphs[0]
            p.text = h
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(10)
            p.font.name = self.theme.FONT_BOLD
            p.alignment = PP_ALIGN.CENTER

    def _fill_cell(self, cell: Any, text: str, align: Any = PP_ALIGN.CENTER, is_bold: bool = False):
        p = cell.text_frame.paragraphs[0]
        p.text = str(text)
        p.font.name = self.theme.FONT_MEDIUM if is_bold else self.theme.FONT_LIGHT
        p.font.size = Pt(10)
        p.alignment = align

    def _apply_banding(self, table: Any, row_idx: int):
        for c in range(len(table.columns)):
            cell = table.cell(row_idx, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme.get_rgb_by_name("brand_glass_blue")
