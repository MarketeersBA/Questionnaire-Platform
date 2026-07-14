from __future__ import annotations

import logging
from typing import Any, Dict, List

from pptx.util import Inches, Pt

from .layout import PPTXLayout
from .narrative_pagination import chunk_sequence, split_text_blocks
from .template_adapter import TemplateAdapter
from .theme import PPTXTheme

logger = logging.getLogger(__name__)


class PPTXExecutiveSummary:
    """Builder for executive summary and key findings narrative slides."""

    def __init__(self, theme: PPTXTheme, layout: PPTXLayout):
        self.theme = theme
        self.layout = layout

    def build_deck(self, presentation, template_adapter: TemplateAdapter, data: Dict[str, Any]) -> int:
        return self.build_summary_deck(presentation, template_adapter, data) + self.build_findings_deck(
            presentation,
            template_adapter,
            data,
        )

    def build_summary_deck(
        self,
        presentation,
        template_adapter: TemplateAdapter,
        data: Dict[str, Any],
    ) -> int:
        return self._build_summary_slides(presentation, template_adapter, data)

    def build_findings_deck(
        self,
        presentation,
        template_adapter: TemplateAdapter,
        data: Dict[str, Any],
    ) -> int:
        return self._build_findings_slides(presentation, template_adapter, data)

    def build(self, slide, data: Dict[str, Any]):
        """Backward-compatible single-slide entrypoint for tests."""
        summary = data.get("executive_summary") or data.get("summary", "")
        findings = data.get("key_findings") or data.get("findings", [])
        if slide.placeholders:
            slide.placeholders[0].text = "EXECUTIVE SUMMARY"
        if summary:
            self._add_summary_text(slide, summary)
        if findings:
            self._add_findings_list(slide, findings[:3])

    def _build_summary_slides(
        self,
        presentation,
        template_adapter: TemplateAdapter,
        data: Dict[str, Any],
    ) -> int:
        summary = data.get("executive_summary") or data.get("summary", "")
        pages = split_text_blocks(summary, max_chars=900, max_paragraphs=3)
        if not pages:
            return 0

        rendered = 0
        for index, page in enumerate(pages, start=1):
            slide = presentation.slides.add_slide(
                template_adapter.get_layout(presentation, "ai_narrative")
            )
            if slide.placeholders:
                title = "EXECUTIVE SUMMARY"
                if len(pages) > 1:
                    title = f"{title} ({index}/{len(pages)})"
                slide.placeholders[0].text = title
            self._add_summary_text(slide, page)
            rendered += 1
        return rendered

    def _build_findings_slides(
        self,
        presentation,
        template_adapter: TemplateAdapter,
        data: Dict[str, Any],
    ) -> int:
        findings = data.get("key_findings") or data.get("findings", [])
        if not findings:
            return 0

        rendered = 0
        for index, chunk in enumerate(chunk_sequence(findings, 3), start=1):
            slide = presentation.slides.add_slide(
                template_adapter.get_layout(presentation, "ai_narrative")
            )
            if slide.placeholders:
                title = "CRITICAL FINDINGS"
                if len(findings) > 3:
                    title = f"{title} ({index})"
                slide.placeholders[0].text = title
            self._add_findings_list(slide, chunk)
            rendered += 1
        return rendered

    def _add_summary_text(self, slide, text: str):
        left = self.layout.CHART_LEFT
        top = Inches(1.5)
        width = self.layout.WIDTH - (left * 2)

        tx_box = slide.shapes.add_textbox(left, top, width, Inches(4.5))
        tf = tx_box.text_frame
        tf.word_wrap = True
        paragraph = tf.paragraphs[0]
        paragraph.text = text
        paragraph.font.name = self.theme.FONT_MEDIUM
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = self.theme.TEXT_COLOR
        paragraph.line_spacing = 1.1

    def _add_findings_list(self, slide, findings: List[Any]):
        left = self.layout.CHART_LEFT
        top = Inches(1.5)
        width = self.layout.WIDTH - (left * 2)

        for index, finding in enumerate(findings):
            y_pos = top + (index * Inches(1.4))
            finding_text = finding.get("finding") if isinstance(finding, dict) else str(finding)
            finding_label = (
                finding.get("label", f"Finding {index + 1}")
                if isinstance(finding, dict)
                else f"Observation {index + 1}"
            )

            box = slide.shapes.add_textbox(left, y_pos, width, Inches(1.2))
            tf = box.text_frame
            tf.word_wrap = True

            label = tf.paragraphs[0]
            label.text = str(finding_label).upper()
            label.font.name = self.theme.FONT_BOLD
            label.font.size = Pt(10)
            label.font.color.rgb = self.theme.SUBTITLE_COLOR

            body = tf.add_paragraph()
            body.text = str(finding_text)
            body.font.name = self.theme.FONT_MEDIUM
            body.font.size = Pt(12)
            body.font.color.rgb = self.theme.TEXT_COLOR
