import logging
from typing import Dict, Any, List
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from .builder_render_status import BuilderEmptyDataError
from .base_builder import BaseChartBuilder

logger = logging.getLogger(__name__)

class PPTXCriteriaTable(BaseChartBuilder):
    """
    Builder for the 'Criteria — Overall' table slide.
    Features:
    - Multi-brand comparison
    - Heatmap-based T2B% background coloring
    - Boolean conditional formatting for Diff (Red/Green)
    - Importance/Significance labeling
    """

    def render(self, slide: Slide, chart_data: Dict[str, Any]) -> None:
        data_payload = chart_data.get("data", {})
        raw_rows = data_payload.get("raw", [])
        brands = data_payload.get("brands", [])
        my_brand = data_payload.get("my_brand")
        
        if not raw_rows:
            raise BuilderEmptyDataError("No criteria table rows to render.")

        # 1. Define Table Structure
        # Columns: [Criteria, Importance] + [Brand 1, Brand 2...] + [Diff]
        cols = ["Criteria", "Importance"] + brands + ["Diff"]
        num_rows = len(raw_rows) + 1 # +1 for header
        num_cols = len(cols)
        
        # 2. Add Table Shape
        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            self.layout.CHART_LEFT, self.layout.CHART_TOP,
            self.layout.CHART_WIDTH, self.layout.CHART_HEIGHT
        )
        table = table_shape.table
        
        # 3. Style & Populate Header
        self._populate_header(table, cols)
        
        # 4. Populate Data Rows with Conditional Formatting
        for row_idx, data_row in enumerate(raw_rows, start=1):
            # Column 0: Criteria Name
            self._fill_cell(table.cell(row_idx, 0), data_row.get("criteria_name", ""), align=PP_ALIGN.LEFT, is_bold=True)
            
            # Column 1: Significance (Importance)
            sig = data_row.get("significance", 0)
            self._fill_cell(table.cell(row_idx, 1), f"{sig:.3f}", align=PP_ALIGN.CENTER)
            
            # Brand Columns: T2B% with Heatmap
            brand_scores = data_row.get("brand_scores", {})
            for b_idx, brand in enumerate(brands):
                score = brand_scores.get(brand, 0)
                cell = table.cell(row_idx, 2 + b_idx)
                self._fill_cell(cell, f"{score:.1f}%", align=PP_ALIGN.CENTER)
                self._apply_heatmap_fill(cell, score)
                
            # Last Column: Diff with Red/Green indicator
            diff = data_row.get("diff", 0)
            diff_cell = table.cell(row_idx, num_cols - 1)
            diff_text = f"+{diff:.1f}" if diff > 0 else f"{diff:.1f}"
            self._fill_cell(diff_cell, diff_text, align=PP_ALIGN.CENTER, is_bold=True)
            self._apply_diff_formatting(diff_cell, diff)

    def _populate_header(self, table: Any, cols: List[str]):
        """Styles and fills the top row of the table."""
        for col_idx, col_name in enumerate(cols):
            cell = table.cell(0, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme.PRIMARY_BRAND
            
            p = cell.text_frame.paragraphs[0]
            p.text = col_name.upper()
            p.font.name = self.theme.FONT_BOLD
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

    def _fill_cell(self, cell: Any, text: str, align: Any = PP_ALIGN.CENTER, is_bold: bool = False):
        """Standard cell text applicator."""
        p = cell.text_frame.paragraphs[0]
        p.text = str(text)
        p.font.name = self.theme.FONT_MEDIUM if is_bold else self.theme.FONT_LIGHT
        p.font.size = Pt(10)
        p.font.color.rgb = self.theme.TEXT_COLOR
        p.alignment = align

    def _apply_heatmap_fill(self, cell: Any, value: float):
        """Applies a blue heatmap background based on score (0-100)."""
        # Linear interpolation for alpha-like effect using primary brand and white
        # We'll use a threshold system for simplicity and better visual clarity in PPT
        cell.fill.solid()
        if value >= 80:
            rgb = RGBColor(0xDB, 0xEA, 0xFE) # Blue 100
        elif value >= 60:
            rgb = RGBColor(0xEF, 0xF6, 0xFF) # Blue 50
        else:
            return # Keep white
            
        cell.fill.fore_color.rgb = rgb

    def _apply_diff_formatting(self, cell: Any, value: float):
        """Applies Green/Red text color based on positive/negative difference."""
        if abs(value) < 0.1: return
        
        green = RGBColor(0x16, 0xA3, 0x4A) # Green 600
        red = RGBColor(0xDC, 0x26, 0x26) # Red 600
        
        cell.text_frame.paragraphs[0].font.color.rgb = green if value > 0 else red
