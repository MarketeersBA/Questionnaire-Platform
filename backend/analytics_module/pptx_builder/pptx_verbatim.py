import logging
from typing import Dict, Any, List
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from .base_builder import BaseChartBuilder
from .builder_render_status import BuilderEmptyDataError

logger = logging.getLogger(__name__)

class PPTXVerbatim(BaseChartBuilder):
    """
    Builder for Qualitative Verbatim Summaries.
    Features:
    - Split-column layout (Themes vs Quote Callouts)
    - Semantic bullet points for AI-synthesized themes
    - High-impact 'Pull Quote' boxes with italicized typography
    - Multi-brand attribution labeling
    """

    def render(self, slide: Slide, chart_data: Dict[str, Any]) -> None:
        data_payload = chart_data.get("data", {})
        themes = data_payload.get("themes", [])
        quotes = data_payload.get("quotes", [])
        
        if not themes and not quotes:
            raise BuilderEmptyDataError("Verbatim slide requires themes and/or quotes payload.")

        # 1. Layout: Two-Column Split
        col_gap = Inches(0.4)
        col_w = (self.layout.CHART_WIDTH - col_gap) / 2
        
        left_col = self.layout.CHART_LEFT
        right_col = self.layout.CHART_LEFT + col_w + col_gap
        
        # 2. Render Themes (Left Column)
        if themes:
            self._render_theme_column(slide, left_col, self.layout.CHART_TOP, col_w, themes)
            
        # 3. Render Quotes (Right Column)
        if quotes:
            self._render_quotes_column(slide, right_col, self.layout.CHART_TOP, col_w, quotes)

    def _render_theme_column(self, slide: Slide, left: float, top: float, width: float, themes: List[dict]):
        """Renders the AI-synthesized themes as structured bullets."""
        # Section Title
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = "KEY QUALITATIVE THEMES"
        tp.font.name = self.theme.FONT_BOLD
        tp.font.size = Pt(11)
        tp.font.color.rgb = self.theme.PRIMARY_BRAND
        
        # Theme List
        curr_y = top + Inches(0.5)
        for theme in themes[:5]: # Max 5 themes for safety
            box = slide.shapes.add_textbox(left, curr_y, width, Inches(0.8))
            tf = box.text_frame
            tf.word_wrap = True
            
            # Title
            p1 = tf.paragraphs[0]
            p1.text = f"• {theme.get('title', '').upper()}"
            p1.font.name = self.theme.FONT_BOLD
            p1.font.size = Pt(10)
            p1.font.color.rgb = self.theme.TEXT_COLOR
            
            # Description
            p2 = tf.add_paragraph()
            p2.text = theme.get('desc', '')
            p2.font.name = self.theme.FONT_LIGHT
            p2.font.size = Pt(9)
            p2.font.color.rgb = self.theme.SUBTITLE_COLOR
            
            curr_y += Inches(0.7)

    def _render_quotes_column(self, slide: Slide, left: float, top: float, width: float, quotes: List[dict]):
        """Renders raw respondent quotes as stylized pull-quote cards."""
        # Section Title
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = "VOICE OF THE CONSUMER"
        tp.font.name = self.theme.FONT_BOLD
        tp.font.size = Pt(11)
        tp.font.color.rgb = self.theme.get_rgb_by_name("brand_slate")
        
        curr_y = top + Inches(0.5)
        for quote in quotes[:3]: # Max 3 pull quotes for impact
            # Quote Bubble/Card
            card_h = Inches(1.2)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, curr_y, width, card_h
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self.theme.get_rgb_by_name("brand_glass_blue")
            card.line.visible = False
            
            # Quote Text
            qt_box = slide.shapes.add_textbox(left + Inches(0.1), curr_y + Inches(0.1), width - Inches(0.2), card_h - Inches(0.4))
            tf = qt_box.text_frame
            tf.word_wrap = True
            qp = tf.paragraphs[0]
            qp.text = f"\"{quote.get('text', '')}\""
            qp.font.name = self.theme.FONT_LIGHT
            qp.font.italic = True
            qp.font.size = Pt(9)
            qp.font.color.rgb = self.theme.TEXT_COLOR
            
            # Brand Tag
            tag_p = tf.add_paragraph()
            tag_p.text = f"— {quote.get('brand', 'Respondent')}"
            tag_p.font.name = self.theme.FONT_BOLD
            tag_p.font.size = Pt(8)
            tag_p.font.color.rgb = self.theme.PRIMARY_BRAND
            tag_p.alignment = PP_ALIGN.RIGHT
            
            curr_y += card_h + Inches(0.2)
