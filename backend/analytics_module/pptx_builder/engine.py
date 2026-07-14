import logging
import io
import os
from pathlib import Path
from typing import Dict, Any, List, Optional, Union, Tuple
from datetime import datetime
from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from .presentation_planner import SlideIntent, SlideType

from .layout import PPTXLayout
from .theme import PPTXTheme
from .template_adapter import TemplateAdapter
from .pptx_criteria_table import PPTXCriteriaTable
from .pptx_profile_chart import PPTXProfileChart
from .pptx_likeness_profile import PPTXLikenessProfile
from .pptx_grouped_bar import PPTXGroupedBar
from .pptx_horizontal_bar import PPTXHorizontalBar
from .pptx_stacked_bar import PPTXStackedBar
from .pptx_snake_line import PPTXSnakeLine
from .pptx_funnel_cards import PPTXFunnelCards
from .pptx_waterfall_bar import PPTXWaterfallBar
from .pptx_radar import PPTXRadar
from .pptx_funnel import PPTXFunnel
from .pptx_importance_combined import PPTXImportanceCombined
from .pptx_importance_matrix import PPTXImportanceMatrix
from .pptx_reference_table import PPTXReferenceTable
from .pptx_scatter import PPTXScatter
from .pptx_sigma_intent import PPTXSigmaIntent
from .pptx_positioning_matrix import PPTXPositioningMatrix
from .pptx_nps_gauge import PPTXNPSGauge
from .pptx_affinity_heatmap import PPTXAffinityHeatmap
from .pptx_brand_comparison import PPTXBrandComparison
from .pptx_scorecard import PPTXScorecard
from .pptx_wordcloud import PPTXWordcloud
from .pptx_verbatim import PPTXVerbatim
from .chart_resolver import PPTXChartResolver, build_builder_registry
from .slide_chrome import ContentSlideChromeApplier
from .pptx_market_position import PPTXMarketPosition
from .pptx_strategic_command import PPTXStrategicCommand
from .pptx_executive_summary import PPTXExecutiveSummary
from .pptx_strategic_narrative import PPTXStrategicNarrative
from .pptx_swot import PPTXSwot
from .pptx_recommendations import PPTXRecommendations
from .pptx_brand_profile_card import PPTXBrandProfileCard

logger = logging.getLogger(__name__)

class PPTXEngine:
    """
    Main Orchestrator for the Native PPTX V2 Engine.
    Coordinates theme loading, layout geometry, and builder delegation.
    """

    def __init__(self, theme: Optional[PPTXTheme] = None, layout: Optional[PPTXLayout] = None, template_path: Optional[str] = None):
        self.theme = theme or PPTXTheme()
        self.layout = layout or PPTXLayout.for_reference()
        self.template_adapter = TemplateAdapter()
        
        # Resolve Template Path: Docker-Aware Fallback
        if template_path:
            self.template_path = Path(template_path)
        else:
            res_dir_env = os.environ.get("ANALYTICS_RESOURCES_DIR")
            if res_dir_env:
                base_dir = Path(res_dir_env)
            # Dynamic fallback for Docker volume mount discrepancies
            elif Path("resources/analytics").exists():
                 base_dir = Path("resources/analytics")
            else:
                 base_dir = Path("backend/resources/analytics")
                 
            self.template_path = base_dir / "marketeers_template.pptx"
            
        self._builder_registry = {}
        self.chart_resolver = PPTXChartResolver()
        self.render_journal: List[Dict[str, Any]] = []
        self.narrative_render_journal: List[Dict[str, Any]] = []
        self.automation_notes: List[Dict[str, Any]] = []
        self.layout_geometry = self.layout.geometry_manifest()
        self._content_chrome = ContentSlideChromeApplier(self.theme, self.layout)
        self._register_default_builders()

    def _resolve_logo_path(self) -> Optional[Path]:
        base = Path(os.environ.get("ANALYTICS_RESOURCES_DIR", "backend/resources/analytics"))
        logo_path = base / "logo.png"
        return logo_path if logo_path.exists() else None

    def _apply_brand_logo(self, slide: Slide) -> None:
        logo_path = self._resolve_logo_path()
        if not logo_path:
            return

        for ph in slide.placeholders:
            try:
                if hasattr(ph, "insert_picture"):
                    ph.insert_picture(str(logo_path))
                    return
            except Exception:
                continue

        try:
            slide.shapes.add_picture(
                str(logo_path),
                self.layout.WIDTH - Inches(1.4),
                Inches(0.2),
                width=Inches(1.0),
            )
        except Exception as exc:
            logger.warning("[PPTXEngine] Failed to place logo: %s", exc)

    def _resolve_report_date_text(self, metadata: Dict[str, Any]) -> str:
        explicit = metadata.get("date")
        if isinstance(explicit, str) and explicit.strip():
            return explicit.strip()
        generated = metadata.get("report_generated_at")
        if generated:
            try:
                if hasattr(generated, "strftime"):
                    return generated.strftime("%B %Y")
                parsed = datetime.fromisoformat(str(generated).replace("Z", "+00:00"))
                return parsed.strftime("%B %Y")
            except Exception:
                pass
        return datetime.now().strftime("%B %Y")

    def _register_default_builders(self):
        """Initializes the resolver registry with all specialized chart builders."""
        self._builder_registry = build_builder_registry()
        self.chart_resolver = PPTXChartResolver(self._builder_registry)

    def generate_presentation(self, intents: List[SlideIntent]) -> Tuple[io.BytesIO, int]:
        """
        The Master Build Pipeline.
        Now a pure renderer: consumes a list of SlideIntents and executes them.
        """
        # 1. Initialize from Brand Template
        if not self.template_path.exists():
            logger.warning(f"[PPTXEngine] Template not found at {self.template_path}, falling back to blank presentation.")
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
        else:
            logger.info(f"[PPTXEngine] Initializing presentation from template: {self.template_path}")
            prs = Presentation(str(self.template_path))

        self.layout = PPTXLayout.from_presentation(prs)
        self.layout_geometry = self.layout.geometry_manifest()
        self._content_chrome = ContentSlideChromeApplier(self.theme, self.layout)
        if not self.layout_geometry.get("chart_frame_fits_slide"):
            logger.warning(
                "[PPTXEngine] Chart frame exceeds slide canvas: %s",
                self.layout_geometry,
            )
        logger.info(
            "[PPTXEngine] Canvas bound to %.3f\" x %.3f\"",
            self.layout.slide_width_emu / 914400,
            self.layout.slide_height_emu / 914400,
        )

        self.render_journal = []
        self.narrative_render_journal = []
        self.automation_notes = []

        # 2. Template Audit (Phase 7 - Hardening)
        audit_report = self.template_adapter.audit(prs)
        if audit_report.get("missing"):
            logger.warning(
                "[PPTXEngine] Template audit missing roles: %s",
                ", ".join(audit_report["missing"]),
            )

        # 3. Iterate and Dispatch Intents with Template Enforcement
        for intent_index, intent in enumerate(intents, start=1):
            try:
                with TemplateEnforcementGuard(self, prs, intent):
                    self._dispatch_intent(prs, intent)
            except Exception as e:
                self.automation_notes.append(
                    {
                        "type": "intent_failure",
                        "severity": "warning",
                        "intent_index": intent_index,
                        "intent_type": str(getattr(intent.type, "value", intent.type)),
                        "message": str(e),
                    }
                )
                logger.error(f"[PPTXEngine] Failed to execute intent {intent.type}: {e}", exc_info=True)

        # 3. Finalize Stream and Metadata
        stream = io.BytesIO()
        prs.save(stream)
        stream.seek(0)
        
        slide_count = len(prs.slides)
        logger.info(f"[PPTXEngine] Generation complete. Total slides: {slide_count}")
        
        return stream, slide_count

    def _dispatch_intent(self, prs: Presentation, intent: SlideIntent):
        """Pure-render dispatcher for slide intents."""
        if intent.type == SlideType.COVER:
            self._add_cover_slide(prs, intent.data)
        
        elif intent.type == SlideType.SURVEY_OVERVIEW:
            self._add_survey_info_slide(prs, intent.data)
            
        elif intent.type == SlideType.EXECUTIVE_SUMMARY:
            self._add_executive_summary_slide(prs, intent.data)
            
        elif intent.type == SlideType.STRATEGIC_NARRATIVE:
            self._add_strategic_narrative_slide(prs, intent.data)
            
        elif intent.type == SlideType.STRATEGIC_INTELLIGENCE:
            # Multi-part dispatcher for strategic content
            self._add_strategic_intelligence_slides(prs, intent.data)
            
        elif intent.type == SlideType.SECTION_DIVIDER:
            self._add_section_divider(prs, intent.title)
            
        elif intent.type == SlideType.CONTENT_SLIDE:
            self._add_content_slide(prs, intent.data)
            
        elif intent.type == SlideType.SWOT:
            self._add_swot_slide(prs, intent.title, intent.data)
            
        elif intent.type == SlideType.RECOMMENDATIONS_4P:
            self._add_recommendations_slide(prs, intent.data)
            
        elif intent.type == SlideType.BRAND_PROFILE:
            self._add_brand_profile_slide(prs, intent.data)
            
        elif intent.type == SlideType.CLOSING:
            self._add_closing_slide(prs, intent.data)

    def _record_narrative_render(
        self,
        section_id: str,
        *,
        title: str,
        rendered_slides: int,
        status: str = "rendered",
        message: str = "",
    ) -> None:
        self.narrative_render_journal.append(
            {
                "section_id": section_id,
                "title": title,
                "rendered_slides": rendered_slides,
                "status": status,
                "message": message,
            }
        )

    def _add_executive_summary_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Render executive summary, findings, and opportunity insights in screen order."""
        executive_builder = PPTXExecutiveSummary(self.theme, self.layout)
        opportunity_builder = PPTXStrategicCommand(self.theme, self.layout)

        try:
            summary_slides = executive_builder.build_summary_deck(prs, self.template_adapter, data)
            if summary_slides:
                self._record_narrative_render(
                    "executive_summary",
                    title="Executive Summary",
                    rendered_slides=summary_slides,
                )

            findings_slides = executive_builder.build_findings_deck(prs, self.template_adapter, data)
            if findings_slides:
                self._record_narrative_render(
                    "key_findings",
                    title="Critical Findings",
                    rendered_slides=findings_slides,
                )

            opportunity_slides = opportunity_builder.build_deck(
                prs,
                self.template_adapter,
                data.get("opportunity_insights") or data.get("opportunities", []),
            )
            if opportunity_slides:
                self._record_narrative_render(
                    "opportunity_insights",
                    title="Opportunity Insights",
                    rendered_slides=opportunity_slides,
                )
        except Exception as e:
            logger.error(f"[PPTXEngine] Failed to build Executive Summary deck: {e}")
            self._record_narrative_render(
                "executive_summary",
                title="Executive Summary",
                rendered_slides=0,
                status="failed",
                message=str(e),
            )

    def _add_strategic_narrative_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Render the premium Strategic Narrative slide."""
        builder = PPTXStrategicNarrative(self.theme, self.layout)
        try:
            slide = prs.slides.add_slide(self.template_adapter.get_layout(prs, "ai_narrative"))
            builder.build(slide, data)
            self._record_narrative_render(
                "strategic_narrative",
                title="Strategic Narrative",
                rendered_slides=1,
            )
        except Exception as e:
            logger.error(f"[PPTXEngine] Failed to build Strategic Narrative slide: {e}")
            self._record_narrative_render(
                "strategic_narrative",
                title="Strategic Narrative",
                rendered_slides=0,
                status="failed",
                message=str(e),
            )

    def _add_swot_slide(self, prs: Presentation, title: str, data: Dict[str, Any]):
        """Inject paginated SWOT slides for a single brand."""
        builder = PPTXSwot(self.theme, self.layout)
        try:
            rendered = builder.build_deck(prs, self.template_adapter, data)
            self._record_narrative_render(
                f"swot::{data.get('brand', 'Brand')}",
                title=title,
                rendered_slides=rendered,
            )
        except Exception as e:
            logger.error(f"[PPTXEngine] Failed to build SWOT slide: {e}")
            self._record_narrative_render(
                f"swot::{data.get('brand', 'Brand')}",
                title=title,
                rendered_slides=0,
                status="failed",
                message=str(e),
            )

    def _add_recommendations_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Inject paginated 4P recommendation roadmap slides."""
        builder = PPTXRecommendations(self.theme, self.layout)
        try:
            rendered = builder.build_deck(prs, self.template_adapter, data)
            self._record_narrative_render(
                "recommendations_4p",
                title="4P Recommendations",
                rendered_slides=rendered,
            )
        except Exception as e:
            logger.error(f"[PPTXEngine] Failed to build Recommendations slide: {e}")
            self._record_narrative_render(
                "recommendations_4p",
                title="4P Recommendations",
                rendered_slides=0,
                status="failed",
                message=str(e),
            )


    def _add_brand_profile_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Renders the premium dark-mode Brand Profile card."""
        builder = PPTXBrandProfileCard(self.theme, self.layout)
        try:
            # We use the blank layout for full-bleed dark mode designs
            slide = prs.slides.add_slide(self.template_adapter.get_layout(prs, "blank"))
            builder.render(slide, data)
            
            brand_name = data.get("brand_data", {}).get("title", "Unknown Brand")
            self._record_narrative_render(
                f"brand_profile::{brand_name}",
                title=f"Brand Profile: {brand_name}",
                rendered_slides=1,
            )
            logger.info(f"[PPTXEngine] Brand Profile slide injected for {brand_name}.")
        except Exception as e:
            logger.error(f"[PPTXEngine] Failed to build Brand Profile slide: {e}", exc_info=True)
            self._record_narrative_render(
                "brand_profile",
                title="Brand Profile",
                rendered_slides=0,
                status="failed",
                message=str(e),
            )

    def _add_strategic_intelligence_slides(self, prs: Presentation, data: Dict[str, Any]):
        """Render market position narrative sections and strategic evidence charts."""
        market_report = data.get("market_position_report") or data.get("report")
        charts = data.get("charts", [])

        try:
            if market_report:
                self._add_section_divider(prs, "Strategic Positioning")
                rendered = PPTXMarketPosition(self.theme, self.layout).build_deck(
                    prs,
                    self.template_adapter,
                    market_report,
                )
                self._record_narrative_render(
                    "market_position_report",
                    title="Market Position Report",
                    rendered_slides=rendered,
                )

            if charts:
                self._add_section_divider(prs, "Positioning Visualizations")
                for chart in charts:
                    self._add_content_slide(prs, chart)
        except Exception as e:
            logger.error(f"[PPTXEngine] Failed to build Strategic Intelligence deck: {e}")

    def _add_cover_slide(self, prs: Presentation, metadata: dict):
        """
        Creates a cinematic branded cover slide using the named 'cover' role.
        Inherits all freeform brand decorations from the template.
        """
        slide_layout = self.template_adapter.get_layout(prs, "cover")
        slide = prs.slides.add_slide(slide_layout)
        
        # 1. Populate Title Placeholder (idx 0)
        title_ph = slide.placeholders[0]
        project_name = (
            metadata.get("title")
            or metadata.get("project_name")
            or "Market Research Insights"
        )
        title_ph.text = str(project_name).upper()
        
        # 2. Populate Subtitle Placeholder (idx 1)
        # We combine Target Brand and Date for a professional look
        company = (
            metadata.get("company_name")
            or metadata.get("brand")
            or "CONFIDENTIAL CLIENT"
        )
        date_str = self._resolve_report_date_text(metadata).upper()
        subtitle_ph = slide.placeholders[1]
        subtitle_ph.text = f"PREPARED FOR: {str(company).upper()}\n{date_str}"
        self._apply_brand_logo(slide)

    def _add_survey_info_slide(self, prs: Presentation, metadata: dict):
        """
        Injects a premium research methodology slide.
        Uses TemplateAdapter to resolve the methodology layout by role.
        """
        from .pptx_survey_info import PPTXSurveyInfo
        slide_layout = self.template_adapter.get_layout(prs, "survey_overview")
        slide = prs.slides.add_slide(slide_layout)
        builder = PPTXSurveyInfo(self.theme, self.layout)
        
        # We populate the native title placeholder [0] in Layout 3
        if slide.placeholders[0]:
            slide.placeholders[0].text = "SURVEY OVERVIEW & METHODOLOGY"
            
        try:
            builder.build(slide, metadata)
            logger.info("[PPTXEngine] Survey Info slide injected successfully on Layout 3.")
        except Exception as e:
            logger.error(f"[PPTXEngine] Failed to build Survey Info slide: {e}")

    def _add_section_divider(self, prs: Presentation, title: str):
        """
        Creates an impactful section separator.
        Uses TemplateAdapter to resolve the divider layout by role.
        """
        slide_layout = self.template_adapter.get_layout(prs, "section_divider")
        slide = prs.slides.add_slide(slide_layout)
        
        # Populate Section Title (idx 0)
        try:
            ph = slide.placeholders[0]
            ph.text = title.upper()
        except Exception as e:
            logger.warning(f"[PPTXEngine] Layout 6 placeholder missing for title: {e}")
        self._apply_brand_logo(slide)

    def _add_content_slide(self, prs: Presentation, chart_data: Dict[str, Any]):
        """
        Delegates specific chart building to registered builders using Layout 1.
        Template chrome is owned by the engine; builders render chart content only.
        """
        chart_type = chart_data.get("chart_type")
        chart_title = chart_data.get("title", "Insight Analysis").upper()

        slide_layout = self.template_adapter.get_layout(prs, "content")
        slide = prs.slides.add_slide(slide_layout)
        self._apply_brand_logo(slide)

        chrome_spec = self._content_chrome.apply(slide, chart_data)

        from .hybrid_export.render_routing import resolve_content_slide_builder

        resolution, builder_class = resolve_content_slide_builder(chart_data, self.chart_resolver)
        logger.info(
            "[PPTXEngine] Resolved chart '%s' via %s -> %s",
            chart_title,
            resolution.source,
            resolution.registry_key,
        )

        journal_entry = {
            "pptx_slide_id": chart_data.get("_pptx_slide_id")
            or f"{chart_data.get('chart_id')}::{chart_data.get('chart_type')}",
            "slide_index": len(prs.slides),
            "chart_id": chart_data.get("chart_id"),
            "chart_type": chart_data.get("chart_type"),
            "title": chart_data.get("title"),
            "registry_key": resolution.registry_key,
            "resolution_source": resolution.source,
            "render_mode": "image_capture" if resolution.registry_key == "image_capture" else "native",
        }
        self.render_journal.append(journal_entry)

        try:
            builder = builder_class(self.theme, self.layout)
            render_result = builder.build(slide, chart_data, chrome_owned_by_engine=True)
            journal_entry["render_status"] = render_result.status.value
            journal_entry["render_message"] = render_result.message
        except Exception as e:
            journal_entry["render_status"] = "failed"
            journal_entry["render_message"] = str(e)
            self.automation_notes.append(
                {
                    "type": "content_slide_render_failure",
                    "severity": "warning",
                    "chart_id": chart_data.get("chart_id"),
                    "chart_type": chart_type,
                    "registry_key": resolution.registry_key,
                    "message": str(e),
                }
            )
            logger.error(
                "[PPTXEngine] Rendering failed for %s/%s ('%s'): %s",
                chart_type,
                resolution.registry_key,
                chart_title,
                e,
                exc_info=True,
            )
            from .base_builder import BaseChartBuilder

            BaseChartBuilder.inject_error_placeholder(
                slide,
                self.layout,
                self.theme,
                chart_title,
                f"Rendering Error: {str(e)}",
            )

        if self._content_chrome.analysis_requires_followup_slide(chrome_spec.ai_deep_analysis):
            self._add_deep_analysis_followup_slide(prs, chart_data, chrome_spec)

    def _add_deep_analysis_followup_slide(
        self,
        prs: Presentation,
        chart_data: Dict[str, Any],
        chrome_spec,
    ) -> None:
        analysis_points = self._content_chrome.normalized_analysis_points(chrome_spec.ai_deep_analysis)
        if not analysis_points:
            return

        chunk_size = 3
        chunks = [
            analysis_points[i : i + chunk_size]
            for i in range(0, len(analysis_points), chunk_size)
        ]
        for chunk_index, chunk in enumerate(chunks, start=1):
            slide_layout = self.template_adapter.get_layout(prs, "content")
            slide = prs.slides.add_slide(slide_layout)
            self._apply_brand_logo(slide)

            title = f"{(chart_data.get('title') or 'Insight Analysis').upper()} - AI DEEP ANALYSIS"
            if len(chunks) > 1:
                title = f"{title} ({chunk_index}/{len(chunks)})"
            payload = {
                "title": title,
                "subtitle": chart_data.get("subtitle"),
                "ai_headline": chart_data.get("ai_headline"),
                "insight": chart_data.get("insight"),
                "ai_deep_analysis": [],
            }
            self._content_chrome.apply(slide, payload)

            card_left = self.layout.CHART_LEFT
            card_top = self.layout.CHART_TOP
            card_width = self.layout.CHART_WIDTH
            card_height = self.layout.CHART_HEIGHT / 3
            gap = Inches(0.2)

            for idx, point in enumerate(chunk):
                top = card_top + (idx * (card_height + gap))
                box = slide.shapes.add_textbox(
                    int(card_left),
                    int(top),
                    int(card_width),
                    int(card_height),
                )
                tf = box.text_frame
                tf.clear()
                sentiment = point.get("sentiment", "neutral").upper()
                p0 = tf.paragraphs[0]
                p0.text = f"[{sentiment}] {point.get('title', 'Insight')}"
                p0.font.name = self.theme.FONT_BOLD
                p0.font.size = Pt(12)
                p0.font.color.rgb = self.theme.PRIMARY_BRAND

                p1 = tf.add_paragraph()
                p1.text = point.get("body", "")
                p1.font.name = self.theme.FONT_MEDIUM
                p1.font.size = Pt(10)
                p1.font.color.rgb = self.theme.TEXT_COLOR

                recommended_action = point.get("recommended_action")
                if recommended_action:
                    p2 = tf.add_paragraph()
                    p2.text = f"Action: {recommended_action}"
                    p2.font.name = self.theme.FONT_LIGHT
                    p2.font.size = Pt(9)
                    p2.font.color.rgb = self.theme.get_rgb_by_name("brand_slate")

    def _inject_unsupported_placeholder(self, slide: Slide, title: str, chart_type: str):
        """
        Injects a visible, branded 'Unsupported' placeholder.
        Prevents silent data loss and provides immediate visual feedback.
        """
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt

        msg_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(12), Inches(2))
        tf = msg_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{title or 'Analysis'} (Preview Unavailable)"
        p.font.name = self.theme.FONT_BOLD
        p.font.size = Pt(16)
        p.font.color.rgb = self.theme.get_rgb_by_name("brand_slate")
        p.alignment = PP_ALIGN.CENTER

        icon = slide.shapes.add_shape(
            MSO_SHAPE.HEART,
            Inches(9.25), Inches(4.0), Inches(1.5), Inches(1.5)
        )
        icon.fill.solid()
        icon.fill.fore_color.rgb = self.theme.get_rgb_by_name("brand_slate")
        icon.fill.transparency = 0.8
        icon.line.visible = False

        detail_box = slide.shapes.add_textbox(Inches(4), Inches(5.5), Inches(12), Inches(2))
        detail_tf = detail_box.text_frame
        detail_p = detail_tf.paragraphs[0]
        detail_p.text = f"DEVELOPMENT IN PROGRESS: NATIVE '{chart_type.upper()}' RENDERER"
        detail_p.font.name = self.theme.FONT_BOLD
        detail_p.font.size = Pt(14)
        detail_p.font.color.rgb = self.theme.get_rgb_by_name("brand_slate")
        detail_p.alignment = PP_ALIGN.CENTER

        detail_p2 = detail_tf.add_paragraph()
        detail_p2.text = "This component is visible in the web report but requires specialized PPTX native mapping."
        detail_p2.font.name = self.theme.FONT_LIGHT
        detail_p2.font.size = Pt(11)
        detail_p2.font.color.rgb = self.theme.get_rgb_by_name("brand_slate")
        detail_p2.alignment = PP_ALIGN.CENTER

    def _add_closing_slide(self, prs: Presentation, metadata: dict):
        """
        Creates a professional closing slide using the named 'closing' role.
        Inherits the branded 'THANK YOU' designs and decorative shapes.
        """
        slide_layout = self.template_adapter.get_layout(prs, "closing")
        slide = prs.slides.add_slide(slide_layout)
        self._apply_brand_logo(slide)
        
        # Populate Placeholder (idx 10 - Text Placeholder 4)
        # This typically contains contact info or final call to action
        try:
            ph = slide.placeholders[10]
            project_name = (
                metadata.get("title")
                or metadata.get("project_name")
                or "MARKET RESEARCH INSIGHTS"
            )
            company_name = metadata.get("company_name") or metadata.get("brand") or "CONFIDENTIAL CLIENT"
            date_str = self._resolve_report_date_text(metadata)
            ph.text = (
                f"{str(project_name).upper()}\n"
                f"CLIENT: {str(company_name).upper()}\n"
                f"DATE: {str(date_str).upper()}\n"
                "QUESTIONS & NEXT STEPS\ncontact@marketeers.com"
            )
        except Exception as e:
            logger.warning(f"[PPTXEngine] Failed to populate closing slide placeholder: {e}")

    def register_builder(self, chart_type: str, builder_class: Any):
        self._builder_registry[chart_type] = builder_class
        self.chart_resolver.register(chart_type, builder_class)

class TemplateEnforcementGuard:
    """
    Active security monitor for PPTX corporate identity.
    Tracks slide creation during intent dispatch and verifies that all new slides
    use trusted layouts from the TemplateAdapter.
    """
    def __init__(self, engine: 'PPTXEngine', prs: Presentation, intent: SlideIntent):
        self.engine = engine
        self.prs = prs
        self.intent = intent
        self._start_count = 0
        self._trusted_indices = set()

    def __enter__(self):
        self._start_count = len(self.prs.slides)
        # Refresh trusted indices per intent if needed, or cache on engine
        self._trusted_indices = self.engine.template_adapter.get_trusted_layout_indices(self.prs)
        return self

    def __exit__(self, exc_type, exc_val, exc_tb):
        if exc_type is not None:
            return  # Let the engine handle the actual error

        end_count = len(self.prs.slides)
        if end_count > self._start_count:
            # Slides were added. Verify their layouts.
            for i in range(self._start_count, end_count):
                slide = self.prs.slides[i]
                # We check the layout index in the slide_layouts collection
                layout_found = False
                for idx, layout in enumerate(self.prs.slide_layouts):
                    if layout == slide.slide_layout:
                        if idx not in self._trusted_indices:
                            logger.warning(
                                "[TemplateEnforcement] SECURITY ALERT: Slide %d in intent '%s' "
                                "uses UNTRUSTED layout '%s' (idx %d).",
                                i + 1, self.intent.type, slide.slide_layout.name, idx
                            )
                        layout_found = True
                        break
                
                if not layout_found:
                    logger.error(
                        "[TemplateEnforcement] CORRUPTION: Slide %d uses a layout "
                        "not found in the presentation masters.", i + 1
                    )
