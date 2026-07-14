from __future__ import annotations

import logging
from typing import Dict, List

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .layout import PPTXLayout
from .narrative_pagination import split_text_blocks
from .template_adapter import TemplateAdapter
from .theme import PPTXTheme

logger = logging.getLogger(__name__)


class PPTXRecommendations:
    """Builder for 4P strategic roadmap slides with paginated commentary."""

    PILLARS = ("product", "price", "place", "promotion")

    def __init__(self, theme: PPTXTheme, layout: PPTXLayout):
        self.theme = theme
        self.layout = layout
        self.COLORS = self.theme.get_recommendation_palette()

    def build_deck(
        self,
        presentation,
        template_adapter: TemplateAdapter,
        recommendations: Dict[str, str],
    ) -> int:
        pages = self._paginate_recommendations(recommendations)
        rendered = 0

        for index, page in enumerate(pages, start=1):
            slide = presentation.slides.add_slide(
                template_adapter.get_layout(presentation, "ai_narrative")
            )
            title = "STRATEGIC ROADMAP: 4P RECOMMENDATIONS"
            if len(pages) > 1:
                title = f"{title} ({index}/{len(pages)})"
            if slide.placeholders:
                slide.placeholders[0].text = title
            self._render_page(slide, page)
            rendered += 1
        return rendered

    def build(self, slide, recommendations: Dict[str, str]):
        self._render_page(slide, recommendations)

    def _paginate_recommendations(self, recommendations: Dict[str, str]) -> List[Dict[str, str]]:
        split_by_pillar = {
            pillar: split_text_blocks(recommendations.get(pillar, "Continue current strategy."), max_chars=420)
            for pillar in self.PILLARS
        }
        page_count = max((len(blocks) for blocks in split_by_pillar.values()), default=1)
        pages: List[Dict[str, str]] = []
        for page_index in range(page_count):
            pages.append(
                {
                    pillar: split_by_pillar[pillar][page_index]
                    if page_index < len(split_by_pillar[pillar])
                    else ""
                    for pillar in self.PILLARS
                }
            )
        return pages

    def _render_page(self, slide, recommendations: Dict[str, str]) -> None:
        margin = self.layout.CHART_LEFT
        spacing = Inches(0.4)
        col_w = (self.layout.WIDTH - (margin * 2) - (spacing * 3)) / 4
        col_h = Inches(5.0)
        start_top = Inches(1.8)

        for index, pillar in enumerate(self.PILLARS):
            left = margin + (index * (col_w + spacing))
            advice = recommendations.get(pillar, "Continue current strategy.")
            self._draw_recommendation_column(slide, left, start_top, col_w, col_h, pillar, advice)

    def _draw_recommendation_column(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        label: str,
        text: str,
    ):
        color = self.COLORS.get(label, self.theme.PRIMARY_BRAND)

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left + (width / 2) - Inches(0.4),
            top - Inches(0.6),
            Inches(0.8),
            Inches(0.8),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.visible = False

        header = slide.shapes.add_textbox(left, top + Inches(0.3), width, Inches(0.4))
        header.text_frame.paragraphs[0].text = label.upper()
        header.text_frame.paragraphs[0].font.name = self.theme.FONT_BOLD
        header.text_frame.paragraphs[0].font.size = Pt(14)
        header.text_frame.paragraphs[0].font.color.rgb = self.theme.PRIMARY_BRAND
        header.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        content = slide.shapes.add_textbox(left, top + Inches(1.0), width, height - Inches(1.0))
        tf = content.text_frame
        tf.word_wrap = True
        paragraph = tf.paragraphs[0]
        paragraph.text = text
        paragraph.font.name = self.theme.FONT_MEDIUM
        paragraph.font.size = Pt(13)
        paragraph.font.color.rgb = self.theme.TEXT_COLOR
        paragraph.line_spacing = 1.1
        paragraph.alignment = PP_ALIGN.CENTER
