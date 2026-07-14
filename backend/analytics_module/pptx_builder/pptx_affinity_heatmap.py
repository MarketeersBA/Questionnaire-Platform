import logging
from typing import Dict, Any, List
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from .base_builder import BaseChartBuilder
from .builder_render_status import BuilderEmptyDataError

logger = logging.getLogger(__name__)

class PPTXAffinityHeatmap(BaseChartBuilder):
    """
    Builder for Affinity Heatmaps (Demographic vs Brand).
    Features:
    - High-density native table architecture
    - 5-Tier Semantic Color Mapping (Blue -> White -> Orange)
    - Precision categorical alignment
    - Auto-scaling font logic for complex matrices
    """

    def render(self, slide: Slide, chart_data: Dict[str, Any]) -> None:
        data_payload = chart_data.get("data", {})
        brands = data_payload.get("brands", [])
        demographics = data_payload.get("demographics", [])
        matrix = data_payload.get("matrix", []) # [[row...]] where row matches demographics[i]
        
        if not demographics or not brands or not matrix:
            raise BuilderEmptyDataError("Missing affinity heatmap matrix data.")

        # 1. Table Dimensions
        num_cols = len(brands) + 1
        num_rows = len(demographics) + 1
        
        # 2. Add Table
        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            self.layout.CHART_LEFT, self.layout.CHART_TOP,
            self.layout.CHART_WIDTH, self.layout.CHART_HEIGHT
        )
        table = table_shape.table
        
        # 3. Populate Header (Brands)
        cols = ["SEGMENT"] + [str(b).upper() for b in brands]
        self._populate_header(table, cols)
        
        # 4. Populate Matrix with Heatmap Coloring
        for r_idx, segment in enumerate(demographics):
            row_idx = r_idx + 1
            row_data = matrix[r_idx] if r_idx < len(matrix) else []
            
            # Col 0: Segment Name
            self._fill_cell(table.cell(row_idx, 0), segment, align=PP_ALIGN.LEFT, is_bold=True)
            
            # Data Cells
            for c_idx, val in enumerate(row_data):
                if c_idx + 1 >= num_cols: break
                cell = table.cell(row_idx, c_idx + 1)
                
                # Format: Assume Index data (100 is neutral) or %
                text = f"{val:.0f}" if val > 10 else f"{val:.1f}"
                self._fill_cell(cell, text, align=PP_ALIGN.CENTER)
                
                # Apply 5-Scale Heatmap
                self._apply_5scale_fill(cell, val)

    def _populate_header(self, table: Any, cols: List[str]):
        """Styles the top row with dark navy for contrast."""
        for j, col_name in enumerate(cols):
            cell = table.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme.PRIMARY_BRAND
            
            p = cell.text_frame.paragraphs[0]
            p.text = col_name
            p.font.name = self.theme.FONT_BOLD
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

    def _fill_cell(self, cell: Any, text: str, align: Any = PP_ALIGN.CENTER, is_bold: bool = False):
        p = cell.text_frame.paragraphs[0]
        p.text = str(text)
        p.font.name = self.theme.FONT_MEDIUM if is_bold else self.theme.FONT_LIGHT
        p.font.size = Pt(9)
        p.font.color.rgb = self.theme.TEXT_COLOR
        p.alignment = align

    def _apply_5scale_fill(self, cell: Any, value: float):
        """
        Maps value to 5-tier heatmap:
        Blue (Strong +) -> Lt Blue -> White (Neutral) -> Lt Orange -> Deep Orange (Strong -)
        Assumes Index values where ~100 is neutral.
        """
        cell.fill.solid()
        
        # Color Definitions (Premium Hues)
        DEEP_BLUE = RGBColor(0x1E, 0x3A, 0x8A) # Blue 900
        LT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)   # Blue 100
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        LT_ORANGE = RGBColor(0xFF, 0xED, 0xD5) # Orange 100
        DEEP_ORANGE = RGBColor(0xF9, 0x73, 0x16) # Orange 500
        
        if value >= 120:
            color = DEEP_BLUE
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        elif value >= 105:
            color = LT_BLUE
        elif value >= 95:
            color = WHITE
        elif value >= 80:
            color = LT_ORANGE
        else: # < 80
            color = DEEP_ORANGE
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            
        cell.fill.fore_color.rgb = color
