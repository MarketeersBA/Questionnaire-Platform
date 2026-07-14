from __future__ import annotations

import logging
from abc import ABC, abstractmethod
from typing import Any, Dict, Optional

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Pt

from .builder_render_status import BuilderEmptyDataError, ChartRenderResult
from .layout import PPTXLayout
from .slide_chrome import ContentSlideChromeApplier
from .theme import PPTXTheme

logger = logging.getLogger(__name__)


class BaseChartBuilder(ABC):
    """Base class for native chart builders. Chrome is owned by the engine in production."""

    def __init__(self, theme: PPTXTheme, layout: PPTXLayout):
        self.theme = theme
        self.layout = layout
        self._brand_color_map: Dict[str, Any] = {}

    def build(
        self,
        slide: Slide,
        chart_data: Dict[str, Any],
        metadata: Optional[Dict[str, Any]] = None,
        *,
        chrome_owned_by_engine: bool = False,
    ) -> ChartRenderResult:
        if not chrome_owned_by_engine:
            self._apply_standalone_chrome(slide, chart_data, metadata)

        title = chart_data.get("title") or (metadata.get("title") if metadata else None)
        try:
            logger.info("[Builder] Invoking %s for '%s'", self.__class__.__name__, title)
            self.render(slide, chart_data)
            return ChartRenderResult.rendered()
        except BuilderEmptyDataError as exc:
            logger.warning("[Builder] Empty payload for %s: %s", self.__class__.__name__, exc)
            return ChartRenderResult.skipped_empty_data(str(exc))
        except Exception as exc:
            logger.error("[Builder] Rendering failed for %s: %s", self.__class__.__name__, exc, exc_info=True)
            BaseChartBuilder.inject_error_placeholder(
                slide,
                self.layout,
                self.theme,
                str(title or "Analysis"),
                str(exc),
            )
            return ChartRenderResult.failed(str(exc))

    def _apply_standalone_chrome(
        self,
        slide: Slide,
        chart_data: Dict[str, Any],
        metadata: Optional[Dict[str, Any]],
    ) -> None:
        """Test-only path that reuses the same chrome contract as the engine."""
        payload = dict(chart_data)
        if metadata:
            payload.setdefault("title", metadata.get("title"))
            payload.setdefault("subtitle", metadata.get("subtitle"))
        ContentSlideChromeApplier(self.theme, self.layout).apply(slide, payload)

    @staticmethod
    def inject_error_placeholder(
        slide: Slide,
        layout: PPTXLayout,
        theme: PPTXTheme,
        title: str,
        error_msg: str,
    ) -> None:
        left, top, width, height = layout.chart_body_bounds()
        icon_size = min(int(width * 0.08), int(height * 0.12))
        icon_left = int(left) + int((int(width) - icon_size) / 2)
        icon_top = int(top) + int(height * 0.25)

        icon_box = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, icon_left, icon_top, icon_size, icon_size)
        icon_box.fill.solid()
        icon_box.fill.fore_color.rgb = theme.get_rgb_by_name("brand_crimson")

        msg_box = slide.shapes.add_textbox(int(left), icon_top + icon_size, int(width), int(height * 0.2))
        paragraph = msg_box.text_frame.paragraphs[0]
        paragraph.text = f"ANALYSIS INTERRUPTED\nReason: {error_msg[:100]}..."
        paragraph.font.name = theme.FONT_BOLD
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = theme.get_rgb_by_name("brand_crimson")
        paragraph.alignment = PP_ALIGN.CENTER

    def _apply_brand_colors(self, series_name: str) -> Any:
        brand_key = series_name.strip().lower()

        if brand_key in self._brand_color_map:
            return self._brand_color_map[brand_key]

        hero_keywords = ["hero", "primary", "total"]
        if any(keyword in brand_key for keyword in hero_keywords):
            color = self.theme.PRIMARY_BRAND
        elif "comp" in brand_key or "other" in brand_key:
            color = self.theme.get_rgb_by_name("brand_slate")
        else:
            color = self.theme.get_color(len(self._brand_color_map))

        self._brand_color_map[brand_key] = color
        return color

    def apply_series_data_labels(
        self,
        series: Any,
        *,
        position: Any,
        number_format: str,
        font_size_pt: int = 9,
        font_name: Optional[str] = None,
        font_color: Optional[Any] = None,
    ) -> None:
        """
        Shared data-label styling contract for all chart builders.
        Keeps label rendering consistent across native visual types.
        """
        series.has_data_labels = True
        labels = series.data_labels
        labels.position = position
        labels.number_format = number_format
        labels.font.size = Pt(font_size_pt)
        labels.font.name = font_name or self.theme.FONT_BOLD
        labels.font.color.rgb = font_color or self.theme.TEXT_COLOR

    @abstractmethod
    def render(self, slide: Slide, chart_data: Dict[str, Any]) -> None:
        pass
