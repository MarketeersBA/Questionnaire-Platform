import logging
from typing import Dict, Any, List
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from .base_builder import BaseChartBuilder
from .builder_render_status import BuilderEmptyDataError

logger = logging.getLogger(__name__)

class PPTXReferenceTable(BaseChartBuilder):
    """
    Builder for Funnel Reference Tables.
    Provides a high-density tabular view of conversion stages across all brands.
    Features:
    - Native PPTX Table
    - Banded row styling for legibility
    - All-caps headers with brand-aligned typography
    """

    def render(self, slide: Slide, chart_data: Dict[str, Any]) -> None:
        data_payload = chart_data.get("data", {})
        labels = data_payload.get("labels", []) # The Stages
        datasets = data_payload.get("datasets", []) # The Brands + their data
        
        if not labels or not datasets:
            raise BuilderEmptyDataError("Reference table requires labels and datasets.")

        # 1. Define Structure
        # Columns: [Brand] + Funnel Stages
        cols = ["BRAND"] + [str(l).upper() for l in labels]
        num_rows = len(datasets) + 1
        num_cols = len(cols)
        
        # 2. Add Table
        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            self.layout.CHART_LEFT, self.layout.CHART_TOP,
            self.layout.CHART_WIDTH, self.layout.CHART_HEIGHT
        )
        table = table_shape.table
        
        # 3. Populate Header
        self._populate_header(table, cols)
        
        # 4. Populate Brand Rows
        for i, ds in enumerate(datasets):
            row_idx = i + 1
            brand_name = ds.get("label", "Brand")
            scores = ds.get("data", [])
            
            # Col 0: Brand Name
            self._fill_cell(table.cell(row_idx, 0), brand_name, align=PP_ALIGN.LEFT, is_bold=True)
            
            # Data Cols
            for j, score in enumerate(scores):
                if j + 1 >= num_cols: break
                cell = table.cell(row_idx, j + 1)
                self._fill_cell(cell, f"{score:.1f}%", align=PP_ALIGN.CENTER)
            
            # Apply Banding (Alternating row colors)
            if i % 2 == 1:
                self._apply_row_banding(table, row_idx)

    def _populate_header(self, table: Any, cols: List[str]):
        """Styles the header row with primary brand colors."""
        for j, col_name in enumerate(cols):
            cell = table.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme.PRIMARY_BRAND
            
            p = cell.text_frame.paragraphs[0]
            p.text = col_name
            p.font.name = self.theme.FONT_BOLD
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

    def _fill_cell(self, cell: Any, text: str, align: Any = PP_ALIGN.CENTER, is_bold: bool = False):
        p = cell.text_frame.paragraphs[0]
        p.text = str(text)
        p.font.name = self.theme.FONT_MEDIUM if is_bold else self.theme.FONT_LIGHT
        p.font.size = Pt(10)
        p.font.color.rgb = self.theme.TEXT_COLOR
        p.alignment = align

    def _apply_row_banding(self, table: Any, row_idx: int):
        """Applies a subtle slate fill to even rows for better visual tracking."""
        for col_idx in range(len(table.columns)):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme.get_rgb_by_name("brand_glass_blue") # Very light slate
