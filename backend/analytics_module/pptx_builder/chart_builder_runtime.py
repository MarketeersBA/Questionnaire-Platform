from __future__ import annotations

from typing import Any, Dict, Iterable, List, Optional

from pptx.enum.chart import XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION, XL_MARKER_STYLE
from pptx.util import Pt


def ensure_chart_legend(chart: Any, position: XL_LEGEND_POSITION = XL_LEGEND_POSITION.BOTTOM):
    chart.has_legend = True
    legend = chart.legend
    legend.position = position
    return legend


def set_marker_style(marker: Any, style: XL_MARKER_STYLE) -> None:
    marker.style = style


def apply_xy_point_labels(chart: Any, datasets: Iterable[Dict[str, Any]]) -> None:
    """
    Apply point labels on XY charts using the Desktop-PowerPoint-safe activation path.
    
    CRITICAL: We use `point.has_data_label = True` (safe) instead of
    `data_label.has_text_frame = True` (causes OXML corruption on desktop).
    """
    for series, dataset in zip(chart.series, datasets):
        points_meta = dataset.get("data", []) if isinstance(dataset, dict) else []
        for point, meta in zip(series.points, points_meta):
            if not isinstance(meta, dict):
                continue
            label = meta.get("label") or meta.get("name") or meta.get("attribute")
            if not label:
                continue
            try:
                point.has_data_label = True
                dl = point.data_label
                dl.position = XL_DATA_LABEL_POSITION.RIGHT
                dl.font.size = Pt(8)
                if hasattr(dl, "text_frame"):
                    dl.text_frame.text = str(label)
            except Exception:
                pass  # Silently skip malformed points


def count_chart_shapes(slide: Any) -> int:
    count = 0
    for shape in slide.shapes:
        try:
            _ = shape.chart
            count += 1
        except Exception:
            continue
    return count
