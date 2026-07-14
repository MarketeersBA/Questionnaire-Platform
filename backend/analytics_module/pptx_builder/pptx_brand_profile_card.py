import logging
from typing import Any, Dict, List, NamedTuple, Tuple
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from .base_builder import BaseChartBuilder

logger = logging.getLogger(__name__)


class KpiTile(NamedTuple):
    """Resolved KPI tile content for the brand profile scorecard grid."""

    label: str
    value: str
    descriptor: str


class PPTXBrandProfileCard(BaseChartBuilder):
    """
    Premium N-Slide Brand Profile Builder.
    Features:
    - Full-bleed Dark Navy aesthetic (#0F172A)
    - KPI Grid with glassmorphism-inspired tiles
    - Dynamic Top Strengths sidebar with score bars
    - AI Insight headline with Signal Rose accents
    """

    DARK_BG = RGBColor(15, 23, 42)    # #0F172A
    GLOW_ROSE = RGBColor(244, 63, 94)  # #F43F5E
    INDIGO = RGBColor(99, 102, 241)     # #6366F1
    KPI_CARD_BG = RGBColor(30, 41, 59) # #1E293B

    KPI_GRID_START_LEFT_IN = 0.5
    KPI_GRID_RIGHT_BOUND_IN = 9.8
    KPI_GRID_SPACING_IN = 0.3
    KPI_GRID_TOP_IN = 2.8
    KPI_CARD_HEIGHT_IN = 3.2

    _KPI_PROFILE_SPECS: Tuple[Tuple[str, str, str, str], ...] = (
        ("Overall Score", "OVERALL SCORE", "7.42 baseline", "score"),
        ("T2B %", "T2B %", "Top 2 Boxes", "percent"),
        ("NPS", "NPS", "Net Promoter Score", "signed_nps"),
        ("Evaluations", "EVALUATIONS", "Respondent Base", "integer"),
    )

    @staticmethod
    def _format_signed_nps(value: Any) -> str:
        """Format NPS as a signed integer score (e.g. +30, -10, 0)."""
        try:
            score = int(round(float(value)))
        except (TypeError, ValueError):
            return str(value)
        if score > 0:
            return f"+{score}"
        return str(score)

    @classmethod
    def _format_kpi_value(cls, profile_key: str, value: Any, value_kind: str) -> str:
        if value is None:
            return "N/A"
        if value_kind == "percent":
            return f"{value}%"
        if value_kind == "signed_nps":
            return cls._format_signed_nps(value)
        if value_kind == "integer":
            try:
                return f"{int(value):,}"
            except (TypeError, ValueError):
                return str(value)
        return str(value)

    @classmethod
    def _resolve_kpi_tiles(cls, profile: Dict[str, Any]) -> List[KpiTile]:
        """Build ordered KPI tiles from present profile keys (Brand excluded)."""
        tiles: List[KpiTile] = []
        for profile_key, label, descriptor, value_kind in cls._KPI_PROFILE_SPECS:
            if profile_key not in profile:
                continue
            tiles.append(
                KpiTile(
                    label=label,
                    value=cls._format_kpi_value(profile_key, profile.get(profile_key), value_kind),
                    descriptor=descriptor,
                )
            )
        return tiles

    @classmethod
    def _compute_kpi_grid_layout(cls, tile_count: int) -> Tuple[float, float, float, float]:
        """
        Compute responsive KPI card geometry so 3–4 tiles fit left of the sidebar.

        Returns (start_left_in, card_width_in, spacing_in, card_height_in).
        """
        count = max(tile_count, 1)
        available = cls.KPI_GRID_RIGHT_BOUND_IN - cls.KPI_GRID_START_LEFT_IN
        spacing = cls.KPI_GRID_SPACING_IN
        card_width = (available - spacing * (count - 1)) / count
        return cls.KPI_GRID_START_LEFT_IN, card_width, spacing, cls.KPI_CARD_HEIGHT_IN

    def render(self, slide: Slide, intent_data: Dict[str, Any]) -> None:
        """
        Renders a premium brand profile slide.
        intent_data structure:
        {
            "brand_data": {...},
            "ai_insight": "...",
            "brand_index": 1,
            "total_brands": 4
        }
        """
        brand_payload = intent_data.get("brand_data", {})
        data = brand_payload.get("data", {})
        profile = data.get("profile", {})
        strengths = data.get("strengths", [])
        
        brand_name = profile.get("Brand", "N/A")
        ai_insight = intent_data.get("ai_insight", "Performance analysis and strategic positioning pending further optimization.")
        brand_idx = intent_data.get("brand_index", 1)
        
        # 1. Background Reinforcement (Full-bleed Dark Navy)
        self._set_premium_background(slide)

        # 2. Section Header
        self._add_section_badge(slide, brand_idx)

        # 3. Brand Identity & AI Banner
        self._add_identity_banner(slide, brand_name, ai_insight)

        # 4. KPI Card Grid (Main Content Left)
        self._add_kpi_grid(slide, profile)

        # 5. Top Strengths Sidebar (Right)
        self._add_strengths_sidebar(slide, strengths)

    def _set_premium_background(self, slide: Slide):
        """Ensures the background is the premium corporate dark mode."""
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.layout.WIDTH, self.layout.HEIGHT
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.DARK_BG
        bg.line.visible = False
        
        # Add a subtle Rose Glow at the very top
        glow = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.layout.WIDTH, Inches(0.08)
        )
        glow.fill.solid()
        glow.fill.fore_color.rgb = self.GLOW_ROSE
        glow.line.visible = False

    def _add_section_badge(self, slide: Slide, idx: int):
        """Adds 'SECTION 01 — BRAND PROFILES' identifier."""
        label = f"SECTION {idx:02d} — BRAND PROFILES"
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(4), Inches(0.4))
        p = box.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.name = self.theme.FONT_BOLD
        p.font.color.rgb = RGBColor(148, 163, 184) # Slate 400
        p.font.bold = True

    def _add_identity_banner(self, slide: Slide, name: str, insight: str):
        """Adds Big Brand Name and the AI Insight Headline."""
        # Brand Name
        name_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(5), Inches(1.0))
        p = name_box.text_frame.paragraphs[0]
        p.text = name.upper()
        p.font.size = Pt(48)
        p.font.name = self.theme.FONT_BOLD
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # AI Insight Banner (Italic Rose)
        # Use a text frame with specific width to wrap
        insight_box = slide.shapes.add_textbox(Inches(5.0), Inches(1.15), Inches(7.8), Inches(1.0))
        p = insight_box.text_frame.paragraphs[0]
        p.text = f"\"{insight}\""
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.name = self.theme.FONT_MEDIUM
        p.font.color.rgb = self.GLOW_ROSE
        insight_box.text_frame.word_wrap = True

    def _add_kpi_grid(self, slide: Slide, profile: Dict[str, Any]):
        """Render premium KPI cards for all present profile metrics (3 or 4 tiles)."""
        tiles = self._resolve_kpi_tiles(profile)
        if not tiles:
            return

        start_left_in, card_w_in, spacing_in, card_h_in = self._compute_kpi_grid_layout(len(tiles))
        start_left = Inches(start_left_in)
        card_w = Inches(card_w_in)
        spacing = Inches(spacing_in)
        start_top = Inches(self.KPI_GRID_TOP_IN)
        card_h = Inches(card_h_in)

        for index, tile in enumerate(tiles):
            left = start_left + (index * (card_w + spacing))

            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, start_top, card_w, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = self.KPI_CARD_BG
            card.line.color.rgb = self.INDIGO
            card.line.width = Pt(1.5)

            lbl_box = slide.shapes.add_textbox(left, start_top + Inches(0.3), card_w, Inches(0.5))
            lp = lbl_box.text_frame.paragraphs[0]
            lp.text = tile.label
            lp.font.size = Pt(14)
            lp.font.name = self.theme.FONT_BOLD
            lp.font.color.rgb = RGBColor(255, 255, 255)
            lp.alignment = PP_ALIGN.CENTER

            val_box = slide.shapes.add_textbox(left, start_top + Inches(0.8), card_w, Inches(1.2))
            vp = val_box.text_frame.paragraphs[0]
            vp.text = tile.value
            vp.font.size = Pt(54)
            vp.font.bold = True
            vp.font.name = self.theme.FONT_BOLD
            vp.font.color.rgb = RGBColor(255, 255, 255)
            vp.alignment = PP_ALIGN.CENTER

            desc_box = slide.shapes.add_textbox(left, start_top + Inches(2.3), card_w, Inches(0.4))
            dp = desc_box.text_frame.paragraphs[0]
            dp.text = tile.descriptor
            dp.font.size = Pt(12)
            dp.font.name = self.theme.FONT_MEDIUM
            dp.font.color.rgb = RGBColor(148, 163, 184)
            dp.alignment = PP_ALIGN.CENTER

    def _add_strengths_sidebar(self, slide: Slide, strengths: List[Dict[str, Any]]):
        """Renders the Top Strengths sidebar with visual bars."""
        sidebar_left = Inches(9.8)
        sidebar_top = Inches(2.8)
        sidebar_w = Inches(3.0)
        sidebar_h = Inches(3.2)
        
        # Sidebar BG (Glassmorphism inspired - slightly lighter)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, sidebar_left, sidebar_top, sidebar_w, sidebar_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(51, 65, 85) # Slate 700
        bg.line.visible = False
        
        # Header
        head_box = slide.shapes.add_textbox(sidebar_left + Inches(0.2), sidebar_top + Inches(0.2), sidebar_w - Inches(0.4), Inches(0.4))
        hp = head_box.text_frame.paragraphs[0]
        hp.text = "TOP STRENGTHS"
        hp.font.size = Pt(14)
        hp.font.bold = True
        hp.font.color.rgb = self.GLOW_ROSE
        
        # Strength Items
        for i, item in enumerate(strengths[:3]):
            attr = item.get("attribute", "Attribute")
            score = item.get("score", 0)
            
            y_offset = sidebar_top + Inches(0.8) + (i * Inches(0.8))
            
            # Attribute Name
            attr_box = slide.shapes.add_textbox(sidebar_left + Inches(0.2), y_offset, sidebar_w - Inches(0.4), Inches(0.4))
            ap = attr_box.text_frame.paragraphs[0]
            ap.text = f"{attr} ({score})"
            ap.font.size = Pt(12)
            ap.font.color.rgb = RGBColor(255, 255, 255)
            
            # Score Bar (Base)
            bar_w_max = sidebar_w - Inches(0.4)
            base_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                sidebar_left + Inches(0.2),
                y_offset + Inches(0.45),
                bar_w_max,
                Inches(0.08),
            )
            base_bar.fill.solid()
            base_bar.fill.fore_color.rgb = RGBColor(30, 41, 59)
            base_bar.line.visible = False
            
            # Score Bar (Value)
            score_ratio = min(score / 10.0, 1.0)
            val_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, sidebar_left + Inches(0.2), y_offset + Inches(0.45), bar_w_max * score_ratio, Inches(0.08))
            val_bar.fill.solid()
            val_bar.fill.fore_color.rgb = self.INDIGO
            val_bar.line.visible = False
