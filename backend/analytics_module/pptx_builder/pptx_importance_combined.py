"""
PPTXImportanceCombined — Desktop-Safe Dual-Panel Scatter Builder
================================================================
Renders Main Category Drivers and Micro-Drivers side-by-side on a single slide.

OXML Compliance Strategy (Desktop PowerPoint Safe):
1. All fill operations target Marker-level, not Series-level, for XY_SCATTER charts.
2. Data labels use the safe `point.has_data_label = True` activation path.
3. Connectors use typed `MSO_CONNECTOR.STRAIGHT` enum, not raw integer `1`.
4. All shape coordinates are int-cast to EMU precision.
5. Alignment uses `PP_ALIGN` enum, not raw integers.
6. Native chart titles are disabled to prevent orphan OXML title blocks.
7. Axis titles use the safe sequential: has_title -> text_frame -> text pattern.
"""

import logging
from typing import Any, Dict, List, Tuple

from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_MARKER_STYLE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from .base_builder import BaseChartBuilder
from .builder_render_status import BuilderEmptyDataError
from .chart_builder_runtime import set_marker_style, ensure_chart_legend

logger = logging.getLogger(__name__)


class PPTXImportanceCombined(BaseChartBuilder):
    """
    Advanced Dual-Scatter Builder for Importance vs Performance Analysis.
    Renders Main Category Drivers and Micro-Drivers side-by-side.
    """

    def render(self, slide: Slide, chart_data: Dict[str, Any]) -> None:
        data_payload = chart_data.get("data", {})
        main_data = data_payload.get("main_scatter", {})
        sub_data = data_payload.get("sub_scatter", {})

        if not main_data:
            raise BuilderEmptyDataError("Combined Importance: Missing main_scatter data.")

        # 1. Resolve Geometry: side-by-side panels (55/45 split)
        left_panel, right_panel = self.layout.dual_chart_bounds(gap_in=0.5, split_ratio=0.55)

        # 2. Render Left Panel: Main Attributes
        main_datasets = main_data.get("datasets", [])
        if main_datasets:
            self._render_scatter_panel(
                slide,
                main_data,
                left_panel,
                panel_label="Overall Scatter",
                show_legend=True,
            )
        else:
            logger.warning("Combined Importance: No main_scatter datasets found.")

        # 3. Render Right Panel: Sub-Attributes (Micro-Drivers)
        if sub_data and sub_data.get("datasets"):
            self._render_scatter_panel(
                slide,
                sub_data,
                right_panel,
                panel_label="Sub Scatter",
                show_legend=False,
            )
        else:
            logger.info("Combined Importance: No sub-attribute data for right panel.")

        # 4. Global Footnote
        self._add_importance_footnote(slide)

    # ──────────────────────────────────────────────────────────────────────
    #  Panel Renderer
    # ──────────────────────────────────────────────────────────────────────

    def _render_scatter_panel(
        self,
        slide: Slide,
        panel_data: Dict[str, Any],
        bounds: Tuple[Any, Any, Any, Any],
        panel_label: str,
        show_legend: bool = True,
    ) -> None:
        left, top, width, height = bounds
        datasets = panel_data.get("datasets", [])

        # Build XyChartData
        chart_data_obj = XyChartData()
        all_x: List[float] = []
        all_y: List[float] = []

        for ds in datasets:
            label = ds.get("brand") or ds.get("label") or "Unknown"
            series = chart_data_obj.add_series(label)
            for p in ds.get("data", []):
                x = float(p.get("x", 0))
                y = float(p.get("y", 0))
                series.add_data_point(x, y)
                all_x.append(x)
                all_y.append(y)

        # Guard: Skip chart creation if no data points exist
        if not all_x:
            logger.warning("Panel '%s' has no data points — skipping chart shape.", panel_label)
            self._add_panel_title(slide, left, top, width, panel_label)
            return

        # 1. Add Aesthetic Overlays (Behind charts)
        mean_x = float(sum(all_x) / len(all_x))
        mean_y = float(sum(all_y) / len(all_y))
        self._add_quadrant_aesthetics(slide, bounds, mean_x, mean_y)

        # 2. Add Chart Shape — all coordinates cast to int (EMU)
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER,
            int(left), int(top), int(width), int(height),
            chart_data_obj,
        )
        chart = chart_shape.chart

        # Disable native chart title to avoid orphan OXML blocks
        chart.has_title = False

        # Style axes
        self._style_panel_axes(chart)

        # Style series markers (marker-level, not series-level)
        self._style_panel_series(chart, datasets)

        # Crosshair quadrant lines
        if all_x and all_y:
            self._add_crosshairs(slide, chart_shape, mean_x, mean_y)

        # Hero-brand data labels
        self._apply_hero_labels(chart, datasets)

        # Legend
        if show_legend:
            legend = ensure_chart_legend(chart, XL_LEGEND_POSITION.BOTTOM)
            legend.font.size = Pt(9)
            legend.font.name = self.theme.FONT_LIGHT
        else:
            chart.has_legend = False

        # Panel title (textbox above chart)
        self._add_panel_title(slide, left, top, width, panel_label)

    # ──────────────────────────────────────────────────────────────────────
    #  Axes
    # ──────────────────────────────────────────────────────────────────────

    def _style_panel_axes(self, chart: Any) -> None:
        # X-Axis: Importance (Correlation × 100)
        x_axis = chart.value_axis
        x_axis.minimum_scale = 0
        x_axis.maximum_scale = 100
        x_axis.has_title = True
        x_axis.axis_title.text_frame.text = "Importance*"
        x_axis.has_major_gridlines = False
        x_axis.tick_labels.font.size = Pt(8)
        x_axis.tick_labels.font.name = self.theme.FONT_LIGHT

        # Y-Axis: Performance (T2B%)
        y_axis = chart.category_axis
        y_axis.minimum_scale = 0
        y_axis.maximum_scale = 100
        y_axis.has_title = True
        y_axis.axis_title.text_frame.text = "Performance (T2B%)"
        y_axis.has_major_gridlines = False
        y_axis.tick_labels.font.size = Pt(8)
        y_axis.tick_labels.font.name = self.theme.FONT_LIGHT

    # ──────────────────────────────────────────────────────────────────────
    #  Series Styling (Marker-Level — Desktop Safe)
    # ──────────────────────────────────────────────────────────────────────

    def _style_panel_series(self, chart: Any, datasets: List[Dict[str, Any]]) -> None:
        """All fill/color is applied at Marker level for OXML stability."""
        for i, series in enumerate(chart.series):
            brand_name = (series.name or "").lower()
            is_hero = (
                any(kw in brand_name for kw in ("hero", "primary", "total"))
                or (i == 0 and len(chart.series) > 1)
            )

            # Marker size
            series.marker.size = 12 if is_hero else 9

            # Marker shape
            if is_hero:
                set_marker_style(series.marker, XL_MARKER_STYLE.TRIANGLE)
            else:
                shape = XL_MARKER_STYLE.SQUARE if i % 2 != 0 else XL_MARKER_STYLE.CIRCLE
                set_marker_style(series.marker, shape)

            # Marker fill (MUST be marker-level for XY_SCATTER OXML compliance)
            series.marker.format.fill.solid()
            color = self.theme.PRIMARY_BRAND if is_hero else self.theme.get_color(i)
            series.marker.format.fill.fore_color.rgb = color

            # Marker outline off
            series.marker.format.line.visible = False

            # Series connector line off (scatter should show points only)
            series.format.line.visible = False

    # ──────────────────────────────────────────────────────────────────────
    #  Hero Labels (Safe Activation Path)
    # ──────────────────────────────────────────────────────────────────────

    def _apply_hero_labels(self, chart: Any, datasets: List[Dict[str, Any]]) -> None:
        """Attribute labels shown only on the hero (first) series for clarity."""
        for series, ds in zip(chart.series, datasets):
            brand_name = (ds.get("brand") or ds.get("label") or "").lower()
            is_hero = any(kw in brand_name for kw in ("hero", "primary", "total"))
            if not is_hero:
                continue

            points_meta = ds.get("data", [])
            for point, meta in zip(series.points, points_meta):
                if not isinstance(meta, dict):
                    continue
                label_text = meta.get("sub_attribute") or meta.get("attribute")
                if not label_text:
                    continue
                # Safe activation: flag first, then format
                try:
                    point.has_data_label = True
                    dl = point.data_label
                    dl.font.size = Pt(7)
                    dl.font.name = self.theme.FONT_LIGHT
                    dl.font.color.rgb = self.theme.get_rgb_by_name("brand_slate")
                    if hasattr(dl, "text_frame"):
                        dl.text_frame.text = str(label_text)
                except Exception:
                    logger.debug("Label application skipped for point.", exc_info=True)

    def _add_quadrant_aesthetics(self, slide: Slide, bounds: Tuple[int, int, int, int], mean_x: float, mean_y: float) -> None:
        """Adds subtle quadrant overlays and labels for strategic context."""
        left, top, width, height = [int(v) for v in bounds]
        
        # Calculate visual center based on mean (0-100 scale normalization)
        frac_x = max(0.1, min(0.9, mean_x / 100.0))
        frac_y = max(0.1, min(0.9, (100.0 - mean_y) / 100.0))

        cw = int(width * frac_x)
        ch = int(height * frac_y)
        rem_w = width - cw
        rem_h = height - ch

        quads = [
            # TR: Pillars
            (left + cw, top, rem_w, ch, "brand_glass_blue", "CORE STRENGTHS"),
            # TL: Hidden Gems
            (left, top, cw, ch, "brand_light_gray", "HIDDEN GEMS"),
            # BR: Risks
            (left + cw, top + ch, rem_w, rem_h, "brand_light_gray", "VULNERABILITIES"),
            # BL: Low Priority
            (left, top + ch, cw, rem_h, "brand_glass_blue", "SECONDARY"),
        ]
        
        for qx, qy, qw, qh, color_key, label in quads:
            if qw <= 0 or qh <= 0: continue
            # 1. Background Rectangle
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, qx, qy, qw, qh)
            rect.fill.solid()
            rect.fill.fore_color.rgb = self.theme.get_rgb_by_name(color_key)
            rect.line.visible = False
            
            # 2. Quadrant Label (Small, subtle text in corner)
            tx = slide.shapes.add_textbox(qx + int(Inches(0.05)), qy + int(Inches(0.05)), qw, Pt(12)) 
            p = tx.text_frame.paragraphs[0]
            p.text = label
            p.font.size = Pt(7)
            p.font.name = self.theme.FONT_BOLD
            p.font.color.rgb = self.theme.get_rgb_by_name("brand_slate")
            p.alignment = PP_ALIGN.LEFT

    # ──────────────────────────────────────────────────────────────────────
    #  Crosshairs (Typed Connectors)
    # ──────────────────────────────────────────────────────────────────────

    def _add_crosshairs(self, slide: Slide, chart_shape: Any, mean_x: float, mean_y: float) -> None:
        """Quadrant crosshairs using MSO_CONNECTOR.STRAIGHT with int-cast EMU coords."""
        left = int(chart_shape.left)
        top = int(chart_shape.top)
        width = int(chart_shape.width)
        height = int(chart_shape.height)

        frac_x = max(0.0, min(1.0, mean_x / 100.0))
        frac_y = max(0.0, min(1.0, (100.0 - mean_y) / 100.0))

        cx = int(left + width * frac_x)
        cy = int(top + height * frac_y)

        # Vertical
        v = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, cx, top, cx, top + height)
        v.line.color.rgb = self.theme.get_rgb_by_name("brand_slate")
        v.line.width = Pt(1.5)
        v.line.dash_style = 1 # Dashed

        # Horizontal
        h = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, cy, left + width, cy)
        h.line.color.rgb = self.theme.get_rgb_by_name("brand_slate")
        h.line.width = Pt(1.5)
        h.line.dash_style = 1 # Dashed

    # ──────────────────────────────────────────────────────────────────────
    #  Panel Title (PP_ALIGN enum — not raw int)
    # ──────────────────────────────────────────────────────────────────────

    def _add_panel_title(self, slide: Slide, left: Any, top: Any, width: Any, text: str) -> None:
        title_box = slide.shapes.add_textbox(
            int(left),
            int(int(top) - int(Inches(0.4))),
            int(width),
            int(Inches(0.3)),
        )
        tf = title_box.text_frame
        tf.text = text
        p = tf.paragraphs[0]
        p.font.name = self.theme.FONT_BOLD
        p.font.size = Pt(11)
        p.font.color.rgb = self.theme.SUBTITLE_COLOR
        p.alignment = PP_ALIGN.CENTER

    # ──────────────────────────────────────────────────────────────────────
    #  Footnote
    # ──────────────────────────────────────────────────────────────────────

    def _add_importance_footnote(self, slide: Slide) -> None:
        left = int(self.layout.FOOTNOTE_LEFT)
        top = int(self.layout.FOOTNOTE_TOP) + int(self.layout.FOOTNOTE_HEIGHT) - int(Inches(0.1))
        width = int(self.layout.FOOTNOTE_WIDTH)
        box_height = int(Inches(0.2))

        shape = slide.shapes.add_textbox(left, top, width, box_height)
        shape.text_frame.text = "* Importance was derived using correlation analysis against overall likeness."
        p = shape.text_frame.paragraphs[0]
        p.font.size = Pt(7)
        p.font.italic = True
        p.font.name = self.theme.FONT_LIGHT
        p.font.color.rgb = self.theme.get_rgb_by_name("brand_slate")
