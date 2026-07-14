import logging
from typing import Dict, Any, List
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from .theme import PPTXTheme
from .layout import PPTXLayout
from .narrative_pagination import chunk_sequence

logger = logging.getLogger(__name__)

class PPTXStrategicCommand:
    """
    Advanced Strategic Intelligence Slide Builder.
    Replicates the 'Decision-Maker's Command Center' UI from the web report.
    Custom full-bleed layout with dark-mode aesthetics.
    """

    def __init__(self, theme: PPTXTheme, layout: PPTXLayout):
        self.theme = theme
        self.layout = layout
        
        # Strategic Color Palette (Tailored for high-impact dark mode)
        self.DARK_BG = RGBColor(15, 23, 42)      # #0F172A (Deep Navy/Slate)
        self.ROSE = RGBColor(244, 63, 94)         # #F43F5E (Signal Rose)
        self.INDIGO = RGBColor(99, 102, 241)      # #6366F1 (Strategic Indigo)
        self.SLATE_TEXT = RGBColor(148, 163, 184) # #94A3B8 (Metadata Slate)
        self.WHITE = RGBColor(255, 255, 255)
        self.CARD_BG = RGBColor(30, 41, 59)       # #1E293B (Slightly lighter slate)

    def build(self, slide: Slide, insights: List[Dict[str, Any]]) -> Slide:
        """
        Orchestrates the creation of the Strategic Command Center slide.
        Expects one or two opportunity insight payloads per slide.
        """
        self._add_background(slide)
        self._add_accent_glow(slide)
        self._add_slide_header(slide)

        import os
        from pathlib import Path

        logo_path = Path(os.environ.get("ANALYTICS_RESOURCES_DIR", "backend/resources/analytics")) / "logo.png"
        if logo_path.exists():
            slide.shapes.add_picture(
                str(logo_path),
                self.layout.WIDTH - Inches(2.5),
                Inches(0.5),
                height=Inches(0.6),
            )

        current_top = Inches(2.7)
        card_height = Inches(3.9)
        spacing = Inches(0.3)

        for index, insight in enumerate(insights[:2]):
            top_pos = current_top + (index * (card_height + spacing))
            self._render_insight_card(slide, top_pos, card_height, insight)

        return slide

    def build_deck(self, presentation, template_adapter, insights: List[Dict[str, Any]]) -> int:
        if not insights:
            return 0

        rendered = 0
        for chunk in chunk_sequence(insights, 2):
            slide = presentation.slides.add_slide(
                template_adapter.get_layout(presentation, "ai_narrative")
            )
            self.build(slide, chunk)
            rendered += 1
        return rendered

    def _add_background(self, slide: Slide):
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.layout.WIDTH, self.layout.HEIGHT
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.DARK_BG
        bg.line.visible = False

    def _add_accent_glow(self, slide: Slide):
        # Top gradient-like line
        glow = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.layout.WIDTH, Inches(0.05)
        )
        glow.fill.solid()
        glow.fill.fore_color.rgb = self.ROSE # We use rose as the primary signal
        glow.line.visible = False

    def _add_slide_header(self, slide: Slide):
        # 1. Icon Placeholder (Zap/Bolt equivalent) - Positioned relative to CHART_LEFT
        icon_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, self.layout.CHART_LEFT, Inches(0.75), Inches(0.6), Inches(0.6)
        )
        icon_box.fill.solid()
        icon_box.fill.fore_color.rgb = self.ROSE
        icon_box.line.visible = False
        
        # 2. Main Title Group
        # "STRATEGIC INTELLIGENCE" (Small & uppercase)
        t1 = slide.shapes.add_textbox(self.layout.CHART_LEFT + Inches(0.75), Inches(0.67), Inches(6), Inches(0.45))
        p1 = t1.text_frame.paragraphs[0]
        p1.text = "STRATEGIC INTELLIGENCE"
        p1.font.name = self.theme.FONT_BOLD
        p1.font.size = Pt(15)
        p1.font.color.rgb = self.ROSE
        
        # "Decision-Maker's Command Center"
        t2 = slide.shapes.add_textbox(self.layout.CHART_LEFT + Inches(0.75), Inches(1.05), Inches(12), Inches(0.9))
        p2 = t2.text_frame.paragraphs[0]
        p2.text = "Decision-Maker's Command Center"
        p2.font.name = self.theme.FONT_BOLD
        p2.font.size = Pt(42)
        p2.font.color.rgb = self.WHITE
        
        # 3. Subtext Description
        t3 = slide.shapes.add_textbox(self.layout.CHART_LEFT + Inches(0.75), Inches(1.87), Inches(12), Inches(0.45))
        p3 = t3.text_frame.paragraphs[0]
        p3.text = "High-impact strategic playbooks derived from performance gaps and consumer undercurrents."
        p3.font.name = self.theme.FONT_MEDIUM
        p3.font.size = Pt(18)
        p3.font.color.rgb = self.SLATE_TEXT

    def _render_insight_card(self, slide: Slide, top: float, height: float, insight: Dict[str, Any]):
        """Renders one horizontal instruction block (Left Signal + Right Playbook)."""
        left_margin = self.layout.CHART_LEFT
        width = self.layout.WIDTH - (left_margin * 2)
        
        # 1. Card Container (Slightly lighter than background)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left_margin, top, width, height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.CARD_BG
        card.line.color.rgb = self.INDIGO
        card.line.width = Pt(1.5)
        
        # Set radius safely
        if hasattr(card, "adjustments"):
            card.adjustments[0] = 5000 
        
        # --- COLUMN GEOMETRY ---
        # We use a 40/60 split for Signal vs Playbook
        col_gap = Inches(0.6)
        left_col_w = (width * 0.42) - col_gap
        right_col_w = (width * 0.58) - col_gap
        
        # --- LEFT COLUMN: SIGNAL INTELLIGENCE ---
        self._add_left_signal_content(slide, left_margin + Inches(0.45), top + Inches(0.4), 
                                    left_col_w, height - Inches(0.75), insight)
        
        # --- RIGHT COLUMN: TACTICAL PLAYBOOK ---
        right_start = left_margin + left_col_w + col_gap
        self._add_right_playbook_content(slide, right_start, top + Inches(0.4), 
                                       right_col_w, height - Inches(0.75), insight)

    def _add_left_signal_content(self, slide: Slide, left: float, top: float, width: float, height: float, insight: Dict[str, Any]):
        # 1. Category Badge
        cat = str(insight.get("strategic_category", "Product")).upper()
        badge_color = self.INDIGO
        if any(kw in cat for kw in ["MARKET", "GROWTH", "EXPANSION"]): 
            badge_color = RGBColor(16, 185, 129) # Emerald
        
        self._draw_badge(slide, left, top, f"{cat} FOCUS", badge_color)
        
        # 2. Title (With Wrap)
        title_box = slide.shapes.add_textbox(left, top + Inches(0.67), width, Inches(0.9))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(insight.get("title", "Insight Summary")).upper()
        p.font.name = self.theme.FONT_BOLD
        p.font.size = Pt(21)
        p.font.color.rgb = self.WHITE
        p.font.line_spacing = 0.95
        
        # 3. Insight Text (With Wrap)
        text_box = slide.shapes.add_textbox(left, top + Inches(1.57), width, Inches(1.2))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(insight.get("insight", ""))
        p.font.name = self.theme.FONT_MEDIUM
        p.font.size = Pt(14)
        p.font.color.rgb = self.SLATE_TEXT
        
        # 4. Decision Matrix Grid
        grid_top = top + height - Inches(1.05)
        mini_w = (width / 2) - Inches(0.15)
        
        # Gap Magnitude
        self._draw_mini_metric(slide, left, grid_top, mini_w, Inches(1.0), 
                             "PERFORMANCE GAP", f"-{insight.get('gap_magnitude', 0):.1f}", self.ROSE)
        
        # Business Impact
        impact_label = str(insight.get("impact", "Medium")).upper()
        impact_color = self.ROSE if "HIGH" in impact_label else RGBColor(245, 158, 11)
        self._draw_mini_metric(slide, left + mini_w + Inches(0.3), grid_top, mini_w, Inches(1.0), 
                             "BUSINESS IMPACT", impact_label, impact_color)

    def _add_right_playbook_content(self, slide: Slide, left: float, top: float, width: float, height: float, insight: Dict[str, Any]):
        # 1. Header
        header_box = slide.shapes.add_textbox(left, top, width, Inches(0.6))
        p = header_box.text_frame.paragraphs[0]
        p.text = "EXECUTION PLAYBOOK"
        p.font.name = self.theme.FONT_BOLD
        p.font.size = Pt(21)
        p.font.color.rgb = self.WHITE
        
        # 2. Action Items (Numbered list)
        actions = insight.get("actions", [])[:3]
        action_top_start = top + Inches(0.75)
        item_height = Inches(0.82) # Scaled spacing
        
        for i, action_obj in enumerate(actions):
            item_top = action_top_start + (i * item_height)
            action_text = action_obj.get("action") if isinstance(action_obj, dict) else str(action_obj)
            
            # Number circle identity
            num_box = slide.shapes.add_textbox(left, item_top, Inches(0.6), item_height)
            p_num = num_box.text_frame.paragraphs[0]
            p_num.text = f"0{i+1}"
            p_num.font.name = self.theme.FONT_BOLD
            p_num.font.size = Pt(15)
            p_num.font.color.rgb = self.INDIGO
            
            # Action content (Crucial: Wrap enabled)
            act_box = slide.shapes.add_textbox(left + Inches(0.75), item_top, width - Inches(0.75), item_height)
            tf = act_box.text_frame
            tf.word_wrap = True
            p_act = tf.paragraphs[0]
            p_act.text = action_text
            p_act.font.name = self.theme.FONT_MEDIUM
            p_act.font.size = Pt(14)
            p_act.font.color.rgb = self.WHITE
            p_act.font.line_spacing = 1.0

        # 3. Recommendation Banner (Stick to bottom of height)
        banner_top = top + height - Inches(0.675)
        banner = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, banner_top, width, Inches(0.675)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = RGBColor(67, 56, 202) # Deeper indigo
        banner.line.visible = False
        
        t_banner = slide.shapes.add_textbox(left, banner_top + Inches(0.12), width, Inches(0.45))
        tf = t_banner.text_frame
        p_b = tf.paragraphs[0]
        p_b.text = f"PRIORITY RECOMMENDATION: EXECUTE STEP 01 FOR MAXIMUM {insight.get('attribute', 'BRAND').upper()} LIFT."
        p_b.font.name = self.theme.FONT_BOLD
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = self.WHITE
        p_b.alignment = PP_ALIGN.CENTER


    def _draw_badge(self, slide: Slide, left: float, top: float, text: str, color: RGBColor):
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.8), Inches(0.375)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.visible = False
        
        t = slide.shapes.add_textbox(left, top, Inches(1.8), Inches(0.375))
        p = t.text_frame.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.theme.FONT_BOLD
        p.font.size = Pt(10)
        p.font.color.rgb = self.WHITE
        
    def _draw_mini_metric(self, slide: Slide, left: float, top: float, width: float, height: float, label: str, value: str, color: RGBColor):
        # Card
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = self.DARK_BG # Inside card use dark back
        box.line.color.rgb = self.SLATE_TEXT
        box.line.width = Pt(0.3)
        
        # Label
        t_label = slide.shapes.add_textbox(left, top + Inches(0.15), width, Inches(0.375))
        p_l = t_label.text_frame.paragraphs[0]
        p_l.text = label
        p_l.alignment = PP_ALIGN.CENTER
        p_l.font.name = self.theme.FONT_BOLD
        p_l.font.size = Pt(10)
        p_l.font.color.rgb = self.SLATE_TEXT
        
        # Value
        t_val = slide.shapes.add_textbox(left, top + Inches(0.45), width, Inches(0.6))
        p_v = t_val.text_frame.paragraphs[0]
        p_v.text = value
        p_v.alignment = PP_ALIGN.CENTER
        p_v.font.name = self.theme.FONT_BOLD
        p_v.font.size = Pt(30)
        p_v.font.color.rgb = color
