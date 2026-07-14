from __future__ import annotations

from collections import Counter
from typing import Any, Dict, List

from pptx import Presentation

from .validation_gating import issue


BOUNDS_TOLERANCE_EMU = 150


def audit_presentation_geometry(presentation: Presentation) -> Dict[str, Any]:
    """Inspect rendered slide geometry and title duplication."""
    layout_warnings: List[Dict[str, Any]] = []
    issues: List[Dict[str, Any]] = []
    duplicate_title_count = 0

    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    title_counts: Counter[str] = Counter()

    for slide_index, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        if title_shape is not None and getattr(title_shape, "text", ""):
            normalized_title = title_shape.text.strip().upper()
            if normalized_title:
                title_counts[normalized_title] += 1

        for shape_index, shape in enumerate(slide.shapes):
            try:
                left = int(shape.left or 0)
                top = int(shape.top or 0)
                width = int(shape.width or 0)
                height = int(shape.height or 0)
            except Exception:
                continue

            right = left + width
            bottom = top + height
            if (
                left < -BOUNDS_TOLERANCE_EMU
                or top < -BOUNDS_TOLERANCE_EMU
                or right > slide_width + BOUNDS_TOLERANCE_EMU
                or bottom > slide_height + BOUNDS_TOLERANCE_EMU
            ):
                warning = {
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                    "slide_width": slide_width,
                    "slide_height": slide_height,
                }
                layout_warnings.append(warning)
                issues.append(
                    issue(
                        "layout_out_of_bounds",
                        (
                            f"Shape {shape_index} on slide {slide_index} extends outside the slide canvas "
                            f"({right}x{bottom} vs {slide_width}x{slide_height})."
                        ),
                        severity="warning",
                        slide_index=slide_index,
                    )
                )

    for title, count in title_counts.items():
        if count > 1:
            duplicate_title_count += count - 1
            issues.append(
                issue(
                    "duplicate_title",
                    f"Duplicate slide title detected: '{title}' appears {count} times.",
                    severity="warning",
                )
            )

    return {
        "layout_warnings": layout_warnings,
        "layout_warning_count": len(layout_warnings),
        "duplicate_title_count": duplicate_title_count,
        "issues": issues,
    }
