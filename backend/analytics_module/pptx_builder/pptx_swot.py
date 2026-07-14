from __future__ import annotations

import logging
from typing import Any, Dict, List

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .layout import PPTXLayout
from .narrative_pagination import chunk_sequence
from .template_adapter import TemplateAdapter
from .theme import PPTXTheme

logger = logging.getLogger(__name__)


class PPTXSwot:
    """Builder for competitive SWOT analysis slides with pagination."""

    QUADRANT_KEYS = ("strengths", "weaknesses", "opportunities", "threats")
    ITEMS_PER_QUADRANT = 4

    def __init__(self, theme: PPTXTheme, layout: PPTXLayout):
        self.theme = theme
        self.layout = layout
        self.COLORS = self.theme.get_swot_palette()

    def build_deck(self, presentation, template_adapter: TemplateAdapter, data: Dict[str, Any]) -> int:
        brand = data.get("brand", "TARGET BRAND")
        swot = data.get("swot", {})
        pages = self._paginate_swot(swot)
        rendered = 0

        for index, page in enumerate(pages, start=1):
            slide = presentation.slides.add_slide(
                template_adapter.get_layout(presentation, "ai_narrative")
            )
            title = f"COMPETITIVE SWOT: {brand.upper()}"
            if len(pages) > 1:
                title = f"{title} ({index}/{len(pages)})"
            if slide.placeholders:
                slide.placeholders[0].text = title
            self._render_page(slide, page)
            rendered += 1
        return rendered

    def build(self, slide, data: Dict[str, Any]):
        swot = data.get("swot", {})
        self._render_page(slide, swot)

    def _paginate_swot(self, swot: Dict[str, List[str]]) -> List[Dict[str, List[str]]]:
        chunked = {
            key: chunk_sequence([str(item) for item in swot.get(key, [])], self.ITEMS_PER_QUADRANT)
            for key in self.QUADRANT_KEYS
        }
        page_count = max((len(chunks) for chunks in chunked.values()), default=1)
        if page_count == 0:
            page_count = 1

        pages: List[Dict[str, List[str]]] = []
        for page_index in range(page_count):
            pages.append(
                {
                    key: chunked[key][page_index] if page_index < len(chunked[key]) else []
                    for key in self.QUADRANT_KEYS
                }
            )
        return pages

    def _render_page(self, slide, swot: Dict[str, List[str]]) -> None:
        margin = self.layout.CHART_LEFT
        spacing = Inches(0.25)
        quad_w = (self.layout.WIDTH - (margin * 2) - spacing) / 2
        quad_h = (self.layout.HEIGHT - Inches(2.0) - spacing) / 2
        start_top = Inches(1.5)

        coords = [
            ("strengths", margin, start_top),
            ("weaknesses", margin + quad_w + spacing, start_top),
            ("opportunities", margin, start_top + quad_h + spacing),
            ("threats", margin + quad_w + spacing, start_top + quad_h + spacing),
        ]

        for key, x, y in coords:
            self._draw_swot_quadrant(slide, x, y, quad_w, quad_h, key, swot.get(key, []))

    def _draw_swot_quadrant(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        key: str,
        items: List[str],
    ):
        color = self.COLORS.get(key, self.theme.PRIMARY_BRAND)

        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        frame.fill.solid()
        frame.fill.fore_color.rgb = self.theme.get_rgb_by_name("white")
        frame.line.color.rgb = color
        frame.line.width = Pt(2.0)

        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.5))
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.visible = False

        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
        title_box.text_frame.paragraphs[0].text = key.upper()
        title_box.text_frame.paragraphs[0].font.name = self.theme.FONT_BOLD
        title_box.text_frame.paragraphs[0].font.size = Pt(14)
        title_box.text_frame.paragraphs[0].font.color.rgb = self.theme.get_rgb_by_name("white")
        title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        content = slide.shapes.add_textbox(
            left + Inches(0.2),
            top + Inches(0.6),
            width - Inches(0.4),
            height - Inches(0.7),
        )
        tf = content.text_frame
        tf.word_wrap = True

        if not items:
            tf.paragraphs[0].text = "No critical data identified for this segment."
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.italic = True
            tf.paragraphs[0].font.color.rgb = self.theme.SUBTITLE_COLOR
            return

        for index, item in enumerate(items):
            paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
            paragraph.text = f"• {item}"
            paragraph.font.name = self.theme.FONT_MEDIUM
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = self.theme.TEXT_COLOR
            paragraph.space_after = Pt(2)
