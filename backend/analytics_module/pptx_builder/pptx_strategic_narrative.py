import logging
from typing import Dict, Any, List, Optional
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from .theme import PPTXTheme
from .layout import PPTXLayout

logger = logging.getLogger(__name__)

class PPTXStrategicNarrative:
    """
    Builder for the premium Strategic Narrative slide.
    
    Layout: Unified 3-Section Slide
    1. Strategic Narrative (Phase framing)
    2. Executive Summary (Core insights)
    3. Broad Observations & Business Question (Strategic roadmap footer)
    
    Aesthetics: High-impact dark mode with rose & indigo accents.
    """

    def __init__(self, theme: PPTXTheme, layout: PPTXLayout):
        self.theme = theme
        self.layout = layout
        
        # Design Tokens (Alignment with Strategic Command)
        self.DARK_BG = RGBColor(15, 23, 42)      # Deep Navy/Slate
        self.ROSE = RGBColor(244, 63, 94)         # Signal Rose (Accent)
        self.INDIGO = RGBColor(99, 102, 241)     # Strategic Indigo
        self.EMERALD = RGBColor(16, 185, 129)    # Emerald (Success/Growth)
        self.WHITE = RGBColor(255, 255, 255)
        self.SLATE_TEXT = RGBColor(148, 163, 184) # Metadata Slate
        self.CARD_BG = RGBColor(30, 41, 59)       # Card background

    def build(self, slide: Slide, data: Dict[str, Any]) -> Slide:
        """Entry point for rendering the unified slide."""
        insights = data.get("insights", {})
        metadata = data.get("metadata", {})
        
        # 1. Apply Branded Canvas
        self._apply_background(slide)
        self._apply_accent_glow(slide)
        self._add_slide_header(slide, metadata)

        # 2. Section 1: Strategic Narrative (Top)
        narrative_top = Inches(2.3)
        narrative_height = Inches(1.2)
        self._render_strategic_narrative(slide, narrative_top, narrative_height, insights, metadata)

        # 3. Section 2: Executive Summary (Middle)
        summary_top = narrative_top + narrative_height + Inches(0.4)
        summary_height = Inches(3.2)
        self._render_executive_summary(slide, summary_top, summary_height, insights)

        # 4. Section 3: Broad Observations & Business Question (Bottom)
        footer_top = summary_top + summary_height + Inches(0.4)
        footer_height = Inches(3.2)
        self._render_footer_strategic_block(slide, footer_top, footer_height, insights, metadata)

        return slide

    def _apply_background(self, slide: Slide):
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.layout.WIDTH, self.layout.HEIGHT
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.DARK_BG
        bg.line.visible = False

    def _apply_accent_glow(self, slide: Slide):
        glow = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.layout.WIDTH, Inches(0.06)
        )
        glow.fill.solid()
        glow.fill.fore_color.rgb = self.ROSE
        glow.line.visible = False

    def _add_slide_header(self, slide: Slide, metadata: Dict[str, Any]):
        left = self.layout.CHART_LEFT
        
        # Branded Icon
        icon_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, Inches(0.75), Inches(0.6), Inches(0.6)
        )
        icon_box.fill.solid()
        icon_box.fill.fore_color.rgb = self.ROSE
        icon_box.line.visible = False

        # Main Title
        t1 = slide.shapes.add_textbox(left + Inches(0.75), Inches(0.65), Inches(8), Inches(0.5))
        p1 = t1.text_frame.paragraphs[0]
        p1.text = "STRATEGIC ARCHITECTURE & NARRATIVE"
        p1.font.name = self.theme.FONT_BOLD
        p1.font.size = Pt(16)
        p1.font.color.rgb = self.ROSE
        
        t2 = slide.shapes.add_textbox(left + Inches(0.75), Inches(1.05), Inches(12), Inches(1.0))
        p2 = t2.text_frame.paragraphs[0]
        p2.text = f"Executive Summary: {metadata.get('category', 'Market Analysis')}"
        p2.font.name = self.theme.FONT_BOLD
        p2.font.size = Pt(40)
        p2.font.color.rgb = self.WHITE

    def _render_strategic_narrative(self, slide: Slide, top: float, height: float, insights: Dict[str, Any], metadata: Dict[str, Any]):
        """Renders the framing section: Strategic Narrative."""
        left = self.layout.CHART_LEFT
        width = self.layout.WIDTH - (left * 2)

        # Fallback synthesis
        narrative_text = insights.get("strategic_narrative")
        if not narrative_text:
            brand = metadata.get("brand", "the client brand")
            category = metadata.get("category", "the category")
            research_type = metadata.get("research_type", "strategic")
            narrative_text = f"This {research_type} engagement synthesizes consumer sentiment and competitive market signals to architect a growth roadmap for {brand} within {category}."

        # Container
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.CARD_BG
        card.line.color.rgb = self.INDIGO
        card.line.width = Pt(1.5)

        # Narrative Header Badge
        badge_w = Inches(2.2)
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.3), top + Inches(0.25), badge_w, Inches(0.35)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = self.INDIGO
        badge.line.visible = False
        
        t_badge = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.25), badge_w, Inches(0.35))
        p_badge = t_badge.text_frame.paragraphs[0]
        p_badge.text = "STRATEGIC NARRATIVE"
        p_badge.font.name = self.theme.FONT_BOLD
        p_badge.font.size = Pt(10)
        p_badge.font.color.rgb = self.WHITE
        p_badge.alignment = PP_ALIGN.CENTER

        # Content Text
        t_content = slide.shapes.add_textbox(left + badge_w + Inches(0.6), top + Inches(0.25), width - badge_w - Inches(0.9), height - Inches(0.5))
        tf = t_content.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = narrative_text
        p.font.name = self.theme.FONT_MEDIUM
        p.font.size = Pt(16)
        p.font.color.rgb = self.WHITE
        p.font.italic = True

    def _render_executive_summary(self, slide: Slide, top: float, height: float, insights: Dict[str, Any]):
        """Renders the core Executive Summary block."""
        left = self.layout.CHART_LEFT
        width = self.layout.WIDTH - (left * 2)

        summary_text = insights.get("executive_summary") or insights.get("summary", "Analysis synthesis in progress. Data indicates significant market shifts and consumer alignment opportunities.")

        # Container (No border, just clear typography)
        t_box = slide.shapes.add_textbox(left, top, width, height)
        tf = t_box.text_frame
        tf.word_wrap = True
        
        # Label
        p_label = tf.paragraphs[0]
        p_label.text = "EXECUTIVE SUMMARY"
        p_label.font.name = self.theme.FONT_BOLD
        p_label.font.size = Pt(12)
        p_label.font.color.rgb = self.SLATE_TEXT
        
        # Body
        p_body = tf.add_paragraph()
        p_body.text = summary_text
        p_body.font.name = self.theme.FONT_MEDIUM
        p_body.font.size = Pt(20)
        p_body.font.color.rgb = self.WHITE
        p_body.line_spacing = 1.1

    def _render_footer_strategic_block(self, slide: Slide, top: float, height: float, insights: Dict[str, Any], metadata: Dict[str, Any]):
        """Renders phase 3: Phased Strategic Plan & Business Question."""
        left = self.layout.CHART_LEFT
        width = self.layout.WIDTH - (left * 2)

        # Split width: 70% Plan, 30% Business Question
        plan_w = width * 0.68
        bq_w = width * 0.28
        spacing = width * 0.04

        # --- A. PHASED STRATEGIC PLAN ---
        # Header with systematic numbering
        t_h = slide.shapes.add_textbox(left, top - Inches(0.1), plan_w, Inches(0.4))
        p_h = t_h.text_frame.paragraphs[0]
        p_h.text = "PHASED STRATEGIC ROADMAP"
        p_h.font.name = self.theme.FONT_BOLD
        p_h.font.size = Pt(13)
        p_h.font.color.rgb = self.ROSE

        # Phase Data Extraction
        raw_observations = insights.get("broad_observations") or insights.get("key_findings") or []
        phases = self._derive_roadmap_phases(raw_observations, metadata)
        
        card_w = (plan_w - Inches(0.4)) / 3
        for i, phase in enumerate(phases):
            p_left = left + (i * (card_w + Inches(0.2)))
            
            # Phase Container (Subtle Glass Effect)
            c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p_left, top + Inches(0.4), card_w, height - Inches(0.5))
            c.fill.solid()
            c.fill.fore_color.rgb = self.CARD_BG
            c.line.color.rgb = self.INDIGO
            c.line.width = Pt(1.2)

            # Phase Tab / Badge
            phase_num = f"PHASE 0{i+1}"
            tb_p = slide.shapes.add_textbox(p_left + Inches(0.15), top + Inches(0.55), card_w - Inches(0.3), Inches(0.3))
            pb_p = tb_p.text_frame.paragraphs[0]
            pb_p.text = phase_num
            pb_p.font.name = self.theme.FONT_BOLD
            pb_p.font.size = Pt(9)
            pb_p.font.color.rgb = self.ROSE

            # Phase Title
            tt_p = slide.shapes.add_textbox(p_left + Inches(0.15), top + Inches(0.85), card_w - Inches(0.3), Inches(0.5))
            tf_t = tt_p.text_frame
            tf_t.word_wrap = True
            pt_p = tf_t.paragraphs[0]
            pt_p.text = phase.get("title", "Initiative").upper()
            pt_p.font.name = self.theme.FONT_BOLD
            pt_p.font.size = Pt(13)
            pt_p.font.color.rgb = self.WHITE

            # Phase Description
            td_p = slide.shapes.add_textbox(p_left + Inches(0.15), top + Inches(1.4), card_w - Inches(0.3), height - Inches(1.6))
            tfd_p = td_p.text_frame
            tfd_p.word_wrap = True
            pd_p = tfd_p.paragraphs[0]
            pd_p.text = phase.get("description", "Awaiting data.")
            pd_p.font.name = self.theme.FONT_MEDIUM
            pd_p.font.size = Pt(10.5)
            pd_p.font.color.rgb = self.SLATE_TEXT
            pd_p.line_spacing = 1.05

        # --- B. BUSINESS QUESTION & PROJECT GOAL ---
        bq_left = left + plan_w + spacing
        
        # Header
        t_bq_h = slide.shapes.add_textbox(bq_left, top - Inches(0.1), bq_w, Inches(0.4))
        p_bq_h = t_bq_h.text_frame.paragraphs[0]
        p_bq_h.text = "CORE PROJECT GOAL"
        p_bq_h.font.name = self.theme.FONT_BOLD
        p_bq_h.font.size = Pt(13)
        p_bq_h.font.color.rgb = self.EMERALD

        # BQ Box (Branded Accent)
        bq_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bq_left, top + Inches(0.4), bq_w, height - Inches(0.5))
        bq_box.fill.solid()
        bq_box.fill.fore_color.rgb = self.CARD_BG
        bq_box.line.color.rgb = self.EMERALD
        bq_box.line.width = Pt(2)
        
        # BQ Decorative Icon / Question Mark
        q_mark = slide.shapes.add_textbox(bq_left + bq_w - Inches(0.5), top + Inches(0.45), Inches(0.4), Inches(0.4))
        pq = q_mark.text_frame.paragraphs[0]
        pq.text = "?"
        pq.font.size = Pt(28)
        pq.font.bold = True
        pq.font.color.rgb = self.EMERALD
        pq.alignment = PP_ALIGN.RIGHT

        # BQ Text
        bq_text = insights.get("business_question") or insights.get("project_goal")
        if not bq_text:
            brand = metadata.get("brand", "Brand")
            category = metadata.get("category", "Category")
            bq_text = f"How can {brand} unlock incremental volume and brand stickiness within the {category} landscape?"

        t_bq = slide.shapes.add_textbox(bq_left + Inches(0.2), top + Inches(0.7), bq_w - Inches(0.5), height - Inches(0.9))
        tf_bq = t_bq.text_frame
        tf_bq.word_wrap = True
        p_bq = tf_bq.paragraphs[0]
        p_bq.text = bq_text
        p_bq.font.name = self.theme.FONT_BOLD
        p_bq.font.size = Pt(17)
        p_bq.font.color.rgb = self.WHITE
        p_bq.alignment = PP_ALIGN.CENTER

    def _derive_roadmap_phases(self, raw_data: List[Any], metadata: Dict[str, Any]) -> List[Dict[str, str]]:
        """Systematically derive three strategic phases from available findings."""
        base_phases = [
            {"title": "Foundational Audit", "description": "Identify core sentiment drivers and competitive gaps."},
            {"title": "Strategic Activation", "description": "Deploy targeted messaging to address unmet consumer needs."},
            {"title": "Leadership & Scale", "description": "Optimize operations for long-term category dominance."}
        ]
        
        # Attempt to populate with real findings
        for i in range(min(3, len(raw_data))):
            item = raw_data[i]
            description = ""
            if isinstance(item, dict):
                description = item.get("finding") or item.get("description") or str(item)
            else:
                description = str(item)
            
            # Smart shortening for card fit
            if len(description) > 120:
                description = description[:117] + "..."
            
            base_phases[i]["description"] = description
            
            # Dynamic titles based on category/research
            if i == 0: base_phases[0]["title"] = f"{metadata.get('brand', 'Brand')} Diagnosis"
            if i == 1: base_phases[1]["title"] = "Market Penetration"
            if i == 2: base_phases[2]["title"] = "Growth Optimization"

        return base_phases
