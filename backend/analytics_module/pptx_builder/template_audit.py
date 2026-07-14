import sys
import os
import json
import logging
from pathlib import Path
from pptx import Presentation

# Add parent directory to path to allow absolute imports when run as module
sys.path.append(str(Path(__file__).resolve().parent.parent.parent.parent))

try:
    from backend.analytics_module.pptx_builder.template_adapter import TemplateAdapter
except ImportError:
    # Fallback for direct script execution in the same directory
    from template_adapter import TemplateAdapter

def run_audit(template_path: str):
    print(f"\n{'='*60}")
    print(f" PPTX TEMPLATE AUDIT: {template_path}")
    print(f"{'='*60}")
    
    if not os.path.exists(template_path):
        print(f"❌ ERROR: Template file not found at: {template_path}")
        return

    try:
        prs = Presentation(template_path)
        adapter = TemplateAdapter()
        report = adapter.audit(prs)

        print("\n[VITAL ROLES STATUS]")
        for item in report["found"]:
            print(f"  ✅ {item}")
        for item in report["missing"]:
            print(f"  ❌ MISSING: {item}")

        if report["missing"]:
            print("\n⚠️ WARNING: Some roles are missing. Generation will use hard-coded index fallbacks.")

        print("\n[TEMPLATE STRUCTURE ENUMERATION]")
        for layout in report["all_layouts"]:
            print(f"\nIndex {layout['index']:2}: \"{layout['name']}\"")
            if not layout['placeholders']:
                print("   (No native placeholders)")
            for ph in layout['placeholders']:
                print(f"   • Ph [{ph['idx']}] - {ph['name']} ({ph['type']})")

        print(f"\n{'='*60}")
        print(f" Audit Complete: {len(report['found'])} roles found, {len(report['missing'])} roles missing.")
        print(f"{'='*60}\n")

    except Exception as e:
        print(f"\n☢️ FATAL: Audit pipeline failed: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    # Default path relative to project root
    default_template = "backend/resources/analytics/Marketeers_Template_2025.pptx"
    
    # Try to resolve relative to script location if root-relative fails
    if not os.path.exists(default_template):
        alt_path = Path(__file__).parent.parent / "resources" / "analytics" / "Marketeers_Template_2025.pptx"
        if alt_path.exists():
            default_template = str(alt_path)

    path = sys.argv[1] if len(sys.argv) > 1 else default_template
    run_audit(path)
