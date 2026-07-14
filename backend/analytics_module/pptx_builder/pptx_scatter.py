import logging
from typing import Any, Dict, List

from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_MARKER_STYLE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.slide import Slide
from pptx.util import Inches, Pt

from .base_builder import BaseChartBuilder
from .builder_render_status import BuilderEmptyDataError
from .chart_builder_runtime import apply_xy_point_labels, set_marker_style

logger = logging.getLogger(__name__)


class PPTXScatter(BaseChartBuilder):
    """Builder for competitive XY scatter maps."""

    def render(self, slide: Slide, chart_data: Dict[str, Any]) -> None:
        data_payload = chart_data.get("data", {})
        datasets = data_payload.get("datasets", [])

        if not datasets:
            raise BuilderEmptyDataError("No datasets found for scatter chart.")

        chart_data_obj = XyChartData()
        all_x: List[float] = []
        all_y: List[float] = []

        for index, dataset in enumerate(datasets):
            if isinstance(dataset, str):
                label = dataset
                points: List[Any] = []
            else:
                label = dataset.get("label", f"Point {index}")
                points = dataset.get("data", [])

            series = chart_data_obj.add_series(label)
            for point in points:
                if isinstance(point, dict):
                    x, y = point.get("x", 0), point.get("y", 0)
                else:
                    x, y = 0, 0
                series.add_data_point(x, y)
                all_x.append(x)
                all_y.append(y)

        if not all_x or not all_y:
            raise BuilderEmptyDataError("Scatter chart has no plotted points.")

        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER,
            self.layout.CHART_LEFT,
            self.layout.CHART_TOP,
            self.layout.CHART_WIDTH,
            self.layout.CHART_HEIGHT,
            chart_data_obj,
        )
        chart = chart_shape.chart

        self._style_axes(chart, all_x, all_y)
        self._style_series(chart)
        apply_xy_point_labels(chart, datasets)

        mean_x = sum(all_x) / len(all_x)
        mean_y = sum(all_y) / len(all_y)
        self._add_quadrant_lines(slide, chart, mean_x, mean_y)

    def _style_axes(self, chart: Any, all_x: list, all_y: list):
        x_axis = chart.value_axis
        x_axis.minimum_scale = 0
        x_axis.maximum_scale = 10
        x_axis.has_major_gridlines = False
        x_axis.tick_labels.font.size = Pt(9)
        x_axis.tick_labels.font.name = self.theme.FONT_LIGHT

        y_axis = chart.category_axis
        y_axis.minimum_scale = 0
        max_y = max(all_y) if all_y else 1.0
        y_axis.maximum_scale = max_y * 1.2
        y_axis.visible = True
        y_axis.tick_labels.font.size = Pt(9)

    def _style_series(self, chart: Any):
        for i, series in enumerate(chart.series):
            set_marker_style(series.marker, XL_MARKER_STYLE.CIRCLE)
            series.marker.size = 8
            # Marker-level fill for OXML compliance
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = self.theme.get_color(i)
            series.marker.format.line.visible = False
            series.format.line.visible = False

    def _add_quadrant_lines(self, slide: Slide, chart: Any, mean_x: float, mean_y: float):
        ratio_x = mean_x / 10
        line_x = self.layout.CHART_LEFT + (self.layout.CHART_WIDTH * ratio_x)

        v_line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            int(line_x),
            int(self.layout.CHART_TOP),
            int(line_x),
            int(self.layout.CHART_TOP + self.layout.CHART_HEIGHT),
        )
        v_line.line.color.rgb = self.theme.get_rgb_by_name("brand_slate")
        v_line.line.width = Pt(1.5)

        y_max = chart.category_axis.maximum_scale
        ratio_y = mean_y / y_max if y_max else 0
        line_y = (self.layout.CHART_TOP + self.layout.CHART_HEIGHT) - (self.layout.CHART_HEIGHT * ratio_y)

        h_line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            int(self.layout.CHART_LEFT),
            int(line_y),
            int(self.layout.CHART_LEFT + self.layout.CHART_WIDTH),
            int(line_y),
        )
        h_line.line.color.rgb = self.theme.get_rgb_by_name("brand_slate")
        h_line.line.width = Pt(1.5)

        self._add_quadrant_desc(slide, line_x + Inches(0.1), self.layout.CHART_TOP + Inches(0.1), "KEY DRIVERS")
        self._add_quadrant_desc(slide, self.layout.CHART_LEFT + Inches(0.1), self.layout.CHART_TOP + Inches(0.1), "CRITICAL GAP")

    def _add_quadrant_desc(self, slide: Slide, x: float, y: float, text: str):
        box = slide.shapes.add_textbox(x, y, Inches(1.5), Inches(0.3))
        paragraph = box.text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.size = Pt(8)
        paragraph.font.name = self.theme.FONT_BOLD
        paragraph.font.color.rgb = self.theme.SUBTITLE_COLOR
