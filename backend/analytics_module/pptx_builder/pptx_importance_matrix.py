import logging
from typing import Dict, Any
from pptx.slide import Slide
from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from .base_builder import BaseChartBuilder
from .builder_render_status import BuilderEmptyDataError
from .chart_builder_runtime import apply_xy_point_labels

logger = logging.getLogger(__name__)

class PPTXImportanceMatrix(BaseChartBuilder):
    """
    Builder for Importance vs Performance Matrices (Heatmaps).
    Implements a 4-quadrant scatter chart.
    """

    def render(self, slide: Slide, chart_data: Dict[str, Any]) -> None:
        data_payload = chart_data.get("data", {})
        datasets = data_payload.get("datasets", [])
        
        if not datasets:
            raise BuilderEmptyDataError("No datasets found for importance matrix chart.")

        # 1. Populate XY Data
        chart_data_obj = XyChartData()
        for ds in datasets:
            series = chart_data_obj.add_series(ds.get("label", ""))
            for point in ds.get("data", []):
                series.add_data_point(point.get("x", 0), point.get("y", 0))

        # 2. Add Scatter Chart
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER,
            self.layout.CHART_LEFT, self.layout.CHART_TOP,
            self.layout.CHART_WIDTH, self.layout.CHART_HEIGHT,
            chart_data_obj
        )
        chart = chart_shape.chart
        
        # 3. Apply Premium Matrix Styling
        self._style_matrix_axes(chart)
        self._style_quadrants(slide, chart_shape)
        self._style_points(chart)
        apply_xy_point_labels(chart, datasets)

    def _style_matrix_axes(self, chart: Any):
        # Value Axis (X) - Importance
        x_axis = chart.value_axis
        x_axis.minimum_scale = 0
        x_axis.maximum_scale = 1.0
        x_axis.has_title = True
        x_axis.axis_title.text_frame.text = "Relative Importance"
        
        # Category Axis (Y) - Performance
        y_axis = chart.category_axis
        y_axis.minimum_scale = 0
        y_axis.maximum_scale = 10.0
        y_axis.has_title = True
        y_axis.axis_title.text_frame.text = "Performance (Mean)"

    def _style_quadrants(self, slide: Slide, chart_shape: Any):
        """Adds subtle quadrant overlays (Background colors)."""
        # We can add 4 colored rectangles behind the plot area
        # For now, we'll just add crosshair lines
        
        mid_x = int(chart_shape.left + (chart_shape.width / 2))
        mid_y = int(chart_shape.top + (chart_shape.height / 2))
        
        # Horizontal Divider
        h_line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, 
            int(chart_shape.left), mid_y,
            int(chart_shape.left + chart_shape.width), mid_y
        )
        h_line.line.color.rgb = self.theme.get_rgb_by_name("brand_light_gray")
        h_line.line.width = Pt(0.75)
        
        # Vertical Divider
        v_line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, 
            mid_x, int(chart_shape.top),
            mid_x, int(chart_shape.top + chart_shape.height)
        )
        v_line.line.color.rgb = self.theme.get_rgb_by_name("brand_light_gray")
        v_line.line.width = Pt(0.75)

    def _style_points(self, chart: Any):
        """Styles scatter points with brand colors."""
        for i, series in enumerate(chart.series):
            series.marker.size = 12
            series.format.line.visible = False # No lines between points
            
            # Use theme color
            color = self.theme.get_color(i)
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = color
            series.marker.format.line.visible = False
            
            # Add labels if available (e.g. attribute names)
            # labels.position = XL_DATA_LABEL_POSITION.RIGHT
