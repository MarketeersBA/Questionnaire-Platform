from pptx import Presentation
import sys

try:
    pres = Presentation('backend/resources/analytics/template.pptx')
    for i, slide in enumerate(pres.slides):
        title = slide.shapes.title.text if slide.shapes.title else "No Title"
        print(f"Slide {i}: {title}")
        if i > 20: break # Just a sample
except Exception as e:
    print(f"Error: {e}")
