from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

def inspect_charts():
    pres = Presentation('backend/resources/analytics/template.pptx')
    for i, slide in enumerate(pres.slides):
        for shape in slide.shapes:
            if shape.has_chart:
                chart = shape.chart
                print(f"Slide {i} has chart: {chart.chart_type} ({chart.chart_title.text_frame.text if chart.has_title else 'No Title'})")

if __name__ == "__main__":
    inspect_charts()
