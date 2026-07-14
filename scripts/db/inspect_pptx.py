from pptx import Presentation
import os

path = 'backend/reports/V2_Report_69d3aeb468f389e465151c53_20260513_085648.pptx'
if os.path.exists(path):
    p = Presentation(path)
    print(f"Total Slides: {len(p.slides)}")
    print(f"Presentation Dimensions: {p.slide_width.inches} x {p.slide_height.inches}")
    for i, s in enumerate(p.slides):
        title = s.shapes.title.text if s.shapes.title else "NO TITLE"
        layout = s.slide_layout.name if s.slide_layout else "UNKNOWN"
        shape_count = len(s.shapes)
        print(f"Slide {i:02d}: [{layout}] {title} ({shape_count} shapes)")
        
        # Check for Overlaps (Simple Bounding Box Collision)
        boxes = []
        for sh in s.shapes:
            if hasattr(sh, "left"):
                boxes.append({
                    "name": sh.name,
                    "x1": sh.left.inches, "y1": sh.top.inches,
                    "x2": sh.left.inches + sh.width.inches, "y2": sh.top.inches + sh.height.inches
                })
        
        overlap_count = 0
        for idx1 in range(len(boxes)):
            for idx2 in range(idx1 + 1, len(boxes)):
                b1 = boxes[idx1]
                b2 = boxes[idx2]
                # Check intersection
                if not (b1["x2"] <= b2["x1"] or b1["x1"] >= b2["x2"] or b1["y2"] <= b2["y1"] or b1["y1"] >= b2["y2"]):
                    overlap_count += 1
                    if overlap_count < 5:
                        print(f"  [OVERLAP] {b1['name']} and {b2['name']}")
        if overlap_count > 0:
            print(f"  TOTAL OVERLAPS: {overlap_count}")
else:
    print(f"File not found: {path}")
